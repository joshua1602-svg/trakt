#!/usr/bin/env python3
"""tests/test_deck_generation_route.py

The React "generate investor deck" action, proven end to end.

Until now a deck existed only as a side effect of a nightly pipeline run: the UI
could discover and download one, but nobody could ask for one. This exercises the
whole chain the button drives —

    POST /mi/decks/generate
      → artefact.investor_pack.generate  (scope, tenant, portfolio, period)
      → pptx_stage.generate_investor_pptx   (the SAME stage the pipeline calls)
      → mi_agent_pptx generator + publication gates
      → durable deck store
    GET  /mi/decks/generate/{jobId}       → completed
    GET  /mi/decks                        → the new deck is discoverable
    GET  /mi/decks/download               → the real PPTX comes back

— against a filesystem-backed blob store, with the real generator. Nothing about
the deck is mocked: the file the test downloads is a PowerPoint file that
python-pptx opens and that passed its publication gates.

The governance properties are asserted alongside, because a route that produces
an artefact is a bigger target than one that reads it: the tenant comes from the
authenticated context, another tenant's job is invisible, concurrent requests for
the same scope do not race, and a blocked deck is reported as blocked rather than
as a success or a failure.
"""

from __future__ import annotations

import json
import sys
import time
from pathlib import Path

import pytest

_REPO_ROOT = Path(__file__).resolve().parents[1]
if str(_REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(_REPO_ROOT))

pd = pytest.importorskip("pandas")

from trakt_core.context import ExecutionContext  # noqa: E402
from trakt_core.envelope import STATUS_BLOCKED, STATUS_SUCCESS  # noqa: E402
from trakt_core.errors import ErrorCode  # noqa: E402
from trakt_core.tenancy import TenantRecord, TenantRegistry  # noqa: E402

CLIENT = "acme"
_CENTRAL = "18_central_lender_tape.csv"


# --------------------------------------------------------------------------- #
# A governed two-period book, in the layout the run discovery actually walks.
# --------------------------------------------------------------------------- #

def _loan(i, pid, ptype, balance, *, cut, region, ltv=45.0):
    return {
        "unique_identifier": f"{pid}_L{i:04d}",
        "source_portfolio_id": pid,
        "source_portfolio_type": ptype,
        "source_portfolio_label": pid.replace("_", " ").title(),
        "current_outstanding_balance": balance,
        "current_principal_balance": balance,
        "original_principal_balance": balance * 1.04,
        "current_valuation_amount": balance / (ltv / 100.0),
        "original_valuation_amount": balance / (ltv / 100.0),
        "current_loan_to_value": ltv,
        "original_loan_to_value": ltv / 100 - 0.02,
        "current_interest_rate": 7.1,
        "youngest_borrower_age": 72,
        "geographic_region_collateral": region,
        "collateral_geography": region,
        "origination_channel": "Direct",
        "origination_date": "2021-04-01",
        "data_cut_off_date": cut,
    }


_PRIOR = [("direct_001", "direct", 300_000.0, "London"),
          ("direct_001", "direct", 200_000.0, "South East"),
          ("acquired_001", "acquired", 250_000.0, "Wales")]
_CURRENT = [("direct_001", "direct", 300_000.0, "London"),
            ("direct_001", "direct", 260_000.0, "South East"),
            ("direct_001", "direct", 120_000.0, "South East"),
            ("acquired_001", "acquired", 230_000.0, "Wales")]

_PERIODS = (("mi_2026_05", "2026-05-31", _PRIOR),
            ("mi_2026_06", "2026-06-30", _CURRENT))


def _write_runs(root: Path, client: str = CLIENT) -> Path:
    for run_id, date, spec in _PERIODS:
        central = root / client / run_id / "central"
        central.mkdir(parents=True, exist_ok=True)
        rows = [_loan(i, pid, ptype, bal, cut=date, region=region)
                for i, (pid, ptype, bal, region) in enumerate(spec, start=1)]
        pd.DataFrame(rows).to_csv(central / _CENTRAL, index=False)
    return root


@pytest.fixture()
def deployment(tmp_path, monkeypatch):
    """One tenant, two completed runs, a filesystem deck store, auth off."""
    root = _write_runs(tmp_path / "runs")
    monkeypatch.setenv("MI_AGENT_ONBOARDING_OUTPUT_ROOT", str(root))
    monkeypatch.setenv("MI_AGENT_CLIENT_ID", CLIENT)
    monkeypatch.setenv("MI_AGENT_PIPELINE_ROOT", str(root))
    monkeypatch.setenv("TRAKT_LOCAL_BLOB_ROOT", str(tmp_path / "blobstore"))
    monkeypatch.setenv("TRAKT_INVESTOR_PPTX_PERSIST", "true")
    monkeypatch.setenv("TRAKT_INVESTOR_PPTX_ON_DEMAND", "true")
    monkeypatch.setenv("MI_AGENT_AUTH_ENABLED", "false")
    # Deck DISCOVERY must read the durable store the publish writes to, not a
    # local override — otherwise the download half of this test proves nothing.
    monkeypatch.delenv("MI_AGENT_DECK_ROOT", raising=False)
    monkeypatch.delenv("AZURE_STORAGE_CONNECTION_STRING", raising=False)
    monkeypatch.delenv("TRAKT_BLOB_CONNECTION", raising=False)

    from mi_agent_api import deck_generation
    deck_generation.reset_jobs()
    yield tmp_path
    deck_generation.reset_jobs()


def _context(tenant=CLIENT, **kw):
    return ExecutionContext(tenant_id=tenant, actor_id="user-1", **kw)


def _registry(*tenant_ids):
    return TenantRegistry(
        {t: TenantRecord(tenant_id=t, default_portfolio_id=t) for t in tenant_ids},
        configured=True)


def _deps(*tenants):
    from mi_agent_api.dependencies import build_dependencies
    return build_dependencies(tenant_registry=_registry(*(tenants or (CLIENT,))))


def _inline(work):
    """Run the generation on the calling thread, so a test asserts an outcome
    rather than a race."""
    work()


# --------------------------------------------------------------------------- #
# The whole chain, through HTTP.
# --------------------------------------------------------------------------- #

def test_react_button_generates_a_real_downloadable_deck(deployment):
    """The end-to-end proof: request → generate → discover → download."""
    from fastapi.testclient import TestClient
    from mi_agent_api.app import app

    client = TestClient(app)

    # Nothing published yet.
    assert client.get("/mi/decks").json()["available"] is False

    accepted = client.post("/mi/decks/generate", json={})
    assert accepted.status_code == 202, accepted.text
    job = accepted.json()
    assert job["ok"] is True
    assert job["state"] in ("queued", "generating")
    job_id = job["jobId"]

    deadline = time.time() + 240
    while time.time() < deadline:
        status = client.get(f"/mi/decks/generate/{job_id}")
        assert status.status_code in (200, 202), status.text
        body = status.json()
        if body["state"] in ("completed", "blocked", "failed"):
            break
        time.sleep(0.5)
    else:
        pytest.fail("generation did not finish within the timeout")

    assert body["state"] == "completed", body
    assert body["failedGates"] == []
    assert body["gateCount"] and body["gateCount"] >= 18
    assert body["reportingPeriod"] == "2026-06-30"
    assert body["runId"] == "mi_2026_06"

    # Discoverable...
    index = client.get("/mi/decks").json()
    assert index["available"] is True
    assert index["latest"]["period"] == "2026-06"

    # ...and downloadable as a real PowerPoint file.
    got = client.get("/mi/decks/download")
    assert got.status_code == 200
    assert got.content[:4] == b"PK\x03\x04"
    assert "investor_deck" in got.headers.get("content-disposition", "")

    import io

    from pptx import Presentation
    deck = Presentation(io.BytesIO(got.content))
    assert len(deck.slides.__iter__.__self__._sldIdLst) >= 10
    text = "\n".join(sh.text_frame.text for s in deck.slides
                     for sh in s.shapes if sh.has_text_frame)
    assert "Executive Summary" in text


# --------------------------------------------------------------------------- #
# Governance. These run the capability directly with a synchronous executor —
# no HTTP, no thread, no generation.
# --------------------------------------------------------------------------- #

def _request(deployment, *, context=None, deps=None, executor=lambda w: None, **kw):
    from mi_agent_api import deck_generation
    return deck_generation.request_generation(
        context or _context(), dependencies=deps or _deps(),
        executor=executor, **kw)


def test_tenant_comes_from_the_context_not_the_request(deployment):
    """Naming another tenant's portfolio cannot produce its pack."""
    result = _request(deployment, deps=_deps(CLIENT, "other_tenant"),
                      portfolio_id="other_tenant")
    assert result.status == STATUS_BLOCKED
    assert result.error.code == ErrorCode.TENANT_MISMATCH
    assert result.result is None


def test_generation_requires_the_generate_scope(deployment):
    """Reading artefacts does not imply producing them."""
    from trakt_core.context import SCOPE_ARTEFACT_READ, SCOPE_PORTFOLIO_READ

    reader = _context(scopes=frozenset({SCOPE_PORTFOLIO_READ, SCOPE_ARTEFACT_READ}))
    result = _request(deployment, context=reader)
    assert result.status == STATUS_BLOCKED
    assert result.error.code == ErrorCode.SCOPE_MISSING


def test_a_deployment_can_switch_on_demand_generation_off(deployment, monkeypatch):
    monkeypatch.setenv("TRAKT_INVESTOR_PPTX_ON_DEMAND", "false")
    result = _request(deployment)
    assert result.status == STATUS_BLOCKED
    assert result.error.code == ErrorCode.PERMISSION_DENIED


def test_an_unknown_period_is_refused_before_a_job_is_accepted(deployment):
    """"There is no data for that month" is an answer to the request, not a job
    that fails a moment later."""
    result = _request(deployment, period="2019-01")
    assert result.status != STATUS_SUCCESS
    assert result.error.code == ErrorCode.DATA_SOURCE_UNAVAILABLE


@pytest.mark.parametrize("period", ["../../etc", "2026-13-99-", "latest; rm", "**"])
def test_a_malformed_period_is_rejected(deployment, period):
    result = _request(deployment, period=period)
    assert result.error.code == ErrorCode.INVALID_INPUT


@pytest.mark.parametrize("context", ["../secrets", "direct;drop", "a" * 65])
def test_a_malformed_portfolio_context_is_rejected(deployment, context):
    result = _request(deployment, portfolio_context=context)
    assert result.error.code == ErrorCode.INVALID_INPUT


def test_a_dated_period_resolves_the_run_for_that_month(deployment):
    result = _request(deployment, period="2026-05")
    assert result.status == STATUS_SUCCESS
    assert result.result.reporting_period == "2026-05-31"
    assert result.result.run_id == "mi_2026_05"


def test_concurrent_requests_for_the_same_scope_join_one_job(deployment):
    """Two identical requests must not generate the same deck twice — they would
    race for the same output file."""
    first = _request(deployment)
    second = _request(deployment)
    assert first.status == second.status == STATUS_SUCCESS
    assert first.result.job_id == second.result.job_id


def test_a_different_scope_gets_its_own_job(deployment):
    total = _request(deployment)
    direct = _request(deployment, portfolio_context="direct")
    assert total.result.job_id != direct.result.job_id


def test_an_idempotency_key_replays_a_finished_job(deployment):
    """Without a key, asking again after the first finished generates again —
    new data is the usual reason to ask."""
    from mi_agent_api import deck_generation

    first = _request(deployment, idempotency_key="req-1")
    first.result.state = deck_generation.STATE_COMPLETED
    replay = _request(deployment, idempotency_key="req-1")
    assert replay.result.job_id == first.result.job_id

    fresh = _request(deployment, idempotency_key="req-2")
    assert fresh.result.job_id != first.result.job_id


def test_another_tenants_job_is_not_visible(deployment):
    """Reported as unknown rather than forbidden: whether another tenant has a
    generation in flight is itself information."""
    from mi_agent_api import deck_generation

    mine = _request(deployment)
    theirs = deck_generation.get_generation(
        _context(tenant="other_tenant"), mine.result.job_id,
        dependencies=_deps("other_tenant"))
    assert theirs.error.code == ErrorCode.ARTEFACT_NOT_FOUND


def test_an_unknown_job_id_is_not_found(deployment):
    from mi_agent_api import deck_generation

    result = deck_generation.get_generation(_context(), "deckgen_nope",
                                            dependencies=_deps())
    assert result.error.code == ErrorCode.ARTEFACT_NOT_FOUND


# --------------------------------------------------------------------------- #
# Outcome mapping — the three ways a generation ends.
# --------------------------------------------------------------------------- #

def _finish(deployment, artifact, monkeypatch):
    from apps.blob_trigger_app import pptx_stage

    monkeypatch.setattr(pptx_stage, "generate_investor_pptx",
                        lambda *a, **k: artifact)
    return _request(deployment, executor=_inline).result


def test_a_gate_failure_is_reported_as_blocked_not_failed(deployment, monkeypatch):
    """The pack exists for an operator to diagnose, and nothing was published —
    that is a third outcome, and the UI must be able to say so."""
    from mi_agent_api import deck_generation

    job = _finish(deployment, {
        "status": "generated_not_published",
        "preflight": {"publishable": False, "gates": [{}] * 18,
                      "failed_gates": ["no_empty_slides", "totals_reconcile"]},
    }, monkeypatch)
    assert job.state == deck_generation.STATE_BLOCKED
    assert job.failed_gates == ["no_empty_slides", "totals_reconcile"]
    assert "not been released" in job.message
    assert "previously published deck is unchanged" in job.message


def test_a_generator_fault_is_reported_without_leaking_internals(deployment,
                                                                monkeypatch):
    from apps.blob_trigger_app import pptx_stage
    from mi_agent_api import deck_generation

    def _boom(*a, **k):
        raise RuntimeError("/srv/secret/path/run_42 exploded")

    monkeypatch.setattr(pptx_stage, "generate_investor_pptx", _boom)
    job = _request(deployment, executor=_inline).result
    assert job.state == deck_generation.STATE_FAILED
    assert "/srv/secret" not in (job.message or "")
    assert job.error_code == ErrorCode.INTERNAL_ERROR


def test_the_job_payload_carries_no_paths(deployment, monkeypatch):
    job = _finish(deployment, {"status": "available",
                               "preflight": {"publishable": True, "gates": []}},
                  monkeypatch)
    blob = json.dumps(job.to_payload())
    assert "/" not in blob.replace("2026-06-30", "").replace("2026-05-31", "")
    assert "run_dir" not in blob and "local_path" not in blob


# --------------------------------------------------------------------------- #
# The scope actually reaches the generator.
# --------------------------------------------------------------------------- #

def test_the_requested_scope_is_passed_to_the_one_generation_stage(deployment,
                                                                  monkeypatch):
    """Both channels reach the generator through pptx_stage — the on-demand path
    must not assemble its own argv."""
    from apps.blob_trigger_app import pptx_stage

    seen = {}

    def _capture(run_dir, **kw):
        seen["run_dir"] = Path(run_dir)
        seen.update(kw)
        return {"status": "available", "preflight": {"publishable": True, "gates": []}}

    monkeypatch.setattr(pptx_stage, "generate_investor_pptx", _capture)
    _request(deployment, portfolio_context="acquired", period="2026-05",
             executor=_inline)

    assert seen["portfolio_context"] == "acquired"
    assert seen["client_id"] == CLIENT and seen["tenant_id"] == CLIENT
    assert seen["as_of_date"] == "2026-05-31"
    assert seen["source_run_id"] == "mi_2026_05"
    assert seen["run_dir"].name == "mi_2026_05"
    # The multi-period surfaces need the client's run ROOT, not this run's parent.
    assert Path(seen["output_root"]).name == "runs"


def test_the_stage_forwards_the_scope_to_the_generator_cli(monkeypatch, tmp_path):
    """One level lower: pptx_stage translates it into the generator's own flag."""
    from apps.blob_trigger_app import pptx_stage

    argv_seen = []

    def _fake(argv):
        argv_seen.extend(argv)
        out = Path(argv[argv.index("--output") + 1])
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_bytes(b"PK\x03\x04")
        return 0

    monkeypatch.setattr(pptx_stage, "_invoke_generator", _fake)
    monkeypatch.setenv("TRAKT_INVESTOR_PPTX_PERSIST", "false")
    pptx_stage.generate_investor_pptx(
        tmp_path / "run", client_name=CLIENT, as_of_date="2026-06-30",
        portfolio_context="direct", client_id=CLIENT, tenant_id=CLIENT,
        output_root=str(tmp_path))

    assert "--portfolio-context" in argv_seen
    assert argv_seen[argv_seen.index("--portfolio-context") + 1] == "direct"
    assert argv_seen[argv_seen.index("--output-root") + 1] == str(tmp_path)


def test_the_pipeline_path_is_unchanged_when_no_scope_is_given(monkeypatch,
                                                              tmp_path):
    """The nightly run passes none of the new arguments, and its argv must look
    exactly as it did."""
    from apps.blob_trigger_app import pptx_stage

    argv_seen = []

    def _fake(argv):
        argv_seen.extend(argv)
        out = Path(argv[argv.index("--output") + 1])
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_bytes(b"PK\x03\x04")
        return 0

    monkeypatch.setattr(pptx_stage, "_invoke_generator", _fake)
    monkeypatch.setenv("TRAKT_INVESTOR_PPTX_PERSIST", "false")
    pptx_stage.generate_investor_pptx(tmp_path / "run", client_name=CLIENT,
                                      as_of_date="2026-06-30")

    for flag in ("--portfolio-context", "--client-id", "--tenant-id", "--output-root"):
        assert flag not in argv_seen
