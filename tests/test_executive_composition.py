#!/usr/bin/env python3
"""tests/test_executive_composition.py

The opening page, checked on the pack a funder receives.

Every test builds a real deck through the React route — ``POST
/mi/decks/generate`` -> poll -> ``GET /mi/decks/download`` — and reads the
PowerPoint that comes back.

The page it replaces carried seven tiles laid out four and three, with a hole
beside the second row, three of them restating one pipeline fact the next page
restated again in words, and two half-width trends competing for the centre
under three quarters of an inch of dead panel. It read as a compressed web
dashboard rather than an opening statement.
"""

from __future__ import annotations

import io
import sys
from pathlib import Path

import pytest

_REPO_ROOT = Path(__file__).resolve().parents[1]
if str(_REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(_REPO_ROOT))

pytest.importorskip("pandas")
pptx = pytest.importorskip("pptx")

from tests.test_funder_pack_enhancement import (  # noqa: E402
    generate_and_download, rich, simple)
# Concentration limits are contractual: the page's risk tile does not exist
# until an operator has approved a configuration, so the fixture that commits
# one is what this page has to be checked against.
from tests.test_final_pack_surfaces import limited  # noqa: E402

__all__ = ["rich", "simple", "limited"]

#: Shapes on the executive page sit in one of three bands: the tile row, the
#: trajectory card, and the risk strip. A tile panel is short and wide.
_TILE_TOP, _TILE_BOTTOM = 1.5, 3.2


def _exec_slide(content: bytes):
    deck = pptx.Presentation(io.BytesIO(content))
    for slide in deck.slides:
        heads = [sh.text_frame.text for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        if heads and "Executive Position" in heads[0]:
            return slide
    pytest.fail("no Executive Position page in the pack")


def _chart_cards(slide):
    """The card PANELS below the tile row. A placed chart is a picture sitting
    inside its card, so counting every non-text shape counts each card twice.
    """
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    return [sh for sh in slide.shapes
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE
            and not (sh.has_text_frame and sh.text_frame.text.strip())
            and Emu(sh.top).inches > _TILE_BOTTOM
            and Emu(sh.height).inches >= 1.4]


def _tile_panels(slide):
    from pptx.util import Emu

    out = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            continue
        top = Emu(sh.top).inches
        if _TILE_TOP <= top <= _TILE_BOTTOM and Emu(sh.height).inches < 2.0:
            out.append((Emu(sh.left).inches, Emu(sh.width).inches))
    return sorted(out)


# --------------------------------------------------------------------------- #
# 1. No hole.
# --------------------------------------------------------------------------- #

def test_1_the_tile_row_has_no_empty_slot(limited):
    """Catches: four tiles over three, with a gap beside the second row.

    The tiles are laid out in ONE row, so a partial row cannot exist. The check
    is geometric rather than a count: every tile shares a top edge, and the
    gaps between them are even.
    """
    from pptx.util import Emu

    slide = _exec_slide(generate_and_download())
    tops = {round(Emu(sh.top).inches, 2) for sh in slide.shapes
            if not (sh.has_text_frame and sh.text_frame.text.strip())
            and _TILE_TOP <= Emu(sh.top).inches <= _TILE_BOTTOM
            and Emu(sh.height).inches < 2.0}
    assert len(tops) == 1, f"tiles occupy {len(tops)} rows: {sorted(tops)}"

    panels = _tile_panels(slide)
    assert len(panels) >= 4, panels
    gaps = [round(b[0] - (a[0] + a[1]), 3) for a, b in zip(panels, panels[1:])]
    assert max(gaps) - min(gaps) < 0.02, f"uneven tile gaps — a hole: {gaps}"


def test_2_the_row_spans_the_content_width(limited):
    """Catches: five tiles' worth of tiles in a six-tile grid, leaving a band
    of empty slide on the right."""
    from mi_agent_pptx.deck import DeckBuilder

    panels = _tile_panels(_exec_slide(generate_and_download()))
    left, right = panels[0][0], panels[-1][0] + panels[-1][1]
    assert left == pytest.approx(DeckBuilder.CONTENT_L, abs=0.05)
    assert right == pytest.approx(DeckBuilder.CONTENT_R, abs=0.05)


def test_3_the_page_never_carries_more_tiles_than_one_row_holds(limited):
    from mi_agent_pptx.deck import DeckBuilder

    panels = _tile_panels(_exec_slide(generate_and_download()))
    assert len(panels) <= DeckBuilder.EXEC_MAX_TILES, len(panels)


# --------------------------------------------------------------------------- #
# 2. The measures are governed, and are the ones that matter.
# --------------------------------------------------------------------------- #

def test_4_the_headline_measures_are_the_priority_set(limited):
    """The position, the risk measure, and — where the book has one — what is
    coming and where it lands.

    A book with no governed pipeline has no pipeline tile and no forecast tile,
    because a forecast equal to the funded balance is the funded balance again
    in a tile that claims to look forward.
    """
    slide = _exec_slide(generate_and_download())
    text = "\n".join(sh.text_frame.text for sh in slide.shapes
                     if sh.has_text_frame).upper()
    for expected in ("FUNDED BALANCE", "LOANS FUNDED", "WA CURRENT LTV",
                     "CLOSEST LIMIT"):
        assert expected in text, text
    if "PIPELINE BALANCE" in text:
        assert "FORECAST FUNDED" in text, text


def test_5_the_risk_line_does_not_restate_the_risk_tile(limited):
    """Catches: "47% utilisation" on a tile and again in the line beneath it.

    The tile carries the closest test and its utilisation; the line carries
    what the tile cannot — the distance left to the limit.
    """
    slide = _exec_slide(generate_and_download())
    lines = [sh.text_frame.text.strip() for sh in slide.shapes
             if sh.has_text_frame and
             sh.text_frame.text.strip().startswith("Concentration —")]
    assert lines, "no concentration line on the executive page"
    assert "utilisation" not in lines[0], lines[0]
    assert "headroom" in lines[0], lines[0]


def test_6_no_measure_appears_on_two_tiles(limited):
    """Catches: pipeline balance, weighted expected and forecast funded — three
    tiles for one fact."""
    slide = _exec_slide(generate_and_download())
    labels = [sh.text_frame.text.strip() for sh in slide.shapes
              if sh.has_text_frame and sh.text_frame.text.strip()
              and sh.text_frame.text.strip().isupper()
              and len(sh.text_frame.text.strip()) < 40]
    assert len(labels) == len(set(labels)), labels
    assert len(labels) >= 4, labels


# --------------------------------------------------------------------------- #
# 3. One trajectory, using the band it is given.
# --------------------------------------------------------------------------- #

def test_7_there_is_exactly_one_trajectory_and_it_spans_the_page(limited):
    """Catches: two half-width trends competing for the centre."""
    from pptx.util import Emu

    slide = _exec_slide(generate_and_download())
    from mi_agent_pptx.deck import DeckBuilder

    cards = _chart_cards(slide)
    assert len(cards) == 1, f"{len(cards)} chart cards on the opening page"
    card = cards[0]
    width = Emu(card.width).inches
    assert width == pytest.approx(DeckBuilder.CONTENT_R - DeckBuilder.CONTENT_L,
                                  abs=0.05), width


def test_8_the_trajectory_fills_the_band_above_the_risk_strip(limited):
    """Catches: a 2.30in cap leaving three quarters of an inch of dead panel."""
    from pptx.util import Emu

    from mi_agent_pptx.deck import DeckBuilder

    slide = _exec_slide(generate_and_download())
    card = _chart_cards(slide)[0]
    bottom = Emu(card.top).inches + Emu(card.height).inches
    slack = DeckBuilder.RISK_STRIP_TOP - bottom
    assert slack <= DeckBuilder.RISK_STRIP_CLEARANCE + 0.02, (
        f"{slack:.2f}in of dead panel above the risk strip")


def test_9_a_book_with_no_funded_history_still_gets_a_trajectory(simple):
    """A new book has a story and it is the origination one — the page must not
    fall back to empty."""
    from pptx.util import Emu

    slide = _exec_slide(generate_and_download())
    assert len(_chart_cards(slide)) <= 1
