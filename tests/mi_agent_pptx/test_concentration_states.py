"""Concentration reporting, end to end, across the four representative states.

The capability shipped in v2.1 but had only ever been exercised against a book
with no pipeline, so the forward-looking states — the ones a warehouse lender
scrutinises hardest — were proven only in their negative form ("absent when no
pipeline exists").

This fixture supplies a funded book AND an active governed pipeline, with
thresholds chosen so the governed evaluator produces, in one configuration:

  * a comfortably compliant test;
  * an amber / near-limit test;
  * an expected-forecast deterioration into breach;
  * an all-pipeline-converts stress breach;
  * and NO current funded breach.

Thresholds were derived by observing the evaluator's own output on this book
(London 30.00 / 38.66 / 37.04, South East 40.00 / 33.61 / 29.63, Wales 24.00 /
22.69 / 28.89 for funded / expected / stress) and setting limits around it. The
shares are the evaluator's; only the limits are the fixture's.
"""

from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
import pytest
import yaml
from pptx import Presentation

from mi_agent_pptx import concentration as C
from mi_agent_pptx.cli import preflight_path, run as cli_run
from mi_agent_pptx.mi_api import build_dashboard_data

DECK_CONFIG = Path("configs/pptx/investor_pack.yaml")
_CENTRAL = "18_central_lender_tape.csv"

#: Funded composition, chosen so three regions land at usefully different shares.
FUNDED = (("London", 3_000_000.0), ("South East", 4_000_000.0),
          ("Wales", 2_400_000.0), ("Scotland", 600_000.0))
#: Active pipeline. London converts at high probability (so it moves the
#: EXPECTED state); Wales at low probability (so it moves only the STRESS state).
PIPELINE = (("London", 2_000_000.0, 0.8), ("Wales", 1_500_000.0, 0.2))

#: Limits placed around the evaluator's own output to produce each state.
LIMITS = (
    ("Scotland concentration", 50.0, ["Scotland"]),       # 6%  → comfortably compliant
    ("South East concentration", 60.0, ["South East"]),   # 40% → compliant
    ("London concentration", 35.0, ["London"]),           # 30% now, 38.7% expected
    ("Wales concentration", 26.0, ["Wales"]),             # 24% now, 28.9% stress
)


def _funded(i, region, bal, cut, pid="direct_001", ptype="direct"):
    return {"unique_identifier": f"F{i:04d}", "source_portfolio_id": pid,
            "source_portfolio_type": ptype,
            "source_portfolio_label": "Direct Book",
            "current_outstanding_balance": bal, "current_principal_balance": bal,
            "original_principal_balance": bal,
            "current_valuation_amount": bal / 0.45,
            "original_valuation_amount": bal / 0.45,
            "current_loan_to_value": 45.0, "original_loan_to_value": 0.43,
            "current_interest_rate": 7.0, "youngest_borrower_age": 72,
            "collateral_geography": region, "geographic_region_obligor": region,
            "origination_date": "2021-04-01", "data_cut_off_date": cut}


def _pipeline(i, region, bal, prob):
    return {"unique_identifier": f"P{i:04d}",
            "current_outstanding_balance": bal, "current_principal_balance": bal,
            "current_valuation_amount": bal / 0.45,
            "collateral_geography": region, "geographic_region_obligor": region,
            "pipeline_status": "pipeline", "completion_probability": prob,
            "pipeline_stage": "OFFER", "current_loan_to_value": 45.0}


@pytest.fixture
def storage_env(tmp_path, monkeypatch):
    monkeypatch.setenv("TRAKT_STORAGE_BACKEND", "file")
    monkeypatch.setenv("TRAKT_LOCAL_BLOB_ROOT", str(tmp_path / "blob"))
    monkeypatch.setenv("TRAKT_RUNTIME_MODE", "test")
    return tmp_path


@pytest.fixture
def approved_limits(storage_env):
    """An operator-approved configuration, committed through the real store."""
    from apps.blob_trigger_app.storage import open_storage
    from mi_agent.concentration_tests.models import (
        ActiveConfiguration, ActiveTest, ApprovalRecord, SourceEvidence)
    from mi_agent.concentration_tests.store import ConcentrationStore

    tests = [
        ActiveTest(metric_id="geo_region_share", threshold=thr, operator="max",
                   display_name=name, parameters={"regions": regions},
                   evidence=SourceEvidence(
                       source_reference="facility_agreement.pdf",
                       source_text="Schedule 8, clause 3.1"),
                   approval=ApprovalRecord(decision="approved", operator="Operator",
                                           decided_at="2026-01-10T09:00:00+00:00"))
        for name, thr, regions in LIMITS]
    config = ActiveConfiguration(client_id="acme", version=1,
                                 activated_by="Operator",
                                 library_version="1.0.0", tests=tests)
    ConcentrationStore(open_storage(),
                       container="operations-control").commit_configuration(
        "acme", config)
    return config


@pytest.fixture
def run_dir(storage_env):
    """Funded run in the governed layout, plus an active pipeline extract."""
    root = storage_env / "root"
    for run_id, cut in (("mi_2026_05", "2026-05-31"), ("mi_2026_06", "2026-06-30")):
        central = root / "acme" / run_id / "central"
        central.mkdir(parents=True, exist_ok=True)
        rows = [_funded(i, r, b, cut) for i, (r, b) in enumerate(FUNDED, 1)]
        pd.DataFrame(rows).to_csv(central / _CENTRAL, index=False)
        dated = root / cut
        dated.mkdir(parents=True, exist_ok=True)
        pd.DataFrame(rows).to_csv(dated / "platform_canonical_typed.csv", index=False)

    # The governed pipeline source, discoverable under its own root.
    pipe = storage_env / "pipe" / "2026-06-30"
    pipe.mkdir(parents=True, exist_ok=True)
    pd.DataFrame([_pipeline(i, r, b, p)
                  for i, (r, b, p) in enumerate(PIPELINE, 1)]).to_csv(
        pipe / "M2L_KFI_and_Pipeline_2026_06_30.csv", index=False)

    rd = root / "orun_conc"
    (rd / "out_platform").mkdir(parents=True, exist_ok=True)
    pd.DataFrame([_funded(i, r, b, "2026-06-30")
                  for i, (r, b) in enumerate(FUNDED, 1)]).to_csv(
        rd / "out_platform" / "platform_canonical_typed.csv", index=False)
    (rd / "run_state.json").write_text(json.dumps({
        "run_id": "orun_conc", "client_id": "acme",
        "reporting_date": "2026-06-30", "out_root": str(root)}), encoding="utf-8")
    return rd


def _data(run_dir):
    return build_dashboard_data(
        run_dir, client_id="acme", as_of="2026-06-30",
        output_root=str(Path(run_dir).parent),
        pipeline_root=str(Path(run_dir).parent.parent / "pipe"))


def _rows(run_dir):
    return {r["label"]: r for r in C.adapt_tests(_data(run_dir).concentration)}


# --------------------------------------------------------------------------- #
# The four representative states.
# --------------------------------------------------------------------------- #

def test_governed_configuration_resolves_as_approved(run_dir, approved_limits):
    env = _data(run_dir).concentration
    assert env.get("source") == C.SOURCE_APPROVED
    assert len(env.get("tests") or []) == len(LIMITS)


def test_state_1_comfortably_compliant(run_dir, approved_limits):
    row = _rows(run_dir)["Scotland concentration"]
    assert row["status"] == C.STATUS_PASS
    assert row["utilisation"] < 50.0, "should sit well clear of its limit"
    assert row["headroom"] and row["headroom"] > 0


def test_state_2_amber_near_limit(run_dir, approved_limits):
    row = _rows(run_dir)["Wales concentration"]
    assert row["status"] == C.STATUS_WARNING
    assert 85.0 <= row["utilisation"] < 100.0


def test_state_3_expected_forecast_deterioration(run_dir, approved_limits):
    """Within limit today; the governed EXPECTED forecast crosses it."""
    row = _rows(run_dir)["London concentration"]
    assert row["status"] != C.STATUS_BREACH, "must not already be in breach"
    assert row["expected_breach"] is True
    assert row["expected_value"] > row["limit"]
    assert row["expected_utilisation"] > 100.0


def test_state_4_stress_only_breach(run_dir, approved_limits):
    """Within limit today AND on the expected forecast; only the
    all-pipeline-converts stress breaches."""
    row = _rows(run_dir)["Wales concentration"]
    assert row["status"] != C.STATUS_BREACH
    assert row["expected_breach"] is False, "the expected case must not breach"
    assert row["stress_breach"] is True
    assert row["stress_value"] > row["limit"]


def test_state_5_no_current_breach(run_dir, approved_limits):
    """The fixture deliberately contains no current funded breach."""
    assert not [r for r in _rows(run_dir).values()
                if r["status"] == C.STATUS_BREACH]


def test_forward_states_are_evaluated(run_dir, approved_limits):
    assert C.forward_states_available(_data(run_dir).concentration) is True


# --------------------------------------------------------------------------- #
# Parity — the deck must show the governed numbers, not its own.
# --------------------------------------------------------------------------- #

def test_deck_values_equal_the_governed_service(run_dir, approved_limits):
    """Every displayed figure is the evaluator's, to the underlying value."""
    env = _data(run_dir).concentration
    governed = {t["displayName"]: t for t in env["tests"]}
    for row in C.adapt_tests(env):
        src = governed[row["label"]]
        assert row["value"] == src["currentValue"]
        assert row["limit"] == src["threshold"]
        assert row["utilisation"] == src["utilization"]
        assert row["headroom"] == src["headroom"]
        assert row["reporting_date"] == src["reportingDate"]
        expected = src.get("expected") or {}
        assert row["expected_value"] == expected.get("value")
        stress = src.get("fullPipeline") or {}
        assert row["stress_value"] == stress.get("value")


# --------------------------------------------------------------------------- #
# Rendering.
# --------------------------------------------------------------------------- #

def _build(run_dir, out):
    return cli_run([
        "--run-dir", str(run_dir), "--deck-config", str(DECK_CONFIG),
        "--client-name", "Test Client", "--as-of-date", "2026-06-30",
        "--output-root", str(Path(run_dir).parent),
        "--pipeline-root", str(Path(run_dir).parent.parent / "pipe"),
        "--output", str(out)])


def _text(path):
    return "\n".join(sh.text_frame.text for s in Presentation(str(path)).slides
                     for sh in s.shapes if sh.has_text_frame)


def test_concentration_slide_renders_all_three_states(run_dir, approved_limits,
                                                      tmp_path):
    out = tmp_path / "conc.pptx"
    assert _build(run_dir, out) == 0
    text = _text(out)
    assert "Concentration Tests and Headroom" in text
    for label, _thr, _r in LIMITS[-2:]:                # the two interesting ones
        assert label in text
    # The three states must be distinguishable on the page.
    assert "FORECAST BREACH" in text and "STRESS BREACH" in text
    assert "Expected" in text
    assert "Operator-approved" in text


def test_stress_is_never_presented_as_the_expected_outcome(run_dir,
                                                           approved_limits,
                                                           tmp_path):
    out = tmp_path / "stress.pptx"
    _build(run_dir, out)
    low = _text(out).lower()
    if "all-pipeline-converts" in low:
        assert "stress, not the expected outcome" in low


def test_takeaway_states_forecast_breach_without_inventing_actions(
        run_dir, approved_limits, tmp_path):
    out = tmp_path / "takeaway.pptx"
    _build(run_dir, out)
    text = _text(out)
    assert "expected forecast breaches" in text.lower()
    for imperative in ("should ", "must ", "recommend", "we suggest"):
        assert imperative not in text.lower()


def test_watch_list_raises_the_forecast_and_stress_items(run_dir, approved_limits):
    from mi_agent_pptx import watchlist as WL
    types = {i.insight_type for i in _data(run_dir).watchlist["watch"]}
    assert WL.CONCENTRATION_FORECAST in types
    assert WL.CONCENTRATION_STRESS in types


def test_all_gates_pass_with_concentration_enabled(run_dir, approved_limits,
                                                   tmp_path):
    out = tmp_path / "gates.pptx"
    rc = _build(run_dir, out)
    verdict = json.loads(preflight_path(out).read_text(encoding="utf-8"))
    assert rc == 0, verdict["preflight"]["failed_gates"]
    gate = next(g for g in verdict["preflight"]["gates"]
                if g["gate"] == "concentration_reconciles")
    assert gate["passed"] is True
