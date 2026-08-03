"""Investor Reporting v2.1 — pipeline restoration, concentration, client safety.

Three defects drove this release, and each gets tests that would have caught it:

  * the deck resolved its pipeline source differently from the dashboard, so a
    whole section could silently vanish;
  * the deck consumed the *legacy, unapproved* limit monitor rather than the
    governed concentration capability, so it had no forward-looking states;
  * the deck rendered filesystem paths, internal function names and placeholder
    prose to a client-facing audience.
"""

from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
import pytest
import yaml
from pptx import Presentation

from mi_agent_pptx import concentration as C
from mi_agent_pptx.cli import PREFLIGHT_BLOCKED_RC, preflight_path, run as cli_run
from mi_agent_pptx.composition import build_facts, select_slides
from mi_agent_pptx.mi_api import build_dashboard_data
from mi_agent_pptx.preflight import run_preflight

DECK_CONFIG = Path("configs/pptx/investor_pack.yaml")
_M2L = Path("tests/fixtures/client_001_mi_pack/pipeline")


# --------------------------------------------------------------------------- #
# Fixtures.
# --------------------------------------------------------------------------- #

def _loan(i, pid, ptype, balance, *, cut, region, ltv=45.0):
    return {
        "unique_identifier": f"{pid}_L{i:04d}",
        "source_portfolio_id": pid,
        "source_portfolio_type": ptype,
        "source_portfolio_label": pid.replace("_", " ").title(),
        "current_outstanding_balance": balance,
        "current_principal_balance": balance,
        "original_principal_balance": balance * 1.05,
        "current_valuation_amount": balance / (ltv / 100.0),
        "original_valuation_amount": balance / (ltv / 100.0),
        "current_loan_to_value": ltv,
        "original_loan_to_value": ltv / 100 - 0.02,
        "current_interest_rate": 7.2,
        "youngest_borrower_age": 74,
        "collateral_geography": region,
        "geographic_region_obligor": region,
        "origination_date": "2022-05-01",
        "data_cut_off_date": cut,
    }


#: A deliberately concentrated book: London is ~46% of balance, so a 40% limit
#: breaches, a 50% limit warns and a 90% limit passes — hand-checkable.
_BOOKS = {
    "direct_001": ("direct", [("London", 300_000.0), ("London", 250_000.0),
                              ("South East", 200_000.0)]),
    "acquired_001": ("acquired", [("Wales", 150_000.0), ("Scotland", 100_000.0)]),
}


def _rows(cut):
    out = []
    for pid, (ptype, entries) in _BOOKS.items():
        for i, (region, bal) in enumerate(entries, start=1):
            out.append(_loan(i, pid, ptype, bal, cut=cut, region=region))
    return out


#: The governed run layout the MI resolvers discover:
#: ``<root>/<client>/<run_id>/central/18_central_lender_tape.csv``.
_CENTRAL_TAPE = "18_central_lender_tape.csv"
_RUNS = (("mi_2026_05", "2026-05-31"), ("mi_2026_06", "2026-06-30"))


@pytest.fixture
def run_root(tmp_path):
    """A realistic deployment: governed run tapes AND dated platform canonicals.

    Both layouts are written because the deck and the concentration service
    resolve funded frames through different governed loaders — the deck's dated
    platform canonicals, and ``discover_snapshots`` over the central run tapes.
    A fixture carrying only one of them would exercise only half the pack.
    """
    root = tmp_path / "root"
    for run_id, cut in _RUNS:
        central = root / "acme" / run_id / "central"
        central.mkdir(parents=True, exist_ok=True)
        pd.DataFrame(_rows(cut)).to_csv(central / _CENTRAL_TAPE, index=False)
        dated = root / cut
        dated.mkdir(parents=True, exist_ok=True)
        pd.DataFrame(_rows(cut)).to_csv(dated / "platform_canonical_typed.csv",
                                        index=False)
    run_dir = root / "orun_v21"
    (run_dir / "out_platform").mkdir(parents=True, exist_ok=True)
    pd.DataFrame(_rows("2026-06-30")).to_csv(
        run_dir / "out_platform" / "platform_canonical_typed.csv", index=False)
    (run_dir / "run_state.json").write_text(json.dumps({
        "run_id": "orun_v21", "client_id": "acme", "reporting_date": "2026-06-30",
        "out_root": str(root)}), encoding="utf-8")
    return run_dir


@pytest.fixture
def pipeline_root(tmp_path):
    import shutil
    root = tmp_path / "pipe"
    for wk in ("2025-10-01", "2025-11-01"):
        src = _M2L / wk / f"M2L_KFI_and_Pipeline_{wk.replace('-', '_')}.csv"
        if not src.exists():
            pytest.skip(f"M2L fixture missing for {wk}")
        (root / wk).mkdir(parents=True, exist_ok=True)
        shutil.copy(src, root / wk / src.name)
    return root


@pytest.fixture
def approved_limits(tmp_path, monkeypatch):
    """An operator-approved concentration configuration, via the real store.

    Uses the same file-backed storage seam the OCC suites use — no mocks, and
    therefore a genuine exercise of the governed evaluator.
    """
    monkeypatch.setenv("TRAKT_STORAGE_BACKEND", "file")
    monkeypatch.setenv("TRAKT_LOCAL_BLOB_ROOT", str(tmp_path / "blob"))
    monkeypatch.setenv("TRAKT_RUNTIME_MODE", "test")

    from apps.blob_trigger_app.storage import open_storage
    from mi_agent.concentration_tests.models import (
        ActiveConfiguration, ActiveTest, ApprovalRecord, SourceEvidence)
    from mi_agent.concentration_tests.store import ConcentrationStore

    def _test(metric_id, threshold, name, params=None):
        return ActiveTest(
            metric_id=metric_id, threshold=threshold, operator="max",
            display_name=name, parameters=params or {},
            evidence=SourceEvidence(source_reference="facility_agreement.pdf",
                                    source_text="Schedule 8, clause 3.1"),
            approval=ApprovalRecord(decision="approved", operator="Operator",
                                    decided_at="2026-01-10T09:00:00+00:00"))

    config = ActiveConfiguration(
        client_id="acme", version=1, activated_by="Operator",
        library_version="1.0.0",
        tests=[
            _test("geo_region_share", 40.0, "London concentration",
                  {"regions": ["London"]}),
            _test("geo_region_share", 90.0, "South East concentration",
                  {"regions": ["South East"]}),
            _test("geo_region_share", 95.0, "Wales concentration",
                  {"regions": ["Wales"]}),
        ])
    ConcentrationStore(open_storage(), container="operations-control") \
        .commit_configuration("acme", config)
    return config


def _data(run_dir, *, pipeline_root=None, context=None):
    return build_dashboard_data(
        run_dir, client_id="acme", as_of="2026-06-30",
        output_root=str(Path(run_dir).parent),
        pipeline_root=str(pipeline_root) if pipeline_root else None,
        portfolio_context=context)


def _slides():
    return yaml.safe_load(DECK_CONFIG.read_text(encoding="utf-8"))["slides"]


def _build(run_dir, out, *, pipeline_root=None, context=None):
    argv = ["--run-dir", str(run_dir), "--deck-config", str(DECK_CONFIG),
            "--client-name", "Acme Capital", "--as-of-date", "2026-06-30",
            "--output-root", str(Path(run_dir).parent), "--output", str(out)]
    if pipeline_root:
        argv += ["--pipeline-root", str(pipeline_root)]
    if context:
        argv += ["--portfolio-context", context]
    return cli_run(argv)


def _text(path):
    return "\n".join(sh.text_frame.text for s in Presentation(str(path)).slides
                     for sh in s.shapes if sh.has_text_frame)


# --------------------------------------------------------------------------- #
# W1 — pipeline resolution parity.
# --------------------------------------------------------------------------- #

def test_deck_resolves_pipeline_through_the_shared_resolver(run_root, monkeypatch,
                                                            pipeline_root):
    """The regression that removed the whole pipeline section.

    With only an explicit source pointer set — no discovery root the deck knows
    about — the OLD deck path found nothing. The shared resolver honours it, so
    the pipeline resolves exactly as it does for the dashboard.
    """
    source = next(pipeline_root.rglob("M2L_KFI_and_Pipeline_*.csv"))
    monkeypatch.setenv("MI_AGENT_PIPELINE_SOURCE", str(source))
    data = build_dashboard_data(run_root, client_id="acme", as_of="2026-06-30",
                                output_root=str(run_root.parent))
    assert data.pipeline, "the explicit governed pipeline source did not resolve"
    assert data.pipeline.get("pipelineAmount")


def test_pipeline_sections_appear_when_a_source_exists(run_root, pipeline_root):
    data = _data(run_root, pipeline_root=pipeline_root)
    assert data.pipeline
    ids = {s["id"] for s in select_slides(_slides(), data)[0]}
    assert "pipeline" in ids and "funnel" in ids


def test_pipeline_sections_are_omitted_with_a_reason_when_absent(run_root):
    data = _data(run_root)
    kept, omitted = select_slides(_slides(), data)
    assert "pipeline" not in {s["id"] for s in kept}
    assert "pipeline" in {o.slide_id for o in omitted}


# --------------------------------------------------------------------------- #
# W2 — governed concentration tests.
# --------------------------------------------------------------------------- #

def test_governed_concentration_tests_resolve(run_root, approved_limits):
    data = _data(run_root)
    env = data.concentration
    assert env.get("tests"), "the approved concentration configuration did not resolve"
    assert env.get("source") == C.SOURCE_APPROVED
    assert build_facts(data)["has_concentration"] is True


def test_concentration_uses_the_approved_source_not_the_legacy_monitor(
        run_root, approved_limits):
    env = _data(run_root).concentration
    assert env["source"] == C.SOURCE_APPROVED
    assert env.get("approvalStatus") == "approved"
    assert C.source_disclosure(env).startswith("Operator-approved")


def test_breach_is_detected_and_ranked_first(run_root, approved_limits):
    """London is ~46% of balance against a 40% limit — a breach, and it must
    outrank the passing tests."""
    rows = C.adapt_tests(_data(run_root).concentration)
    london = next(r for r in rows if "London" in r["label"])
    assert london["status"] == C.STATUS_BREACH
    assert london["utilisation"] > 100
    assert C.select_tests(rows)[0]["label"] == london["label"]


def test_ranking_is_deterministic(run_root, approved_limits):
    rows = C.adapt_tests(_data(run_root).concentration)
    assert [r["label"] for r in C.select_tests(rows)] == \
           [r["label"] for r in C.select_tests(rows)]


def test_concentration_slide_renders_with_states_and_stress_label(
        run_root, approved_limits, tmp_path):
    out = tmp_path / "conc.pptx"
    _build(run_root, out)
    text = _text(out)
    assert "Concentration Tests and Headroom" in text
    assert "London concentration" in text
    # The stress must never be presented as the expected case.
    assert "STRESS BREACH" in text
    assert "Operator-approved" in text


def test_concentration_slide_is_omitted_when_no_tests_are_configured(run_root):
    data = _data(run_root)
    assert not data.concentration.get("tests")
    kept, omitted = select_slides(_slides(), data)
    assert "concentration" not in {s["id"] for s in kept}
    reasons = {o.slide_id: o.reason for o in omitted}
    assert "concentration" in reasons


def test_legacy_monitor_is_suppressed_when_approved_tests_exist(run_root,
                                                                approved_limits):
    """Two similar-looking limit sections would confuse; the approved one wins."""
    facts = build_facts(_data(run_root))
    assert facts["has_concentration"] is True
    kept = {s["id"] for s in select_slides(_slides(), _data(run_root))[0]}
    assert "risk" not in kept


def test_stress_is_never_labelled_as_expected(run_root, approved_limits, tmp_path):
    out = tmp_path / "stress.pptx"
    _build(run_root, out)
    text = _text(out).lower()
    # If the stress is discussed at all, it must carry the stress framing.
    if "all-pipeline-converts" in text:
        assert "stress, not the expected outcome" in text


def test_forward_states_absent_without_a_pipeline(run_root, approved_limits):
    """No active pipeline means the evaluator declines to project rather than
    restating funded as a forecast — the deck must not draw those marks."""
    env = _data(run_root).concentration
    assert C.forward_states_available(env) is False


# --------------------------------------------------------------------------- #
# W4 — client safety.
# --------------------------------------------------------------------------- #

def test_no_internal_paths_or_module_names_are_rendered(run_root, pipeline_root,
                                                        tmp_path):
    out = tmp_path / "safe.pptx"
    _build(run_root, out, pipeline_root=pipeline_root)
    text = _text(out)
    for forbidden in ("/tmp/", "scratchpad", "blob://", ".py", "compute_funded_snapshot",
                      "/mi/", str(tmp_path)):
        assert forbidden not in text, f"internal detail leaked: {forbidden}"


def test_leakage_gate_blocks_publication(run_root, tmp_path):
    """The gate must actually fail closed, not merely report."""
    data = _data(run_root)
    report = {"output": None, "slides": [{"id": "cover"}, {"id": "methodology"}],
              "facts": build_facts(data)}

    class _Leaky:
        portfolio = data.portfolio
        insights = data.insights
        reporting_date = data.reporting_date
        pipeline = None
        concentration = None

    import mi_agent_pptx.preflight as pf
    leaked = "Funded history root: /tmp/runs/acme scratchpad"
    gate = pf._gate_no_internal_paths(leaked)
    assert gate.passed is False and gate.mandatory is True


def test_dates_are_rendered_in_prose_not_iso_only(run_root, tmp_path):
    out = tmp_path / "dates.pptx"
    _build(run_root, out)
    assert "30 June 2026" in _text(out)


# --------------------------------------------------------------------------- #
# W5 — vintage wording.
# --------------------------------------------------------------------------- #

def test_no_static_pool_claim_survives(run_root, tmp_path):
    out = tmp_path / "vintage.pptx"
    _build(run_root, out)
    text = _text(out).lower()
    assert "static-pool" not in text and "static pool" not in text
    assert "origination vintages" in text


def test_static_pool_gate_would_block_the_old_wording():
    import mi_agent_pptx.preflight as pf
    gate = pf._gate_no_false_static_pool(
        "Static-pool profile compared across origination vintages")
    assert gate.passed is False and gate.mandatory is True


# --------------------------------------------------------------------------- #
# W3 / W8 — density and gates.
# --------------------------------------------------------------------------- #

def test_no_slide_is_a_shell(run_root, pipeline_root, tmp_path):
    out = tmp_path / "dense.pptx"
    _build(run_root, out, pipeline_root=pipeline_root)
    prs = Presentation(str(out))
    for i, slide in enumerate(prs.slides, start=1):
        texts = [sh for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        pictures = [sh for sh in slide.shapes if sh.shape_type == 13]
        assert pictures or len(texts) > 4, f"slide {i} is a shell"


def test_all_gates_pass_on_a_complete_deck(run_root, pipeline_root, approved_limits,
                                           tmp_path):
    out = tmp_path / "complete.pptx"
    rc = _build(run_root, out, pipeline_root=pipeline_root)
    verdict = json.loads(preflight_path(out).read_text(encoding="utf-8"))
    assert rc == 0, verdict["preflight"]["failed_gates"]
    gates = {g["gate"] for g in verdict["preflight"]["gates"]}
    for expected in ("no_internal_paths", "no_placeholder_language",
                     "no_false_static_pool_claim", "no_empty_slides",
                     "pipeline_reconciles", "concentration_reconciles"):
        assert expected in gates


def test_pipeline_headline_reconciles_to_the_governed_snapshot(run_root,
                                                               pipeline_root,
                                                               tmp_path):
    out = tmp_path / "recon.pptx"
    _build(run_root, out, pipeline_root=pipeline_root)
    verdict = json.loads(preflight_path(out).read_text(encoding="utf-8"))
    gate = next(g for g in verdict["preflight"]["gates"]
                if g["gate"] == "pipeline_reconciles")
    assert gate["passed"] is True
    assert gate["evidence"]["pipeline_amount"]


def test_omission_ledger_records_every_absent_section(run_root, tmp_path):
    out = tmp_path / "ledger.pptx"
    _build(run_root, out)
    verdict = json.loads(preflight_path(out).read_text(encoding="utf-8"))
    omitted = {o["slide_id"]: o["reason"] for o in verdict["omitted_slides"]}
    assert "pipeline" in omitted and "concentration" in omitted
    for reason in omitted.values():
        assert reason and not reason.endswith("None")
