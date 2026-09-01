"""FINAL SIGN-OFF — the deck's Stage Movement slide serves the dashboard's answer.

`test_stage_transition_parity.py` already proves the two channels carry equal
VALUES. Certification needs more than equality: two independent computations can
be equal today and diverge tomorrow. These tests pin the properties that make
divergence structurally impossible, and close the chain all the way to the
pixels on the slide:

  * the deck does not RECONSTRUCT the result — patch the governed resolver and
    the deck returns the patched object, so there is exactly one producer;
  * the deck's data function only DELEGATES — no aggregation of its own;
  * the deck's payload equals the ACTUAL HTTP RESPONSE BODY the dashboard
    receives, including under a SCOPED request, which is the one call-argument
    asymmetry between the two call sites;
  * the numbers RENDERED on the slide are the numbers in that response body;
  * the deck reaches the capability WITHOUT a web framework — the property that
    lets it run in the Azure Functions PPTX stage, where no HTTP server exists.

The fixture is deliberately local rather than shared with the parity module: a
certification suite that could be broken by an edit to another test file is not
a certification.
"""

from __future__ import annotations

import json
import re
from pathlib import Path

import pandas as pd
import pytest

_REPO = Path(__file__).resolve().parents[2]
PIPELINE = _REPO / "tests" / "fixtures" / "pipeline_transition_2w"
DECK_CONFIG = _REPO / "configs" / "pptx" / "investor_pack.yaml"
CLIENT, AS_OF, RUN_ID = "acme", "2026-06-30", "mi_2026_06"
ROUTE = "/mi/insight/movement-detail"
DETAIL_TYPE = "PIPELINE_STAGE_TRANSITION"

pytestmark = pytest.mark.skipif(not PIPELINE.exists(),
                                reason="transition fixture pack not present")

#: Provenance of the RUN, not of the analysis — it differs between a deck build
#: and an HTTP request by construction and carries no stage-transition value.
_PROVENANCE = ("run_id", "scope", "portfolio_id", "portfolioScope")


def _norm(payload: dict) -> dict:
    out = dict(payload)
    for key in _PROVENANCE:
        out.pop(key, None)
    return out


def _loan(i: int, cut: str) -> dict:
    return {
        "unique_identifier": f"direct_001_L{i:04d}",
        "source_portfolio_id": "direct_001",
        "source_portfolio_type": "direct",
        "source_portfolio_label": "Direct 001",
        "current_outstanding_balance": 300_000.0,
        "current_principal_balance": 300_000.0,
        "original_principal_balance": 312_000.0,
        "current_valuation_amount": 666_666.0,
        "current_loan_to_value": 45.0,
        "current_interest_rate": 7.1,
        "youngest_borrower_age": 72,
        "geographic_region_collateral": "London",
        "collateral_geography": "London",
        "geographic_region_obligor": "TLI3",
        "origination_channel": "Direct",
        "origination_date": "2021-04-01",
        "data_cut_off_date": cut,
    }


@pytest.fixture()
def book(tmp_path, monkeypatch):
    """One governed environment both channels resolve from."""
    root = tmp_path / "root"
    for run_id, date in (("mi_2026_05", "2026-05-31"), (RUN_ID, AS_OF)):
        central = root / CLIENT / run_id / "central"
        central.mkdir(parents=True, exist_ok=True)
        frame = pd.DataFrame([_loan(i, date) for i in range(1, 4)])
        frame.to_csv(central / "18_central_lender_tape.csv", index=False)
        dated = root / date
        dated.mkdir(parents=True, exist_ok=True)
        frame.to_csv(dated / "platform_canonical_typed.csv", index=False)

    run_dir = root / "orun_signoff"
    (run_dir / "out_platform").mkdir(parents=True, exist_ok=True)
    pd.DataFrame([_loan(i, AS_OF) for i in range(1, 4)]).to_csv(
        run_dir / "out_platform" / "platform_canonical_typed.csv", index=False)
    (run_dir / "run_state.json").write_text(json.dumps({
        "run_id": RUN_ID, "client_id": CLIENT, "reporting_date": AS_OF,
        "out_root": str(root)}), encoding="utf-8")

    monkeypatch.setenv("MI_AGENT_ONBOARDING_OUTPUT_ROOT", str(root))
    monkeypatch.setenv("MI_AGENT_CLIENT_ID", CLIENT)
    monkeypatch.setenv("MI_AGENT_PIPELINE_ROOT", str(PIPELINE))
    monkeypatch.setenv("MI_AGENT_AUTH_ENABLED", "false")
    monkeypatch.setenv("TRAKT_MI_ENHANCED_HOVERS", "true")
    return run_dir


def _deck_data(book):
    from mi_agent_pptx.mi_api import build_dashboard_data
    return build_dashboard_data(book, client_id=CLIENT, as_of=AS_OF,
                                output_root=str(Path(book).parent))


def _http(portfolio_context: str | None = None) -> dict:
    """The response body the dashboard actually receives, over the real route."""
    from fastapi.testclient import TestClient
    from mi_agent_api.app import app

    params = {"portfolioId": f"{CLIENT}/{RUN_ID}", "detailType": DETAIL_TYPE}
    if portfolio_context:
        params["portfolioContext"] = portfolio_context
    res = TestClient(app).get(ROUTE, params=params)
    assert res.status_code == 200, f"{res.status_code}: {res.text}"
    return res.json()


# --------------------------------------------------------------------------- #
# 1. One producer — the deck cannot be computing its own answer.
# --------------------------------------------------------------------------- #
class TestSingleProducer:

    def test_the_deck_returns_whatever_the_governed_resolver_returns(
            self, book, monkeypatch):
        """The strongest available proof that the deck does not reconstruct.

        Equal values could be two computations that happen to agree. If the deck
        hands back a sentinel the resolver was patched to produce, it is not
        computing anything — it is relaying the one governed result.
        """
        from mi_agent_api import movement_detail as md

        sentinel = {"detail_type": DETAIL_TYPE, "available": True,
                    "transitions": [{"marker": "from-the-governed-resolver"}]}
        seen: list[tuple] = []

        def fake(root, client_id, **kw):
            seen.append((str(root), client_id, tuple(sorted(kw))))
            return dict(sentinel)

        monkeypatch.setattr(md, "resolve_stage_transition_detail", fake)
        data = _deck_data(book)

        assert data.stage_transitions == sentinel, \
            "the deck produced a payload the governed resolver did not"
        assert seen, "the deck never called the governed resolver at all"

    def test_the_route_uses_that_same_resolver(self, book, monkeypatch):
        """The other half: patching the SAME symbol changes the HTTP answer, so
        route and deck demonstrably share one producer."""
        from mi_agent_api import movement_detail as md

        monkeypatch.setattr(md, "resolve_stage_transition_detail",
                            lambda *a, **k: {"detail_type": DETAIL_TYPE,
                                             "available": True,
                                             "transitions": [{"marker": "patched"}]})
        assert _http()["transitions"] == [{"marker": "patched"}]

    def test_the_deck_data_function_only_delegates(self):
        """Structural: the function is an import and one delegating return.

        Asserted on the parsed syntax tree rather than on substrings, so it
        cannot pass vacuously and cannot fail on a harmless edit like a type
        annotation. A deck that summed, grouped, filtered or reshaped here would
        be a second analytical owner however well its numbers agreed today —
        this makes that a test failure rather than a code review.
        """
        import ast
        import inspect
        import textwrap
        from mi_agent_pptx import mi_api

        tree = ast.parse(textwrap.dedent(inspect.getsource(mi_api._stage_transitions)))
        body = list(tree.body[0].body)
        if (isinstance(body[0], ast.Expr)
                and isinstance(body[0].value, ast.Constant)
                and isinstance(body[0].value.value, str)):
            body.pop(0)                      # the docstring

        assert all(isinstance(n, (ast.Import, ast.ImportFrom)) for n in body[:-1]), \
            "_stage_transitions does work before delegating"
        ret = body[-1]
        assert isinstance(ret, ast.Return) and isinstance(ret.value, ast.Call), \
            "_stage_transitions does not end in a single delegating call"
        assert ret.value.func.attr == "resolve_stage_transition_detail", \
            "_stage_transitions delegates to something other than the governed resolver"
        # Every argument is a name passed straight through — nothing derived.
        for node in [*ret.value.args, *(k.value for k in ret.value.keywords)]:
            assert isinstance(node, ast.Name), \
                f"_stage_transitions computes an argument ({ast.dump(node)[:80]})"


# --------------------------------------------------------------------------- #
# 2. The deck's payload IS the dashboard's response body.
# --------------------------------------------------------------------------- #
class TestSameResponseBody:

    def test_the_deck_payload_equals_the_http_response_body(self, book):
        assert _norm(_deck_data(book).stage_transitions) == _norm(_http())

    def test_they_still_agree_under_a_scoped_dashboard_request(self, book):
        """The one call-argument asymmetry between the two sites: the route
        passes ``scope``, the deck does not. It must not move a number.

        (``scope`` labels the payload; it does not filter the population — the
        same treatment every other deck pipeline payload gets.)
        """
        scoped = _http(portfolio_context="direct")
        assert _norm(_deck_data(book).stage_transitions) == _norm(scoped)

    def test_every_governed_block_is_present_on_both_sides(self, book):
        deck, api = _deck_data(book).stage_transitions, _http()
        for key in ("detail_type", "available", "identifier", "measure",
                    "counts", "transitions", "new_arrivals", "stayers",
                    "departures", "event_totals", "reconciliation",
                    "methodology", "source_dates"):
            assert key in deck and key in api, f"{key} missing from a channel"
            assert deck[key] == api[key], f"{key} differs between the channels"


# --------------------------------------------------------------------------- #
# 3. The chain closes on the rendered slide, not just on the payload.
# --------------------------------------------------------------------------- #
class TestRenderedSlideMatchesTheResponse:

    def _slide_text(self, book, tmp_path) -> str:
        from pptx import Presentation
        from mi_agent_pptx.cli import run as cli_run

        out = tmp_path / "signoff.pptx"
        assert cli_run(["--run-dir", str(book), "--deck-config", str(DECK_CONFIG),
                        "--client-name", "Acme Capital", "--as-of-date", AS_OF,
                        "--output-root", str(Path(book).parent),
                        "--output", str(out)]) == 0
        for slide in Presentation(str(out)).slides:
            text = " ".join(sh.text_frame.text for sh in slide.shapes
                            if sh.has_text_frame)
            if "Pipeline Stage Movement" in text:
                return text
        pytest.fail("the deck contains no Pipeline Stage Movement slide")

    def test_the_slide_states_the_response_bodys_window_and_identifier(
            self, book, tmp_path):
        api, text = _http(), self._slide_text(book, tmp_path)
        assert api["identifier"] in text
        assert "5 June 2026" in text and "12 June 2026" in text

    def test_the_slide_states_the_response_bodys_population_and_residuals(
            self, book, tmp_path):
        api, text = _http(), self._slide_text(book, tmp_path)
        counts, recon = api["counts"], api["reconciliation"]
        assert f"({counts['comparison']} prior, {counts['current']} latest)" in text
        assert (f"residual {recon['count_reconciliation_residual']} cases / "
                f"{recon['amount_reconciliation_residual']} by value") in text

    def test_the_slide_never_presents_a_synthetic_stage_as_a_real_one(
            self, book, tmp_path):
        """A new arrival has no source stage in the response body, and must not
        acquire one on the slide."""
        text = self._slide_text(book, tmp_path)
        for invented in ("NEW →", "New →", "KFI → KFI",
                         "Application → Application"):
            assert invented not in text


# --------------------------------------------------------------------------- #
# 4. Why the transport is in-process, pinned so it cannot regress silently.
# --------------------------------------------------------------------------- #
class TestDeckRunsWhereNoHttpServerExists:
    """The deck ships in the Azure Functions PPTX stage, which runs no FastAPI
    app (root ``requirements.txt``: *"unused by the Function App"*). Reaching
    the capability over HTTP would make deck generation depend on a second
    service being up, reachable and authenticated — a runtime dependency the
    current design deliberately avoids, and one no other deck payload has.

    These pin the property rather than merely asserting it in a report.
    """

    def test_reaching_the_capability_imports_no_web_framework(self):
        import subprocess
        import sys

        probe = (
            "import sys;"
            "import mi_agent_pptx.mi_api;"
            "from mi_agent_api import movement_detail;"
            "movement_detail.resolve_stage_transition_detail;"
            "print(sorted(k for k in sys.modules "
            "if k.split('.')[0] in ('fastapi','starlette','uvicorn')))"
        )
        out = subprocess.run([sys.executable, "-c", probe], cwd=str(_REPO),
                             capture_output=True, text=True, timeout=180)
        assert out.returncode == 0, out.stderr[-2000:]
        assert out.stdout.strip().endswith("[]"), \
            f"the deck now drags in a web framework: {out.stdout.strip()}"

    def test_the_deck_stack_makes_no_http_calls(self):
        """No client, no URL, no transport — in any deck module."""
        offenders = []
        for path in sorted((_REPO / "mi_agent_pptx").glob("*.py")):
            for n, line in enumerate(path.read_text(encoding="utf-8").splitlines(), 1):
                code = line.split("#", 1)[0]
                if re.search(r"\b(requests|httpx|aiohttp|urllib\.request|TestClient)\b",
                             code) or re.search(r"https?://", code):
                    offenders.append(f"{path.name}:{n}: {line.strip()}")
        assert not offenders, "the deck acquired an HTTP data path:\n" + "\n".join(offenders)
