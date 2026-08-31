"""Investor Reporting v2 — portfolio context, composition, insights, gates.

Covers the release's acceptance criteria against real multi-book tapes:

  * a deck always states the portfolio it describes;
  * total, direct-only, acquired-only and multi-acquired scopes each compose
    correctly;
  * direct + acquired reconcile to the total;
  * sections with no data are OMITTED with a reason, never placeholdered;
  * the executive summary is deterministic and reproducible;
  * a mandatory gate failure withholds publication without destroying the deck.
"""

from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
import pytest
import yaml
from pptx import Presentation

from mi_agent_pptx.cli import PREFLIGHT_BLOCKED_RC, preflight_path, run as cli_run
from mi_agent_pptx.composition import build_facts, select_slides
from mi_agent_pptx.mi_api import build_dashboard_data
from mi_agent_pptx.preflight import run_preflight

DECK_CONFIG = Path("configs/pptx/investor_pack.yaml")


# --------------------------------------------------------------------------- #
# Fixtures — governed multi-book tapes with provenance.
# --------------------------------------------------------------------------- #

def _loan(i: int, pid: str, ptype: str, balance: float, *, cut: str,
          ltv: float = 45.0, rate: float = 7.2, age: int = 74,
          orig: str = "2022-05-01") -> dict:
    return {
        "unique_identifier": f"{pid}_L{i:04d}",
        "source_portfolio_id": pid,
        "source_portfolio_type": ptype,
        "source_portfolio_label": pid.replace("_", " ").title(),
        # The canonical balance the MI layer aggregates on.
        "current_outstanding_balance": balance,
        "current_principal_balance": balance,
        "current_valuation_amount": balance / (ltv / 100.0),
        "current_loan_to_value": ltv,
        "original_loan_to_value": ltv / 100 - 0.02,
        "current_interest_rate": rate,
        "youngest_borrower_age": age,
        "geographic_region_obligor": "TLI3" if i % 2 else "TLJ2",
        "origination_date": orig,
        "data_cut_off_date": cut,
    }


def _write_cut(root: Path, date: str, books: dict) -> None:
    """One dated platform canonical containing every book."""
    rows = []
    for pid, (ptype, balances) in books.items():
        for i, bal in enumerate(balances, start=1):
            rows.append(_loan(i, pid, ptype, bal, cut=date))
    d = root / date
    d.mkdir(parents=True, exist_ok=True)
    pd.DataFrame(rows).to_csv(d / "platform_canonical_typed.csv", index=False)


def _make_run(tmp_path: Path, current: dict, prior: dict | None = None,
              *, cut: str = "2026-06-30", prior_cut: str = "2026-05-31") -> Path:
    """A run directory plus dated history, mirroring a real deployment."""
    root = tmp_path / "root"
    if prior is not None:
        _write_cut(root, prior_cut, prior)
    _write_cut(root, cut, current)

    run_dir = root / "orun_v2"
    (run_dir / "out_platform").mkdir(parents=True, exist_ok=True)
    rows = []
    for pid, (ptype, balances) in current.items():
        for i, bal in enumerate(balances, start=1):
            rows.append(_loan(i, pid, ptype, bal, cut=cut))
    pd.DataFrame(rows).to_csv(
        run_dir / "out_platform" / "platform_canonical_typed.csv", index=False)
    (run_dir / "run_state.json").write_text(json.dumps({
        "run_id": "orun_v2", "client_id": "acme", "target": "mi",
        "reporting_date": cut, "out_root": str(root)}), encoding="utf-8")
    return run_dir


#: A growing direct book and a redeeming acquired book — the case the release
#: exists for: the total alone describes neither.
_CURRENT = {
    "direct_001": ("direct", [300_000.0, 280_000.0, 310_000.0]),
    "acquired_001": ("acquired", [250_000.0, 240_000.0]),
}
_PRIOR = {
    "direct_001": ("direct", [290_000.0, 275_000.0]),
    "acquired_001": ("acquired", [260_000.0, 248_000.0]),
}


@pytest.fixture
def mixed_run(tmp_path):
    return _make_run(tmp_path, _CURRENT, _PRIOR)


@pytest.fixture
def direct_only_run(tmp_path):
    return _make_run(tmp_path, {"direct_001": ("direct", [300_000.0, 280_000.0])})


@pytest.fixture
def acquired_only_run(tmp_path):
    return _make_run(tmp_path, {"acquired_001": ("acquired", [250_000.0, 240_000.0])})


@pytest.fixture
def multi_acquired_run(tmp_path):
    return _make_run(tmp_path, {
        "direct_001": ("direct", [300_000.0]),
        "acquired_001": ("acquired", [250_000.0, 240_000.0]),
        "acquired_002": ("acquired", [180_000.0, 175_000.0, 160_000.0]),
    })


def _data(run_dir, context=None, as_of="2026-06-30"):
    return build_dashboard_data(run_dir, client_id="acme", as_of=as_of,
                               output_root=str(Path(run_dir).parent),
                               portfolio_context=context)


def _slides():
    return yaml.safe_load(DECK_CONFIG.read_text(encoding="utf-8"))["slides"]


# --------------------------------------------------------------------------- #
# Phase 1 — portfolio context.
# --------------------------------------------------------------------------- #

def test_context_resolves_books_types_and_dates(mixed_run):
    ctx = _data(mixed_run).portfolio
    assert ctx is not None
    assert ctx.report_scope == "total"
    assert ctx.portfolio_count == 2
    assert set(ctx.portfolio_types) == {"direct", "acquired"}
    assert ctx.is_mixed
    assert set(ctx.portfolio_ids) == {"direct_001", "acquired_001"}
    assert ctx.reporting_dates == {"direct_001": "2026-06-30",
                                   "acquired_001": "2026-06-30"}
    assert ctx.total_balance and ctx.loan_count == 5


@pytest.mark.parametrize("context,expected_types,expected_loans", [
    ("total", {"direct", "acquired"}, 5),
    ("direct", {"direct"}, 3),
    ("acquired", {"acquired"}, 2),
])
def test_scope_narrows_correctly(mixed_run, context, expected_types, expected_loans):
    ctx = _data(mixed_run, context).portfolio
    assert ctx.report_scope == context
    assert set(ctx.portfolio_types) == expected_types
    assert int(ctx.loan_count) == expected_loans


def test_direct_and_acquired_reconcile_to_total(mixed_run):
    """The additivity guarantee: the parts are the same function as the whole."""
    ctx = _data(mixed_run).portfolio
    # Guard against a vacuous pass: 0 == 0 would satisfy the sum but prove nothing.
    assert ctx.total_balance > 0
    assert ctx.slice_for("direct").balance > 0
    assert ctx.slice_for("acquired").balance > 0
    assert ctx.type_balance_total() == pytest.approx(ctx.total_balance, rel=1e-9)
    assert int(ctx.type_loan_total()) == int(ctx.loan_count)
    # And the split is genuinely a split, not one book plus an empty one.
    assert ctx.slice_for("direct").balance != ctx.total_balance


def test_books_at_different_dates_are_detected_and_disclosed(tmp_path):
    """A total that combines books at different dates must say so — an investor
    must never read a mixed-date total as a clean cut."""
    root = tmp_path / "root"
    run_dir = root / "orun_v2"
    (run_dir / "out_platform").mkdir(parents=True)
    rows = [_loan(1, "direct_001", "direct", 300_000.0, cut="2026-06-30"),
            _loan(2, "direct_001", "direct", 280_000.0, cut="2026-06-30"),
            _loan(1, "acquired_001", "acquired", 250_000.0, cut="2026-05-31")]
    pd.DataFrame(rows).to_csv(
        run_dir / "out_platform" / "platform_canonical_typed.csv", index=False)
    (run_dir / "run_state.json").write_text(json.dumps({
        "run_id": "orun_v2", "client_id": "acme", "reporting_date": "2026-06-30",
        "out_root": str(root)}), encoding="utf-8")

    data = _data(run_dir)
    ctx = data.portfolio
    assert ctx.has_mixed_reporting_dates
    assert ctx.type_reporting_dates == {"direct": "2026-06-30",
                                        "acquired": "2026-05-31"}
    # The caveat is raised as an ATTENTION observation, not buried.
    scope_notes = [i for i in data.insights["insights"]
                   if i.insight_type == "REPORTING_SCOPE"]
    assert scope_notes and scope_notes[0].severity == "attention"

    out = tmp_path / "mixed_dates.pptx"
    assert _build(run_dir, out) == 0
    text = "\n".join(sh.text_frame.text for s in Presentation(str(out)).slides
                     for sh in s.shapes if sh.has_text_frame)
    assert "2026-06-30" in text and "2026-05-31" in text
    assert "different dates" in text


def test_an_unscoped_run_is_not_reported_as_a_fallback(mixed_run):
    """A plain client id means Total; it is not a failed scope request."""
    ctx = _data(mixed_run, None).portfolio
    assert ctx.is_total
    assert ctx.fell_back_to_total is False


# --------------------------------------------------------------------------- #
# Phase 2 / 5 — composition.
# --------------------------------------------------------------------------- #

def test_single_type_shows_no_per_type_split(acquired_only_run):
    """With one type there is nothing to compare, and the composition slide
    says so rather than rendering a one-column table."""
    ctx = _data(acquired_only_run).portfolio
    assert ctx.portfolio_types == ("acquired",)
    assert not ctx.is_mixed
    # Per-type snapshots are only computed when there is more than one type —
    # a single-book scope must not pay for a split it cannot show.
    assert ctx.type_slices == ()


def test_mixed_portfolio_carries_the_per_type_comparison(mixed_run, tmp_path):
    """The per-type comparison lives on Portfolio Composition.

    A separate "Direct vs Acquired" slide was consolidated away in the close-out
    pass: its table duplicated this one and its waterfall was identical to the
    one on Portfolio Movement and Drivers.
    """
    data = _data(mixed_run)
    assert data.portfolio.is_mixed
    kept = {s["id"] for s in select_slides(_slides(), data)[0]}
    assert "portfolio_composition" in kept
    assert "portfolio_comparison" not in kept, "the duplicated slide is gone"

    out = tmp_path / "composition.pptx"
    assert _build(mixed_run, out) == 0
    text = "\n".join(sh.text_frame.text for s in Presentation(str(out)).slides
                     for sh in s.shapes if sh.has_text_frame)
    # Every measure the removed slide carried is still reported.
    for measure in ("Funded balance", "Loans", "Average balance",
                    "WA current LTV", "WA interest rate", "WA borrower age",
                    "Period movement", "Loan movement"):
        assert measure in text, f"{measure} lost in consolidation"


def test_no_pipeline_omits_every_pipeline_section(mixed_run):
    data = _data(mixed_run)
    assert not data.pipeline
    kept, omitted = select_slides(_slides(), data)
    ids = {s["id"] for s in kept}
    for sid in ("pipeline", "pipeline_evolution", "funnel", "origination_flow"):
        assert sid not in ids
    assert {"pipeline", "funnel"} <= {o.slide_id for o in omitted}


def test_no_forecast_omits_the_forecast_sections(mixed_run):
    kept, _ = select_slides(_slides(), _data(mixed_run))
    ids = {s["id"] for s in kept}
    assert "forecast_bridge" not in ids
    assert "forecast_projection" not in ids


def test_every_omission_carries_a_reason(mixed_run):
    _kept, omitted = select_slides(_slides(), _data(mixed_run))
    assert omitted
    for o in omitted:
        assert o.reason and o.slide_id and o.title
        assert o.category in ("condition", "no data")


def test_forecast_fact_is_false_without_a_pipeline_contribution(mixed_run):
    """A bridge with no weighted pipeline restates the balance; it is not a
    forecast and must not be presented as one."""
    assert build_facts(_data(mixed_run))["has_forecast"] is False


# --------------------------------------------------------------------------- #
# Pipeline / forecast PRESENT — the mirror of the omission cases above.
# --------------------------------------------------------------------------- #

_M2L = Path("tests/fixtures/client_001_mi_pack/pipeline")


@pytest.fixture
def pipeline_root(tmp_path):
    """Governed weekly M2L extracts laid out as a discovery root."""
    import shutil
    weeks = ["2025-10-01", "2025-11-01"]
    root = tmp_path / "pipe"
    for wk in weeks:
        src = _M2L / wk / f"M2L_KFI_and_Pipeline_{wk.replace('-', '_')}.csv"
        if not src.exists():
            pytest.skip(f"M2L fixture missing for {wk}")
        dst = root / wk
        dst.mkdir(parents=True, exist_ok=True)
        shutil.copy(src, dst / src.name)
    return root


def test_pipeline_present_includes_the_pipeline_sections(mixed_run, pipeline_root):
    data = build_dashboard_data(mixed_run, client_id="acme", as_of="2026-06-30",
                                output_root=str(Path(mixed_run).parent),
                                pipeline_root=str(pipeline_root))
    assert data.pipeline, "the governed pipeline source did not resolve"
    facts = build_facts(data)
    assert facts["has_pipeline"] is True

    kept, _omitted = select_slides(_slides(), data)
    ids = {s["id"] for s in kept}
    assert "pipeline" in ids
    # HOW THE PIPELINE PROGRESSES, by whichever page answers it for this book.
    # Origination Funnel measures conversion from weekly FLOW; Pipeline Stage
    # Movement reconciles the same question case by case and supersedes it
    # where a stable case identifier exists. The durable property is that the
    # question is answered, not which page answers it.
    assert ids & {"funnel", "pipeline_movement"}, ids

    # And the pipeline is narrated as an outlook, labelled direct-origination only.
    outlook = [i for i in data.insights["insights"]
               if i.insight_type == "PIPELINE_OUTLOOK"]
    assert outlook
    assert "acquired" in outlook[0].methodology["basis"].lower()


def test_forecast_present_includes_the_bridge(mixed_run, pipeline_root):
    data = build_dashboard_data(mixed_run, client_id="acme", as_of="2026-06-30",
                                output_root=str(Path(mixed_run).parent),
                                pipeline_root=str(pipeline_root))
    bridge = (data.forecast or {}).get("forecastBridge") or {}
    if not bridge.get("weightedExpectedFundedAmount"):
        pytest.skip("this fixture week contributes no weighted pipeline")
    assert build_facts(data)["has_forecast"] is True
    assert "forecast_bridge" in {s["id"] for s in select_slides(_slides(), data)[0]}
    assert [i for i in data.insights["insights"]
            if i.insight_type == "FORECAST_OUTLOOK"]


# --------------------------------------------------------------------------- #
# Phase 3 — deterministic executive summary.
# --------------------------------------------------------------------------- #

def test_executive_summary_attributes_the_movement(mixed_run):
    brief = _data(mixed_run).insights
    types = {i.insight_type for i in brief["insights"]}
    assert "MOVEMENT_ATTRIBUTION" in types
    attribution = next(i for i in brief["insights"]
                       if i.insight_type == "MOVEMENT_ATTRIBUTION")
    # The attribution names BOTH books, not just the net total.
    assert set(attribution.contributors) == {"direct", "acquired"}


def test_executive_summary_is_reproducible(mixed_run):
    """The same approved run must produce the same words every time."""
    first = [i.headline for i in _data(mixed_run).insights["insights"]]
    second = [i.headline for i in _data(mixed_run).insights["insights"]]
    assert first == second and first


def test_insights_never_exceed_the_cap_and_carry_evidence(mixed_run):
    brief = _data(mixed_run).insights
    assert 0 < brief["count"] <= 8
    for ins in brief["insights"]:
        assert ins.headline and ins.summary
        assert ins.methodology.get("source"), "every claim names its source"
        assert ins.portfolio_context == "total"


def test_immaterial_and_unavailable_observations_are_stated_not_dropped(mixed_run):
    brief = _data(mixed_run).insights
    categories = {o.category for o in brief["omitted"]}
    assert categories & {"immaterial", "unavailable"}
    for o in brief["omitted"]:
        assert o.reason


def test_single_type_scope_omits_attribution(acquired_only_run):
    brief = _data(acquired_only_run).insights
    omitted = {o.insight_type for o in brief["omitted"]}
    assert "MOVEMENT_ATTRIBUTION" in omitted


# --------------------------------------------------------------------------- #
# Phase 4 — multiple acquired books.
# --------------------------------------------------------------------------- #

def test_multiple_acquired_books_consolidate_into_one_type(multi_acquired_run):
    ctx = _data(multi_acquired_run).portfolio
    assert ctx.portfolio_count == 3
    acquired = ctx.slice_for("acquired")
    assert set(acquired.portfolio_ids) == {"acquired_001", "acquired_002"}
    assert int(acquired.loan_count) == 5          # 2 + 3, consolidated
    assert ctx.type_balance_total() == pytest.approx(ctx.total_balance, rel=1e-9)


def test_acquired_only_scope_covers_every_acquired_book(multi_acquired_run):
    ctx = _data(multi_acquired_run, "acquired").portfolio
    assert set(ctx.portfolio_ids) == {"acquired_001", "acquired_002"}
    assert int(ctx.loan_count) == 5


# --------------------------------------------------------------------------- #
# Phase 6 — publication gates.
# --------------------------------------------------------------------------- #

def _build(run_dir, out, context=None):
    argv = ["--run-dir", str(run_dir), "--deck-config", str(DECK_CONFIG),
            "--client-name", "Acme Capital", "--as-of-date", "2026-06-30",
            "--output-root", str(Path(run_dir).parent), "--output", str(out)]
    if context:
        argv += ["--portfolio-context", context]
    return cli_run(argv)


@pytest.mark.parametrize("context", [None, "direct", "acquired"])
def test_decks_generate_and_pass_every_gate(mixed_run, tmp_path, context):
    out = tmp_path / f"deck_{context or 'total'}.pptx"
    assert _build(mixed_run, out, context) == 0
    assert out.exists()
    prs = Presentation(str(out))          # opens as a real presentation
    assert len(list(prs.slides)) >= 3
    verdict = json.loads(preflight_path(out).read_text(encoding="utf-8"))
    assert verdict["preflight"]["publishable"] is True


def test_every_deck_states_its_scope_and_dates(mixed_run, tmp_path):
    """The finding this release exists to fix: a scoped deck must SAY it is
    scoped, so it can never be read as a total-portfolio pack."""
    for context, expected in (("direct", "Direct"), ("acquired", "Acquired")):
        out = tmp_path / f"scope_{context}.pptx"
        assert _build(mixed_run, out, context) == 0
        prs = Presentation(str(out))
        text = "\n".join(sh.text_frame.text for s in prs.slides
                         for sh in s.shapes if sh.has_text_frame)
        assert expected in text, f"{context} deck never names its scope"
        assert "2026-06-30" in text, "reporting date is not disclosed"


def test_no_deck_contains_a_placeholder_page(mixed_run, tmp_path):
    out = tmp_path / "no_placeholders.pptx"
    _build(mixed_run, out)
    verdict = json.loads(preflight_path(out).read_text(encoding="utf-8"))
    assert not [s for s in verdict["slides"] if s["placeholder"]]
    text = "\n".join(sh.text_frame.text for s in Presentation(str(out)).slides
                     for sh in s.shapes if sh.has_text_frame)
    assert "No data for this run" not in text


def test_reconciliation_gate_is_evaluated_on_a_mixed_book(mixed_run, tmp_path):
    data = _data(mixed_run)
    out = tmp_path / "recon.pptx"
    _build(mixed_run, out)
    verdict = json.loads(preflight_path(out).read_text(encoding="utf-8"))
    gates = {g["gate"]: g for g in verdict["preflight"]["gates"]}
    assert gates["totals_reconcile"]["passed"] is True
    assert gates["totals_reconcile"]["mandatory"] is True
    assert gates["loan_counts_reconcile"]["passed"] is True


def test_a_reconciliation_break_blocks_publication(mixed_run, tmp_path):
    """If the parts stopped summing to the whole, the pack must not publish."""
    data = _data(mixed_run)
    build_report = {"output": None, "slides": [], "facts": build_facts(data)}

    class _Broken:
        portfolio = data.portfolio
        insights = data.insights
        reporting_date = data.reporting_date

    # Corrupt the total so the type slices no longer reconcile to it.
    object.__setattr__(_Broken.portfolio, "snapshot",
                       {"kpis": [{"id": "balance", "raw": 1.0, "available": True,
                                  "value": "£1"},
                                 {"id": "loans", "raw": 1, "available": True,
                                  "value": "1"}]})
    report = run_preflight(build_report, _Broken)
    assert report.publishable is False
    assert "totals_reconcile" in [g.gate for g in report.failures]


def test_omission_ledger_is_rendered_in_the_appendix(mixed_run, tmp_path):
    out = tmp_path / "ledger.pptx"
    _build(mixed_run, out)
    text = "\n".join(sh.text_frame.text for s in Presentation(str(out)).slides
                     for sh in s.shapes if sh.has_text_frame)
    assert "SECTIONS NOT INCLUDED" in text
    assert "no governed pipeline source" in text


# --------------------------------------------------------------------------- #
# Phase 7 — artefact metadata.
# --------------------------------------------------------------------------- #

def test_artefact_metadata_carries_the_governed_context(mixed_run, tmp_path):
    out = tmp_path / "meta.pptx"
    _build(mixed_run, out, "direct")
    sidecar = json.loads(preflight_path(out).read_text(encoding="utf-8"))
    ctx = sidecar["portfolio_context"]
    assert ctx["report_scope"] == "direct"
    assert ctx["portfolio_ids"] == ["direct_001"]
    assert ctx["reporting_dates"]
    assert sidecar["template_version"]
    assert sidecar["insights"]["insight_version"]
    assert sidecar["source_runs"]["run_id"]


def test_pointer_distinguishes_two_scopes_of_the_same_client(tmp_path):
    """Without scope on the artefact, a total pack and a direct pack for the same
    client and period are indistinguishable to every consumer."""
    from apps.blob_trigger_app.pptx_stage import build_deck_pointer

    deck = tmp_path / "d.pptx"
    deck.write_bytes(b"PK\x03\x04" + b"0" * 64)
    made = []
    for scope in ("total", "direct"):
        made.append(build_deck_pointer(
            deck, client_id="acme", period="2026-06",
            latest_uri="decks/acme/latest/investor_pack.pptx",
            context={"portfolio_context": {"report_scope": scope,
                                           "portfolio_ids": ["direct_001"]},
                     "template_version": "2.0.0"}))
    assert made[0]["report_scope"] == "total"
    assert made[1]["report_scope"] == "direct"
    assert made[0]["report_scope"] != made[1]["report_scope"]


def test_pointer_without_context_keeps_its_original_shape(tmp_path):
    """Back-compat: a caller that supplies no context still gets a valid pointer."""
    from apps.blob_trigger_app.pptx_stage import build_deck_pointer

    deck = tmp_path / "d.pptx"
    deck.write_bytes(b"PK\x03\x04" + b"0" * 64)
    pointer = build_deck_pointer(deck, client_id="acme", period="2026-06",
                                 latest_uri="decks/acme/latest/investor_pack.pptx")
    assert pointer["client_id"] == "acme"
    assert pointer["checksum"].startswith("sha256:")
    assert "report_scope" not in pointer
