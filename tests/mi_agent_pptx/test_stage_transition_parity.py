"""The deck and the dashboard must show the SAME stage movement.

Not "computed the same way" — the SAME object. This is the parity discipline
`test_channel_parity.py` already applies to the funded surfaces, pointed at the
governed stage-transition capability: drive the REAL React HTTP route and the
REAL deck data path over one fixture, and compare what each produced.

It has its own fixture rather than extending the funded parity book because the
transition capability needs a governed pair of WEEKLY PIPELINE extracts, which
that book deliberately does not have. Adding extracts to it would have changed
what the existing parity tests see.

The defect this is built to catch is specific and has a precedent in this
estate: one surface quietly deriving a number the other reads from the engine.
A deck that recomputed "cases that moved from Offer" would look right on its own
and disagree with the dashboard the moment anything changed. So the assertions
compare the SHARED payload, never two separately recalculated values.
"""

from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
import pytest

_REPO = Path(__file__).resolve().parents[2]
_CENTRAL = "18_central_lender_tape.csv"
#: The deterministic two-snapshot pack from the engine sprint.
PIPELINE = _REPO / "tests" / "fixtures" / "pipeline_transition_2w"
CLIENT = "acme"
AS_OF = "2026-06-30"
RUN_ID = "mi_2026_06"

pytestmark = pytest.mark.skipif(not PIPELINE.exists(),
                                reason="transition fixture pack not present")


def _loan(i: int, balance: float, cut: str) -> dict:
    return {
        "unique_identifier": f"direct_001_L{i:04d}",
        "source_portfolio_id": "direct_001",
        "source_portfolio_type": "direct",
        "source_portfolio_label": "Direct 001",
        "current_outstanding_balance": balance,
        "current_principal_balance": balance,
        "original_principal_balance": balance * 1.04,
        "current_valuation_amount": balance / 0.45,
        "current_loan_to_value": 45.0,
        "current_interest_rate": 7.1,
        "youngest_borrower_age": 72,
        "geographic_region_collateral": "London",
        "collateral_geography": "London",
        "geographic_region_obligor": "TLI3",
        "origination_channel": "Direct",
        "origination_date": "2021-04-01",
        "data_cut_off_date": cut,
    }


@pytest.fixture()
def book(tmp_path, monkeypatch):
    """A governed run both channels resolve from the same environment.

    The funded side is deliberately minimal — this is a PIPELINE parity test —
    but it has to exist, because both channels resolve a portfolio before they
    resolve anything about it.
    """
    root = tmp_path / "root"
    for run_id, date in (("mi_2026_05", "2026-05-31"), (RUN_ID, AS_OF)):
        central = root / CLIENT / run_id / "central"
        central.mkdir(parents=True, exist_ok=True)
        frame = pd.DataFrame([_loan(i, 300_000.0, date) for i in range(1, 4)])
        frame.to_csv(central / _CENTRAL, index=False)
        dated = root / date
        dated.mkdir(parents=True, exist_ok=True)
        frame.to_csv(dated / "platform_canonical_typed.csv", index=False)

    run_dir = root / "orun_stx"
    (run_dir / "out_platform").mkdir(parents=True, exist_ok=True)
    pd.DataFrame([_loan(i, 300_000.0, AS_OF) for i in range(1, 4)]).to_csv(
        run_dir / "out_platform" / "platform_canonical_typed.csv", index=False)
    (run_dir / "run_state.json").write_text(json.dumps({
        "run_id": RUN_ID, "client_id": CLIENT, "reporting_date": AS_OF,
        "out_root": str(root)}), encoding="utf-8")

    monkeypatch.setenv("MI_AGENT_ONBOARDING_OUTPUT_ROOT", str(root))
    monkeypatch.setenv("MI_AGENT_CLIENT_ID", CLIENT)
    # BOTH channels are pointed at the SAME governed pipeline pack. That is the
    # premise of the whole test: same source, same engine, same answer.
    monkeypatch.setenv("MI_AGENT_PIPELINE_ROOT", str(PIPELINE))
    monkeypatch.setenv("MI_AGENT_AUTH_ENABLED", "false")
    monkeypatch.setenv("TRAKT_MI_ENHANCED_HOVERS", "true")
    return run_dir


@pytest.fixture()
def deck(book):
    """The PPTX channel, through its real shared aggregate."""
    from mi_agent_pptx.mi_api import build_dashboard_data
    return build_dashboard_data(book, client_id=CLIENT, as_of=AS_OF,
                                output_root=str(Path(book).parent))


@pytest.fixture()
def react(book):
    """The React channel, through its real HTTP route."""
    from fastapi.testclient import TestClient
    from mi_agent_api.app import app
    from mi_agent_api import movement_detail as md

    client = TestClient(app)

    def get():
        res = client.get("/mi/insight/movement-detail", params={
            "portfolioId": f"{CLIENT}/{RUN_ID}",
            "detailType": md.DETAIL_STAGE_TRANSITION})
        assert res.status_code == 200, f"{res.status_code}: {res.text}"
        body = res.json()
        body.pop("portfolioScope", None)
        return body

    return get


# --------------------------------------------------------------------------- #
# Both channels resolved the capability at all.
# --------------------------------------------------------------------------- #

def test_both_channels_resolved_the_governed_capability(deck, react):
    api = react()
    assert api["available"], f"the dashboard has no transitions: {api.get('reason')}"
    assert deck.stage_transitions, "the deck resolved no stage transitions"
    assert deck.stage_transitions["available"], deck.stage_transitions.get("reason")


def test_the_deck_carries_the_capability_on_the_shared_payload(deck):
    """Not a PPTX-only loader: it is a field on `DashboardData`, the same
    aggregate every other slide reads, resolved by the same governed function."""
    from mi_agent_pptx.mi_api import DashboardData
    assert "stage_transitions" in DashboardData.__dataclass_fields__
    assert deck.stage_transitions["detail_type"] == "PIPELINE_STAGE_TRANSITION"


# --------------------------------------------------------------------------- #
# The parity requirement, block by block.
# --------------------------------------------------------------------------- #

def test_the_whole_governed_object_is_identical(deck, react):
    """The strongest form of the requirement: not field-by-field agreement but
    the same object. Everything below is a readable restatement of this."""
    mine = dict(deck.stage_transitions)
    theirs = dict(react())
    # Resolution provenance differs by construction (the deck knows its run id).
    for k in ("run_id", "scope", "portfolio_id"):
        mine.pop(k, None)
        theirs.pop(k, None)
    assert mine == theirs


def test_the_identifier_matches(deck, react):
    assert deck.stage_transitions["identifier"] == react()["identifier"]
    assert react()["identifier"] == "pipeline_case_identifier"


def test_the_reporting_window_matches(deck, react):
    api = react()
    assert deck.stage_transitions["comparison_date"] == api["comparison_date"]
    assert deck.stage_transitions["as_of_date"] == api["as_of_date"]


def test_the_source_destination_transitions_match(deck, react):
    assert deck.stage_transitions["transitions"] == react()["transitions"]


def test_the_transition_counts_and_amounts_match(deck, react):
    for mine, theirs in zip(deck.stage_transitions["transitions"],
                            react()["transitions"]):
        assert (mine["source_stage"], mine["destination_stage"]) == \
               (theirs["source_stage"], theirs["destination_stage"])
        assert mine["case_count"] == theirs["case_count"]
        assert mine["prior_amount"] == theirs["prior_amount"]
        assert mine["latest_amount"] == theirs["latest_amount"]
        assert mine["amount_change"] == theirs["amount_change"]


def test_the_kfi_to_application_flow_matches(deck, react):
    def find(payload):
        return next(t for t in payload["transitions"]
                    if t["source_stage"] == "KFI"
                    and t["destination_stage"] == "APPLICATION")
    mine, theirs = find(deck.stage_transitions), find(react())
    assert mine == theirs
    assert mine["case_count"] == 2                 # fixture truth
    assert mine["latest_amount"] == 920_000.0


def test_the_application_to_offer_flow_matches(deck, react):
    def find(payload):
        return next(t for t in payload["transitions"]
                    if t["source_stage"] == "APPLICATION"
                    and t["destination_stage"] == "OFFER")
    mine, theirs = find(deck.stage_transitions), find(react())
    assert mine == theirs
    assert mine["case_count"] == 2
    assert mine["amount_change"] == -10_000.0


def test_the_completion_flow_matches(deck, react):
    def find(payload):
        return next(t for t in payload["transitions"]
                    if t["destination_stage"] == "COMPLETED")
    mine, theirs = find(deck.stage_transitions), find(react())
    assert mine == theirs
    assert mine["source_stage"] == "OFFER"
    assert mine["case_count"] == 1


def test_the_arrivals_match(deck, react):
    assert deck.stage_transitions["new_arrivals"] == react()["new_arrivals"]
    # And neither side invented a source stage for them.
    for row in react()["new_arrivals"]:
        assert "source_stage" not in row


def test_the_stayers_and_their_amendments_match(deck, react):
    assert deck.stage_transitions["stayers"] == react()["stayers"]
    changes = {s["stage"]: s["amount_change"] for s in react()["stayers"]}
    assert changes["KFI"] == 20_000.0            # amended up, still stayers
    assert changes["APPLICATION"] == -20_000.0   # amended down


def test_the_departures_and_their_outcomes_match(deck, react):
    assert deck.stage_transitions["departures"] == react()["departures"]
    outcomes = {d["source_stage"]: d["governed_outcome"] for d in react()["departures"]}
    assert outcomes["COMPLETED"] == "COMPLETED"
    assert outcomes["WITHDRAWN"] == "WITHDRAWN"
    # Neither channel resolved these into an outcome the data never showed.
    assert outcomes["OFFER"] == "unclassified_departure"
    assert outcomes["APPLICATION"] == "unclassified_departure"


def test_the_reconciliation_and_residuals_match(deck, react):
    mine = deck.stage_transitions["reconciliation"]
    theirs = react()["reconciliation"]
    assert mine == theirs
    assert mine["count_reconciliation_residual"] == 0
    assert mine["amount_reconciliation_residual"] == 0.0
    assert mine["global"]["residual"] == 0


def test_the_event_totals_match(deck, react):
    assert deck.stage_transitions["event_totals"] == react()["event_totals"]


# --------------------------------------------------------------------------- #
# The slide renders from that payload, through the real composition path.
# --------------------------------------------------------------------------- #

def test_the_slide_is_selected_and_renders(deck, tmp_path):
    import yaml
    from mi_agent_pptx.composition import build_facts, select_slides

    spec = yaml.safe_load(
        (_REPO / "configs/pptx/investor_pack.yaml").read_text(encoding="utf-8"))
    facts = build_facts(deck)
    assert facts["has_stage_transitions"] is True
    kept, _omitted = select_slides(spec["slides"], deck, facts)
    assert "stage_transitions" in {s["id"] for s in kept}


def _real_deck(book, out) -> int:
    """The REAL production deck route — the same CLI a scheduled pack runs."""
    from mi_agent_pptx.cli import run as cli_run
    return cli_run(["--run-dir", str(book),
                    "--deck-config", str(_REPO / "configs/pptx/investor_pack.yaml"),
                    "--client-name", "Acme Capital", "--as-of-date", AS_OF,
                    "--output-root", str(Path(book).parent), "--output", str(out)])


def test_the_real_deck_route_renders_the_slide_from_the_engines_numbers(book, tmp_path):
    """Rendered text out of the REAL CLI, so a deck that stored the payload
    correctly but formatted it into something else is still caught."""
    from pptx import Presentation

    out = tmp_path / "stx.pptx"
    assert _real_deck(book, out) == 0
    assert out.exists()

    text = " ".join(sh.text_frame.text for s in Presentation(str(out)).slides
                    for sh in s.shapes if sh.has_text_frame)
    assert "Pipeline Stage Movement" in text
    assert "Cases that moved stage" in text
    assert "Arrivals, stayers and departures" in text
    assert "Case reconciliation by stage" in text
    assert "Stage value — opening vs closing" in text
    # The identifier, the population and the residuals are stated, never implied.
    assert "pipeline_case_identifier" in text
    assert "Reconciliation residual 0 cases" in text
    # The window the engine resolved, in the deck's own date style — stated on
    # the slide, never left for the reader to assume.
    assert "5 June 2026" in text and "12 June 2026" in text


def test_an_unavailable_capability_omits_the_slide_with_the_engines_reason(deck):
    """The deck never decides availability for itself, and never renders an
    empty matrix — which a reader would take as "nothing moved"."""
    import yaml
    from mi_agent_pptx.composition import build_facts, select_slides

    deck.stage_transitions = {
        "available": False,
        "reason": "2 duplicate pipeline_case_identifier value(s) in the latest "
                  "snapshot prevent deterministic case matching.",
        "reason_code": "duplicate_case_identifiers",
    }
    spec = yaml.safe_load(
        (_REPO / "configs/pptx/investor_pack.yaml").read_text(encoding="utf-8"))
    kept, omitted = select_slides(spec["slides"], deck, build_facts(deck))
    assert "stage_transitions" not in {s["id"] for s in kept}
    reason = next(o.reason for o in omitted if o.slide_id == "stage_transitions")
    assert "duplicate pipeline_case_identifier" in reason
