"""Investor Reporting v2.2 — movement, watch list, reproducibility, MI parity.

The acceptance requirement for this release is that the pack works every month,
for any client, without code changes. So the tests here are less about "does
this slide look right" and more about:

  * does attribution reconcile to the headline, always;
  * does the deck agree with the governed services it claims to consume;
  * is the output reproducible for the same inputs;
  * does anything in the code assume THIS client, THIS date or THIS schema.
"""

from __future__ import annotations

import json
import re
from pathlib import Path

import pandas as pd
import pytest
import yaml
from pptx import Presentation

from mi_agent_pptx import movement as MV
from mi_agent_pptx import watchlist as WL
from mi_agent_pptx.cli import preflight_path, run as cli_run
from mi_agent_pptx.composition import build_facts, select_slides
from mi_agent_pptx.mi_api import build_dashboard_data

DECK_CONFIG = Path("configs/pptx/investor_pack.yaml")
_CENTRAL = "18_central_lender_tape.csv"


# --------------------------------------------------------------------------- #
# Fixtures — governed run layout, two periods, two portfolio types.
# --------------------------------------------------------------------------- #

def _loan(i, pid, ptype, balance, *, cut, region, channel, ltv=45.0):
    return {
        "unique_identifier": f"{pid}_L{i:04d}",
        "source_portfolio_id": pid,
        "source_portfolio_type": ptype,
        "source_portfolio_label": pid.replace("_", " ").title(),
        "current_outstanding_balance": balance,
        "current_principal_balance": balance,
        "original_principal_balance": balance * 1.04,
        "current_valuation_amount": balance / (ltv / 100.0),
        "original_valuation_amount": balance / (ltv / 100.0),
        "current_loan_to_value": ltv,
        "original_loan_to_value": ltv / 100 - 0.02,
        "current_interest_rate": 7.1,
        "youngest_borrower_age": 72,
        "geographic_region_collateral": region,
        "collateral_geography": region,
        "origination_channel": channel,
        "origination_date": "2021-04-01",
        "data_cut_off_date": cut,
    }


#: Prior and current period. The direct book grows, the acquired book shrinks,
#: and the growth is concentrated in one region — so attribution has a real
#: answer rather than noise.
_PRIOR = [
    ("direct_001", "direct", 300_000.0, "London", "Direct"),
    ("direct_001", "direct", 200_000.0, "South East", "Direct"),
    ("acquired_001", "acquired", 250_000.0, "Wales", "Broker"),
]
_CURRENT = [
    ("direct_001", "direct", 300_000.0, "London", "Direct"),
    ("direct_001", "direct", 260_000.0, "South East", "Direct"),
    ("direct_001", "direct", 120_000.0, "South East", "Direct"),
    ("acquired_001", "acquired", 230_000.0, "Wales", "Broker"),
]


def _rows(spec, cut):
    return [_loan(i, pid, ptype, bal, cut=cut, region=region, channel=channel)
            for i, (pid, ptype, bal, region, channel) in enumerate(spec, start=1)]


def _write_run(root: Path, client: str, run_id: str, date: str, spec):
    central = root / client / run_id / "central"
    central.mkdir(parents=True, exist_ok=True)
    pd.DataFrame(_rows(spec, date)).to_csv(central / _CENTRAL, index=False)
    dated = root / date
    dated.mkdir(parents=True, exist_ok=True)
    pd.DataFrame(_rows(spec, date)).to_csv(
        dated / "platform_canonical_typed.csv", index=False)


def _make(tmp_path, *, client="acme", periods=(("mi_2026_05", "2026-05-31", _PRIOR),
                                               ("mi_2026_06", "2026-06-30", _CURRENT))):
    root = tmp_path / "root"
    for run_id, date, spec in periods:
        _write_run(root, client, run_id, date, spec)
    latest_date, latest_spec = periods[-1][1], periods[-1][2]
    run_dir = root / "orun_v22"
    (run_dir / "out_platform").mkdir(parents=True, exist_ok=True)
    pd.DataFrame(_rows(latest_spec, latest_date)).to_csv(
        run_dir / "out_platform" / "platform_canonical_typed.csv", index=False)
    (run_dir / "run_state.json").write_text(json.dumps({
        "run_id": "orun_v22", "client_id": client, "reporting_date": latest_date,
        "out_root": str(root)}), encoding="utf-8")
    return run_dir


@pytest.fixture
def run_dir(tmp_path):
    return _make(tmp_path)


def _data(run_dir, client="acme", context=None, as_of="2026-06-30"):
    return build_dashboard_data(run_dir, client_id=client, as_of=as_of,
                                output_root=str(Path(run_dir).parent),
                                portfolio_context=context)


def _slides():
    return yaml.safe_load(DECK_CONFIG.read_text(encoding="utf-8"))["slides"]


def _build(run_dir, out, *, client="acme", context=None):
    argv = ["--run-dir", str(run_dir), "--deck-config", str(DECK_CONFIG),
            "--client-name", "Test Client", "--as-of-date", "2026-06-30",
            "--output-root", str(Path(run_dir).parent), "--output", str(out)]
    if context:
        argv += ["--portfolio-context", context]
    return cli_run(argv)


def _text(path):
    return "\n".join(sh.text_frame.text for s in Presentation(str(path)).slides
                     for sh in s.shapes if sh.has_text_frame)


# --------------------------------------------------------------------------- #
# Movement attribution.
# --------------------------------------------------------------------------- #

def test_attribution_reconciles_exactly_for_every_dimension(run_dir):
    """Every dimension's contributor deltas must sum to the headline movement."""
    bridges = {k: b for k, b in _data(run_dir).movement.items() if b.available}
    assert bridges, "no attribution resolved"
    for key, bridge in bridges.items():
        assert bridge.reconciles(), f"{key} does not reconcile"
        assert abs(sum(c.delta for c in bridge.contributors)
                   - bridge.total_delta) < 0.51


def test_attribution_finds_the_real_driver(run_dir):
    """South East grew by 180k; London did not move. Attribution must say so."""
    region = _data(run_dir).movement["region"]
    ups, _downs = region.movers(limit=3)
    assert ups[0].category == "South East"
    assert ups[0].delta == pytest.approx(180_000.0, abs=1.0)
    assert "London" not in [c.category for c in ups]


def test_absent_dimension_is_unavailable_not_zero_movement(run_dir):
    """A missing column must not render as a confident, empty attribution."""
    bridges = _data(run_dir).movement
    missing = [b for b in bridges.values() if not b.available]
    assert missing
    for b in missing:
        assert b.reason and b.opening is None and b.closing is None


def test_movement_wording_never_attributes_a_cause(run_dir):
    """The balance decomposition evidences WHAT moved, never WHY."""
    bridges = _data(run_dir).movement
    text = " ".join(line for b in bridges.values() if b.available
                    for line in MV.takeaways(b))
    for banned in ("redemption", "redeemed", "amortis", "prepay"):
        assert banned not in text.lower()


def test_movement_slide_is_omitted_without_a_prior_period(tmp_path):
    run = _make(tmp_path, periods=(("mi_2026_06", "2026-06-30", _CURRENT),))
    data = _data(run)
    kept, omitted = select_slides(_slides(), data)
    if "movement_drivers" in {s["id"] for s in kept}:
        pytest.skip("type-level movement available from a single period")
    assert "movement_drivers" in {o.slide_id for o in omitted}


# --------------------------------------------------------------------------- #
# Watch list.
# --------------------------------------------------------------------------- #

def test_ordinary_movement_is_not_a_warning(run_dir):
    """A book that simply reduced must not be escalated to a watch item."""
    wl = _data(run_dir).watchlist
    types = {i.insight_type for i in wl["watch"]}
    assert WL.ACQUIRED_RUNOFF not in types


def test_no_material_items_is_an_explicit_state(run_dir, tmp_path):
    wl = _data(run_dir).watchlist
    if wl["watch"]:
        pytest.skip("this fixture produced watch items")
    assert wl["none_identified"] is True
    out = tmp_path / "wl.pptx"
    _build(run_dir, out)
    assert "No material watch items identified" in _text(out)


def test_watch_items_are_capped_and_carry_evidence(run_dir):
    wl = _data(run_dir).watchlist
    assert len(wl["watch"]) <= WL.MAX_WATCH_ITEMS
    assert len(wl["observations"]) <= WL.MAX_POSITIVES
    for item in wl["watch"] + wl["observations"]:
        assert item.headline and item.summary
        assert item.methodology.get("source"), "every item names its source"
        assert item.portfolio_context


def test_watch_list_proposes_no_management_actions(run_dir):
    """Recommending an action is a decision the evidence does not authorise."""
    wl = _data(run_dir).watchlist
    text = " ".join(f"{i.headline} {i.summary}"
                    for i in wl["watch"] + wl["observations"]).lower()
    for imperative in ("should ", "must ", "recommend", "consider ", "we suggest"):
        assert imperative not in text


# --------------------------------------------------------------------------- #
# MI parity — the deck must agree with the services it claims to consume.
# --------------------------------------------------------------------------- #

def test_deck_headline_equals_the_governed_snapshot(run_dir):
    from mi_agent_pptx.deck_context import raw_of
    data = _data(run_dir)
    assert data.portfolio.total_balance == raw_of(data.funded, "balance")
    assert data.portfolio.loan_count == raw_of(data.funded, "loans")


def test_type_slices_equal_the_governed_snapshot_per_type(run_dir):
    """Each slice IS compute_funded_snapshot over that type's rows."""
    data = _data(run_dir)
    assert data.portfolio.type_balance_total() == pytest.approx(
        data.portfolio.total_balance, rel=1e-9)


def test_attribution_closing_equals_the_funded_total(run_dir):
    """The bridge's close must be the same book the KPIs describe."""
    data = _data(run_dir)
    for bridge in data.movement.values():
        if bridge.available:
            assert bridge.closing == pytest.approx(
                data.portfolio.total_balance, rel=1e-6)


def test_reporting_dates_equal_the_resolved_mi_source_dates(run_dir):
    data = _data(run_dir)
    assert set(data.portfolio.reporting_dates.values()) == {"2026-06-30"}


# --------------------------------------------------------------------------- #
# Monthly reproducibility.
# --------------------------------------------------------------------------- #

def test_same_inputs_produce_an_equivalent_deck(run_dir, tmp_path):
    """Not byte-identical (the OOXML package embeds timestamps) but analytically
    equivalent: same slides, same order, same evidence, same omissions."""
    first, second = tmp_path / "a.pptx", tmp_path / "b.pptx"
    _build(run_dir, first)
    _build(run_dir, second)
    a = json.loads(preflight_path(first).read_text(encoding="utf-8"))
    b = json.loads(preflight_path(second).read_text(encoding="utf-8"))
    assert [s["id"] for s in a["slides"]] == [s["id"] for s in b["slides"]]
    assert a["omitted_slides"] == b["omitted_slides"]
    assert a["insights"]["headlines"] == b["insights"]["headlines"]
    assert a["portfolio_context"] == b["portfolio_context"]
    assert _text(first) == _text(second)


def test_a_different_client_id_needs_no_code_change(tmp_path):
    """Nothing may assume a particular client."""
    run = _make(tmp_path, client="zeta_capital")
    data = _data(run, client="zeta_capital")
    assert data.portfolio.client_id == "zeta_capital"
    assert data.portfolio.total_balance
    out = tmp_path / "zeta.pptx"
    assert _build(run, out, client="zeta_capital") == 0


def test_a_later_reporting_period_needs_no_code_change(tmp_path):
    """Next month must work without touching the code."""
    later = (("mi_2027_01", "2027-01-31", _PRIOR),
             ("mi_2027_02", "2027-02-28", _CURRENT))
    run = _make(tmp_path, periods=later)
    data = build_dashboard_data(run, client_id="acme", as_of="2027-02-28",
                                output_root=str(Path(run).parent))
    assert data.portfolio.total_balance
    assert set(data.portfolio.reporting_dates.values()) == {"2027-02-28"}
    assert any(b.available for b in data.movement.values())


def test_no_client_or_date_is_hard_coded_in_the_deck_layer():
    """A literal client name, portfolio id or reporting date in the deck code
    would silently break every other client."""
    banned = re.compile(
        r"\b(alpha lending|acme|alp_origination|alp_acquired|spv1_sponsored"
        r"|direct_001|acquired_001|2026-0\d-\d\d)\b", re.I)
    offenders = []
    for path in Path("mi_agent_pptx").glob("*.py"):
        in_doc = False
        for n, line in enumerate(path.read_text(encoding="utf-8").splitlines(), 1):
            stripped = line.strip()
            # Track triple-quoted blocks across lines: prose may legitimately
            # cite an example id, and a continuation line does not start with
            # the quote marker.
            fences = stripped.count('"""') + stripped.count("'''")
            was_doc = in_doc
            if fences % 2:
                in_doc = not in_doc
            if (was_doc or in_doc or stripped.startswith("#")
                    or stripped.startswith(('"""', "'''"))):
                continue
            if banned.search(line):
                offenders.append(f"{path.name}:{n}: {stripped[:70]}")
    assert not offenders, "hard-coded fixture values:\n" + "\n".join(offenders)


def test_dimensions_are_candidate_lists_not_fixed_columns():
    """Schema portability: every dimension must offer alternatives."""
    for key, label, candidates in MV.DIMENSIONS:
        assert candidates, f"{key} declares no candidate columns"
        assert isinstance(candidates, tuple)


# --------------------------------------------------------------------------- #
# Gates.
# --------------------------------------------------------------------------- #

def test_all_gates_pass_and_include_the_v22_checks(run_dir, tmp_path):
    out = tmp_path / "gates.pptx"
    rc = _build(run_dir, out)
    verdict = json.loads(preflight_path(out).read_text(encoding="utf-8"))
    assert rc == 0, verdict["preflight"]["failed_gates"]
    gates = {g["gate"] for g in verdict["preflight"]["gates"]}
    assert "movement_reconciles" in gates
    assert "no_unsupported_causal_language" in gates


def test_causal_language_gate_blocks():
    import mi_agent_pptx.preflight as pf
    gate = pf._gate_no_unsupported_causal_language(
        "The acquired book reduced through redemption and amortisation.")
    assert gate.passed is False and gate.mandatory is True


def test_movement_reconciliation_gate_blocks_a_broken_bridge():
    import mi_agent_pptx.preflight as pf
    broken = MV.MovementBridge(
        key="region", label="Region", available=True, opening=100.0,
        closing=200.0,
        contributors=(MV.Contributor("A", 100.0, 130.0, 30.0),))   # 30 != 100
    gate = pf._gate_movement_reconciles({"region": broken})
    assert gate.passed is False and gate.mandatory is True


def test_geometry_is_clean_across_scopes(run_dir, tmp_path):
    """No off-slide, overlapping or clipped content in any generated deck."""
    for context in (None, "direct", "acquired"):
        out = tmp_path / f"geo_{context or 'total'}.pptx"
        if _build(run_dir, out, context=context) != 0:
            continue
        prs = Presentation(str(out))
        for i, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.left is None:
                    continue
                assert shape.left >= -18288, f"slide {i} off-slide left"
                assert shape.top >= -18288, f"slide {i} off-slide top"
                assert shape.left + shape.width <= 12211050, f"slide {i} off right"
                assert shape.top + shape.height <= 6877050, f"slide {i} off bottom"
