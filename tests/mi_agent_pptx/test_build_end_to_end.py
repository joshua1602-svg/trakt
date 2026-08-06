"""End-to-end deck build tests for the portfolio-aware, dashboard-aligned deck.

Every slide renders directly from an MI Agent API payload
(``mi_agent_pptx.mi_api``). These tests assert the deck always builds, stays
16:9, keeps the pipeline lens separate from the funded lens, composes itself to
the portfolio it is describing (rather than emitting placeholder pages), and
fails CLOSED — a deck that cannot state its own scope is generated but withheld
from publication.
"""

from __future__ import annotations

import json
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Emu

from mi_agent_pptx.cli import PREFLIGHT_BLOCKED_RC, preflight_path, run as cli_run
from mi_agent_pptx.composition import ConditionError, evaluate_condition
from mi_agent_pptx.deck import DeckBuilder, DeckContext
from mi_agent_pptx.mi_api import build_dashboard_data

#: Slides the pack must always carry — they are its disclosure spine.
#: The pack's disclosure spine. The former separate appendix is now folded
#: into the single investor-safe Data and Methodology page.
MANDATORY_SLIDE_IDS = {"cover", "methodology"}


def _slides(deck_config_path):
    cfg = yaml.safe_load(Path(deck_config_path).read_text(encoding="utf-8")) or {}
    return cfg.get("slides", [])


def _facts(deck_config_path):
    from mi_agent_pptx.composition import build_facts

    class _Empty:
        portfolio = None

    return build_facts(_Empty())


def test_deck_config_is_dispatchable_and_well_formed(deck_config_path):
    """Every configured slide maps to a handler and declares a sane condition.

    The deck no longer has a fixed length — composition decides that — so the old
    12-to-20 bound is replaced by structural assertions that still hold.
    """
    slides = _slides(deck_config_path)
    assert slides, "deck config declares no slides"

    ids = [s.get("id") for s in slides]
    assert len(ids) == len(set(ids)), "duplicate slide ids in the deck config"
    assert MANDATORY_SLIDE_IDS.issubset(set(ids))

    for spec in slides:
        assert spec.get("type") in DeckBuilder._DISPATCH, spec.get("type")
        # Every slide states the investor question it answers.
        assert spec.get("question"), f"{spec.get('id')} declares no question"

    # Every `when:` must parse against the real fact vocabulary — a typo in the
    # config must fail here, not silently include a slide in production.
    facts = _facts(deck_config_path)
    for spec in slides:
        if spec.get("when"):
            evaluate_condition(str(spec["when"]), facts)


def test_unknown_fact_in_a_condition_is_rejected():
    """A condition may only refer to governed facts."""
    try:
        evaluate_condition("has_unicorns", {"has_direct": True})
    except ConditionError:
        pass
    else:  # pragma: no cover - the assertion below reports the failure
        raise AssertionError("an unknown fact must not evaluate")


def test_a_condition_cannot_execute_code():
    """`when:` is not an eval hatch."""
    for hostile in ("__import__('os').system('true')", "open('/etc/passwd')",
                    "has_direct.__class__"):
        try:
            evaluate_condition(hostile, {"has_direct": True})
        except ConditionError:
            continue
        raise AssertionError(f"hostile condition was not rejected: {hostile}")


def test_cli_builds_valid_pptx(run_dir, deck_config_path, tmp_path):
    out = tmp_path / "deck.pptx"
    rc = cli_run([
        "--run-dir", str(run_dir),
        "--deck-config", str(deck_config_path),
        "--client-name", "Test Client",
        "--as-of-date", "2026-01-31",
        "--output", str(out),
    ])
    assert rc == 0
    assert out.exists()

    prs = Presentation(str(out))
    slides = list(prs.slides)
    assert len(slides) >= len(MANDATORY_SLIDE_IDS)
    # 16:9 widescreen.
    assert round(Emu(prs.slide_width).inches, 2) == 13.33
    # Every slide carries at least a title + one more text frame.
    for s in slides:
        texts = [sh for sh in s.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        assert len(texts) >= 2


def test_missing_artefacts_still_build_a_valid_deck(empty_run_dir, deck_config_path,
                                                    tmp_path):
    """A run with no artefacts must still yield a valid, composed deck."""
    out = tmp_path / "empty_deck.pptx"
    cli_run([
        "--run-dir", str(empty_run_dir),
        "--deck-config", str(deck_config_path),
        "--client-name", "No Data Co",
        "--output", str(out),
    ])
    assert out.exists()
    ids = {r["id"] for r in json.loads(
        preflight_path(out).read_text(encoding="utf-8"))["slides"]}
    assert MANDATORY_SLIDE_IDS.issubset(ids)


def test_publication_fails_closed_when_a_mandatory_gate_fails(
        run_dir, deck_config_path, tmp_path):
    """A deck missing its disclosure spine is generated but NOT published.

    The cover is what states the reporting scope, so a deck without one is
    exactly the artefact this release exists to prevent shipping. It is still
    written — an operator needs to see it to diagnose the run — but publication
    is withheld and the reason travels with it.
    """
    broken = tmp_path / "no_cover.yaml"
    cfg = yaml.safe_load(Path(deck_config_path).read_text(encoding="utf-8"))
    cfg["slides"] = [s for s in cfg["slides"] if s.get("id") != "cover"]
    broken.write_text(yaml.safe_dump(cfg), encoding="utf-8")

    out = tmp_path / "blocked.pptx"
    rc = cli_run([
        "--run-dir", str(run_dir),
        "--deck-config", str(broken),
        "--client-name", "Test Client",
        "--as-of-date", "2026-01-31",
        "--output", str(out),
    ])
    assert rc == PREFLIGHT_BLOCKED_RC
    assert out.exists(), "the deck must still be generated for diagnosis"

    sidecar = json.loads(preflight_path(out).read_text(encoding="utf-8"))
    assert sidecar["preflight"]["publishable"] is False
    assert "mandatory_slides" in sidecar["preflight"]["failed_gates"]


def test_funded_kpis_come_from_the_dashboard_snapshot(run_dir):
    """The funded tiles are the funded snapshot KPIs verbatim — the same payload
    the dashboard's Funded header renders."""
    data = build_dashboard_data(run_dir, as_of="2026-01-31")
    kpis = {k.get("id"): k for k in data.funded.get("kpis", [])}
    assert "balance" in kpis
    assert kpis["balance"].get("value")           # rendered, formatted value
    assert data.cohorts.get("cohorts")            # vintage composition resolves


def test_pipeline_lens_is_separate_from_funded(run_dir):
    """The funded run has no pipeline source, so the pipeline payload must be
    empty rather than borrowing the funded total — the lens separation the
    dashboard enforces."""
    data = build_dashboard_data(run_dir, as_of="2026-01-31")
    assert data.funded.get("kpis")                # funded resolved
    assert not data.pipeline                       # pipeline did NOT resolve
    # And a coverage note explains the omission.
    assert any("pipeline" in n.lower() for n in data.notes)


def test_absent_sections_are_omitted_with_a_reason_not_placeholdered(
        run_dir, deck_config_path, tmp_path):
    """A section with no data is left OUT and explained — never rendered as a
    'no data available' page."""
    data = build_dashboard_data(run_dir, as_of="2026-01-31")
    ctx = DeckContext(client_name="X", as_of_date="2026-01-31",
                      run_dir=str(run_dir), work_dir=str(tmp_path / "charts"))
    report = DeckBuilder(data, ctx).build(_slides(deck_config_path),
                                          tmp_path / "d.pptx")

    for rec in report["slides"]:
        assert rec.get("title")
    # No placeholder page reaches an investor deck.
    assert not [r for r in report["slides"] if r.get("placeholder")]

    # This run has no pipeline source, so the pipeline section is omitted — and
    # the omission carries a reason a reader can act on.
    omitted = {o["slide_id"]: o for o in report["omitted_slides"]}
    assert "pipeline" in omitted
    assert omitted["pipeline"]["reason"]
    assert "pipeline" not in {r["id"] for r in report["slides"]}
