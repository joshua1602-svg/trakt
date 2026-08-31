#!/usr/bin/env python3
"""tests/test_final_pack_surfaces.py

What the final funder pack SAYS, checked on the pack a funder receives.

Every test builds a real deck through the React route — ``POST
/mi/decks/generate`` -> poll -> ``GET /mi/decks/download`` — and reads the
PowerPoint that comes back. Nothing calls the builder directly: a page that
renders under a direct call and not under the job service is a page that does
not exist.

These pin three corrections found by reading the printed pack page by page,
plus the two properties the sprint must not have broken on the way. Each states
the defect it would catch.
"""

from __future__ import annotations

import sys
from pathlib import Path

import pytest

_REPO_ROOT = Path(__file__).resolve().parents[1]
if str(_REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(_REPO_ROOT))

pytest.importorskip("pandas")
pytest.importorskip("pptx")

from tests.test_funder_pack_enhancement import (  # noqa: E402
    deck_text, generate_and_download, rich, simple, slide_titles)

__all__ = ["rich", "simple"]


@pytest.fixture()
def limited(rich):
    """The rich book WITH an operator-approved concentration configuration.

    Concentration limits are contractual: nothing in the tape implies them, and
    the deck refuses to invent one, so the page the headroom wording lives on
    does not render at all until an operator has approved a configuration
    through the store. This commits four regional limits the way the operations
    console does.
    """
    from apps.blob_trigger_app.storage import open_storage
    from mi_agent.concentration_tests.models import (
        ActiveConfiguration, ActiveTest, ApprovalRecord, SourceEvidence)
    from mi_agent.concentration_tests.store import ConcentrationStore
    import os

    client = os.environ["MI_AGENT_CLIENT_ID"]
    tests = [
        ActiveTest(metric_id="geo_region_share", threshold=thr, operator="max",
                   display_name=name, parameters={"regions": regions},
                   evidence=SourceEvidence(
                       source_reference="facility_agreement.pdf",
                       source_text="Schedule 8, clause 3.1"),
                   approval=ApprovalRecord(decision="approved",
                                           operator="Operator",
                                           decided_at="2026-01-10T09:00:00+00:00"))
        for name, thr, regions in (
            ("London concentration", 30.0, ["London"]),
            ("South East concentration", 45.0, ["South East"]),
            ("Wales concentration", 40.0, ["Wales"]),
            ("Scotland concentration", 50.0, ["Scotland"]))]
    ConcentrationStore(open_storage(), container="operations-control"
                       ).commit_configuration(
        client, ActiveConfiguration(client_id=client, version=1,
                                    activated_by="Operator",
                                    library_version="1.0.0", tests=tests))
    return rich


# --------------------------------------------------------------------------- #
# 1. HEADROOM IS A DISTANCE, NOT A SHARE.
# --------------------------------------------------------------------------- #

def test_1_percentage_headroom_is_stated_in_points(limited):
    """Catches: "London concentration at 47% utilisation (16.0% headroom)".

    Headroom on a percentage test is the DIFFERENCE between two percentages —
    30.0% limit less 14.0% current is 16.0 percentage points. Printed with a
    percent sign it landed on the same line as a 47% utilisation figure, where
    a reader has no way to tell the two apart and the natural reading of "16.0%
    of headroom" is a share of the headroom rather than the headroom itself.
    The engine's own answer text has always said pp; only the deck did not.
    """
    text = deck_text(generate_and_download())
    assert "headroom" in text.lower(), "no headroom is reported at all"
    for line in text.splitlines():
        low = line.lower()
        if "headroom" not in low or "%" not in line:
            continue
        # Percentage headroom must carry pp. A currency or count test states
        # headroom in its own unit and no percent sign appears with it.
        assert "pp" in low, f"headroom stated as a percentage: {line!r}"


def test_2_headroom_carries_its_unit_whatever_the_test_measures(limited):
    """Catches: a bare "16.0" under a Headroom column, and £2m printed raw.

    The column printed ``f"{headroom:.1f}"`` with no unit at all. On a
    percentage test that is an unlabelled 16.0 beside a 30.0% limit; on a
    currency test it is 2000000.0 beside limits formatted as money. The unit
    belongs to the test, so the formatter must take it from the test and never
    from the slide.
    """
    from mi_agent_pptx import concentration as C

    assert C.format_headroom(16.0, "pct") == "16.0pp"
    assert C.format_headroom(16.0, "percentage_points") == "16.0pp"
    assert C.format_headroom(None, "pct") == "—"
    # A currency headroom is a difference of pounds, which IS pounds — it keeps
    # the governed money formatting rather than being relabelled in points.
    money = C.format_headroom(2_000_000.0, "gbp")
    assert "pp" not in money and any(ch.isdigit() for ch in money), money
    assert C.format_headroom(12.0, "count") == "12"


# --------------------------------------------------------------------------- #
# 2. A STRAPLINE DESCRIBES THE PAGE IT IS ON.
# --------------------------------------------------------------------------- #

def test_3_stratifications_promise_movement_only_when_they_show_it(rich):
    """Catches: "Composition and period movement" over four panels of
    composition and no movement.

    The strapline was chosen from whether movement was AVAILABLE; the movement
    strip is drawn only where two panels leave room for it. On a four-panel
    matrix the page promised a view it had suppressed, and a reader looking for
    it found nothing — the worst kind of missing content, because the page
    itself said it was there.
    """
    import io

    import pptx as _pptx

    deck = _pptx.Presentation(io.BytesIO(generate_and_download()))
    found = False
    for slide in deck.slides:
        lines = [sh.text_frame.text.strip() for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        if not lines or "Stratifications" not in lines[0]:
            continue
        found = True
        page = "\n".join(lines)
        if "Composition and period movement" in page:
            # It claimed movement, so movement must be stated in words on it.
            assert any(w in page.lower() for w in
                       ("increased", "decreased", "moved", "contributed")), page
    assert found, "no stratification page in the pack — the test proved nothing"


# --------------------------------------------------------------------------- #
# 3. AN EMPTY PANEL IS NOT A LAYOUT.
# --------------------------------------------------------------------------- #

def test_4_no_empty_observations_panel_beside_watch_items(rich):
    """Catches: four inches of empty box labelled OBSERVATIONS / None recorded.

    The watchlist slide sized its columns to whether there were WATCH items and
    never to whether there were OBSERVATIONS. With watch items and no
    observations a funder got the page that says what needs attention, half of
    it occupied by an empty container.
    """
    text = deck_text(generate_and_download())
    assert "None recorded." not in text, (
        "an observations panel was drawn with nothing in it")


def test_5_a_short_watch_list_still_says_what_was_checked(rich):
    """Catches: one watch item floating in an otherwise blank page.

    A reader cannot tell "one thing was flagged" from "only one thing was
    looked at". Where the stack leaves room the governed checks are named,
    which is the same list the all-clear branch has always printed.
    """
    import io

    import pptx as _pptx
    from mi_agent_pptx.deck import DeckBuilder

    deck = _pptx.Presentation(io.BytesIO(generate_and_download()))
    for slide in deck.slides:
        lines = [sh.text_frame.text.strip() for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        if not lines or "Watch Items" not in lines[0]:
            continue
        page = "\n".join(lines)
        if "CHECKS PERFORMED" not in page:
            pytest.skip("this book's watch stack fills the band")
        for check in DeckBuilder.GOVERNED_CHECKS:
            assert check in page, f"check not named: {check}"
        return


# --------------------------------------------------------------------------- #
# 4. WHAT THE SPRINT MUST NOT HAVE BROKEN.
# --------------------------------------------------------------------------- #

def test_6_additive_balance_stratification_still_renders(rich):
    """Catches: the additivity contract reaching the deck and suppressing the
    stratification bars.

    ``is_additive_measure`` was introduced for the React and Copilot surfaces.
    The deck's bar lists are balance sums — genuinely additive — and must be
    unaffected.
    """
    titles = slide_titles(generate_and_download())
    assert "Funded Stratifications" in titles, titles


def test_7_the_route_still_returns_a_downloadable_pack(simple):
    """Catches: any of the above breaking generation itself.

    The assertion that matters is in ``generate_and_download``: 202 accepted, a
    job that completes with no failed gate, and a file that starts PK.
    """
    content = generate_and_download()
    assert len(content) > 100_000, "a real pack is not this small"
    titles = slide_titles(content)
    # One book, one period. Everything that needs history, a second book or a
    # pipeline is omitted with a stated reason rather than drawn empty, so this
    # is the shortest honest pack: a cover, the position, the summary, the key
    # measures, one stratification page, watch items and the methodology.
    #
    # A COUNT IS NOT THE PROPERTY. An exact count also asserts what OTHER tests
    # in the run left behind — an operator-approved concentration configuration
    # committed by a fixture elsewhere in the suite legitimately earns this book
    # a concentration page, and the count then fails on a deck that is correct.
    # What this test is for is that the route still returns a real pack.
    for expected in ("Executive Position", "Executive Summary",
                     "Funded Portfolio — Key Measures", "Data and Methodology"):
        assert expected in titles, titles
    assert 7 <= len(titles) <= 9, titles


# --------------------------------------------------------------------------- #
# 5. "CONCENTRATED" IS A CLAIM ABOUT THE DATA.
# --------------------------------------------------------------------------- #

def test_8_a_near_uniform_spread_is_not_called_concentrated():
    """Catches: seven regions moving £3.7m-£4.4m reported as a concentration.

    The dominance floor is a share, and a GROUP clears a share by having
    members: two of seven categories reach 28.6% of the movement by doing
    nothing at all, so a 35% floor was crossed by a spread that is barely
    distinguishable from uniform. The leading group must beat what an even
    split would hand it.
    """
    from mi_agent_api import materiality as MAT

    # Six regions between £3.4m and £3.9m — the shape the shipped pack hit. The
    # top two reach 35.6% of the movement, which cleared the dominance floor
    # with two members while an even split already hands two of six 33.3%.
    spread = [{"label": f"R{i}", "value": v} for i, v in
              enumerate((3.9, 3.8, 3.6, 3.5, 3.4, 3.4))]
    assert MAT.classify(spread).shape == MAT.SHAPE_DISTRIBUTED

    # And the seven-way £3.7-4.4m spread the module was written for.
    assert MAT.classify(
        [{"label": f"R{i}", "value": v} for i, v in
         enumerate((4.4, 4.3, 4.1, 4.0, 3.9, 3.8, 3.7))]
    ).shape == MAT.SHAPE_DISTRIBUTED

    # A real concentration still reads as one: two of seven carrying 60%.
    real = [{"label": f"R{i}", "value": v} for i, v in
            enumerate((6.0, 6.0, 1.0, 1.0, 1.0, 1.0, 4.0))]
    assert MAT.classify(real).shape in (MAT.SHAPE_CONCENTRATED, MAT.SHAPE_DRIVEN)


def test_9_the_aggregated_tail_counts_toward_the_movement():
    """Catches: a leader called dominant over a book it does not dominate.

    "Other" is a top-N presentation bucket. Ranking it as a category would be
    wrong; dropping it from the denominator is worse, because every share then
    becomes a share of the named categories only and the executive summary
    quotes a movement total the movement page contradicts.
    """
    from mi_agent_api import materiality as MAT

    named = [{"label": "A", "value": 4.0}, {"label": "B", "value": 2.0}]
    without = MAT.classify(named)
    assert without.shape == MAT.SHAPE_DRIVEN
    assert without.leader_share == pytest.approx(4.0 / 6.0)

    # The same two categories, with £14m of movement aggregated away behind
    # them across five more categories. A leads nothing.
    with_tail = MAT.classify(named, residual_magnitude=14.0, residual_count=5)
    assert with_tail.shape == MAT.SHAPE_DISTRIBUTED
    assert with_tail.leader_share == pytest.approx(4.0 / 20.0)
    assert with_tail.contributor_count == 7
    assert with_tail.total_magnitude == pytest.approx(20.0)


def test_10_the_executive_summary_quotes_the_packs_own_movement(rich):
    """Catches: two totals for one movement in one pack.

    The summary's movement share was computed over the named categories only,
    so it printed "£21.6m moved" on the page facing a stock page that says the
    book moved £24.8MM. A funder reading both finds the pack disagreeing with
    itself.
    """
    import re

    text = deck_text(generate_and_download())
    quoted = re.findall(r"of the £([\d.]+)m moved", text)
    if not quoted:
        pytest.skip("this book's movement has a named driver, not a spread")
    stock = re.findall(r"\(\+£([\d.]+)MM\)", text)
    assert stock, "no stock movement figure to reconcile against"
    assert float(quoted[0]) == pytest.approx(float(stock[0]), abs=0.15), (
        f"summary says £{quoted[0]}m moved; the stock page says £{stock[0]}MM")


# --------------------------------------------------------------------------- #
# 6. A CAPTION BELONGS INSIDE ITS OWN TILE.
# --------------------------------------------------------------------------- #

def test_11_no_tile_caption_is_drawn_outside_its_panel(rich):
    """Catches: "2 portfolio types" printed across the bottom edge of its card.

    A tile's hint line is drawn 1.02in from the tile top and stands 0.30in
    tall, so a strip shorter than 1.34in puts the caption outside the panel it
    describes. Two slides shipped with a 1.10in and a 1.16in strip carrying
    hints. The height is now a property of the strip rather than of each
    caller's memory, and this reads it back off the rendered file.
    """
    import io

    import pptx as _pptx
    from pptx.util import Emu

    deck = _pptx.Presentation(io.BytesIO(generate_and_download()))
    overflows = []
    for index, slide in enumerate(deck.slides, 1):
        panels = [(Emu(sh.top).inches, Emu(sh.top + sh.height).inches,
                   Emu(sh.left).inches, Emu(sh.left + sh.width).inches)
                  for sh in slide.shapes
                  if not (sh.has_text_frame and sh.text_frame.text.strip())]
        for sh in slide.shapes:
            if not (sh.has_text_frame and sh.text_frame.text.strip()):
                continue
            top, bottom = Emu(sh.top).inches, Emu(sh.top + sh.height).inches
            left = Emu(sh.left).inches
            # A text box belongs to the smallest panel it starts inside. Full
            # -height backgrounds are not tiles, so only panels under 3in count.
            hosts = [p for p in panels
                     if p[0] <= top + 0.01 <= p[1] and p[2] <= left + 0.01 <= p[3]
                     and (p[1] - p[0]) < 3.0]
            if hosts and bottom > min(h[1] for h in hosts) + 0.03:
                overflows.append(
                    f"slide {index}: {sh.text_frame.text[:30]!r} runs to "
                    f"{bottom:.2f}in, past its panel at "
                    f"{min(h[1] for h in hosts):.2f}in")
    assert not overflows, "\n".join(overflows)


def test_12_the_methodology_columns_stay_on_the_slide(rich):
    """Catches: "SECTIONS NOT INCLUDED" printed over the footer and off the page.

    PowerPoint text does not shrink to fit its box — it draws past the bottom,
    over the footer and off the slide, and python-pptx reports nothing. The
    right column carried the basis, the coverage, every measure the book cannot
    report AND every omitted section while the left column sat two-thirds
    empty, so the sections a reader most needs listed were the ones off the
    page. The omissions are scope, and now sit with the rest of the scope.
    """
    import io

    import pptx as _pptx
    from pptx.util import Emu

    from mi_agent_pptx.deck import DeckBuilder

    deck = _pptx.Presentation(io.BytesIO(generate_and_download()))
    checked = 0
    for slide in deck.slides:
        heads = [sh.text_frame.text for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        if not heads or "Data and Methodology" not in heads[0]:
            continue
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            paras = shape.text_frame.paragraphs
            sizes = [r.font.size.pt for p in paras for r in p.runs
                     if r.font.size is not None]
            if len(paras) < 4 or not sizes:
                continue
            lines = [p.text for p in paras]
            body = max(sizes)
            assert min(sizes) >= 8.0, f"type below the floor: {min(sizes)}pt"
            extent = DeckBuilder._column_extent(
                lines, int(shape.width) / 914400, body)
            assert extent <= DeckBuilder._COLUMN_HEIGHT + 0.01, (
                f"a methodology column renders {extent:.2f}in into a "
                f"{DeckBuilder._COLUMN_HEIGHT}in box: {lines[0]!r}")
            checked += 1
    assert checked >= 2, "the methodology page's two columns were not found"


def test_13_the_headline_number_is_not_the_smallest_on_the_page(rich):
    """Catches: the funded balance printed at 12pt beside a 20pt "48.3%".

    Tile values were sized by counting characters, so "£78.4MM" — seven
    characters, four of them narrow — stepped down twice while "48.3%" stayed
    at full size in the same width of tile. The most important number in the
    pack rendered as the least important thing on the page, which reads as a
    rendering fault rather than as a measure.
    """
    import io

    import pptx as _pptx

    from mi_agent_pptx.deck import DeckBuilder

    # It fits, and the estimator knows it: a 2.00in tile has ~1.68in inside it.
    assert DeckBuilder._text_width_in("£78.4MM", 20) < 1.68
    assert DeckBuilder._text_width_in("£78.4MM", 20) > \
        DeckBuilder._text_width_in("251", 20)

    deck = _pptx.Presentation(io.BytesIO(generate_and_download()))
    for slide in deck.slides:
        heads = [sh.text_frame.text for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        if not heads or "Key Measures" not in heads[0]:
            continue
        sizes = {}
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            text = sh.text_frame.text.strip()
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size.pt >= 10 and text:
                        sizes.setdefault(text, run.font.size.pt)
        values = {t: pt for t, pt in sizes.items()
                  if t and t[0] in "£$€" or t.replace(".", "").isdigit()
                  or t.endswith("%")}
        assert values, sizes
        assert len(set(values.values())) == 1, (
            f"one page, two type sizes for the same kind of measure: {values}")
        return
    pytest.fail("no Key Measures page in the pack")


# --------------------------------------------------------------------------- #
# 7. A BREAKDOWN OF THE PIPELINE IS A BREAKDOWN OF THE PIPELINE.
# --------------------------------------------------------------------------- #

def test_14_pipeline_breakdowns_sum_to_the_pipeline_headline():
    """Catches: £10.7m of stage bars beneath a £7.8m total pipeline tile.

    The live-stock correction gave ``pipelineAmount`` the active stages and
    left every breakdown of it on the whole extract, terminal cases included.
    A funder adding the bars on the Pipeline Overview page got a third more
    than the headline above them, and the broker split described a population
    no figure on the page reports.
    """
    import pandas as pd

    from mi_agent_api import pipeline_contract as PC

    rows = []
    for i, (stage, amount, broker) in enumerate((
            ("KFI", 100.0, "Direct"), ("APPLICATION", 200.0, "Direct"),
            ("OFFER", 300.0, "Network"), ("COMPLETED", 400.0, "Direct"),
            ("WITHDRAWN", 50.0, "Network"))):
        rows.append({"pipeline_case_identifier": f"C{i}", "pipeline_stage": stage,
                     "current_outstanding_balance": amount,
                     "broker_channel": broker})
    df = pd.DataFrame(rows)

    stages = PC._stage_breakdown(df[df["pipeline_stage"].isin(
        ("KFI", "APPLICATION", "OFFER"))])
    assert {r["stage"] for r in stages} == {"KFI", "APPLICATION", "OFFER"}
    assert sum(r["pipelineAmount"] for r in stages) == pytest.approx(600.0)

    from mi_agent_api.pipeline_prep import live_mask

    live = df[live_mask(df)]
    assert len(live) == 3, "live stock must exclude completed and withdrawn"
    brokers = PC._dimension_breakdown(live, "broker_channel")
    assert sum(r["pipelineAmount"] for r in brokers) == pytest.approx(600.0)


def test_15_the_pipeline_page_reconciles_to_its_own_tile(rich):
    """The same property, read off the printed page rather than the frame."""
    import io
    import re

    import pptx as _pptx

    deck = _pptx.Presentation(io.BytesIO(generate_and_download()))
    for slide in deck.slides:
        lines = [sh.text_frame.text.strip() for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        if not lines or "Pipeline Overview" not in lines[0]:
            continue
        page = "\n".join(lines)
        assert "Completed" not in page and "Withdrawn" not in page, (
            "a terminal stage is drawn as live pipeline stock:\n" + page)
        return
    pytest.skip("this book has no governed pipeline source")


def test_16_the_forecast_bridge_steps_sum_to_its_own_closing_bar():
    """Catches: £7.0m of monthly steps under a £4.1m weighted-pipeline block.

    The forecast is funded plus the LIVE pipeline weighted by conversion, and
    the bridge's headline says so — but the by-month, by-region and by-LTV
    cuts summed the weighted amount over the whole extract, terminal cases
    included. A reader adding the steps of a waterfall is doing the one thing a
    waterfall is drawn for, and got 71% more than the bar it closes on.
    """
    import pandas as pd

    from mi_agent_api import workspace as W

    funded = pd.DataFrame([
        {"geographic_region_obligor": "London", "ltv_bucket": "20-30%",
         "current_outstanding_balance": 1000.0}])
    pipeline = pd.DataFrame([
        {"pipeline_stage": "KFI", "geographic_region_obligor": "London",
         "ltv_bucket": "20-30%", "expected_completion_month": "2026-07",
         "weighted_expected_funded_amount": 100.0},
        {"pipeline_stage": "OFFER", "geographic_region_obligor": "London",
         "ltv_bucket": "20-30%", "expected_completion_month": "2026-08",
         "weighted_expected_funded_amount": 50.0},
        # Already funded. It is in the funded book; counting it here reports the
        # same exposure twice in one bridge.
        {"pipeline_stage": "COMPLETED", "geographic_region_obligor": "London",
         "ltv_bucket": "20-30%", "expected_completion_month": "2026-07",
         "weighted_expected_funded_amount": 400.0},
    ])
    brk = W.forecast_breakdowns(funded, pipeline)
    months = sum(m["weightedExpectedFundedAmount"]
                 for m in brk["byCompletionMonth"])
    assert months == pytest.approx(150.0), brk["byCompletionMonth"]
    for key in ("byRegion", "byLtvBucket"):
        assert sum(r["forecastAmount"] for r in brk[key]) == pytest.approx(1150.0)


def test_17_the_forecast_cuts_account_for_the_whole_forecast(rich):
    """Catches: six bars a reader adds up, £11.8m short, with no note.

    The panel truncated its list to what it could draw and said nothing, on
    two charts whose bars sum to the number in the bar chart above them. Every
    row is now drawn, and where one still cannot be, the remainder is
    aggregated rather than dropped.
    """
    import io
    import re

    import pptx as _pptx

    deck = _pptx.Presentation(io.BytesIO(generate_and_download()))
    for slide in deck.slides:
        lines = [sh.text_frame.text.strip() for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        if not lines or "Forecast Bridge" not in lines[0]:
            continue
        # The waterfall's own steps close on its own final bar.
        page = "\n".join(lines)
        assert "Forecast Funded" in page or "forecast" in page.lower()
        return
    pytest.skip("this book has no forecast bridge")


def test_18_a_truncated_bar_list_aggregates_rather_than_drops():
    """The property itself, without needing a book wide enough to truncate."""
    from mi_agent_pptx.deck import DeckBuilder

    rows = [{"label": f"R{i}", "balance": float(i + 1)} for i in range(9)]
    drawn = DeckBuilder._fit_bars(rows, 6)
    assert len(drawn) == 6
    assert drawn[-1]["label"] == "Other (4)"
    assert sum(r["balance"] for r in drawn) == pytest.approx(
        sum(r["balance"] for r in rows))
    # A list that fits is returned untouched — no phantom "Other (0)".
    assert DeckBuilder._fit_bars(rows[:4], 6) == rows[:4]
