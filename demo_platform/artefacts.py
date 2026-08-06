"""demo_platform.artefacts — the governed outputs the closing scenes show.

SYNTHETIC DEMONSTRATION DATA — NOT A REAL CUSTOMER.

Produces, from the assembled platform canonical the recurring run just published:

  * the consolidated canonical tape (already produced by the assembler — recorded
    here with its row count, balance and content hash);
  * a consolidated validation report, aggregated from the per-portfolio Gate 2 /
    Gate 3 / Gate 3b outputs;
  * the investor PowerPoint pack, via ``mi_agent_pptx.cli`` — the real deck
    builder, which reads the same ``/mi/*`` computations the dashboard uses;
  * the risk-monitoring output, via the ``/mi/risk-limits`` service;
  * an audit manifest: every input file, every gate status, every output, with
    provenance and content hashes.

Each artefact records whether it was actually produced. Anything that could not
be produced is reported as unavailable with the reason — the film only shows what
exists.
"""

from __future__ import annotations

import hashlib
import json
import os
import shutil
import subprocess
import sys
import warnings
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence

import pandas as pd

from . import config as cfg

warnings.simplefilter("ignore")

DECK_NAME = "investor_pack.pptx"
POINTER_NAME = "latest_investor_pack.json"


def _sha256(path: Path) -> str:
    return "sha256:" + hashlib.sha256(Path(path).read_bytes()).hexdigest()


def _dated_platform_dir(prd: cfg.PeriodSpec) -> Path:
    return (cfg.local_blob_root() / "processed" / "platform"
            / cfg.CLIENT_ID / prd.reporting_date)


# --------------------------------------------------------------------------- #
# 1. Consolidated canonical tape
# --------------------------------------------------------------------------- #
def canonical_tape() -> Dict[str, Any]:
    """Record the assembled consolidated tape for the current period."""
    path = _dated_platform_dir(cfg.CURRENT_PERIOD) / "platform_canonical_typed.csv"
    manifest_path = path.parent / "platform_canonical_manifest.json"
    if not path.exists():
        return {"available": False, "reason": "the platform canonical has not been assembled"}
    manifest = json.loads(manifest_path.read_text(encoding="utf-8")) if manifest_path.exists() else {}
    df = pd.read_csv(path, low_memory=False)
    return {
        "available": True,
        "kind": "consolidated_canonical_tape",
        "title": "Consolidated canonical tape",
        "fileName": path.name,
        "format": "CSV",
        "rows": int(len(df)),
        "columns": int(len(df.columns)),
        "sizeBytes": path.stat().st_size,
        "totalBalance": manifest.get("output_total_balance"),
        "portfolioCount": manifest.get("portfolio_count"),
        "compositeKey": manifest.get("composite_key"),
        "contentSha256": manifest.get("content_sha256"),
        "reportingDate": cfg.CURRENT_PERIOD.reporting_date,
        "producedBy": "engine/platform_assembler.py",
    }


# --------------------------------------------------------------------------- #
# 2. Consolidated validation report
# --------------------------------------------------------------------------- #
def validation_report() -> Dict[str, Any]:
    """Aggregate the current period's per-portfolio validation into one report."""
    per_portfolio: List[Dict[str, Any]] = []
    rule_totals: Dict[str, int] = {}
    total_rows = 0
    total_exceptions = 0
    gate2_ok = True

    # PLATFORM scope: the validation report covers the assembled platform canonical.
    for portfolio in cfg.PORTFOLIOS:
        val_dir = cfg.run_validation_dir(portfolio.source_portfolio_id,
                                         cfg.CURRENT_PERIOD.run_id)
        dashboard_path = next(iter(sorted(val_dir.glob("*_dashboard.json"))), None)
        canonical_v = next(iter(sorted(val_dir.glob("*_canonical_typed_canonical_violations.csv"))), None)
        business_v = next(iter(sorted(val_dir.glob("*_business_rules_violations.csv"))), None)

        dashboard = json.loads(dashboard_path.read_text(encoding="utf-8")) if dashboard_path else {}
        canonical_errors = 0
        canonical_warnings = 0
        if canonical_v is not None and canonical_v.exists():
            cv = pd.read_csv(canonical_v)
            if "severity" in cv.columns:
                canonical_errors = int((cv["severity"] == "error").sum())
                canonical_warnings = int((cv["severity"] == "warn").sum())
        business_rows: List[Dict[str, Any]] = []
        if business_v is not None and business_v.exists():
            bv = pd.read_csv(business_v)
            if "rule_id" in bv.columns:
                for rule, group in bv.groupby("rule_id"):
                    count = int(len(group))
                    rule_totals[str(rule)] = rule_totals.get(str(rule), 0) + count
                    business_rows.append({
                        "rule": str(rule),
                        "severity": str(group["severity"].iloc[0]) if "severity" in group else "",
                        "description": str(group["description"].iloc[0]) if "description" in group else "",
                        "rows": count,
                    })
                business_rows.sort(key=lambda r: -r["rows"])
                total_exceptions += int(len(bv))
        # The Gate 3b dashboard carries the exception summary, not a row count, so
        # take the row count from the portfolio's own canonical output — the tape
        # that was actually validated.
        canonical = (cfg.run_output_dir(portfolio.source_portfolio_id,
                                        cfg.CURRENT_PERIOD.run_id)
                     / f"{cfg.source_file(portfolio, cfg.CURRENT_PERIOD).stem}"
                       f"_canonical_typed.csv")
        rows = int(len(pd.read_csv(canonical, low_memory=False))) if canonical.exists() else 0
        total_rows += rows
        if canonical_errors:
            gate2_ok = False
        per_portfolio.append({
            "portfolio": portfolio.display_id,
            "label": portfolio.label,
            "rows": rows,
            "gate3bSummary": dashboard.get("summary") or {},
            "canonicalErrors": canonical_errors,
            "canonicalWarnings": canonical_warnings,
            "businessRuleExceptions": sum(r["rows"] for r in business_rows),
            "businessRules": business_rows,
        })

    return {
        "available": bool(per_portfolio),
        "kind": "validation_report",
        "title": "Validation report",
        "format": "CSV + JSON",
        "reportingDate": cfg.CURRENT_PERIOD.reporting_date,
        "rowsValidated": total_rows,
        "canonicalValidation": "Passed" if gate2_ok else "Exceptions raised",
        "businessRuleExceptions": total_exceptions,
        "exceptionRatePct": (round(total_exceptions / total_rows * 100, 3)
                            if total_rows else None),
        "rulesTriggered": sorted(
            ({"rule": r, "rows": n} for r, n in rule_totals.items()),
            key=lambda r: -r["rows"]),
        "perPortfolio": per_portfolio,
        "producedBy": "engine/gate_3_validation/* (Gates 2, 3, 3b)",
    }


# --------------------------------------------------------------------------- #
# 3. Investor deck
# --------------------------------------------------------------------------- #
def investor_deck(*, verbose: bool = True) -> Dict[str, Any]:
    """Generate the investor PPTX pack with the production deck builder.

    Published into the ``MI_AGENT_DECK_ROOT`` layout the MI API discovers
    (``{client}/latest|{period}/investor_pack.pptx`` plus the latest pointer), so
    the Copilot ``getLatestInvestorDeck`` action resolves it exactly as it would
    in production.
    """
    run_dir = _dated_platform_dir(cfg.CURRENT_PERIOD)
    prior_dir = _dated_platform_dir(cfg.PRIOR_PERIOD)
    if not (run_dir / "platform_canonical_typed.csv").exists():
        return {"available": False, "reason": "no assembled platform canonical for the current period"}

    work = cfg.artefact_dir() / "investor_deck"
    work.mkdir(parents=True, exist_ok=True)
    output = work / DECK_NAME

    env = dict(os.environ)
    env.update(cfg.mi_env())
    env["MI_AGENT_ONBOARDING_OUTPUT_ROOT"] = cfg.platform_prefix()

    cmd = [
        sys.executable, "-m", "mi_agent_pptx.cli",
        "--run-dir", str(run_dir),
        "--prior-run-dir", str(prior_dir),
        "--deck-config", str(cfg.REPO_ROOT / "configs" / "pptx" / "investor_pack.yaml"),
        "--client-name", cfg.CLIENT_DISPLAY_NAME,
        "--client-id", cfg.CLIENT_ID,
        "--run-id", cfg.CURRENT_PERIOD.reporting_date,
        "--as-of-date", cfg.CURRENT_PERIOD.reporting_date,
        "--output", str(output),
        "--consolidated",
        "--work-dir", str(work / "scratch"),
        "--repo-root", str(cfg.REPO_ROOT),
    ]
    proc = subprocess.run(cmd, cwd=str(cfg.REPO_ROOT), capture_output=True,
                          text=True, env=env)
    if proc.returncode != 0 or not output.exists():
        tail = "\n".join((proc.stdout or "").splitlines()[-12:])
        err = "\n".join((proc.stderr or "").splitlines()[-12:])
        return {
            "available": False,
            "reason": f"the deck builder exited {proc.returncode}",
            "stdoutTail": tail, "stderrTail": err,
        }

    # Publish into the deck root the MI API + Copilot action discover.
    deck_base = cfg.deck_root() / cfg.CLIENT_ID
    for folder in ("latest", cfg.CURRENT_PERIOD.period):
        target = deck_base / folder
        target.mkdir(parents=True, exist_ok=True)
        shutil.copyfile(output, target / DECK_NAME)
    pointer = {
        "deck_name": DECK_NAME,
        "client_id": cfg.CLIENT_ID,
        "reporting_period": cfg.CURRENT_PERIOD.period,
        "generated_at": f"{cfg.CURRENT_PERIOD.reporting_date}T23:59:00Z",
        "synthetic_notice": cfg.SYNTHETIC_BANNER,
    }
    (deck_base / "latest" / POINTER_NAME).write_text(
        json.dumps(pointer, indent=2), encoding="utf-8")

    slides = _count_slides(output)
    if verbose:
        print(f"  [artefacts] investor deck: {slides} slides, "
              f"{output.stat().st_size / 1024:,.0f} KB")
    return {
        "available": True,
        "kind": "investor_deck",
        "title": "Investor pack",
        "fileName": DECK_NAME,
        "format": "PowerPoint (PPTX)",
        "slides": slides,
        "sizeBytes": output.stat().st_size,
        "contentSha256": _sha256(output),
        "reportingPeriod": cfg.CURRENT_PERIOD.period,
        "path": str(output),
        "publishedTo": str(deck_base / "latest" / DECK_NAME),
        "producedBy": "mi_agent_pptx/cli.py",
        "slideTitles": _slide_titles(output),
    }


def _count_slides(path: Path) -> Optional[int]:
    try:
        from pptx import Presentation
        return len(Presentation(str(path)).slides)
    except Exception:  # noqa: BLE001 - a missing count never blocks the artefact
        return None


def _slide_titles(path: Path) -> List[str]:
    try:
        from pptx import Presentation
        titles: List[str] = []
        for slide in Presentation(str(path)).slides:
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    title = shape.text_frame.text.strip().splitlines()[0]
                    break
            titles.append(title[:80])
        return titles
    except Exception:  # noqa: BLE001
        return []


# --------------------------------------------------------------------------- #
# 4. Risk monitoring
# --------------------------------------------------------------------------- #
def risk_monitor() -> Dict[str, Any]:
    """Capture the risk-limit monitoring output from its production service."""
    os.environ.update(cfg.mi_env())
    try:
        from fastapi.testclient import TestClient
        from mi_agent_api.app import app
        resp = TestClient(app).get("/mi/risk-limits", params={
            "portfolioId": f"{cfg.CLIENT_ID}/{cfg.CURRENT_PERIOD.reporting_date}"})
    except Exception as exc:  # noqa: BLE001
        return {"available": False, "reason": f"risk service unavailable: {exc}"}
    if resp.status_code != 200:
        return {"available": False, "reason": f"/mi/risk-limits returned HTTP {resp.status_code}"}
    payload = resp.json()
    # The monitor returns one entry per limit TEST (`tests`), each with a RAG
    # status and the observed concentration; the earlier `limits` key does not
    # exist on this envelope.
    tests = payload.get("tests") or []
    summary = payload.get("summary") or {}
    statuses: Dict[str, int] = {}
    for row in tests if isinstance(tests, list) else []:
        status = str((row or {}).get("status") or "").lower()
        if status:
            statuses[status] = statuses.get(status, 0) + 1
    available = bool(payload.get("available")) and bool(tests)
    reason = None
    if not available:
        reason = (payload.get("limitsReason")
                  or "no concentration limits are configured for this client")
    elif not payload.get("fundedDataAvailable"):
        # Limits extracted but no funded data to test them against: report the
        # limits, and say plainly that the observed concentrations are absent.
        reason = "limits extracted; observed concentrations unavailable"
    return {
        "available": available,
        "kind": "risk_monitoring",
        "title": "Concentration risk monitor",
        "format": "JSON",
        "limitCount": len(tests) if isinstance(tests, list) else None,
        "testsPassed": summary.get("testsPassed"),
        "breaches": summary.get("breaches"),
        "unavailableTests": summary.get("unavailable"),
        "fundedDataAvailable": bool(payload.get("fundedDataAvailable")),
        "limitsSource": payload.get("limitsSource"),
        "largestConcentration": summary.get("largestConcentration"),
        "closestHeadroom": summary.get("closestHeadroom"),
        "statusCounts": statuses,
        "reportingDate": cfg.CURRENT_PERIOD.reporting_date,
        "producedBy": "mi_agent_api/risk_limits.py (mi_agent/risk_monitor)",
        "payload": payload,
        "reason": reason,
    }


# --------------------------------------------------------------------------- #
# 5. Regulatory output (attempted honestly; reported unavailable if it cannot run)
# --------------------------------------------------------------------------- #
def regulatory_output(*, verbose: bool = True) -> Dict[str, Any]:
    """Produce the ESMA Annex 2 exposure-level submission for the current period.

    Runs the real regulatory path end to end for the consolidated platform
    canonical: Gate 4b regime projection, delivery normalisation, then Gate 5 XML
    generation with XSD validation against the in-repo
    ``DRAFT1auth.099.001.04_1.3.0.xsd``.

    The outcome is a measured fact. If any stage cannot complete, the artefact is
    reported as unavailable with the reason and the film does not show it.
    """
    canonical = (_dated_platform_dir(cfg.CURRENT_PERIOD)
                 / "platform_canonical_typed.csv")
    if not canonical.exists():
        return {"available": False, "reason": "no assembled platform canonical to project"}

    out_dir = cfg.artefact_dir() / "regulatory"
    out_dir.mkdir(parents=True, exist_ok=True)
    # The projector derives the output stem from the input filename; give it a
    # name that reads as a submission rather than as an internal canonical.
    staged = out_dir / f"{cfg.CLIENT_ID}_annex2_{cfg.CURRENT_PERIOD.period}.csv"
    shutil.copyfile(canonical, staged)

    stages: List[Dict[str, Any]] = []

    def _stage(name: str, cmd: List[Any]) -> subprocess.CompletedProcess:
        proc = subprocess.run([str(c) for c in cmd], cwd=str(cfg.REPO_ROOT),
                              capture_output=True, text=True)
        stages.append({"stage": name, "returncode": proc.returncode})
        return proc

    config_root = cfg.REPO_ROOT / "config"
    proj = _stage("gate_4b_projection", [
        sys.executable, cfg.REPO_ROOT / "engine" / "gate_4_projection" / "regime_projector.py",
        staged, "--regime", "ESMA_Annex2",
        "--registry", config_root / "system" / "fields_registry.yaml",
        "--enum-mapping", _demo_enum_mapping(out_dir),
        # Layer 2 — the equity-release asset pack, merged UNDER the client
        # config. Product no-data treatment is owned here, once per asset
        # class, rather than restated in every client configuration.
        "--product-defaults", config_root / "asset" / "product_defaults_ERM.yaml",
        "--config", cfg.demo_client_config(),
        "--template-order", config_root / "system" / "esma_code_order.yaml",
        "--portfolio-type", "equity_release",
        "--output-dir", out_dir,
    ])
    projected = next(iter(sorted(out_dir.glob("*_ESMA_Annex2_projected.csv"))), None)
    if proj.returncode != 0 or projected is None:
        return _regulatory_unavailable(proj, stages, "projection", verbose)

    norm = _stage("gate_4b_delivery_normalisation", [
        sys.executable, cfg.REPO_ROOT / "engine" / "gate_4b_delivery" / "annex2_delivery_normalizer.py",
        "--input", projected,
        "--rules", _demo_delivery_rules(out_dir),
        "--output-dir", out_dir,
    ])
    delivery = next(iter(sorted(out_dir.glob("*_delivery_ready.csv"))), None)
    if norm.returncode != 0 or delivery is None:
        return _regulatory_unavailable(norm, stages, "delivery normalisation", verbose)

    xml_out = out_dir / "annex2_submission.xml"
    xml = _stage("gate_5_xml_and_xsd", [
        sys.executable, cfg.REPO_ROOT / "engine" / "gate_5_delivery" / "xml_builder_annex2.py",
        "--input", delivery,
        "--output", xml_out,
        "--mapping-workbook", cfg.REPO_ROOT
        / "DRAFT1auth.099.001.04_non-ABCP Underlying Exposure Report_Version_1.3.1.xlsx",
        "--sheet", "DRAFT1auth.099.001.04",
        "--code-order-yaml", config_root / "system" / "esma_code_order.yaml",
        "--xsd", config_root / "system" / "DRAFT1auth.099.001.04_1.3.0.xsd",
    ])
    if xml.returncode != 0 or not xml_out.exists():
        return _regulatory_unavailable(xml, stages, "XML generation / XSD validation", verbose)

    projected_df = pd.read_csv(projected, low_memory=False)
    delivery_df = pd.read_csv(delivery, low_memory=False)
    # A bare field count treats a governed "not applicable" the same as a field
    # carrying client data. Gate 4b reports the split; carry it through rather
    # than letting one number stand for both.
    provenance: Dict[str, Any] = {}
    report_json = next(iter(sorted(out_dir.glob("*_delivery_report.json"))), None)
    if report_json is not None:
        try:
            provenance = (json.loads(report_json.read_text(encoding="utf-8"))
                          .get("field_provenance") or {})
        except (OSError, ValueError):
            provenance = {}
    # The builder prints "XSD Validation: PASSED" and exits non-zero on failure;
    # match its own wording rather than inferring from the return code alone.
    xsd_validated = "XSD Validation: PASSED" in (xml.stdout or "")
    exposure_records = _count_xml_records(xml_out)

    from_source = provenance.get("populated_from_projected_source")
    by_rule = provenance.get("populated_by_declared_delivery_rule")

    if verbose:
        fields = (f"{len(delivery_df.columns)} fields" if from_source is None else
                  f"{len(delivery_df.columns)} fields ({from_source} from source, "
                  f"{by_rule} by declared delivery rule)")
        print(f"  [artefacts] ESMA Annex 2: {exposure_records:,} exposure records, "
              f"{fields}, {xml_out.stat().st_size / 1e6:,.1f} MB XML"
              + (", XSD valid" if xsd_validated else ""))

    return {
        "available": True,
        "kind": "regulatory_output",
        "title": "ESMA Annex 2 exposure report",
        "fileName": xml_out.name,
        "format": "XML (ESMA auth.099.001.04)",
        "rows": int(len(projected_df)),
        "exposureRecords": exposure_records,
        "fields": int(len(delivery_df.columns)),
        # A field carrying a governed "not applicable" is not a field carrying
        # client data. `fields` is the total; these two say what it is made of.
        "fieldsFromProjectedSource": from_source,
        "fieldsFromDeliveryRule": by_rule,
        "sizeBytes": xml_out.stat().st_size,
        "xsdValidated": xsd_validated,
        "contentSha256": _sha256(xml_out),
        "reportingDate": cfg.CURRENT_PERIOD.reporting_date,
        "stages": stages,
        "producedBy": "engine/gate_4_projection/regime_projector.py → "
                      "engine/gate_4b_delivery/annex2_delivery_normalizer.py → "
                      "engine/gate_5_delivery/xml_builder_annex2.py",
    }


#: The collateral-type codes the ESMA auth.099.001.04 schema itself defines for
#: ``CollTp``, read off the XSD enumeration. Residential property is ``RBLD``.
_ANNEX2_COLLATERAL_TYPE_CODES = (
    "RBLD", "RALV", "CBLD", "CMTR", "IBLD", "INDE", "MIXD", "OFEQ", "OTRE",
    "NACM", "NALV", "INDV", "ITEQ", "MCHT", "CARX", "AERO", "ENEQ", "GUAR",
    "MDEQ", "OTFA", "OTGI", "OTHV", "OTHE", "SECU", "OTHR",
)


def _demo_enum_mapping(out_dir: Path) -> Path:
    """Write a demo-scoped copy of the regime enum mapping and return its path.

    ``config/system/enum_mapping.yaml`` is production configuration and is left
    untouched. Its ``ESMA_Annex2.collateral_type`` table maps dwelling forms
    (HOUSE, FLAT, LAND …) onto ``R1``/``R2``/``C1``/``C2``, which are the codes the
    Annex 2 *property type* field takes — not the codes the auth.099.001.04 schema
    accepts for ``CollTp``. The repository's own enum library already states the
    canonical answer for this book (``config/system/enum_synonyms.yaml`` maps
    "residential property" to ``RBLD``), but the projector's enum resolver
    discards a synonym whose target is absent from the regime table, so the value
    reaches Gate 5 unmapped and the XSD rejects it.

    The overlay adds identity entries for the codes the XSD enumerates, so the
    canonical value survives projection. It invents no code and remaps nothing
    that was already mapped. Passed to the projector through its own
    ``--enum-mapping`` argument, which is a first-class per-run input.
    """
    import yaml

    source = cfg.REPO_ROOT / "config" / "system" / "enum_mapping.yaml"
    mapping = yaml.safe_load(source.read_text(encoding="utf-8"))

    table = (mapping.setdefault("ESMA_Annex2", {})
                    .setdefault("collateral_type", {}))
    for code in _ANNEX2_COLLATERAL_TYPE_CODES:
        table.setdefault(code, code)

    path = out_dir / "enum_mapping_demo.yaml"
    path.write_text(
        "# GENERATED — demo-scoped copy of config/system/enum_mapping.yaml.\n"
        "# SYNTHETIC DEMONSTRATION DATA — NOT A REAL CUSTOMER.\n"
        "# The only difference from the production mapping is identity entries for\n"
        "# the ESMA auth.099.001.04 CollTp codes under ESMA_Annex2.collateral_type.\n"
        + yaml.safe_dump(mapping, sort_keys=False, allow_unicode=True),
        encoding="utf-8",
    )
    return path


def _demo_delivery_rules(out_dir: Path) -> Path:
    """Write a demo-scoped copy of the Annex 2 delivery rules and return its path.

    ``config/regime/annex2_delivery_rules.yaml`` is production configuration and is
    left untouched. Its ``RREL83`` (originator LEI) rule is a strict controlled
    vocabulary: an LEI that is not a key in ``transform.enum_map`` is rejected, so
    the synthetic originator LEI would fail delivery normalisation for every row.

    The overlay adds exactly one entry — an identity mapping for this
    demonstration's own LEI, read from the demo client config. It changes no rule,
    relaxes no validator, and adds no field: the LEI still has to satisfy the
    ``lei`` validator and the ISO 17442 pattern downstream. This is client
    onboarding configuration, which is per-client by design.
    """
    import yaml

    source = cfg.REPO_ROOT / "config" / "regime" / "annex2_delivery_rules.yaml"
    rules = yaml.safe_load(source.read_text(encoding="utf-8"))

    client_cfg = yaml.safe_load(cfg.demo_client_config().read_text(encoding="utf-8"))
    lei = _find_key(client_cfg, "originator_legal_entity_identifier")
    if not lei:
        raise RuntimeError("demo client config carries no originator LEI")

    rule = rules.setdefault("field_rules", {}).setdefault("RREL83", {})
    # REPLACE rather than extend: the production vocabulary lists other clients'
    # LEIs, and a demonstration artefact must not carry them. A per-client
    # onboarding configuration would only ever hold that client's own LEI.
    rule.setdefault("transform", {})["enum_map"] = {str(lei): str(lei)}

    path = out_dir / "annex2_delivery_rules_demo.yaml"
    path.write_text(
        "# GENERATED — demo-scoped copy of config/regime/annex2_delivery_rules.yaml.\n"
        "# SYNTHETIC DEMONSTRATION DATA — NOT A REAL CUSTOMER.\n"
        "# The only difference from the production rules is one RREL83 enum_map\n"
        "# entry: an identity mapping for this demonstration's originator LEI.\n"
        + yaml.safe_dump(rules, sort_keys=False, allow_unicode=True),
        encoding="utf-8",
    )
    return path


def _find_key(node: Any, key: str) -> Optional[Any]:
    """Depth-first search for ``key`` in a nested mapping/sequence structure."""
    if isinstance(node, dict):
        if key in node:
            return node[key]
        for value in node.values():
            found = _find_key(value, key)
            if found is not None:
                return found
    elif isinstance(node, list):
        for value in node:
            found = _find_key(value, key)
            if found is not None:
                return found
    return None


def _regulatory_unavailable(proc: subprocess.CompletedProcess,
                            stages: List[Dict[str, Any]], stage: str,
                            verbose: bool) -> Dict[str, Any]:
    """Report the regulatory artefact as unavailable, naming the failing stage."""
    reason = next((ln.strip() for ln in reversed((proc.stderr or "").splitlines())
                   if ln.strip()), f"{stage} exited {proc.returncode}")
    if verbose:
        print(f"  [artefacts] ESMA Annex 2 unavailable at {stage}: {reason[:150]}")
    return {
        "available": False,
        "kind": "regulatory_output",
        "title": "ESMA Annex 2 exposure report",
        "failedStage": stage,
        "reason": reason[:400],
        "stages": stages,
        "note": "Excluded from the film: the regulatory submission did not "
                "complete, so the demonstration does not show it.",
    }


def _count_xml_records(path: Path) -> int:
    """Count exposure records in the generated XML without loading it all."""
    try:
        import re as _re
        text = path.read_text(encoding="utf-8", errors="replace")
        # The exposure record element in auth.099.001.04, opening tags only. A
        # looser pattern matches the sixteen other elements whose names share the
        # "Undrlyg" stem and reports sixteen times the true count.
        return len(_re.findall(r"<(?:\w+:)?UndrlygXpsrRcrd\b[^>/]*>", text))
    except Exception:  # noqa: BLE001 - a count is not worth failing the artefact
        return 0


# --------------------------------------------------------------------------- #
# 6. Audit manifest
# --------------------------------------------------------------------------- #
def orchestration_from_disk() -> Dict[str, Any]:
    """Reconstruct the orchestration result from the manifests already written.

    A selective run (``--artefacts`` without ``--orchestrate``) has no in-process
    orchestration result, and an audit manifest that silently recorded zero
    assembled outputs would be worse than one that failed. The assembler and run
    manifests are on disk, so read them.
    """
    result: Dict[str, Any] = {"runs": [], "assemblies": [], "reconciliation": {}}
    for prd in cfg.PERIODS:
        mpath = _dated_platform_dir(prd) / "platform_canonical_manifest.json"
        if not mpath.exists():
            continue
        m = json.loads(mpath.read_text(encoding="utf-8"))
        result["assemblies"].append({
            "period": prd.period,
            "reporting_date": prd.reporting_date,
            "role": prd.role,
            "total_rows": m.get("total_rows"),
            "total_balance": m.get("output_total_balance"),
            "content_sha256": m.get("content_sha256"),
            "dated_uri": f"{cfg.platform_prefix()}/{prd.reporting_date}",
            "per_portfolio_balance": {
                p["source_portfolio_id"]: p["total_balance"]
                for p in m.get("portfolios", [])
            },
        })
    # The audit trail is ALL scope: every portfolio governed through the gates, the sold
    # deal included. That is what makes it an audit trail rather than a platform report.
    for portfolio in cfg.ALL_PORTFOLIOS:
        for prd in cfg.PERIODS:
            run_manifest = (cfg.run_output_dir(portfolio.source_portfolio_id, prd.run_id)
                            / "run_manifest.json")
            canonical = (cfg.run_output_dir(portfolio.source_portfolio_id, prd.run_id)
                         / f"{cfg.source_file(portfolio, prd).stem}_canonical_typed.csv")
            if not canonical.exists():
                continue
            df = pd.read_csv(canonical, low_memory=False)
            result["runs"].append({
                "portfolio": portfolio.display_id,
                "period": prd.period,
                "rows": int(len(df)),
                "total_outstanding_balance": round(float(pd.to_numeric(
                    df.get("current_outstanding_balance", pd.Series(dtype=float)),
                    errors="coerce").fillna(0.0).sum()), 2),
                "gate_summary": (
                    json.loads(run_manifest.read_text(encoding="utf-8")).get("gates", [])
                    if run_manifest.exists() else []),
            })
    by_period = {a["period"]: a for a in result["assemblies"]}
    cur = by_period.get(cfg.CURRENT_PERIOD.period)
    pri = by_period.get(cfg.PRIOR_PERIOD.period)
    if cur and pri:
        result["reconciliation"] = {
            "current_period": cur["period"],
            "prior_period": pri["period"],
            "current_total_balance": cur["total_balance"],
            "prior_total_balance": pri["total_balance"],
            "consolidated_movement": round(
                (cur["total_balance"] or 0.0) - (pri["total_balance"] or 0.0), 2),
            "per_portfolio_movement": {
                pid: round(cur["per_portfolio_balance"].get(pid, 0.0)
                           - pri["per_portfolio_balance"].get(pid, 0.0), 2)
                for pid in sorted(cur["per_portfolio_balance"])
            },
        }
    return result


def audit_manifest(orchestration_result: Dict[str, Any]) -> Dict[str, Any]:
    """A single provenance record covering the whole reporting cycle."""
    # A caller that has no in-process orchestration result still gets a complete
    # manifest, rather than one that quietly reports nothing was assembled.
    if not orchestration_result.get("assemblies"):
        orchestration_result = orchestration_from_disk()
    inputs: List[Dict[str, Any]] = []
    for portfolio in cfg.ALL_PORTFOLIOS:
        for prd in cfg.PERIODS:
            path = cfg.source_file(portfolio, prd)
            if path.exists():
                inputs.append({
                    "portfolio": portfolio.display_id,
                    "period": prd.period,
                    "fileName": path.name,
                    "sizeBytes": path.stat().st_size,
                    "contentSha256": _sha256(path),
                })
    gates: List[Dict[str, Any]] = []
    for run in orchestration_result.get("runs", []):
        gates.append({
            "portfolio": run["portfolio"],
            "period": run["period"],
            "rows": run["rows"],
            "balance": run["total_outstanding_balance"],
            "gates": run.get("gate_summary", []),
            "elapsedSeconds": run.get("elapsed_seconds"),
        })
    outputs = [
        {
            "period": a["period"],
            "reportingDate": a["reporting_date"],
            "rows": a["total_rows"],
            "balance": a["total_balance"],
            "contentSha256": a.get("content_sha256"),
            "uri": a["dated_uri"],
        }
        for a in orchestration_result.get("assemblies", [])
    ]
    return {
        "available": True,
        "kind": "audit_manifest",
        "title": "Audit manifest",
        "format": "JSON",
        "synthetic_notice": cfg.SYNTHETIC_BANNER,
        "clientId": cfg.CLIENT_ID,
        "clientDisplayName": cfg.CLIENT_DISPLAY_NAME,
        "reportingDate": cfg.CURRENT_PERIOD.reporting_date,
        "sourceFiles": inputs,
        "gateExecution": gates,
        "assembledOutputs": outputs,
        "reconciliation": orchestration_result.get("reconciliation", {}),
        "provenanceFields": [
            "source_portfolio_id", "source_portfolio_type", "source_portfolio_label",
            "acquisition_date", "seller_name", "portfolio_cohort",
        ],
        "producedBy": "demo_platform/artefacts.py from the run and assembler manifests",
    }


# --------------------------------------------------------------------------- #
# Build all
# --------------------------------------------------------------------------- #
def build(orchestration_result: Dict[str, Any], *, verbose: bool = True) -> Dict[str, Any]:
    """Produce every governed output and return the artefact catalogue."""
    cfg.artefact_dir().mkdir(parents=True, exist_ok=True)
    catalogue = {
        "synthetic_notice": cfg.SYNTHETIC_BANNER,
        "reportingDate": cfg.CURRENT_PERIOD.reporting_date,
        "canonicalTape": canonical_tape(),
        "validationReport": validation_report(),
        "investorDeck": investor_deck(verbose=verbose),
        "riskMonitor": risk_monitor(),
        "regulatoryOutput": regulatory_output(verbose=verbose),
        "auditManifest": audit_manifest(orchestration_result),
    }
    out = cfg.artefact_dir() / "artefact_catalogue.json"
    out.write_text(json.dumps(catalogue, indent=2, default=str), encoding="utf-8")
    if verbose:
        for key, art in catalogue.items():
            if isinstance(art, dict) and "available" in art:
                mark = "OK  " if art["available"] else "n/a "
                print(f"  [artefacts] {mark}{key}"
                      + ("" if art["available"] else f" — {str(art.get('reason'))[:100]}"))
    return catalogue


def main(argv: Optional[Sequence[str]] = None) -> int:  # pragma: no cover - CLI
    import argparse
    ap = argparse.ArgumentParser(description="Build the demo governed outputs.")
    ap.add_argument("--quiet", action="store_true")
    args = ap.parse_args(argv)
    # Re-derive the orchestration result from the on-disk manifests.
    stub = orchestration_from_disk()
    build(stub, verbose=not args.quiet)
    return 0


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(main())
