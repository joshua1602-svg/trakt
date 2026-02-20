#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
generate_pptx_erm.py

PPTX generator for Equity Release Mortgage analytics.
Mirrors the structure and metrics from streamlit_app_erm.py.

Usage:
    python generate_pptx_erm.py --input portfolio_typed.csv --output report.pptx
    python generate_pptx_erm.py  # Uses default file
"""

import os
import sys
import argparse
from pathlib import Path
from datetime import datetime
from typing import Tuple, Optional
import math
import traceback

# Fix Windows console encoding for emojis
if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding='utf-8')
    except:
        # Fallback for older Python versions
        import codecs
        sys.stdout = codecs.getwriter('utf-8')(sys.stdout.buffer, 'strict')

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.ticker import FuncFormatter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Optional treemap support
try:
    import squarify  # type: ignore
    HAS_SQUARIFY = True
except ImportError:
    HAS_SQUARIFY = False

# Local modules
from config import (
    PRIMARY_COLOR,
    SECONDARY_COLOR,
    ACCENT_COLOR,
    TEXT_DARK,
    TEXT_LIGHT,
    BACKGROUND_LIGHT,
    CHART_COLORS,
    MAX_FILE_SIZE_MB,
    MAX_PPTX_ROWS,
)

# ============================
# CLIENT CONFIG (YAML-driven)
# ============================

import yaml as _yaml

def _load_client_config() -> dict:
    """Load client YAML config, searching common locations."""
    candidates = [
        Path(__file__).resolve().parent.parent / "config" / "client" / "config_client_ERM_UK.yaml",
        Path.cwd() / "config" / "client" / "config_client_ERM_UK.yaml",
        Path("config_client_ERM_UK.yaml"),
    ]
    for p in candidates:
        if p.exists():
            try:
                with open(p, "r", encoding="utf-8") as f:
                    return _yaml.safe_load(f) or {}
            except Exception:
                return {}
    return {}

_CLIENT_CFG = _load_client_config()

# Resolve display names from YAML, fall back to safe defaults
ENTITY_NAME: str = (
    _CLIENT_CFG.get("defaults", {}).get("originator_name")
    or _CLIENT_CFG.get("client", {}).get("display_name")
    or "Portfolio Analytics"
)
CLIENT_DISPLAY_NAME: str = (
    _CLIENT_CFG.get("client", {}).get("display_name")
    or ENTITY_NAME
)
REPORT_FOOTER: str = f"Confidential \u2022 {ENTITY_NAME}"
COVER_TITLE: str = f"{CLIENT_DISPLAY_NAME} \u2014 Portfolio Overview"
STATIC_REPORTING_DATE: str = (
    _CLIENT_CFG.get("portfolio", {}).get("static_reporting_date") or ""
)

# Risk monitoring (optional)
try:
    from risk_monitor import RiskMonitor, LimitCheck
    from risk_limits_config import ALL_LIMITS, LIMIT_CATEGORIES, CONCENTRATION_LIMITS
    RISK_MONITORING_AVAILABLE = True
except ImportError:
    RISK_MONITORING_AVAILABLE = False
    print("Warning: Risk monitoring modules not found. Risk slide will be skipped.")


# Scenario analysis (optional)
try:
    from scenario_charts_pptx import (
        save_scenario_balance_projection,
        save_scenario_nneg_losses,
        save_dual_scenario_charts,
    )
    from scenario_engine import ScenarioAssumptions
    SCENARIO_CHARTS_AVAILABLE = True
except ImportError:
    SCENARIO_CHARTS_AVAILABLE = False
    print("Warning: Scenario charts module not found. Scenario slides will be skipped.")



# ============================
# MATPLOTLIB THEME
# ============================

FONT_TITLE = "Lucida Sans"
FONT_BODY = "Lucida Sans"
FONT_NUMBERS = "Lucida Sans"

plt.style.use("seaborn-v0_8-whitegrid")
plt.rcParams.update(
    {
        "font.family": FONT_BODY,
        "font.size": 12,  # ✓ IMPROVED from 11
        "axes.labelsize": 14,  # ✓ IMPROVED from 12
        "axes.titlesize": 16,  # ✓ IMPROVED from 15
        "axes.titleweight": "bold",
        "axes.labelcolor": TEXT_DARK,
        "axes.edgecolor": "#E0E0E0",
        "axes.linewidth": 1.0,
        "axes.spines.top": False,
        "axes.spines.right": False,
        "xtick.labelsize": 12,  # ✓ IMPROVED from 10
        "ytick.labelsize": 12,  # ✓ IMPROVED from 10
        "xtick.color": "#666666",
        "ytick.color": "#666666",
        "legend.fontsize": 12,  # ✓ IMPROVED from 10
        "legend.frameon": False,
        "figure.facecolor": "white",
        "axes.facecolor": "#FAFAFA",
        "grid.color": "#E0E0E0",
        "grid.alpha": 0.30,  # ✓ CONSISTENT with charts
        "grid.linewidth": 0.8,
    }
)


# ============================
# HELPERS
# ============================

def hex_to_rgb(hex_str: str) -> Tuple[int, int, int]:
    hex_str = hex_str.lstrip("#")
    return tuple(int(hex_str[i : i + 2], 16) for i in (0, 2, 4))


def millions_formatter(x, pos):
    if x >= 1_000_000:
        return f"£{x/1_000_000:.1f}M"
    if x >= 1_000:
        return f"£{x/1_000:.0f}K"
    return f"£{x:.0f}"


def save_dual_bar_chart(
    df: pd.DataFrame,
    bucket_col: str,
    out_path: str,
    chart_title: str = "",
    orientation: str = "vertical"
) -> bool:
    """Generate dual bar charts: balance (left) and count (right) with divider."""
    if bucket_col not in df.columns or "total_balance" not in df.columns:
        return False

    # Aggregate data
    grouped = df.groupby(bucket_col, as_index=False).agg({"total_balance": "sum"})
    grouped["loan_count"] = df.groupby(bucket_col).size().values

    # Remove null / zero buckets
    grouped = grouped.dropna(subset=[bucket_col])
    grouped = grouped[grouped["total_balance"].notna()]
    grouped = grouped[grouped["total_balance"] != 0]

    if grouped.empty:
        return False

    # Respect categorical bucket ordering where applicable
    if pd.api.types.is_categorical_dtype(grouped[bucket_col]):
        grouped = grouped.sort_values(bucket_col)
    else:
        # Otherwise sort by balance descending for readability
        grouped = grouped.sort_values("total_balance", ascending=False)

    # Create figure with two subplots
    fig, (ax1, ax2) = plt.subplots(
        1,
        2,
        figsize=(24, 12),
        gridspec_kw={"wspace": 0.3}
    )

    idx = list(range(len(grouped)))

    is_horizontal = (orientation or "vertical").lower().startswith("h")

    # ===== LEFT CHART: BALANCE =====
    if is_horizontal:
        bars1 = ax1.barh(
            idx,
            grouped["total_balance"],
            color=PRIMARY_COLOR,
            alpha=1.0,
            edgecolor="none",
        )

        # Category labels on Y axis
        ax1.set_yticks(idx)
        ax1.set_yticklabels(grouped[bucket_col])

        # Highest at top (broker charts read better this way)
        ax1.invert_yaxis()

        # Value labels to the right of bars
        max_bal = float(grouped["total_balance"].max()) if len(grouped) else 0.0
        pad = max_bal * 0.01 if max_bal > 0 else 0.0

        for i, balance in enumerate(grouped["total_balance"]):
            label = (
                f"£{balance/1_000_000:.1f}M" if balance >= 1_000_000
                else f"£{balance/1_000:.0f}K"
            )
            ax1.text(
                float(balance) + pad,
                i,
                label,
                ha="left",
                va="center",
                fontsize=14,
                fontweight="bold",
                color=TEXT_DARK,
            )

        ax1.set_title(
            "Outstanding Balance (£)",
            pad=15,
            fontweight="bold",
            color=TEXT_DARK,
            fontsize=16,
        )
        ax1.set_xlabel("Balance (£)", fontweight="bold", color=TEXT_DARK, fontsize=14)
        ax1.xaxis.set_major_formatter(FuncFormatter(millions_formatter))
        ax1.grid(True, linestyle="--", alpha=0.25, linewidth=0.8, axis="x")
        ax1.set_axisbelow(True)

    else:
        bars1 = ax1.bar(
            idx,
            grouped["total_balance"],
            color=PRIMARY_COLOR,
            alpha=1.0,
            edgecolor="none",
        )

        for bar, balance in zip(bars1, grouped["total_balance"]):
            height = bar.get_height()
            label = (
                f"£{balance/1_000_000:.1f}M" if balance >= 1_000_000
                else f"£{balance/1_000:.0f}K"
            )
            ax1.text(
                bar.get_x() + bar.get_width() / 2,
                height,
                label,
                ha="center",
                va="bottom",
                fontsize=14,
                fontweight="bold",
                color=TEXT_DARK,
            )

        ax1.set_xticks(idx)
        ax1.set_xticklabels(grouped[bucket_col], rotation=45, ha="right")
        ax1.set_title(
            "Outstanding Balance (£)",
            pad=15,
            fontweight="bold",
            color=TEXT_DARK,
            fontsize=16,
        )
        ax1.set_ylabel("Balance (£)", fontweight="bold", color=TEXT_DARK, fontsize=14)
        ax1.yaxis.set_major_formatter(FuncFormatter(millions_formatter))
        ax1.grid(True, linestyle="--", alpha=0.25, linewidth=0.8, axis="y")
        ax1.set_axisbelow(True)

    ax1.tick_params(axis="both", labelsize=12)

    # ===== RIGHT CHART: COUNT =====
    if is_horizontal:
        bars2 = ax2.barh(
            idx,
            grouped["loan_count"],
            color=SECONDARY_COLOR,
            alpha=1.0,
            edgecolor="none",
        )

        ax2.set_yticks(idx)
        ax2.set_yticklabels(grouped[bucket_col])
        ax2.invert_yaxis()

        max_cnt = float(grouped["loan_count"].max()) if len(grouped) else 0.0
        pad = max_cnt * 0.01 if max_cnt > 0 else 0.0

        for i, count in enumerate(grouped["loan_count"]):
            ax2.text(
                float(count) + pad,
                i,
                f"{int(count):,}",
                ha="left",
                va="center",
                fontsize=14,
                fontweight="bold",
                color=TEXT_DARK,
            )

        ax2.set_title(
            "Loan Count",
            pad=15,
            fontweight="bold",
            color=TEXT_DARK,
            fontsize=16,
        )
        ax2.set_xlabel("Number of Loans", fontweight="bold", color=TEXT_DARK, fontsize=14)
        ax2.grid(True, linestyle="--", alpha=0.25, linewidth=0.8, axis="x")
        ax2.set_axisbelow(True)

    else:
        bars2 = ax2.bar(
            idx,
            grouped["loan_count"],
            color=SECONDARY_COLOR,
            alpha=1.0,
            edgecolor="none",
        )

        for bar, count in zip(bars2, grouped["loan_count"]):
            height = bar.get_height()
            ax2.text(
                bar.get_x() + bar.get_width() / 2,
                height,
                f"{int(count):,}",
                ha="center",
                va="bottom",
                fontsize=14,
                fontweight="bold",
                color=TEXT_DARK,
            )

        ax2.set_xticks(idx)
        ax2.set_xticklabels(grouped[bucket_col], rotation=45, ha="right")
        ax2.set_title(
            "Loan Count",
            pad=15,
            fontweight="bold",
            color=TEXT_DARK,
            fontsize=16,
        )
        ax2.set_ylabel("Number of Loans", fontweight="bold", color=TEXT_DARK, fontsize=14)
        ax2.grid(True, linestyle="--", alpha=0.25, linewidth=0.8, axis="y")
        ax2.set_axisbelow(True)

    ax2.tick_params(axis="both", labelsize=12)

    # ===== DIVIDER LINE =====
    skip_divider = is_horizontal
    
    if not skip_divider:
        line = plt.Line2D(
            [0.515, 0.515],
            [0.12, 0.8],
            transform=fig.transFigure,
            color="#BFBFBF",
            linewidth=2,
            linestyle="--",
            alpha=0.6,
        )
        fig.add_artist(line)

    if chart_title:
        fig.suptitle(
            chart_title,
            fontsize=18,
            fontweight="bold",
            color=TEXT_DARK,
            y=0.98,
        )

    plt.tight_layout(rect=[0.02, 0.16, 0.98, 0.94] if chart_title else [0, 0, 1, 1])
    fig.savefig(out_path, dpi=300, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return True


def percentage_formatter(x, pos):
    return f"{x:.1f}%"


def disable_shadow(shape) -> None:
    """Disable shadow for a pptx shape (avoids default/theme shadows)."""
    try:
        shape.shadow.inherit = False
        shape.shadow.visible = False
    except Exception:
        pass


def validate_file_path(path_str: str, must_exist: bool = True) -> Path:
    """Sanity-check and normalise file path."""
    try:
        p = Path(path_str).expanduser().resolve()
    except Exception as e:
        raise ValueError(f"Invalid path: {path_str}. Error: {e}")

    if must_exist and not p.exists():
        raise ValueError(f"File does not exist: {path_str}")

    if p.exists() and p.is_file():
        size_mb = p.stat().st_size / (1024 * 1024)
        if size_mb > MAX_FILE_SIZE_MB:
            raise ValueError(
                f"File too large: {size_mb:.1f}MB (max {MAX_FILE_SIZE_MB}MB)"
            )

    if p.exists() and not p.is_file():
        raise ValueError(f"Path is not a file: {path_str}")

    return p


def weighted_average(series: pd.Series, weights: pd.Series) -> float:
    """Calculate weighted average, properly handling NaN values and zero weights.
    
    Special handling: If series name is 'youngest_borrower_age', also filter out
    ages that are 0 or less than 18 (invalid ages).
    """
    mask = series.notna() & weights.notna() & (weights > 0)
    
    # Special handling for age columns - filter out invalid ages
    if series.name == 'youngest_borrower_age' or 'age' in str(series.name).lower():
        mask = mask & (series > 0) & (series >= 18)
    
    if not mask.any():
        return np.nan
    return np.average(series[mask], weights=weights[mask])


# ============================
# DATA LAYER
# ============================

def _normalise_interest_rate(series: pd.Series) -> pd.Series:
    """Normalise interest rates to decimal (0.05 = 5%)."""
    s = pd.to_numeric(series, errors="coerce")
    non_null = s.dropna()
    if not non_null.empty and non_null.median() > 1:
        s = s / 100.0
    return s


def load_data(path: str) -> pd.DataFrame:
    """
    Load canonical CSV for ERM PPTX generation.
    Mirrors the Streamlit app logic from streamlit_app_erm.py.
    """
    file_path = validate_file_path(path, must_exist=True)

    df = pd.read_csv(file_path, nrows=MAX_PPTX_ROWS)

    # Normalise column names (avoid trailing spaces / mixed casing issues)
    df.columns = [str(c).strip() for c in df.columns]

    if df.empty:
        raise ValueError("CSV is empty")

    # Date conversions
    for col in ["origination_date", "maturity_date"]:
        if col in df.columns:
            df[col] = pd.to_datetime(df[col], errors="coerce")

    # Numeric conversions
    numeric_cols = [
        "current_principal_balance", "original_principal_balance",
        "current_interest_rate", "current_valuation_amount",
        "original_valuation_amount", "current_loan_to_value",
        "original_loan_to_value", "accrued_interest", "nneg_ratio",
        "borrower_1_age", "borrower_2_age", "youngest_borrower_age"
    ]
    
    for col in numeric_cols:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce")

    # Interest rate normalization
    if "current_interest_rate" in df.columns:
        df["current_interest_rate"] = _normalise_interest_rate(df["current_interest_rate"])

    # LTV normalization (if LTV is 0-100, convert to 0-1)
    for ltv_col in ["current_loan_to_value", "original_loan_to_value"]:
        if ltv_col in df.columns:
            non_null = df[ltv_col].dropna()
            if not non_null.empty and non_null.median() > 2:
                df[ltv_col] = df[ltv_col] / 100.0

    # NNEG ratio normalization (if NNEG is 0-100, convert to 0-1)
    if "nneg_ratio" in df.columns:
        non_null = df["nneg_ratio"].dropna()
        if not non_null.empty and non_null.median() > 2:
            df["nneg_ratio"] = df["nneg_ratio"] / 100.0

    # Create total_balance (principal + accrued interest)
    if "current_principal_balance" in df.columns:
        df["total_balance"] = df["current_principal_balance"].copy()
        if "accrued_interest" in df.columns:
            df["total_balance"] = df["total_balance"].fillna(0) + df["accrued_interest"].fillna(0)

    # Youngest borrower age
    # - Prefer the already-typed column if present (your loan tape has it)
    # - Only derive from borrower_1_age / borrower_2_age if the typed column is missing/empty
    if "youngest_borrower_age" in df.columns and df["youngest_borrower_age"].notna().any():
        pass  # keep as-is
    elif "borrower_1_age" in df.columns and "borrower_2_age" in df.columns:
        df["youngest_borrower_age"] = df[["borrower_1_age", "borrower_2_age"]].min(axis=1)
    elif "borrower_1_age" in df.columns:
        df["youngest_borrower_age"] = df["borrower_1_age"]
    else:
        df["youngest_borrower_age"] = np.nan

    # Origination year for stratifications
    if "origination_date" in df.columns:
        df["origination_year"] = df["origination_date"].dt.year

    print(f"Loaded {len(df):,} loans for PPTX generation")
    return df


# ============================
# CHART GENERATORS
# ============================

def create_bucket_column(df: pd.DataFrame, col: str, bucket_type: str) -> pd.Series:
    """Create bucketed columns for stratification analysis."""
    if bucket_type == "ltv":
        bins = [0, 0.20, 0.30, 0.40, 0.50, 0.60, 0.70, 0.80, 1.50]
        labels = ["0-20%", "20-30%", "30-40%", "40-50%", "50-60%", "60-70%", "70-80%", "80%+"]
        return pd.cut(df[col], bins=bins, labels=labels, include_lowest=True)
    
    elif bucket_type == "age":
        bins = [0, 55, 60, 65, 70, 75, 80, 85, 150]
        labels = ["<55", "55-60", "60-65", "65-70", "70-75", "75-80", "80-85", "85+"]
        return pd.cut(df[col], bins=bins, labels=labels, include_lowest=True)
    
    elif bucket_type == "rate":
        # Rate is already in decimal (0.05 = 5%)
        bins = [0, 0.03, 0.04, 0.05, 0.06, 0.07, 0.08, 1.0]
        labels = ["<3%", "3-4%", "4-5%", "5-6%", "6-7%", "7-8%", "8%+"]
        return pd.cut(df[col], bins=bins, labels=labels, include_lowest=True)
    
    elif bucket_type == "nneg":
        # NNEG is 0-1 (e.g., 0.60 = 60%)
        bins = [0, 0.50, 0.60, 0.70, 0.80, 0.90, 1.50]
        labels = ["<50%", "50-60%", "60-70%", "70-80%", "80-90%", "90%+"]
        return pd.cut(df[col], bins=bins, labels=labels, include_lowest=True)
    
    else:
        return pd.Series([None] * len(df), index=df.index)


# ------------------------------------------------------------------
# Dual-chart wrappers (balance + count) for PPTX slides
# ------------------------------------------------------------------

def save_ltv_dual_chart(df: pd.DataFrame, out_path: str) -> bool:
    """Generate dual LTV distribution chart (balance + count)."""
    if "current_loan_to_value" not in df.columns or "total_balance" not in df.columns:
        return False

    df_work = df.copy()
    df_work["ltv_bucket"] = create_bucket_column(df_work, "current_loan_to_value", "ltv")
    return save_dual_bar_chart(df_work, "ltv_bucket", out_path, "Current LTV Distribution")


def save_age_dual_chart(df: pd.DataFrame, out_path: str) -> bool:
    """Generate dual age distribution chart (balance + count)."""
    if "youngest_borrower_age" not in df.columns or "total_balance" not in df.columns:
        return False

    # Keep parity with existing age filtering logic (valid adult ages only)
    df_work = df[
        df["youngest_borrower_age"].notna()
        & (df["youngest_borrower_age"] > 0)
        & (df["youngest_borrower_age"] >= 18)
    ].copy()
    if df_work.empty:
        return False

    df_work["age_bucket"] = create_bucket_column(df_work, "youngest_borrower_age", "age")
    return save_dual_bar_chart(df_work, "age_bucket", out_path, "Borrower Age Distribution")


def save_broker_dual_chart(df: pd.DataFrame, out_path: str) -> bool:
    """Generate dual broker channel distribution chart (balance + count)."""
    if "broker_channel" not in df.columns or "total_balance" not in df.columns:
        return False

    broker_totals = df.groupby("broker_channel")["total_balance"].sum()
    top_brokers = broker_totals.nlargest(10).index
    df_filtered = df[df["broker_channel"].isin(top_brokers)].copy()
    if df_filtered.empty:
        return False

    return save_dual_bar_chart(
        df_filtered,
        "broker_channel",
        out_path,
        "Broker Channel Distribution (Top 10)",
        orientation="horizontal"
    )

def save_ticket_size_dual_chart(df: pd.DataFrame, out_path: str) -> bool:
    """Generate dual ticket size distribution chart (balance + count)."""
    if "total_balance" not in df.columns:
        return False
    
    # Create ticket size buckets
    df_work = df.copy()
    df_work["ticket_bucket"] = pd.cut(
        df_work["total_balance"],
        bins=[0, 75000, 100000, 175000, 250000, float('inf')],
        labels=['<£75K', '£75K-£100K', '£100K-£175K', '£175K-£250K', '>£250K'],
        include_lowest=True
    )
    
    # Filter out null buckets
    df_work = df_work[df_work["ticket_bucket"].notna()].copy()
    if df_work.empty:
        return False
    
    # Ensure categorical ordering
    df_work["ticket_bucket"] = pd.Categorical(
        df_work["ticket_bucket"],
        categories=['<£75K', '£75K-£100K', '£100K-£175K', '£175K-£250K', '>£250K'],
        ordered=True
    )
    
    return save_dual_bar_chart(
        df_work,
        "ticket_bucket",
        out_path,
        "Balance by Ticket Size",
        orientation="vertical"
    )


def save_geographic_region_dual_chart(df: pd.DataFrame, out_path: str, top_n: int = 10) -> bool:
    """Generate dual geographic region distribution chart (balance + count).

    Intended for use as a fallback when treemap rendering is unavailable.
    """
    if "geographic_region" not in df.columns or "total_balance" not in df.columns:
        return False

    region_totals = (
        df.groupby("geographic_region")["total_balance"]
        .sum()
        .sort_values(ascending=False)
    )
    top_regions = region_totals.head(top_n).index
    df_filtered = df[df["geographic_region"].isin(top_regions)].copy()
    if df_filtered.empty:
        return False

    return save_dual_bar_chart(
        df_filtered,
        "geographic_region",
        out_path,
        f"Geographic Distribution (Top {top_n})",
    )


def save_product_type_dual_chart(df: pd.DataFrame, out_path: str) -> bool:
    """Generate dual product type distribution chart (balance + count).

    Mirrors the 'Product Type Distribution' section in the Stratifications tab.
    """
    col = next(
        (c for c in ["erm_product_type", "product_type"] if c in df.columns),
        None,
    )
    if col is None or "total_balance" not in df.columns:
        return False

    return save_dual_bar_chart(df, col, out_path, "Product Type Distribution")


def save_geographic_count_treemap(df: pd.DataFrame, out_path: str) -> bool:
    """Generate geographic distribution treemap coloured by loan count.

    Mirrors the right-hand 'Loan Count by Region' treemap in the dashboard.
    Falls back to a dual bar chart if squarify is not installed.
    """
    if "geographic_region" not in df.columns:
        return False

    geo_data = (
        df.groupby("geographic_region")
        .size()
        .reset_index(name="loan_count")
        .sort_values("loan_count", ascending=False)
        .head(12)
    )
    geo_data = geo_data[geo_data["loan_count"] > 0]

    if geo_data.empty:
        return False

    if not HAS_SQUARIFY:
        # Fallback: dual bar chart showing count
        top_regions = geo_data["geographic_region"].tolist()
        df_filtered = df[df["geographic_region"].isin(top_regions)].copy()
        return save_dual_bar_chart(
            df_filtered,
            "geographic_region",
            out_path,
            "Geographic Distribution: Loans by Region (Top 10)",
        )

    fig, ax = plt.subplots(figsize=(14, 8))

    colors = [SECONDARY_COLOR, PRIMARY_COLOR, "#7EBAB5", "#5AA9A3",
              "#A3CCC9", "#FFB84D", "#FF9933", "#919DD1"] * 2

    labels = [
        f"{region}\n{count:,} loans"
        for region, count in zip(
            geo_data["geographic_region"].values,
            geo_data["loan_count"].values,
        )
    ]
    squarify.plot(
        sizes=geo_data["loan_count"].values,
        label=labels,
        color=colors[: len(geo_data)],
        alpha=0.85,
        ax=ax,
        text_kwargs={"fontsize": 9, "weight": "bold", "color": "white"},
        edgecolor="white",
        linewidth=2,
    )
    ax.axis("off")
    ax.set_title(
        "Geographic Concentration: Loan Count by Region",
        pad=20,
        fontweight="bold",
        color=TEXT_DARK,
        fontsize=20,
        loc="center",
    )

    plt.tight_layout(pad=1.5)
    fig.savefig(out_path, dpi=300, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return True


def save_ltv_distribution(df: pd.DataFrame, out_path: str) -> bool:
    """Generate LTV distribution bar chart."""
    if "current_loan_to_value" not in df.columns or "total_balance" not in df.columns:
        return False

    df_work = df.copy()
    df_work["ltv_bucket"] = create_bucket_column(df_work, "current_loan_to_value", "ltv")
    
    # FIX: Use DataFrame index for counting
    grouped = df_work.groupby("ltv_bucket", as_index=False).agg({
        "total_balance": "sum"
    })
    grouped["loan_count"] = df_work.groupby("ltv_bucket").size().values
    
    if grouped.empty:
        return False

    fig, ax = plt.subplots(figsize=(14, 8))
    
    bars = ax.bar(
        range(len(grouped)),
        grouped["total_balance"],
        color=PRIMARY_COLOR,
        alpha=1.0,  # Solid, not transparent
        edgecolor='none',  # No outline
        linewidth=0,
    )
    
    # Add value labels on top of bars
    for i, (bar, balance) in enumerate(zip(bars, grouped["total_balance"])):
        height = bar.get_height()
        label = f"£{balance/1_000_000:.1f}M" if balance >= 1_000_000 else f"£{balance/1_000:.0f}K"
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            height,
            label,
            ha='center',
            va='bottom',
            fontsize=9,
            fontweight='bold',
            color=TEXT_DARK
        )
    
    ax.set_xticks(range(len(grouped)))
    ax.set_xticklabels(grouped["ltv_bucket"], rotation=45, ha="right")
    
    ax.set_title(
        "Current LTV Distribution",
        pad=20,
        fontweight="bold",
        color=TEXT_DARK,
        fontsize=20,
        loc='center'
    )
    # X-axis label removed - tick labels are self-explanatory
    ax.set_ylabel("Outstanding Balance (£)", fontweight="bold", color=TEXT_DARK, fontsize=18)
    ax.yaxis.set_major_formatter(FuncFormatter(millions_formatter))
    
    ax.grid(True, linestyle="--", alpha=0.25, linewidth=0.8, axis='y')
    ax.set_axisbelow(True)
    
    plt.tight_layout(pad=1.5)
    fig.savefig(out_path, dpi=300, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return True


def save_age_distribution(df: pd.DataFrame, out_path: str) -> bool:
    """Generate borrower age distribution bar chart."""
    if "youngest_borrower_age" not in df.columns or "total_balance" not in df.columns:
        return False

    # Filter: Valid ages (not NaN, > 0, >= 18)
    df_work = df[
        df["youngest_borrower_age"].notna() &
        (df["youngest_borrower_age"] > 0) &
        (df["youngest_borrower_age"] >= 18)
    ].copy()

    if df_work.empty:
        return False
    
    df_work["age_bucket"] = create_bucket_column(df_work, "youngest_borrower_age", "age")
    
    # FIX: Use DataFrame index for counting instead of assuming unique_identifier exists
    grouped = df_work.groupby("age_bucket", as_index=False).agg({
        "total_balance": "sum"
    })
    grouped["loan_count"] = df_work.groupby("age_bucket").size().values
    
    if grouped.empty or grouped["total_balance"].sum() == 0:
        return False

    fig, ax = plt.subplots(figsize=(14, 8))
    
    bars = ax.bar(
        range(len(grouped)),
        grouped["total_balance"],
        color=PRIMARY_COLOR,
        alpha=1.0,  # Solid
        edgecolor='none',  # No outline
        linewidth=0,
    )
    
    # Add value labels on top of bars
    for i, (bar, balance) in enumerate(zip(bars, grouped["total_balance"])):
        height = bar.get_height()
        label = f"£{balance/1_000_000:.1f}M" if balance >= 1_000_000 else f"£{balance/1_000:.0f}K"
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            height,
            label,
            ha='center',
            va='bottom',
            fontsize=9,
            fontweight='bold',
            color=TEXT_DARK
        )
    
    ax.set_xticks(range(len(grouped)))
    ax.set_xticklabels(grouped["age_bucket"], rotation=45, ha="right")
    
    ax.set_title(
        "Borrower Age Distribution (Youngest)",
        pad=20,
        fontweight="bold",
        color=TEXT_DARK,
        fontsize=20,
        loc='center'
    )
    # X-axis label removed - tick labels are self-explanatory
    ax.set_ylabel("Outstanding Balance (£)", fontweight="bold", color=TEXT_DARK, fontsize=18)
    ax.yaxis.set_major_formatter(FuncFormatter(millions_formatter))
    
    ax.grid(True, linestyle="--", alpha=0.25, linewidth=0.8, axis='y')
    ax.set_axisbelow(True)
    
    plt.tight_layout(pad=1.5)
    fig.savefig(out_path, dpi=300, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return True


def save_nneg_distribution(df: pd.DataFrame, out_path: str) -> bool:
    """Generate NNEG ratio distribution bar chart."""
    if "nneg_ratio" not in df.columns or "total_balance" not in df.columns:
        return False

    df_work = df.copy()
    df_work["nneg_bucket"] = create_bucket_column(df_work, "nneg_ratio", "nneg")
    
    # FIX: Use DataFrame index for counting
    grouped = df_work.groupby("nneg_bucket", as_index=False).agg({
        "total_balance": "sum"
    })
    grouped["loan_count"] = df_work.groupby("nneg_bucket").size().values
    
    if grouped.empty:
        return False

    fig, ax = plt.subplots(figsize=(14, 8))
    
    # Color gradient for NNEG risk
    colors = ["#449C95", "#5AA9A3", "#7EBAB5", "#A3CCC9", "#FFB84D", "#FF9933"]
    bar_colors = colors[:len(grouped)]
    
    bars = ax.bar(
        range(len(grouped)),
        grouped["total_balance"],
        color=bar_colors,
        alpha=1.0,  # Solid
        edgecolor='none',  # No outline
        linewidth=0,
    )
    
    # Add value labels on top of bars
    for i, (bar, balance) in enumerate(zip(bars, grouped["total_balance"])):
        height = bar.get_height()
        label = f"£{balance/1_000_000:.1f}M" if balance >= 1_000_000 else f"£{balance/1_000:.0f}K"
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            height,
            label,
            ha='center',
            va='bottom',
            fontsize=9,
            fontweight='bold',
            color=TEXT_DARK
        )
    
    ax.set_xticks(range(len(grouped)))
    ax.set_xticklabels(grouped["nneg_bucket"], rotation=45, ha="right")
    
    ax.set_title(
        "NNEG Ratio Distribution (No Negative Equity Guarantee)",
        pad=20,
        fontweight="bold",
        color=TEXT_DARK,
        fontsize=15,
        loc='center'
    )
    # X-axis label removed - tick labels are self-explanatory
    ax.set_ylabel("Outstanding Balance (£)", fontweight="bold", color=TEXT_DARK, fontsize=12)
    ax.yaxis.set_major_formatter(FuncFormatter(millions_formatter))
    
    ax.grid(True, linestyle="--", alpha=0.25, linewidth=0.8, axis='y')
    ax.set_axisbelow(True)
    
    plt.tight_layout(pad=1.5)
    fig.savefig(out_path, dpi=300, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return True


def save_geographic_treemap(df: pd.DataFrame, out_path: str) -> bool:
    """Generate geographic distribution treemap (ALWAYS treemap, no fallback to bars)."""
    if "geographic_region" not in df.columns or "total_balance" not in df.columns:
        return False

    geo_data = (
        df.groupby("geographic_region", as_index=False)["total_balance"]
        .sum()
        .sort_values("total_balance", ascending=False)
        .head(12)
    )
    geo_data = geo_data[geo_data["total_balance"] > 0]

    if geo_data.empty:
        return False

    fig, ax = plt.subplots(figsize=(14, 8))

    # ALWAYS use treemap (install squarify if needed)
    if not HAS_SQUARIFY:
        # Fallback to a dual bar chart (balance + count), consistent with other bar-chart slides
        print("WARNING: squarify not installed. Using dual bar chart fallback for geographic distribution.")

        # Limit to top regions by balance (avoids unreadable labels)
        region_totals = df.groupby("geographic_region")["total_balance"].sum().sort_values(ascending=False)
        top_regions = region_totals.head(10).index
        df_filtered = df[df["geographic_region"].isin(top_regions)].copy()
        if df_filtered.empty:
            plt.close(fig)
            return False

        plt.close(fig)
        return save_dual_bar_chart(
            df_filtered,
            "geographic_region",
            out_path,
            "Geographic Distribution (Top 10)",
        )

    # Primary path: treemap
    # Primary path: treemap
    colors = [PRIMARY_COLOR, ACCENT_COLOR, "#7EBAB5", "#5AA9A3", 
              "#A3CCC9", "#FFB84D", "#FF9933", "#919DD1"] * 2
    
    labels = [
        f"{region}\n£{val/1_000_000:.1f}M"
        for region, val in zip(
            geo_data["geographic_region"].values,
            geo_data["total_balance"].values,
        )
    ]
    squarify.plot(
        sizes=geo_data["total_balance"].values,
        label=labels,
        color=colors[:len(geo_data)],
        alpha=0.85,
        ax=ax,
        text_kwargs={"fontsize": 9, "weight": "bold", "color": "white"},
        edgecolor="white",
        linewidth=2,
    )
    ax.axis("off")
    ax.set_title(
        "Geographic Concentration: Top Regions",
        pad=20,
        fontweight="bold",
        color=TEXT_DARK,
        fontsize=20,
        loc='center'
    )

    plt.tight_layout(pad=1.5)
    fig.savefig(out_path, dpi=300, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return True

def save_broker_channel_distribution(df: pd.DataFrame, out_path: str) -> bool:
    """Generate broker channel distribution bar chart."""
    if "broker_channel" not in df.columns or "total_balance" not in df.columns:
        return False

    # FIX: Use DataFrame index for counting
    broker_data = (
        df.groupby("broker_channel", as_index=False)
        .agg({"total_balance": "sum"})
        .sort_values("total_balance", ascending=False)
        .head(10)
    )
    broker_data["loan_count"] = df.groupby("broker_channel").size().sort_values(ascending=False).head(10).values

    if broker_data.empty:
        return False

    fig, ax = plt.subplots(figsize=(14, 8))
    
    bars = ax.barh(
        broker_data["broker_channel"].values[::-1],
        broker_data["total_balance"].values[::-1],
        color=PRIMARY_COLOR,
        alpha=1.0,  # Solid
        edgecolor='none',  # No outline
        linewidth=0,
    )
    
    ax.set_title(
        "Distribution by Broker Channel",
        pad=20,
        fontweight="bold",
        color=TEXT_DARK,
        fontsize=20,
        loc='center'
    )
    ax.set_xlabel("Outstanding Balance (£)", fontweight="bold", color=TEXT_DARK, fontsize=18, labelpad=10)
    ax.xaxis.set_major_formatter(FuncFormatter(millions_formatter))
    
    # Add value labels
    for bar in bars:
        width = bar.get_width()
        ax.text(
            width,
            bar.get_y() + bar.get_height() / 2,
            f" £{width/1_000_000:.1f}M",
            va="center",
            ha="left",
            fontsize=9,
            color=TEXT_DARK,
            weight="bold",
        )
    
    ax.grid(True, linestyle="--", alpha=0.25, linewidth=0.8, axis='x')
    ax.set_axisbelow(True)
    
    plt.tight_layout(pad=1.5)
    fig.savefig(out_path, dpi=300, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return True

def save_bubble_balance_vs_value(df: pd.DataFrame, out_path: str) -> bool:
    """Generate bubble chart: Current Outstanding Balance vs. Current Property Value.

    Mirrors the dashboard 'Balance vs Current Property Value' scatter chart.
    Points are coloured by geographic region (categorical) when available,
    matching the dashboard's colour dimension.
    """
    required_cols = ["current_principal_balance", "current_valuation_amount", "current_loan_to_value"]
    if not all(col in df.columns for col in required_cols):
        return False

    work_cols = required_cols + (["geographic_region"] if "geographic_region" in df.columns else [])
    df_plot = df[work_cols].dropna(subset=required_cols)
    if df_plot.empty or len(df_plot) < 10:
        return False

    # Sample if too many loans
    if len(df_plot) > 1000:
        df_plot = df_plot.sample(1000, random_state=42)

    fig, ax = plt.subplots(figsize=(14, 8))

    has_region = "geographic_region" in df_plot.columns and df_plot["geographic_region"].notna().any()

    if has_region:
        # Categorical coloring by geographic_region — matches dashboard
        regions = df_plot["geographic_region"].fillna("Unknown").unique()
        color_cycle = CHART_COLORS + [
            "#7EBAB5", "#5AA9A3", "#A3CCC9", "#FFB84D", "#FF9933",
            "#449C95", "#E07B54", "#6C5B7B",
        ]
        region_color = {r: color_cycle[i % len(color_cycle)] for i, r in enumerate(regions)}
        colors_array = [region_color[r] for r in df_plot["geographic_region"].fillna("Unknown")]

        scatter = ax.scatter(
            df_plot["current_valuation_amount"],
            df_plot["current_principal_balance"],
            s=80,
            c=colors_array,
            alpha=0.7,
            edgecolor="white",
            linewidth=0.5,
        )

        # Legend patches (top 10 regions by count)
        top_regions = (
            df_plot["geographic_region"].fillna("Unknown")
            .value_counts()
            .head(10)
            .index
        )
        handles = [
            mpatches.Patch(color=region_color[r], label=r)
            for r in top_regions
            if r in region_color
        ]
        ax.legend(
            handles=handles,
            title="Region",
            loc="upper left",
            fontsize=8,
            title_fontsize=9,
            frameon=False,
            ncol=2,
        )
    else:
        # Fallback: continuous colormap by LTV
        from matplotlib.colors import LinearSegmentedColormap
        brand_cmap = LinearSegmentedColormap.from_list(
            "brand_blue",
            ["#E8F4F8", ACCENT_COLOR, PRIMARY_COLOR],
            N=256,
        )
        scatter = ax.scatter(
            df_plot["current_valuation_amount"],
            df_plot["current_principal_balance"],
            s=80,
            c=df_plot["current_loan_to_value"] * 100,
            cmap=brand_cmap,
            alpha=0.7,
            edgecolor="white",
            linewidth=0.5,
        )
        cbar = plt.colorbar(scatter, ax=ax)
        cbar.set_label("Current LTV (%)", rotation=270, labelpad=20, fontweight="bold")

    # Diagonal reference line (balance = value)
    max_val = max(
        df_plot["current_valuation_amount"].max(),
        df_plot["current_principal_balance"].max(),
    )
    ax.plot([0, max_val], [0, max_val], "k--", alpha=0.3, linewidth=1, label="Balance = Value")

    ax.set_title(
        "Current Outstanding Balance vs. Current Property Value",
        pad=20,
        fontweight="bold",
        color=TEXT_DARK,
        fontsize=14,
        loc="center",
    )
    ax.set_xlabel("Current Property Value (£)", fontweight="bold", color=TEXT_DARK, fontsize=12, labelpad=10)
    ax.set_ylabel("Current Outstanding Balance (£)", fontweight="bold", color=TEXT_DARK, fontsize=12)

    ax.xaxis.set_major_formatter(FuncFormatter(millions_formatter))
    ax.yaxis.set_major_formatter(FuncFormatter(millions_formatter))

    ax.grid(True, linestyle="--", alpha=0.25, linewidth=0.8)
    ax.set_axisbelow(True)

    plt.tight_layout(pad=1.5)
    fig.savefig(out_path, dpi=300, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return True


def save_bubble_ltv_vs_age(df: pd.DataFrame, out_path: str) -> bool:
    """Generate bubble chart: LTV vs Borrower Age.

    Mirrors the dashboard 'LTV vs Youngest Borrower Age' scatter chart.
    Points are coloured by product type (categorical) when available,
    matching the dashboard's colour dimension. Bubble size is proportional
    to loan balance.
    """
    required_cols = ["current_loan_to_value", "youngest_borrower_age", "total_balance"]
    if not all(col in df.columns for col in required_cols):
        return False

    product_col = next(
        (c for c in ["erm_product_type", "product_type"] if c in df.columns),
        None,
    )
    work_cols = required_cols + ([product_col] if product_col else [])
    df_plot = df[work_cols].copy()
    df_plot = df_plot.dropna(subset=required_cols)
    df_plot = df_plot[
        (df_plot["youngest_borrower_age"] > 0) &
        (df_plot["youngest_borrower_age"] >= 18)
    ]

    if df_plot.empty or len(df_plot) < 10:
        return False

    # Sample if too many loans
    if len(df_plot) > 2000:
        df_plot = df_plot.sample(2000, random_state=42)

    fig, ax = plt.subplots(figsize=(14, 8))

    # Size by balance (normalized)
    sizes = (df_plot["total_balance"] / df_plot["total_balance"].max() * 300).clip(20, 300)

    has_product = product_col and product_col in df_plot.columns and df_plot[product_col].notna().any()

    if has_product:
        # Categorical coloring by product type — matches dashboard
        products = df_plot[product_col].fillna("Unknown").unique()
        color_cycle = CHART_COLORS + [
            "#7EBAB5", "#5AA9A3", "#A3CCC9", "#FFB84D", "#FF9933",
        ]
        product_color = {p: color_cycle[i % len(color_cycle)] for i, p in enumerate(products)}
        colors_array = [product_color[p] for p in df_plot[product_col].fillna("Unknown")]

        ax.scatter(
            df_plot["youngest_borrower_age"],
            df_plot["current_loan_to_value"] * 100,
            s=sizes,
            c=colors_array,
            alpha=0.65,
            edgecolor="white",
            linewidth=0.4,
        )

        handles = [
            mpatches.Patch(color=product_color[p], label=p)
            for p in products
            if p in product_color
        ]
        ax.legend(
            handles=handles,
            title="Product Type",
            loc="upper right",
            fontsize=9,
            title_fontsize=10,
            frameon=False,
        )
    else:
        # Fallback: continuous colormap by balance
        from matplotlib.colors import LinearSegmentedColormap
        brand_cmap = LinearSegmentedColormap.from_list(
            "brand_accent",
            ["#E8F4F8", ACCENT_COLOR, PRIMARY_COLOR],
            N=256,
        )
        scatter = ax.scatter(
            df_plot["youngest_borrower_age"],
            df_plot["current_loan_to_value"] * 100,
            s=sizes,
            c=df_plot["total_balance"],
            cmap=brand_cmap,
            alpha=0.65,
            edgecolor="white",
            linewidth=0.4,
        )
        cbar = plt.colorbar(scatter, ax=ax)
        cbar.set_label("Loan Balance (£)", rotation=270, labelpad=20, fontweight="bold")
        cbar.ax.yaxis.set_major_formatter(FuncFormatter(millions_formatter))

    ax.set_title(
        "Current LTV vs. Borrower Age (Youngest)",
        pad=20,
        fontweight="bold",
        color=TEXT_DARK,
        fontsize=14,
        loc="center",
    )
    ax.set_xlabel("Borrower Age (Youngest)", fontweight="bold", color=TEXT_DARK, fontsize=12, labelpad=10)
    ax.set_ylabel("Current LTV (%)", fontweight="bold", color=TEXT_DARK, fontsize=12)

    ax.grid(True, linestyle="--", alpha=0.25, linewidth=0.8)
    ax.set_axisbelow(True)

    plt.tight_layout(pad=1.5)
    fig.savefig(out_path, dpi=300, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return True

def save_vintage_distribution(df: pd.DataFrame, out_path: str) -> bool:
    """Generate vintage year distribution bar chart."""
    if "origination_year" not in df.columns or "total_balance" not in df.columns:
        return False

    # FIX: Use DataFrame index for counting
    vintage_data = (
        df.groupby("origination_year", as_index=False)
        .agg({"total_balance": "sum"})
        .sort_values("origination_year")
    )
    vintage_data["loan_count"] = df.groupby("origination_year").size().sort_index().values

    if vintage_data.empty:
        return False

    fig, ax1 = plt.subplots(figsize=(14, 8))

    # Bar chart for balance
    bars = ax1.bar(
        vintage_data["origination_year"],
        vintage_data["total_balance"],
        color=PRIMARY_COLOR,
        alpha=0.7,
        label="Outstanding Balance",
    )
    
    ax1.set_xlabel("Vintage Year", fontweight="bold", color=TEXT_DARK, fontsize=12, labelpad=10)
    ax1.set_ylabel("Outstanding Balance (£)", fontweight="bold", color=TEXT_DARK, fontsize=12)
    ax1.yaxis.set_major_formatter(FuncFormatter(millions_formatter))
    ax1.tick_params(axis='y', labelcolor=TEXT_DARK)
    
    # Line chart for count
    ax2 = ax1.twinx()
    line = ax2.plot(
        vintage_data["origination_year"],
        vintage_data["loan_count"],
        color=PRIMARY_COLOR,
        marker='o',
        linewidth=2.5,
        markersize=6,
        label="Loan Count",
    )
    
    ax2.set_ylabel("Number of Loans", fontweight="bold", color=TEXT_DARK, fontsize=12)
    ax2.tick_params(axis='y', labelcolor=TEXT_DARK)
    
    ax1.set_title(
        "Portfolio by Vintage Year",
        pad=20,
        fontweight="bold",
        color=TEXT_DARK,
        fontsize=15,
        loc="center"
    )
    
    # Combined legend
    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1 + lines2, labels1 + labels2, loc='upper left', frameon=False)
    
    ax1.grid(True, linestyle="--", alpha=0.25, linewidth=0.8, axis='y')
    ax1.set_axisbelow(True)
    
    plt.tight_layout(pad=1.5)
    fig.subplots_adjust(bottom=0.18)
    fig.savefig(out_path, dpi=300, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return True


# ============================
# SLIDE BUILDERS
# ============================

def add_footer(slide, text: str):
    """Add consistent footer to slide."""
    box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(7.1),
        Inches(9.0),
        Inches(0.3),
    )
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.name = FONT_BODY
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(*hex_to_rgb(TEXT_LIGHT))
    p.alignment = PP_ALIGN.CENTER

def add_page_numbers(prs: Presentation):
    """
    Add page numbers to all slides except cover (first) and end (last).
    Numbering starts at 1 on the first content slide.
    """
    if len(prs.slides) <= 2:
        return

    page_num = 1
    for i, slide in enumerate(prs.slides):
        if i == 0 or i == len(prs.slides) - 1:
            continue

        box = slide.shapes.add_textbox(
            prs.slide_width - Inches(1.0),
            Inches(7.08),
            Inches(0.8),
            Inches(0.25),
        )
        tf = box.text_frame
        tf.text = str(page_num)
        p = tf.paragraphs[0]
        p.font.name = FONT_NUMBERS
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(*hex_to_rgb(TEXT_LIGHT))
        p.alignment = PP_ALIGN.RIGHT
        page_num += 1

def add_cover_slide(prs: Presentation, logo_path: Optional[str], title: str):
    """Create branded cover slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Slide background to PRIMARY_COLOR
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*hex_to_rgb(PRIMARY_COLOR))

    # Color band (now WHITE rectangle)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(2.8),
        prs.slide_width,
        Inches(2.0),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255) 
    bg.line.fill.background()
    disable_shadow(bg)

    # Logo
    if logo_path and os.path.exists(logo_path):
        logo_width = Inches(2.5)
        left = (prs.slide_width - logo_width) / 2
        top = Inches(1.2)
        slide.shapes.add_picture(logo_path, left, top, width=logo_width)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(3.0),
        Inches(9.0),
        Inches(1.0),
    )
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.name = FONT_TITLE
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*hex_to_rgb(PRIMARY_COLOR))  
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(4.0),
        Inches(9.0),
        Inches(0.5),
    )
    tf = sub_box.text_frame
    _report_date = (
        datetime.strptime(STATIC_REPORTING_DATE, "%Y-%m-%d").strftime("%B %d, %Y")
        if STATIC_REPORTING_DATE
        else datetime.now().strftime("%B %d, %Y")
    )
    tf.text = f"Confidential \u2022 Data as of {_report_date}"
    p = tf.paragraphs[0]
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(*hex_to_rgb(PRIMARY_COLOR))  
    p.alignment = PP_ALIGN.CENTER


def style_content_title(prs: Presentation, slide, title_text: str, logo_path: Optional[str]):
    """Style content slide with title, logo, and subtitle - FIXED VERSION."""
    
    # Delete default title placeholder to avoid z-order issues
    shapes_to_delete = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            try:
                if shape.placeholder_format.type == 1:  # Title
                    shapes_to_delete.append(shape)
            except:
                pass
    
    for shape in shapes_to_delete:
        sp = shape.element
        sp.getparent().remove(sp)
    
    # 1. Add dark blue bar (0.85" high - proper size)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        prs.slide_width,
        Inches(0.787),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*hex_to_rgb(PRIMARY_COLOR))
    bar.line.fill.background()
    
    disable_shadow(bar)
# 2. Add title text (WHITE, on top of bar)
    title_box = slide.shapes.add_textbox(
        Inches(0.25),
        Inches(0.14),
        Inches(7.0),
        Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.word_wrap = False
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    p = title_frame.paragraphs[0]
    p.font.name = FONT_TITLE
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)  
    p.alignment = PP_ALIGN.LEFT
    
    # 3. Add logo
    if logo_path:
        if os.path.exists(logo_path):
            try:
                logo_width = Inches(1.5)
                left = prs.slide_width - logo_width - Inches(0.5)
                slide.shapes.add_picture(logo_path, left, Inches(0.25), width=logo_width)
            except Exception as e:
                print(f"   ⚠ Failed to add logo: {e}")
        else:
            print(f"   ⚠ Logo file not found: {logo_path}")
    
    # 4. Add subtitle placeholder
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.25),
        Inches(0.95),
        Inches(9.0),
        Inches(0.3)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = ""
    subtitle_frame.word_wrap = False
    
    p = subtitle_frame.paragraphs[0]
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(*hex_to_rgb(TEXT_LIGHT))
    p.alignment = PP_ALIGN.LEFT


def add_kpi_slide(prs: Presentation, df: pd.DataFrame, logo_path: Optional[str]):
    """Create ERM-specific KPI dashboard slide."""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    style_content_title(prs, slide, "Portfolio Overview Dashboard", logo_path)

    total_loans = len(df)
    total_balance = df["total_balance"].sum() if "total_balance" in df.columns else 0
    
    # Average and maximum loan size
    avg_loan_size = total_balance / total_loans if total_loans > 0 else 0.0
    max_loan_size = df["total_balance"].max() if "total_balance" in df.columns else 0.0
    
    # WA calculations (matching streamlit_app_erm.py)
    if total_balance > 0:
        # LTV (convert to percentage for display)
        wa_current_ltv = weighted_average(df["current_loan_to_value"], df["total_balance"])
        if not pd.isna(wa_current_ltv):
            wa_current_ltv *= 100
        else:
            wa_current_ltv = 0
        
        wa_rate = weighted_average(df["current_interest_rate"], df["total_balance"])
        if pd.isna(wa_rate):
            wa_rate = 0
        
        # Age calculation - keep as None if no valid data
        wa_age = weighted_average(df["youngest_borrower_age"], df["total_balance"])
        if pd.isna(wa_age):
            wa_age = None
        
        wa_nneg = weighted_average(df["nneg_ratio"] if "nneg_ratio" in df.columns else pd.Series(np.nan, index=df.index), df["total_balance"])
        if not pd.isna(wa_nneg):
            wa_nneg *= 100  # Convert to percentage
        else:
            wa_nneg = 0
    else:
        wa_current_ltv = 0
        wa_rate = 0
        wa_age = None  # Keep as None
        wa_nneg = 0
    
    # Original LTV
    if "original_loan_to_value" in df.columns and total_balance > 0:
        wa_original_ltv = weighted_average(df["original_loan_to_value"], df["total_balance"])
        if not pd.isna(wa_original_ltv):
            wa_original_ltv *= 100
        else:
            wa_original_ltv = 0
    else:
        wa_original_ltv = 0

    # Largest geographic exposure
    if "geographic_region" in df.columns and total_balance > 0:
        geo_balance = df.groupby("geographic_region")["total_balance"].sum()
        largest_geo = geo_balance.idxmax() if not geo_balance.empty else "N/A"
        largest_geo_pct = (geo_balance.max() / total_balance * 100) if not geo_balance.empty else 0
    else:
        largest_geo = "N/A"
        largest_geo_pct = 0
    
    # Largest broker
    if "broker_channel" in df.columns and total_balance > 0:
        broker_balance = df.groupby("broker_channel")["total_balance"].sum()
        largest_broker = broker_balance.idxmax() if not broker_balance.empty else "N/A"
        largest_broker_pct = (broker_balance.max() / total_balance * 100) if not broker_balance.empty else 0
    else:
        largest_broker = "N/A"
        largest_broker_pct = 0

    # KPIs matching the ERM dashboard (10 KPIs in 2 rows)
    kpis = [
        # First row
        ("Total Loans", f"{total_loans:,}", "Portfolio count"),
        ("Portfolio Balance", f"£{total_balance/1_000_000:.1f}M", "Outstanding + accrued"),
        ("WA Current LTV", f"{wa_current_ltv:.1f}%", "Indexed to current value"),
        ("WA Interest Rate", f"{wa_rate:.2%}", "Balance-weighted"),
        ("WA Borrower Age", f"{wa_age:.0f}" if wa_age is not None and wa_age > 0 else "N/A", "Youngest borrower"),
        # Second row
        ("WA Original LTV", f"{wa_original_ltv:.1f}%", "At origination"),
        ("Avg. Loan Size", f"£{avg_loan_size/1_000:.0f}K", "Balance-weighted average"),
        ("Largest Loan", f"£{max_loan_size/1_000:.0f}K", "Maximum single exposure"),
        ("Largest Geographic", f"{largest_geo_pct:.1f}%", f"{largest_geo}"),
        ("Largest Broker", f"{largest_broker_pct:.1f}%", f"{largest_broker}"),
    ]
    
    # Card dimensions and layout
    card_w = Inches(1.9)
    card_h = Inches(1.25)
    h_gap = Inches(0.05)
    v_gap = Inches(0.15)
    
    # Calculate to center the grid (5 columns)
    total_width = 5 * card_w + 4 * h_gap
    margin_left = (Inches(10) - total_width) / 2
    start_top = Inches(2.0)

    for idx, (label, value, sub) in enumerate(kpis):
        row = idx // 5
        col = idx % 5
        left = margin_left + col * (card_w + h_gap)
        top = start_top + row * (card_h + v_gap)

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left,
            top,
            card_w,
            card_h,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 249, 250)
        card.line.color.rgb = RGBColor(*hex_to_rgb(SECONDARY_COLOR))
        card.line.width = Pt(1)

        # Label
        label_box = slide.shapes.add_textbox(
            left,
            top + Inches(0.15),
            card_w,
            Inches(0.25),
        )
        tf = label_box.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.font.name = FONT_BODY
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(*hex_to_rgb(TEXT_LIGHT))
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Value
        value_box = slide.shapes.add_textbox(
            left,
            top + Inches(0.45),
            card_w,
            Inches(0.40),
        )
        tf = value_box.text_frame
        tf.text = value
        p = tf.paragraphs[0]
        p.font.name = FONT_NUMBERS
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*hex_to_rgb(PRIMARY_COLOR))
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub_box = slide.shapes.add_textbox(
            left,
            top + Inches(0.88),
            card_w,
            Inches(0.25),
        )
        tf = sub_box.text_frame
        tf.text = sub
        p = tf.paragraphs[0]
        p.font.name = FONT_BODY
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(*hex_to_rgb(TEXT_LIGHT))
        p.alignment = PP_ALIGN.CENTER

    # Risk monitoring tiles (if available)
    if RISK_MONITORING_AVAILABLE:
        try:
            monitor = RiskMonitor(ALL_LIMITS)
            results = monitor.run_all_checks(df)
            
            # Count status
            red_count = sum(1 for r in results if r.status == "red")
            amber_count = sum(1 for r in results if r.status == "amber")
            green_count = sum(1 for r in results if r.status == "green")
            
            # Risk tile dimensions
            risk_tile_w = Inches(2.8)
            risk_tile_h = Inches(0.8)
            risk_start_top = start_top + 2 * (card_h + v_gap) + Inches(0.3)
            risk_margin_left = Inches(1.5)
            risk_h_gap = Inches(0.2)
            
            risk_tiles = [
                ("Limits Breached", red_count, "#DC3545" if red_count > 0 else "#28A745"),
                ("Limits in Warning", amber_count, "#FFC107" if amber_count > 0 else "#28A745"),
                ("Limits Clear", green_count, "#28A745"),
            ]
            
            for idx, (label, value, color) in enumerate(risk_tiles):
                left = risk_margin_left + idx * (risk_tile_w + risk_h_gap)
                
                # Tile background
                tile = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left,
                    risk_start_top,
                    risk_tile_w,
                    risk_tile_h,
                )
                tile.fill.solid()
                tile.fill.fore_color.rgb = RGBColor(*hex_to_rgb(color))
                tile.line.color.rgb = RGBColor(*hex_to_rgb(color))
                tile.line.width = Pt(2)
                
                # Label
                label_box = slide.shapes.add_textbox(
                    left + Inches(0.1),
                    risk_start_top + Inches(0.1),
                    risk_tile_w - Inches(0.2),
                    Inches(0.25),
                )
                tf = label_box.text_frame
                tf.text = label
                p = tf.paragraphs[0]
                p.font.name = FONT_BODY
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True
                p.alignment = PP_ALIGN.LEFT
                
                # Value
                value_box = slide.shapes.add_textbox(
                    left + Inches(0.1),
                    risk_start_top + Inches(0.35),
                    risk_tile_w - Inches(0.2),
                    Inches(0.35),
                )
                tf = value_box.text_frame
                tf.text = str(value)
                p = tf.paragraphs[0]
                p.font.name = FONT_NUMBERS
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.LEFT
        except Exception as e:
            print(f"   Warning: Could not add risk tiles: {e}")
    
    add_footer(slide, REPORT_FOOTER)

def add_chart_slide(
    prs: Presentation,
    title: str,
    chart_path: str,
    logo_path: Optional[str],
    caption: str = "",
):
    """Add a slide with a single chart."""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    style_content_title(prs, slide, title, logo_path)

    if os.path.exists(chart_path):
        width = Inches(9.35)
        left = (prs.slide_width - width) / 2
        top = Inches(1.5)
        slide.shapes.add_picture(chart_path, left, top, width=width)

    if caption:
        caption_box = slide.shapes.add_textbox(
            Inches(0.7),
            Inches(6.5),
            Inches(8.6),
            Inches(0.4),
        )
        tf = caption_box.text_frame
        tf.text = caption
        p = tf.paragraphs[0]
        p.font.name = FONT_BODY
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(*hex_to_rgb(TEXT_LIGHT))
        p.alignment = PP_ALIGN.CENTER

    add_footer(slide, REPORT_FOOTER)


# ============================
# RISK MONITORING SLIDE
# ============================


def render_risk_limits_table_png(results, out_path: str, max_rows: int = 20) -> bool:
    """
    Render the Detailed Limit Status table to a PNG using matplotlib (Streamlit-like table).
    This is intended to be placed into PPTX as an image to avoid PowerPoint table layout issues.

    Args:
        results: list of RiskLimitResult objects (from RiskMonitor.check_all_limits()).
        out_path: destination PNG path.
        max_rows: max number of rows to render (keeps slide readable).

    Returns:
        True if rendered, False otherwise.
    """
    try:
        import numpy as np
        import pandas as pd
        import matplotlib.pyplot as plt
        from matplotlib import rcParams

        if not results:
            return False

        # Default Streamlit filters from the dashboard
        status_filter = {"red", "amber", "green"}
        severity_filter = {"critical", "high"}

        filtered = [
            r for r in results
            if getattr(r, "status", "unknown") in status_filter
            and getattr(r, "severity", "unknown") in severity_filter
        ]

        if not filtered:
            return False

        # Sort: red -> amber -> green; within status, critical before high, then by utilization desc
        status_order = {"red": 0, "amber": 1, "green": 2, "unknown": 3}
        severity_order = {"critical": 0, "high": 1, "medium": 2, "low": 3, "unknown": 9}

        filtered.sort(
            key=lambda r: (
                status_order.get(getattr(r, "status", "unknown"), 9),
                severity_order.get(getattr(r, "severity", "unknown"), 9),
                -float(getattr(r, "utilization_pct", 0.0) or 0.0),
            )
        )

        rows = filtered[:max_rows]

        # Build table rows (match Streamlit formatting closely)
        table_rows = []
        for r in rows:
            limit_id = str(getattr(r, "limit_id", "")).lower()
            current_value = getattr(r, "current_value", np.nan)
            limit_value = getattr(r, "limit_value", np.nan)

            if any(k in limit_id for k in ("pct", "concentration", "ltv", "dpd")):
                current_str = f"{current_value:.1f}%" if not np.isnan(current_value) else "N/A"
                limit_str = f"{limit_value:.1f}%" if not np.isnan(limit_value) else "N/A"
            elif "absolute" in limit_id:
                current_str = f"£{current_value:,.0f}" if not np.isnan(current_value) else "N/A"
                limit_str = f"£{limit_value:,.0f}" if not np.isnan(limit_value) else "N/A"
            elif any(k in limit_id for k in ("age", "rate")):
                current_str = f"{current_value:.1f}" if not np.isnan(current_value) else "N/A"
                limit_str = f"{limit_value:.1f}" if not np.isnan(limit_value) else "N/A"
            elif "size" in limit_id:
                current_str = f"{int(current_value)}" if not np.isnan(current_value) else "N/A"
                limit_str = f"{int(limit_value)}" if not np.isnan(limit_value) else "N/A"
            else:
                current_str = f"{current_value:.2f}" if not np.isnan(current_value) else "N/A"
                limit_str = f"{limit_value:.2f}" if not np.isnan(limit_value) else "N/A"

            # Keep description readable
            desc = str(getattr(r, "description", ""))
            if len(desc) > 85:
                desc = desc[:82] + "..."

            table_rows.append({
                "": "●",  # status bullet (colored per-row)
                "Category": str(getattr(r, "category", "")),
                "Description": desc,
                "Current": current_str,
                "Limit": limit_str,
                "Headroom": f"{(getattr(r, 'limit_value', 0.0) or 0.0) - (getattr(r, 'current_value', 0.0) or 0.0):,.2f}" if "pct" not in limit_id else "",
                "Utilization": f"{float(getattr(r, 'utilization_pct', 0.0) or 0.0):.0f}%",
                "Severity": str(getattr(r, "severity", "")).upper(),
                "_status": str(getattr(r, "status", "unknown")).lower(),
            })

        df_tbl = pd.DataFrame(table_rows)
        status_col = df_tbl.pop("_status")

        # Matplotlib table rendering
        rcParams["font.family"] = "DejaVu Sans"  # ensures bullet glyph is available on most installs

        n_rows = len(df_tbl)
        n_cols = len(df_tbl.columns)

        # Figure size tuned for PPTX insertion (wide and tall, minimal dead space)
        fig_w = 20
        fig_h = 1.2 + n_rows * 0.42
        fig_h = max(6.0, min(fig_h, 12.0))

        fig, ax = plt.subplots(figsize=(fig_w, fig_h))
        ax.axis("off")

        col_labels = list(df_tbl.columns)
        cell_text = df_tbl.values.tolist()

        table = ax.table(
            cellText=cell_text,
            colLabels=col_labels,
            cellLoc="left",
            colLoc="left",
            loc="center",
        )

        # Font sizing + scaling
        table.auto_set_font_size(False)
        base_fs = 12 if n_rows <= 12 else 11 if n_rows <= 18 else 10
        table.set_fontsize(base_fs)
        table.scale(1.0, 1.45)

        # Column width weighting (sum doesn't need to equal 1; matplotlib normalizes)
        # Designed to resemble Streamlit proportions.
        widths = {
            0: 0.04,  # status
            1: 0.12,  # category
            2: 0.44,  # description
            3: 0.10,  # current
            4: 0.10,  # limit
            5: 0.08,  # headroom
            6: 0.08,  # utilization
            7: 0.06,  # severity
        }

        # Colors
        header_bg = PRIMARY_COLOR
        header_fg = "#FFFFFF"
        grid_color = "#D9D9D9"
        row_alt = "#F7F9FC"
        text_color = TEXT_DARK

        status_colors = {
            "red": "#DC3545",
            "amber": "#FFC107",
            "green": "#28A745",
            "unknown": "#B0B0B0",
        }

        # Style header row
        for c in range(n_cols):
            cell = table[(0, c)]
            cell.set_facecolor(header_bg)
            cell.set_edgecolor(header_bg)
            cell.get_text().set_color(header_fg)
            cell.get_text().set_weight("bold")
            cell.PAD = 0.02

        # Style body cells
        for r in range(1, n_rows + 1):
            is_alt = (r % 2 == 0)
            for c in range(n_cols):
                cell = table[(r, c)]
                cell.set_facecolor(row_alt if is_alt else "white")
                cell.set_edgecolor(grid_color)
                cell.set_linewidth(0.6)
                cell.PAD = 0.02
                cell.get_text().set_color(text_color)

            # Status bullet color (col 0)
            st = status_col.iloc[r - 1]
            st_color = status_colors.get(st, status_colors["unknown"])
            table[(r, 0)].get_text().set_color(st_color)
            table[(r, 0)].get_text().set_weight("bold")
            table[(r, 0)].get_text().set_horizontalalignment("center")

            # Utilization color emphasis (col index based on col_labels)
            if "Utilization" in col_labels:
                util_idx = col_labels.index("Utilization")
                table[(r, util_idx)].get_text().set_color(st_color)
                table[(r, util_idx)].get_text().set_weight("bold")
                table[(r, util_idx)].get_text().set_horizontalalignment("center")

        # Tighten layout without tight_layout (tables often warn); use subplots_adjust
        fig.subplots_adjust(left=0.01, right=0.99, top=0.98, bottom=0.02)

        out_path = str(out_path)
        Path(out_path).parent.mkdir(parents=True, exist_ok=True)
        fig.savefig(out_path, dpi=220, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        return True

    except Exception as e:
        print(f"   > WARNING: Failed to render risk limits PNG table: {e}")
        import traceback
        traceback.print_exc()
        return False


def add_risk_monitoring_slide(prs, df, logo_path=None):
    """
    Add risk monitoring summary slide with KPIs and breach table.
    Only adds if risk monitoring is available.
    """
    if not RISK_MONITORING_AVAILABLE:
        print("   > Skipping Risk Monitoring slide (module not available)")
        return
    
    try:
        # Run risk monitor
        monitor = RiskMonitor(df)
        results = monitor.check_all_limits()
        summary = monitor.get_summary_stats(results)
        
        # Create slide with CONSISTENT layout
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Apply CONSISTENT styling (blue bar, white title, logo, subtitle)
        style_content_title(prs, slide, "Risk Limit Monitoring", logo_path)

        # Remove any stray placeholder textboxes from the chosen layout (e.g., a lone "Limits")
        for shp in list(slide.shapes):
            try:
                if getattr(shp, "has_text_frame", False):
                    if (shp.text_frame.text or "").strip() == "Limits":
                        shp._element.getparent().remove(shp._element)
            except Exception:
                pass
        
        # KPI row (summary metrics) - positioned below title bar
        kpi_top = 1.65  # Moved down to create more space for subtitle
        kpi_width = 2.1
        kpi_height = 1.2
        kpi_spacing = 0.15
        kpis = [
            ("Total Limits", summary["total_limits"], TEXT_DARK),
            ("Breaches", summary["breaches"], "#DC3545"),
            ("Warnings", summary["warnings"], "#FFC107"),
            ("Compliant", summary["compliant"], "#28A745")
        ]
        
        # Calculate starting position to center the row
        total_width = 4 * kpi_width + 3 * kpi_spacing
        start_left = (Inches(10) - Inches(total_width)) / 2
        
        for i, (label, value, color) in enumerate(kpis):
            left_pos = start_left + Inches(i * (kpi_width + kpi_spacing))
            
            # KPI box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left=left_pos,
                top=Inches(kpi_top),
                width=Inches(kpi_width),
                height=Inches(kpi_height)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*hex_to_rgb(BACKGROUND_LIGHT))
            box.line.color.rgb = RGBColor(*hex_to_rgb(SECONDARY_COLOR))
            box.line.width = Pt(1)
            
            # Value (in top portion of tile)
            value_box = slide.shapes.add_textbox(
                left=left_pos,
                top=Inches(kpi_top) + Inches(0.15),
                width=Inches(kpi_width),
                height=Inches(0.5)
            )
            text_frame = value_box.text_frame
            text_frame.text = str(value)
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.name = FONT_NUMBERS
            p.font.color.rgb = RGBColor(*hex_to_rgb(color))
            
            # Label (in bottom portion of tile)
            label_box = slide.shapes.add_textbox(
                left=left_pos,
                top=Inches(kpi_top) + Inches(0.65),
                width=Inches(kpi_width),
                height=Inches(0.4)
            )
            text_frame = label_box.text_frame
            text_frame.text = label
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(11)
            p.font.bold = False
            p.font.name = FONT_BODY
            p.font.color.rgb = RGBColor(*hex_to_rgb(TEXT_LIGHT))
        
        
        # Detailed limit status table
        table_top = 2.45  # Slightly higher to allow a larger PNG table
        header_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(3.1),
            width=Inches(9),
            height=Inches(0.30)
        )
        header_frame = header_box.text_frame
        header_frame.text = "All Risk Limits Status"
        p = header_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = FONT_TITLE
        p.font.color.rgb = RGBColor(*hex_to_rgb(TEXT_DARK))
        
        # Sort: red first, then amber, then green - all by utilization within each status
        status_rank = {"red": 0, "amber": 1, "green": 2}
        all_results = sorted(
            results,
            key=lambda r: (
                status_rank.get(getattr(r, "status", "green"), 2),
                -float(getattr(r, "utilization_pct", 0.0))
            )
        )

        # Apply Streamlit default filters (status: red/amber/green; severity: critical/high)
        status_filter = {"red", "amber", "green"}
        severity_filter = {"critical", "high"}

        filtered_results = [
            r for r in results
            if getattr(r, "status", "unknown") in status_filter
            and getattr(r, "severity", "unknown") in severity_filter
        ]

        # Sort by status (red first, then amber, then green), then severity, then utilization desc
        status_order = {"red": 0, "amber": 1, "green": 2, "unknown": 3}
        severity_order = {"critical": 0, "high": 1, "medium": 2, "low": 3, "unknown": 9}
        filtered_results.sort(
            key=lambda r: (
                status_order.get(getattr(r, "status", "unknown"), 9),
                severity_order.get(getattr(r, "severity", "unknown"), 9),
                -float(getattr(r, "utilization_pct", 0.0) or 0.0),
            )
        )

        if not filtered_results:
            info = slide.shapes.add_textbox(Inches(0.5), Inches(table_top + 0.6), Inches(9), Inches(0.6))
            tf = info.text_frame
            tf.text = "No limits match the current filters."
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.name = FONT_TITLE
            p.font.color.rgb = RGBColor(*hex_to_rgb(TEXT_DARK))
            p.alignment = PP_ALIGN.CENTER
            add_footer(slide, REPORT_FOOTER)
            print("   > Risk monitoring slide added (no matching limits)")
            return
        # Render the Detailed Limit Status table to PNG (Streamlit-like) and insert into slide
        try:
            chart_dir = Path("_pptx_charts")
            chart_dir.mkdir(exist_ok=True)

            # Use a unique filename per run to avoid Windows file locks / permission issues
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            png_path = chart_dir / f"risk_limits_table_{ts}.png"

            ok = render_risk_limits_table_png(filtered_results, str(png_path), max_rows=20)

            if not ok:
                info = slide.shapes.add_textbox(Inches(0.5), Inches(table_top + 0.6), Inches(9), Inches(0.6))
                tf = info.text_frame
                tf.text = "No limits match the current filters."
                p = tf.paragraphs[0]
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.name = FONT_TITLE
                p.font.color.rgb = RGBColor(*hex_to_rgb(TEXT_DARK))
                p.alignment = PP_ALIGN.CENTER
                add_footer(slide, REPORT_FOOTER)
                print("   > Risk monitoring slide added (no matching limits)")
                return

            # Place the PNG where the PPTX table previously sat
            png_width = Inches(9.65)

            pic = slide.shapes.add_picture(
                str(png_path),
                left=(prs.slide_width - png_width) / 2,
                top=Inches(table_top + 0.38),
                width=png_width,
                height=Inches(4.15),
            )

            # Send picture behind other content (but not behind the slide background)
            spTree = slide.shapes._spTree
            spTree.remove(pic._element)
            spTree.insert(2, pic._element)


        except Exception as e:
            print(f"   > WARNING: Failed to insert PNG table into PPTX: {e}")
            import traceback
            traceback.print_exc()


        add_footer(slide, REPORT_FOOTER)

    except Exception as e:
        print(f"   Warning: Could not create risk monitoring slide: {e}")
        import traceback
        traceback.print_exc()

# ============================
# SCENARIO ANALYSIS SLIDE
# ============================

def add_scenario_analysis_slide(prs, df, logo_path=None):
    """
    Add scenario analysis slide with dual charts (balance projection + NNEG losses).
    Only adds if scenario charts module is available and required columns exist.
    """
    if not SCENARIO_CHARTS_AVAILABLE:
        print("   > Skipping Scenario Analysis slide (module not available)")
        return

    # Check required columns
    required_cols = [
        "current_principal_balance",
        "current_valuation_amount",
        "current_interest_rate",
    ]
    missing_cols = [c for c in required_cols if c not in df.columns]

    if missing_cols:
        print(f"   > Skipping Scenario Analysis slide (missing columns: {missing_cols})")
        return

    try:
        # Create charts directory
        chart_dir = Path("_pptx_charts")
        chart_dir.mkdir(exist_ok=True)

        # Use Base Case assumptions
        base_assumptions = ScenarioAssumptions(
            hpi_rate=0.02,
            interest_rate_spread=0.0,
            voluntary_prepay_rate=0.02,
            mortality_rate=0.03,
            move_to_care_rate=0.01,
            sale_cost_pct=0.05,
            n_years=25,
        )

        # Generate dual scenario chart (balance + NNEG on one slide)
        scenario_dual_path = chart_dir / "scenario_dual.png"
        if save_dual_scenario_charts(
            df,
            str(scenario_dual_path),
            assumptions=base_assumptions,
            main_title="Base Case Scenario Analysis",
        ):
            # Add the slide
            add_chart_slide(
                prs,
                "Base Case Scenario Analysis",
                str(scenario_dual_path),
                logo_path,
                "25-year projection with 2% HPI, 6% exit rate",
            )
            print("   > Scenario Analysis slide added")
        else:
            print("   > WARNING: Scenario chart generation failed")

    except Exception as e:
        print(f"   Warning: Could not create scenario analysis slide: {e}")
        import traceback
        traceback.print_exc()

# ============================
# MAIN GENERATOR
# ============================

def add_end_slide(prs: Presentation, logo_path: Optional[str]):
    """Create branded end slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Slide background to PRIMARY_COLOR
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*hex_to_rgb(PRIMARY_COLOR))

    # White band (same as cover)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(2.8),
        prs.slide_width,
        Inches(2.0),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    disable_shadow(bg)

    # Logo
    if logo_path and os.path.exists(logo_path):
        logo_width = Inches(2.5)
        left = (prs.slide_width - logo_width) / 2
        top = Inches(0.6)
        slide.shapes.add_picture(logo_path, left, top, width=logo_width)

    # Text
    text_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(3.2),
        Inches(9.0),
        Inches(1.5),
    )
    tf = text_box.text_frame
    tf.text = ENTITY_NAME
    p = tf.paragraphs[0]
    p.font.name = FONT_TITLE
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*hex_to_rgb(PRIMARY_COLOR))
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(5.5),
        Inches(9.0),
        Inches(1.0),
    )
    tf = contact_box.text_frame
    tf.text = "Confidential Portfolio Analytics"
    for p in tf.paragraphs:
        p.font.name = FONT_BODY
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

def generate_pptx(
    input_path: str,
    output_path: str,
    logo_path: Optional[str] = None,
):
    """
    Main PPTX generation function for ERM analytics.
    
    Args:
        input_path: Path to typed canonical CSV
        output_path: Path for output PPTX
        logo_path: Optional path to company logo
    """
    
    print(f"\n{'='*60}")
    print("ERM Portfolio Analytics PPTX Generator")
    print(f"{'='*60}\n")
    
    # Check for required dependencies
    if not HAS_SQUARIFY:
        print("⚠️  WARNING: squarify package not installed!")
        print("   Geographic chart will use bar chart fallback instead of treemap.")
        print("   Install with: pip install squarify\n")
    
    # Load data
    print("[1/7] Loading portfolio data...")
    df = load_data(input_path)
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Look for logo (robust)
    if logo_path is None:
        search_dirs = [
            Path(__file__).resolve().parent,
            Path(input_path).resolve().parent,
            Path.cwd(),
        ]
        candidates = ["ere_logo.png", "ERE_logo.png", "logo.png",
                    "equity_release_europe_logo_highres.png", "buecor_logo.png"]

        for d in search_dirs:
            for name in candidates:
                p = d / name
                if p.exists():
                    logo_path = str(p)
                    break
            if logo_path:
                break

    print(f"   🖼 Logo resolved to: {logo_path}")
    if logo_path:
        print(f"   🖼 Logo exists: {os.path.exists(logo_path)}")

    # Cover slide
    print("[2/9] Creating cover slide...")
    add_cover_slide(
        prs,
        logo_path,
        COVER_TITLE,
    )
    
    # KPI dashboard
    print("[3/9] Creating KPI dashboard...")
    add_kpi_slide(prs, df, logo_path)
    
    # Generate charts
    chart_dir = Path("_pptx_charts")
    chart_dir.mkdir(exist_ok=True)
    
    print("\n[4/9] Generating charts...")
    print(f"   Data: {len(df):,} loans, {len(df.columns)} columns")

    charts = []

    # 1. LTV Distribution (DUAL: balance + count)
    ltv_path = chart_dir / "ltv_distribution_dual.png"
    if save_ltv_dual_chart(df, str(ltv_path)):
        charts.append(("Current LTV Distribution", str(ltv_path), ""))
        print(f"   ✓ LTV Distribution")
    else:
        print(f"   ✗ LTV Distribution FAILED")

    # 2. Product Type Distribution (DUAL: balance + count) — mirrors Stratifications tab
    product_path = chart_dir / "product_type_distribution_dual.png"
    if save_product_type_dual_chart(df, str(product_path)):
        charts.append(("Product Type Distribution", str(product_path), ""))
        print(f"   ✓ Product Type Distribution")
    else:
        print(f"   ✗ Product Type Distribution FAILED (column missing or no data)")

    # 3. Geographic Distribution — Balance treemap (mirrors dashboard left panel)
    geo_path = chart_dir / "geographic_treemap.png"
    if save_geographic_treemap(df, str(geo_path)):
        charts.append(("Geographic Distribution: Balance by Region", str(geo_path), ""))
        print(f"   ✓ Geographic Distribution (balance)")
    else:
        print(f"   ✗ Geographic Distribution (balance) FAILED")

    # 4. Geographic Distribution — Loan Count treemap (mirrors dashboard right panel)
    geo_cnt_path = chart_dir / "geographic_count_treemap.png"
    if save_geographic_count_treemap(df, str(geo_cnt_path)):
        charts.append(("Geographic Distribution: Loans by Region", str(geo_cnt_path), ""))
        print(f"   ✓ Geographic Distribution (count)")
    else:
        print(f"   ✗ Geographic Distribution (count) FAILED")

    # 5. Borrower Age Distribution (DUAL: balance + count)
    age_path = chart_dir / "age_distribution_dual.png"
    if save_age_dual_chart(df, str(age_path)):
        charts.append(("Borrower Age Distribution", str(age_path), ""))
        print(f"   ✓ Borrower Age Distribution")
    else:
        print(f"   ✗ Borrower Age Distribution FAILED")

    # 6. Vintage Distribution (balance bars + count line — mirrors Vintage Distribution section)
    vintage_path = chart_dir / "vintage_distribution.png"
    if save_vintage_distribution(df, str(vintage_path)):
        charts.append(("Portfolio by Vintage Year", str(vintage_path), ""))
        print(f"   ✓ Vintage Distribution")
    else:
        print(f"   ✗ Vintage Distribution FAILED")

    # 7. Broker Channel Distribution (DUAL: balance + count)
    broker_path = chart_dir / "broker_distribution_dual.png"
    if save_broker_dual_chart(df, str(broker_path)):
        charts.append(("Broker Channel Distribution (Top 10)", str(broker_path), ""))
        print(f"   ✓ Broker Channel Distribution")
    else:
        print(f"   ✗ Broker Channel Distribution FAILED")

    # 8. Ticket Size Distribution (DUAL: balance + count)
    ticket_path = chart_dir / "ticket_size_distribution_dual.png"
    if save_ticket_size_dual_chart(df, str(ticket_path)):
        charts.append(("Balance by Ticket Size", str(ticket_path), ""))
        print(f"   ✓ Ticket Size Distribution")
    else:
        print(f"   ✗ Ticket Size Distribution FAILED")

    # 9. Bubble Chart 1: Balance vs Property Value (colored by geographic region)
    bubble1_path = chart_dir / "bubble_balance_vs_value.png"
    if save_bubble_balance_vs_value(df, str(bubble1_path)):
        charts.append(("Outstanding Balance vs. Property Value", str(bubble1_path), ""))
        print(f"   ✓ Balance vs Value bubble chart")
    else:
        print(f"   ✗ Balance vs Value bubble FAILED")

    # 10. Bubble Chart 2: LTV vs Borrower Age (colored by product type)
    bubble2_path = chart_dir / "bubble_ltv_vs_age.png"
    if save_bubble_ltv_vs_age(df, str(bubble2_path)):
        charts.append(("LTV vs. Borrower Age", str(bubble2_path), ""))
        print(f"   ✓ LTV vs Borrower Age bubble chart")
    else:
        print(f"   ✗ LTV vs Borrower Age bubble FAILED")

    # Add chart slides
    n_expected = 10
    print(f"\n   Summary: Generated {len(charts)} of {n_expected} charts")
    print(f"[5/9] Adding chart slides...")
    for title, path, caption in charts:
        add_chart_slide(prs, title, path, logo_path, caption)
    
    # Risk monitoring slide (if available)
    print("\n[6/9] Checking risk monitoring...")
    add_risk_monitoring_slide(prs, df, logo_path)

    # Scenario analysis slides (if available)
    print("\n[7/9] Adding scenario analysis slides...")
    add_scenario_analysis_slide(prs, df, logo_path)

    print("\n[8/9] Adding end slide...")
    add_end_slide(prs, logo_path)
    
    # Save presentation
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    print(f"\n[9/9] Saving presentation to: {output_path}")
    # Add page numbers (exclude cover + end)
    add_page_numbers(prs)
    
    prs.save(str(output_path))
    
    print(f"\n{'='*60}")
    print("SUCCESS! PPTX generation complete")
    print(f"{'='*60}")
    print(f"  Loans processed:  {len(df):,}")
    print(f"  Slides created:   {len(prs.slides)}")
    print(f"  Charts generated: {len(charts)}")
    print(f"  Output file:      {output_path.resolve()}")
    print(f"{'='*60}\n")


# ============================
# CLI
# ============================

def main():
    parser = argparse.ArgumentParser(
        description="Generate ERM analytics PPTX from typed canonical CSV"
    )
    
    parser.add_argument(
        "--input",
        default="synthetic_uk_equity_release_mandatory_full_ESMA_Annex2_canonical_ESMA_Annex2_typed.csv",
        help="Input typed canonical CSV file"
    )
    
    parser.add_argument(
        "--output",
        default=f"erm_portfolio_report_{datetime.now():%Y%m%d_%H%M%S}.pptx",
        help="Output PPTX file path"
    )
    
    parser.add_argument(
        "--logo",
        default=None,
        help="Optional path to company logo (PNG)"
    )
    
    args = parser.parse_args()
    
    try:
        generate_pptx(
            input_path=args.input,
            output_path=args.output,
            logo_path=args.logo
        )
    except Exception as e:
        print(f"\nERROR: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()