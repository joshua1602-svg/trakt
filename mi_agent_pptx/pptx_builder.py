"""mi_agent_pptx.pptx_builder — assemble the institutional investor PPTX pack.

Assembles a 16:9 deck with python-pptx from lens-routed metrics, dashboard-
faithful charts and straplines. The visual system mirrors the MI Agent React
dashboard (navy surfaces, periwinkle accents, mono figures, KPI tiles with
prior-period delta chips).

Two fixes over the first cut drive the quality step-change:

* **No image distortion** — each chart is rendered at the *exact* pixel size of
  the panel it lands in, then placed at that same size, so python-pptx never
  stretches it.
* **Lens routing** — every slide/metric is bound to a lens (funded / pipeline /
  forecast) and draws from that lens's frame; a slide whose lens has no data
  renders a branded placeholder instead of borrowing another lens's numbers.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .chart_resolver import ChartResolver, ChartResult
from .deck_config import DeckConfig, SlideSpec
from .insight_resolver import StraplineResolver
from .metric_resolver import MetricResolver, MetricResult
from .placeholders import AppendixNotes, render_placeholder_png
from .pptx_theme import PptxTheme, THEME

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
EMU_PER_IN = 914400

BROKER_DIMENSIONS = {"broker_channel", "broker"}
# Filled triangles render reliably across all PowerPoint default fonts.
DELTA_ARROW = {"up": "▲", "down": "▼", "flat": "–"}


@dataclass
class BuildContext:
    client_name: str
    as_of_date: str
    run_dir: str
    lens: str = "total"
    consolidated: bool = False
    generated_by: str = "trakt MI Agent"
    logo_path: Optional[str] = None
    prior_label: str = ""
    source_artifacts: List[str] = field(default_factory=list)


class DeckBuilder:
    def __init__(
        self,
        config: DeckConfig,
        context: BuildContext,
        metric_resolver: MetricResolver,
        chart_resolvers: Dict[str, Optional[ChartResolver]],
        strapline_resolver: StraplineResolver,
        appendix: AppendixNotes,
        theme: PptxTheme = THEME,
    ):
        self.config = config
        self.ctx = context
        self.metrics = metric_resolver
        self.chart_resolvers = chart_resolvers or {}
        self.straplines = strapline_resolver
        self.appendix = appendix
        self.theme = theme
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._blank = self.prs.slide_layouts[6]
        self._page = 0
        self.slide_records: List[Dict[str, Any]] = []

    # ------------------------------------------------------------------ colour
    def _rgb(self, hex_str: str) -> RGBColor:
        r, g, b = self.theme.rgb(hex_str)
        return RGBColor(r, g, b)

    # ------------------------------------------------------------------- public
    def build(self, output_path: str | Path) -> Dict[str, Any]:
        for slide_spec in self.config.slides:
            handler = getattr(self, f"_slide_{slide_spec.type}", self._slide_charts)
            self.slide_records.append(handler(slide_spec))
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(out))
        return {"output": str(out), "slides": self.slide_records,
                "coverage_notes": self.appendix.as_list()}

    # ---------------------------------------------------------------- scaffold
    def _new_slide(self):
        slide = self.prs.slides.add_slide(self._blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self._rgb(self.theme.bg_page)
        return slide

    def _text(self, slide, left, top, width, height, text, *, size=14, color=None,
              bold=False, align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP,
              font=None, spacing=None):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, 0)
        p = tf.paragraphs[0]
        p.alignment = align
        if spacing is not None:
            p.line_spacing = spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font or self.theme.font_sans
        run.font.color.rgb = self._rgb(color or self.theme.ink_100)
        return box

    def _panel(self, slide, left, top, width, height, *, fill=None, line=None,
               radius=True, line_w=0.75):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
            left, top, width, height)
        try:
            shape.adjustments[0] = 0.045
        except Exception:
            pass
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(fill or self.theme.bg_panel)
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = self._rgb(line)
            shape.line.width = Pt(line_w)
        shape.shadow.inherit = False
        return shape

    def _header(self, slide, title: str, strapline: str, *, lens: str = "") -> None:
        rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                      Inches(0.12), SLIDE_H)
        rail.fill.solid()
        rail.fill.fore_color.rgb = self._rgb(
            self.theme.mint if lens == "forecast" else self.theme.peri)
        rail.line.fill.background()
        rail.shadow.inherit = False
        self._text(slide, Inches(0.55), Inches(0.34), Inches(10.8), Inches(0.6),
                   title, size=25, bold=True)
        if lens in ("pipeline", "forecast", "funded"):
            self._lens_chip(slide, lens)
        self._text(slide, Inches(0.57), Inches(1.0), Inches(11.4), Inches(0.5),
                   strapline, size=12, color=self.theme.peri, italic=True)

    def _lens_chip(self, slide, lens: str) -> None:
        colors = {"funded": self.theme.peri, "pipeline": self.theme.peri,
                  "forecast": self.theme.mint}
        chip = self._panel(slide, Inches(11.3), Inches(0.42), Inches(1.5),
                           Inches(0.36), fill=self.theme.bg_panel_alt,
                           line=colors.get(lens, self.theme.peri), line_w=1.0)
        chip = self._text(slide, Inches(11.3), Inches(0.47), Inches(1.5),
                          Inches(0.3), lens.upper() + " LENS", size=8.5,
                          color=colors.get(lens, self.theme.peri), bold=True,
                          align=PP_ALIGN.CENTER)

    def _footer(self, slide) -> None:
        self._page += 1
        self._text(slide, Inches(0.55), Inches(7.08), Inches(10.0), Inches(0.3),
                   self.config.footer, size=8, color=self.theme.ink_500)
        self._text(slide, Inches(12.3), Inches(7.08), Inches(0.8), Inches(0.3),
                   str(self._page), size=8, color=self.theme.ink_500,
                   align=PP_ALIGN.RIGHT)

    def _strapline_for(self, spec: SlideSpec) -> str:
        return self.straplines.resolve(spec.id, spec.type)

    def _record(self, spec, strapline, *, placeholder=False, extra=None):
        rec = {"id": spec.id, "type": spec.type, "title": spec.title,
               "strapline": strapline, "mandatory": spec.mandatory,
               "placeholder": placeholder}
        if extra:
            rec.update(extra)
        return rec

    # ------------------------------------------------------------------ COVER
    def _slide_cover(self, spec: SlideSpec) -> Dict[str, Any]:
        slide = self._new_slide()
        # Layered brand backdrop.
        glow = self._panel(slide, Inches(-1), Inches(-2), Inches(9), Inches(6),
                           fill=self.theme.bg_panel, radius=False)
        glow.fill.fore_color.rgb = self._rgb("#0e1430")
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(3.02),
                                     Inches(2.2), Inches(0.07))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(self.theme.peri)
        bar.line.fill.background()
        bar.shadow.inherit = False

        self._text(slide, Inches(0.9), Inches(0.7), Inches(6), Inches(0.4),
                   "TRAKT · MI AGENT", size=12, color=self.theme.peri, bold=True)
        self._text(slide, Inches(0.86), Inches(1.7), Inches(11.5), Inches(1.4),
                   self.ctx.client_name, size=42, bold=True)
        self._text(slide, Inches(0.92), Inches(3.25), Inches(11), Inches(0.6),
                   self.config.name, size=19, color=self.theme.peri)
        strap = self._strapline_for(spec)
        self._text(slide, Inches(0.92), Inches(3.95), Inches(10.5), Inches(0.7),
                   strap, size=13, color=self.theme.ink_300, italic=True,
                   spacing=1.15)
        # Footer facts.
        self._text(slide, Inches(0.92), Inches(5.7), Inches(6), Inches(0.4),
                   f"Data cut-off   {self.ctx.as_of_date or 'n/a'}", size=12.5,
                   color=self.theme.ink_100)
        self._text(slide, Inches(0.92), Inches(6.12), Inches(8), Inches(0.4),
                   f"Generated automatically by {self.ctx.generated_by}", size=11,
                   color=self.theme.ink_400)
        if self.ctx.logo_path and Path(self.ctx.logo_path).exists():
            try:
                slide.shapes.add_picture(self.ctx.logo_path, Inches(11.0),
                                         Inches(0.6), height=Inches(0.85))
            except Exception:
                pass
        return self._record(spec, strap)

    # -------------------------------------------------------------------- KPI
    def _slide_kpi(self, spec: SlideSpec) -> Dict[str, Any]:
        slide = self._new_slide()
        strap = self._strapline_for(spec)
        self._header(slide, spec.title or "Executive Summary", strap)

        specs = self.config.metric_specs(spec.kpis)
        results = self.metrics.resolve_all(specs)

        cols = 5
        rows = max(1, (len(results) + cols - 1) // cols)
        gx, gy = Inches(0.16), Inches(0.22)
        left0, top0 = Inches(0.55), Inches(1.62)
        avail_w = Inches(12.25)
        tile_w = Emu(int((int(avail_w) - (cols - 1) * int(gx)) / cols))
        tile_h = Inches(1.62) if rows <= 2 else Inches(1.32)

        n_missing = 0
        for i, res in enumerate(results):
            r, c = divmod(i, cols)
            left = Emu(int(left0) + c * (int(tile_w) + int(gx)))
            top = Emu(int(top0) + r * (int(tile_h) + int(gy)))
            self._kpi_tile(slide, left, top, tile_w, tile_h, res)
            if not res.ok:
                n_missing += 1
                self.appendix.add(f"Executive KPI '{res.label}' unavailable: {res.note}")
        self._footer(slide)
        return self._record(spec, strap, placeholder=(n_missing == len(results)),
                            extra={"kpis_missing": n_missing})

    def _kpi_tile(self, slide, left, top, width, height, res: MetricResult) -> None:
        self._panel(slide, left, top, width, height, fill=self.theme.bg_panel_alt,
                    line=self.theme.line_soft, line_w=1.0)
        pad = Inches(0.16)
        inner_w = Emu(int(width) - 2 * int(pad))
        # label
        self._text(slide, left + pad, top + Inches(0.14), inner_w, Inches(0.3),
                   res.label.upper(), size=8.5, color=self.theme.ink_400, bold=True)
        # value
        val_color = self.theme.ink_100 if res.ok else self.theme.ink_500
        self._text(slide, left + pad, top + Inches(0.44), inner_w, Inches(0.55),
                   res.display, size=21, bold=True, color=val_color)
        # delta chip or hint
        y = top + Inches(1.02)
        if res.has_delta:
            color = {"up": self.theme.mint, "down": self.theme.rose,
                     "flat": self.theme.ink_400}.get(res.delta_dir, self.theme.ink_400)
            arrow = DELTA_ARROW.get(res.delta_dir, "–")
            prefix = f"{arrow} " if res.delta_dir != "flat" else ""
            self._text(slide, left + pad, y, inner_w, Inches(0.3),
                       f"{prefix}{res.delta_display}", size=9.5, color=color, bold=True)
        elif res.hint:
            self._text(slide, left + pad, y, inner_w, Inches(0.3),
                       res.hint, size=9.5, color=self.theme.peri)

    # --------------------------------------------------------------- CHART SLIDE
    def _slide_charts(self, spec: SlideSpec) -> Dict[str, Any]:
        slide = self._new_slide()
        strap = self._strapline_for(spec)
        lens = spec.lens if spec.lens in ("funded", "pipeline", "forecast") else "funded"
        self._header(slide, spec.title, strap, lens=lens)

        boxes = self._chart_boxes(len(spec.charts))
        rendered: List[ChartResult] = []
        for cspec, box in zip(spec.charts, boxes):
            rendered.append(self._render_into(slide, spec, cspec, lens, box))

        self._footer(slide)
        all_ph = bool(rendered) and all(c.placeholder for c in rendered)
        for c in rendered:
            if c.placeholder:
                self.appendix.add(f"[{spec.title}] {c.title}: {c.note}")
        return self._record(spec, strap, placeholder=all_ph,
                            extra={"charts": len(rendered)})

    def _chart_boxes(self, n: int):
        top = Inches(1.62)
        h = Inches(4.95)
        if n <= 1:
            return [(Inches(0.55), top, Inches(12.25), h)]
        return [(Inches(0.55), top, Inches(6.0), h),
                (Inches(6.78), top, Inches(6.0), h)]

    def _render_into(self, slide, slide_spec, cspec, lens, box) -> ChartResult:
        left, top, width, height = box
        title = cspec.get("title", cspec.get("id", "chart").replace("_", " ").title())
        cid = cspec.get("id", "chart")

        # card + header
        self._panel(slide, left, top, width, height, fill=self.theme.bg_panel,
                    line=self.theme.line, line_w=1.0)
        self._text(slide, left + Inches(0.22), top + Inches(0.16),
                   width - Inches(0.4), Inches(0.34), title, size=12.5, bold=True)

        # broker suppression (consolidated funded)
        is_broker = (cspec.get("dimension") in BROKER_DIMENSIONS) or bool(cspec.get("broker"))
        suppress = (self.config.suppress_broker_consolidated
                    or slide_spec.suppress_broker_consolidated)
        img_left = left + Inches(0.16)
        img_top = top + Inches(0.62)
        img_w_in = (int(width) - 2 * int(Inches(0.16))) / EMU_PER_IN
        img_h_in = (int(height) - int(Inches(0.62)) - int(Inches(0.16))) / EMU_PER_IN

        if is_broker and suppress and self.ctx.consolidated:
            note = ("Broker channel suppressed at consolidated funded level — "
                    "acquired portfolios do not carry broker data.")
            self.appendix.add(f"[{slide_spec.title}] {title}: {note}")
            path = Path(self._charts_out()) / f"{cid}_suppressed.png"
            render_placeholder_png(path, "", "Broker channel suppressed\n"
                                   "(consolidated funded lens)", theme=self.theme,
                                   width_in=img_w_in, height_in=img_h_in)
            self._place(slide, path, img_left, img_top, img_w_in, img_h_in)
            return ChartResult(chart_id=cid, title=title, path=path, available=False,
                               placeholder=True, kind="suppressed", note=note)

        # Forecast bridge waterfall — cross-lens, computed from resolved metrics.
        if cspec.get("type") == "bridge":
            return self._render_bridge(slide, cid, title, img_left, img_top,
                                       img_w_in, img_h_in)

        resolver = self.chart_resolvers.get(lens)
        if resolver is None:
            note = f"{lens.title()} lens data not available for this run."
            path = Path(self._charts_out()) / f"{lens}_{cid}_placeholder.png"
            render_placeholder_png(path, "", note, theme=self.theme,
                                   width_in=img_w_in, height_in=img_h_in)
            self._place(slide, path, img_left, img_top, img_w_in, img_h_in)
            return ChartResult(chart_id=cid, title=title, path=path, available=False,
                               placeholder=True, kind="lens_unavailable", note=note)

        result = resolver.resolve(cspec, img_w_in, img_h_in)
        result.title = title
        if result.path and Path(result.path).exists():
            self._place(slide, result.path, img_left, img_top, img_w_in, img_h_in)
        return result

    def _render_bridge(self, slide, cid, title, img_left, img_top, w_in, h_in):
        """Funded → +weighted pipeline → forecast funded bridge waterfall."""
        from .chart_resolver import render_bridge_waterfall
        funded = self.metrics.resolve(self.config.metric_spec("funded_balance"))
        forecast = self.metrics.resolve(self.config.metric_spec("forecast_funded_balance"))
        out = Path(self._charts_out()) / f"{cid}_bridge.png"
        if not (funded.ok and forecast.ok and forecast.value > funded.value):
            render_placeholder_png(out, "", "Forecast bridge requires funded + "
                                   "pipeline lenses.", theme=self.theme,
                                   width_in=w_in, height_in=h_in)
            self._place(slide, out, img_left, img_top, w_in, h_in)
            return ChartResult(chart_id=cid, title=title, path=out, available=False,
                               placeholder=True, kind="bridge",
                               note="Forecast bridge requires funded + pipeline data.")
        weighted = float(forecast.value) - float(funded.value)
        steps = [("Funded", float(funded.value), "base"),
                 ("+ Weighted Pipeline", weighted, "add"),
                 ("Forecast Funded", float(forecast.value), "total")]
        render_bridge_waterfall(out, steps, w_in, h_in, theme=self.theme)
        self._place(slide, out, img_left, img_top, w_in, h_in)
        return ChartResult(chart_id=cid, title=title, path=out, available=True,
                           placeholder=False, kind="bridge",
                           note="Registry forecast bridge.")

    def _charts_out(self) -> Path:
        for r in self.chart_resolvers.values():
            if r is not None:
                return r.out_dir
        p = Path(self.ctx.run_dir) / "_pptx_charts"
        p.mkdir(parents=True, exist_ok=True)
        return p

    def _place(self, slide, path, left, top, w_in, h_in) -> None:
        try:
            slide.shapes.add_picture(str(path), left, top,
                                     width=Inches(w_in), height=Inches(h_in))
        except Exception:
            pass

    # -------------------------------------------------------------- RISK MONITOR
    def _slide_risk_monitor(self, spec: SlideSpec) -> Dict[str, Any]:
        slide = self._new_slide()
        strap = self._strapline_for(spec)
        self._header(slide, spec.title or "Risk Monitor", strap)
        risk = (self.metrics.analytics or {}).get("risk_monitor")
        placeholder = False
        if isinstance(risk, dict) and risk.get("tests"):
            self._risk_tiles(slide, risk)
        else:
            placeholder = True
            path = Path(self._charts_out()) / "risk_monitor_placeholder.png"
            render_placeholder_png(path, "Risk Monitor",
                                   "Risk-limit artifacts not present for this run — "
                                   "concentration monitor pending", theme=self.theme,
                                   width_in=12.2, height_in=4.9)
            self._place(slide, path, Inches(0.55), Inches(1.62), 12.2, 4.9)
            self.appendix.add("Risk monitor rendered as placeholder: no risk-limit "
                              "artifact in run directory (v1 acceptable).")
        self._footer(slide)
        return self._record(spec, strap, placeholder=placeholder)

    def _risk_tiles(self, slide, risk) -> None:
        summary = risk.get("summary", {}) or {}
        tiles = [("Limits", summary.get("total", "—"), self.theme.peri),
                 ("Breaches", summary.get("breaches", 0), self.theme.rag["red"]),
                 ("Warnings", summary.get("warnings", 0), self.theme.rag["amber"]),
                 ("Within limit", summary.get("testsPassed", 0), self.theme.rag["green"])]
        tw = Inches(2.95)
        for i, (label, value, color) in enumerate(tiles):
            left = Emu(int(Inches(0.55)) + i * int(Inches(3.05)))
            self._panel(slide, left, Inches(1.75), tw, Inches(1.4),
                        fill=self.theme.bg_panel_alt, line=color, line_w=1.2)
            self._text(slide, left + Inches(0.2), Inches(1.98), tw, Inches(0.7),
                       str(value), size=26, bold=True, color=color)
            self._text(slide, left + Inches(0.2), Inches(2.78), tw, Inches(0.3),
                       label.upper(), size=9.5, color=self.theme.ink_400, bold=True)
        y = 3.55
        for line in (risk.get("observations", []) or [])[:5]:
            self._text(slide, Inches(0.6), Inches(y), Inches(12), Inches(0.35),
                       f"•  {line}", size=12, color=self.theme.ink_300)
            y += 0.42

    # -------------------------------------------------------------- METHODOLOGY
    def _slide_methodology(self, spec: SlideSpec) -> Dict[str, Any]:
        slide = self._new_slide()
        strap = self._strapline_for(spec)
        self._header(slide, spec.title or "Methodology & Notes", strap)
        lines = [
            f"Client:  {self.ctx.client_name}",
            f"Data cut-off date:  {self.ctx.as_of_date or 'n/a'}",
            f"Run directory:  {self.ctx.run_dir}",
            f"Lens basis:  {self.ctx.lens}"
            + (" (consolidated funded)" if self.ctx.consolidated else ""),
            "Balance basis:  current outstanding balance (registry-authorised).",
            "Forecast method:  deterministic registry bridge — funded + Σ(weighted pipeline).",
            "Aggregations:  MI Agent semantic registry + analytics_lib (config-driven buckets).",
        ]
        if self.ctx.prior_label:
            lines.append(f"Prior-period comparison:  {self.ctx.prior_label}.")
        if self.ctx.source_artifacts:
            lines.append("Source artifacts:")
            lines += [f"    – {a}" for a in self.ctx.source_artifacts[:8]]
        self._bullets(slide, lines, size=12.5)
        self._footer(slide)
        return self._record(spec, strap)

    # ------------------------------------------------------------------ APPENDIX
    def _slide_appendix(self, spec: SlideSpec) -> Dict[str, Any]:
        slide = self._new_slide()
        strap = self._strapline_for(spec)
        self._header(slide, spec.title or "Appendix — Data Coverage", strap)
        notes = self.appendix.as_list() or [
            "All configured metrics and charts resolved from the current pipeline "
            "run; no coverage gaps."]
        self._bullets(slide, [f"•  {n}" for n in notes[:12]], size=11)
        self._footer(slide)
        return self._record(spec, strap)

    def _bullets(self, slide, lines, *, size=12) -> None:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.62), Inches(12.1), Inches(5.2))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(size)
            run.font.name = self.theme.font_sans
            run.font.color.rgb = self._rgb(self.theme.ink_300)
            p.space_after = Pt(6)
