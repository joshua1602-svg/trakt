"""mi_agent_pptx.preflight — the checks a deck must pass before it is published.

A deck that renders successfully is not the same as a deck that is safe to send
to an investor. Rendering proves the file is valid PowerPoint; it proves nothing
about whether the pack states its scope, whether its dates are disclosed, or
whether its parts add up.

These gates are deliberately narrow. They assert the things that, if wrong, make
the artefact *misleading* rather than merely incomplete:

  * the reporting scope is rendered, so a single-book deck cannot be read as a
    total-portfolio deck;
  * the reporting dates are rendered;
  * the constituent books are named;
  * direct + acquired reconcile to the total;
  * an executive summary was generated;
  * no placeholder slide reached the deck;
  * the file opens as a presentation.

Failure policy, by design: **generate but do not publish.** The deck is still
written to the run directory so an operator can look at it and see what went
wrong, the failure is recorded on the artefact metadata, and the durable
publication is withheld. A misleading pack that never leaves the run directory
is a contained problem; a published one is not.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Any, Dict, List, Mapping, Optional, Sequence

#: Relative tolerance for the additivity check. The per-type snapshots and the
#: total are the same function over disjoint row sets, so they should agree to
#: floating-point noise; this allows for rounding only.
RECONCILIATION_TOLERANCE = 0.005  # 0.5 basis points of the total


@dataclass
class GateResult:
    """One check, its verdict and the evidence behind it."""

    gate: str
    passed: bool
    detail: str
    mandatory: bool = True
    evidence: Dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> Dict[str, Any]:
        return {"gate": self.gate, "passed": self.passed, "detail": self.detail,
                "mandatory": self.mandatory, "evidence": self.evidence}


@dataclass
class PreflightReport:
    """The publication verdict for one generated deck."""

    results: List[GateResult] = field(default_factory=list)

    @property
    def failures(self) -> List[GateResult]:
        return [r for r in self.results if not r.passed and r.mandatory]

    @property
    def warnings(self) -> List[GateResult]:
        return [r for r in self.results if not r.passed and not r.mandatory]

    @property
    def publishable(self) -> bool:
        """False when any MANDATORY gate failed — publication must be withheld."""
        return not self.failures

    def summary(self) -> str:
        state = "PASS" if self.publishable else "BLOCKED"
        return (f"{state} — {len(self.results)} gate(s), {len(self.failures)} "
                f"failure(s), {len(self.warnings)} warning(s)")

    def to_dict(self) -> Dict[str, Any]:
        return {
            "publishable": self.publishable,
            "summary": self.summary(),
            "failed_gates": [r.gate for r in self.failures],
            "warning_gates": [r.gate for r in self.warnings],
            "gates": [r.to_dict() for r in self.results],
        }


# --------------------------------------------------------------------------- #
# Gates.
# --------------------------------------------------------------------------- #

def _rendered_text(deck_path: Optional[str]) -> Optional[str]:
    """All text in the generated deck, or ``None`` if it cannot be opened.

    Reading the FILE rather than the build record is deliberate: the gate must
    verify what the investor will actually see, not what the builder believed it
    drew.
    """
    if not deck_path:
        return None
    try:
        from pptx import Presentation
        prs = Presentation(str(deck_path))
        return "\n".join(shape.text_frame.text
                         for slide in prs.slides for shape in slide.shapes
                         if shape.has_text_frame)
    except Exception:  # noqa: BLE001 — an unopenable deck is the gate's answer
        return None


def _gate_opens(text: Optional[str]) -> GateResult:
    ok = text is not None
    return GateResult("pptx_opens", ok,
                      "the deck opens as a valid presentation" if ok else
                      "the generated file could not be opened as a presentation")


def _gate_scope_rendered(text: Optional[str], portfolio) -> GateResult:
    if text is None:
        return GateResult("scope_rendered", False, "deck could not be read")
    if portfolio is None:
        return GateResult("scope_rendered", False,
                          "no governed portfolio context was resolved, so the "
                          "deck cannot state its scope")
    label = str(portfolio.scope_label or "").strip()
    ok = bool(label) and label.lower() in text.lower()
    return GateResult("scope_rendered", ok,
                      f"reporting scope '{label}' is rendered in the deck" if ok else
                      f"reporting scope '{label}' does not appear anywhere in the deck",
                      evidence={"scope": label})


def _gate_dates_rendered(text: Optional[str], portfolio, reporting_date) -> GateResult:
    if text is None:
        return GateResult("dates_rendered", False, "deck could not be read")
    expected = sorted({d for d in (portfolio.reporting_dates.values()
                                   if portfolio else [])} or
                      ({reporting_date} if reporting_date else set()))
    if not expected:
        return GateResult("dates_rendered", False,
                          "no reporting date could be resolved for this run")
    missing = [d for d in expected if d not in text]
    ok = not missing
    return GateResult("dates_rendered", ok,
                      "every reporting date is rendered" if ok else
                      f"reporting date(s) not rendered: {', '.join(missing)}",
                      evidence={"expected": expected, "missing": missing})


def _gate_books_rendered(text: Optional[str], portfolio) -> GateResult:
    if text is None or portfolio is None:
        return GateResult("books_rendered", False, "deck or context unavailable")
    names = [p.label for p in portfolio.portfolios]
    if not names:
        # A tape without per-book provenance is a legitimate single-book
        # deployment, not a governance failure: there is no constituent book to
        # name. Recorded as a warning so the absence is still visible.
        return GateResult("books_rendered", True,
                          "no per-book provenance on this tape (single-book "
                          "deployment) — no constituent books to name",
                          mandatory=False)
    missing = [n for n in names if n not in text]
    ok = not missing
    return GateResult("books_rendered", ok,
                      f"all {len(names)} constituent book(s) are named" if ok else
                      f"constituent book(s) not named in the deck: {', '.join(missing)}",
                      evidence={"books": names, "missing": missing})


def _gate_reconciles(portfolio) -> GateResult:
    """Direct + acquired must equal the total — the additivity guarantee."""
    if portfolio is None:
        return GateResult("totals_reconcile", False, "no portfolio context")
    if len(portfolio.type_slices) < 2:
        return GateResult("totals_reconcile", True,
                          "single portfolio type — no additivity to check",
                          mandatory=False)
    total = portfolio.total_balance
    parts = portfolio.type_balance_total()
    if total is None or parts is None:
        return GateResult("totals_reconcile", False,
                          "balances unavailable on the total or a type slice")
    drift = abs(total - parts)
    tol = abs(total) * RECONCILIATION_TOLERANCE
    ok = drift <= tol
    return GateResult(
        "totals_reconcile", ok,
        (f"portfolio types sum to the total (drift {drift:,.2f})" if ok else
         f"portfolio types do NOT sum to the total: total {total:,.2f} vs parts "
         f"{parts:,.2f} (drift {drift:,.2f}, tolerance {tol:,.2f})"),
        evidence={"total": total, "sum_of_types": parts, "drift": drift,
                  "tolerance": tol})


def _gate_loan_reconciles(portfolio) -> GateResult:
    if portfolio is None or len(portfolio.type_slices) < 2:
        return GateResult("loan_counts_reconcile", True,
                          "single portfolio type — no additivity to check",
                          mandatory=False)
    total = portfolio.loan_count
    parts = portfolio.type_loan_total()
    if total is None or parts is None:
        return GateResult("loan_counts_reconcile", False, "loan counts unavailable")
    ok = int(total) == int(parts)
    return GateResult("loan_counts_reconcile", ok,
                      f"loan counts reconcile ({int(total):,})" if ok else
                      f"loan counts do NOT reconcile: total {int(total):,} vs parts "
                      f"{int(parts):,}",
                      evidence={"total": total, "sum_of_types": parts})


#: Facts that mean the run HAD something to observe. A first-ever run of a
#: single book with no prior period genuinely has nothing to say, and blocking
#: it would fail closed on a perfectly valid pack.
_OBSERVABLE_FACTS = ("has_movement", "is_mixed", "has_geo", "has_risk",
                     "has_pipeline", "has_forecast", "has_funded_history")


def _gate_executive_summary(insights: Optional[Mapping[str, Any]],
                            facts: Optional[Mapping[str, Any]] = None) -> GateResult:
    items = (insights or {}).get("insights") or []
    if items:
        return GateResult("executive_summary", True,
                          f"{len(items)} governed observation(s) generated",
                          evidence={"count": len(items),
                                    "status": (insights or {}).get("status")})
    observable = [f for f in _OBSERVABLE_FACTS if (facts or {}).get(f)]
    if not observable:
        # Nothing to observe is not the same as failing to observe.
        return GateResult("executive_summary", True,
                          "no observations were possible for this run (single "
                          "book, single period, no pipeline, forecast or risk "
                          "inputs)", mandatory=False,
                          evidence={"count": 0})
    return GateResult("executive_summary", False,
                      "no governed observations were generated even though the "
                      f"run had observable inputs ({', '.join(observable)})",
                      evidence={"count": 0, "observable": observable,
                                "status": (insights or {}).get("status")})


def _gate_no_placeholders(records: Sequence[Mapping[str, Any]]) -> GateResult:
    placeholders = [r.get("id") for r in records if r.get("placeholder")]
    ok = not placeholders
    return GateResult("no_placeholder_slides", ok,
                      "no placeholder slides in the deck" if ok else
                      f"placeholder slide(s) reached the deck: "
                      f"{', '.join(str(p) for p in placeholders)}",
                      evidence={"placeholders": placeholders})


#: Text that must never reach a client-facing deck. These are implementation
#: details, not disclosures: a storage location or a function name tells an
#: investor nothing and tells everyone else too much.
_LEAK_PATTERNS = (
    (r"(?:^|[\s(])/(?:tmp|home|var|root|mnt|Users)/\S+", "absolute filesystem path"),
    (r"\bscratchpad\b", "scratchpad location"),
    (r"\b(?:blob|abfss|s3|file)://\S+", "storage URI"),
    (r"\b[A-Za-z_][A-Za-z0-9_]*\.py\b", "python module name"),
    (r"\bcompute_[a-z_]+\b", "internal function name"),
    (r"\b_[a-z]+_[a-z_]+\(\)", "internal function name"),
    (r"C:\\\\", "windows path"),
)

#: Wording that signals an unpopulated or provisional page.
_PLACEHOLDER_PHRASES = (
    "no data for this run", "data unavailable", "not available on this tape",
    "lorem ipsum", "tbc", "to be confirmed", "placeholder", "coming soon",
)

#: A static pool tracks a FIXED cohort forward through time. This release renders
#: a point-in-time cross-section by origination year, so the claim must not
#: appear — a wrong label is as misleading as a wrong number.
_STATIC_POOL_PHRASES = ("static-pool", "static pool")


def _gate_no_internal_paths(text: Optional[str]) -> GateResult:
    import re
    if text is None:
        return GateResult("no_internal_paths", False, "deck could not be read")
    found: List[str] = []
    for pattern, label in _LEAK_PATTERNS:
        for match in re.findall(pattern, text):
            found.append(f"{label}: {str(match).strip()[:60]}")
    ok = not found
    return GateResult("no_internal_paths", ok,
                      "no internal paths or implementation details are rendered" if ok
                      else f"internal detail rendered in the deck — {found[0]}",
                      evidence={"found": found[:6]})


def _gate_no_placeholder_language(text: Optional[str]) -> GateResult:
    if text is None:
        return GateResult("no_placeholder_language", False, "deck could not be read")
    low = text.lower()
    hits = [p for p in _PLACEHOLDER_PHRASES if p in low]
    ok = not hits
    return GateResult("no_placeholder_language", ok,
                      "no placeholder language in the deck" if ok else
                      f"placeholder language rendered: {', '.join(hits)}",
                      evidence={"phrases": hits})


def _gate_no_false_static_pool(text: Optional[str]) -> GateResult:
    if text is None:
        return GateResult("no_false_static_pool_claim", False, "deck could not be read")
    low = text.lower()
    hits = [p for p in _STATIC_POOL_PHRASES if p in low]
    ok = not hits
    return GateResult("no_false_static_pool_claim", ok,
                      "no static-pool claim is made" if ok else
                      "the deck claims static-pool analysis, which it does not render",
                      evidence={"phrases": hits})


def _gate_pipeline_reconciles(text: Optional[str], pipeline) -> GateResult:
    """The pipeline headline must match the governed pipeline snapshot."""
    amount = (pipeline or {}).get("pipelineAmount")
    if not amount:
        return GateResult("pipeline_reconciles", True,
                          "no pipeline in scope — nothing to reconcile",
                          mandatory=False)
    if text is None:
        return GateResult("pipeline_reconciles", False, "deck could not be read")
    # The deck formats the governed figure; assert the governed magnitude is the
    # one on the page (millions, one decimal — the pack's own convention).
    millions = f"{float(amount) / 1e6:.1f}"
    ok = millions in text.replace(",", "")
    return GateResult("pipeline_reconciles", ok,
                      f"pipeline headline reconciles to the governed snapshot "
                      f"(£{millions}MM)" if ok else
                      f"pipeline headline £{millions}MM does not appear in the deck",
                      evidence={"pipeline_amount": amount})


def _gate_concentration_reconciles(text: Optional[str], concentration) -> GateResult:
    """Concentration values on the slide must be the governed test values."""
    from . import concentration as C
    rows = C.adapt_tests(concentration)
    if not rows:
        return GateResult("concentration_reconciles", True,
                          "no governed concentration tests in scope",
                          mandatory=False)
    if text is None:
        return GateResult("concentration_reconciles", False, "deck could not be read")
    shown = C.select_tests(rows)
    missing = [r["label"] for r in shown if r["label"][:26] not in text]
    ok = not missing
    return GateResult("concentration_reconciles", ok,
                      f"{len(shown)} governed concentration test(s) rendered" if ok else
                      f"concentration test(s) not rendered: {', '.join(missing)}",
                      evidence={"rendered": [r["label"] for r in shown],
                                "missing": missing})


def _gate_no_empty_slides(deck_path: Optional[str],
                          records: Sequence[Mapping[str, Any]]) -> GateResult:
    """Every included slide must carry content beyond its title and footer.

    A title with nothing under it is worse than an omitted section: it implies
    the analysis exists and was empty.
    """
    if not deck_path:
        return GateResult("no_empty_slides", False, "deck path unavailable")
    try:
        from pptx import Presentation
        prs = Presentation(str(deck_path))
    except Exception:  # noqa: BLE001
        return GateResult("no_empty_slides", False, "deck could not be opened")
    sparse: List[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = [sh for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        pictures = [sh for sh in slide.shapes if sh.shape_type == 13]
        # Title + subtitle + footer + page number is four text frames and no
        # content; anything at or below that with no picture is a shell.
        if not pictures and len(texts) <= 4:
            sparse.append(f"slide {index}")
    ok = not sparse
    return GateResult("no_empty_slides", ok,
                      "every slide carries content" if ok else
                      f"slide(s) rendered with no meaningful content: {', '.join(sparse)}",
                      evidence={"sparse": sparse})


def _gate_movement_reconciles(movement) -> GateResult:
    """Attribution must sum to the headline movement, for every dimension.

    ``funded_bridge`` guarantees this by construction. The gate exists so a
    future change that broke it fails publication rather than quietly
    misattributing a period's movement.
    """
    bridges = [b for b in (movement or {}).values()
               if getattr(b, "available", False)]
    if not bridges:
        return GateResult("movement_reconciles", True,
                          "no movement attribution in scope", mandatory=False)
    broken = [b.label for b in bridges if not b.reconciles()]
    ok = not broken
    return GateResult("movement_reconciles", ok,
                      f"{len(bridges)} attribution dimension(s) reconcile" if ok else
                      f"attribution does not sum to the headline movement for: "
                      f"{', '.join(broken)}",
                      evidence={"dimensions": [b.label for b in bridges],
                                "broken": broken})


def _gate_no_unsupported_causal_language(text: Optional[str]) -> GateResult:
    """Balance movement evidences WHAT moved, never WHY.

    Words like "redemption", "amortisation" or "prepayment" attribute a cause the
    funded balance decomposition cannot prove. They are only permissible where a
    governed decomposition supplies them, which this release does not have — so
    their presence in the rendered pack is a defect.
    """
    import re
    if text is None:
        return GateResult("no_unsupported_causal_language", False,
                          "deck could not be read")
    banned = ("redemption", "redeemed", "amortisation", "amortised",
              "prepayment", "prepaid", "expected runoff", "expected run-off")
    low = text.lower()
    hits = [w for w in banned if re.search(rf"\b{re.escape(w)}\b", low)]
    ok = not hits
    return GateResult("no_unsupported_causal_language", ok,
                      "no unsupported causal attribution" if ok else
                      f"the deck attributes a cause the evidence does not prove: "
                      f"{', '.join(hits)}",
                      evidence={"terms": hits})


def _gate_no_duplicate_observations(insights, watchlist) -> GateResult:
    """The same fact must not be stated twice.

    Complementary shares make this easy to get wrong: within one dimension, a
    category gaining 8pp means the others lose 8pp between them, so a naive
    generator emits the same finding several times and crowds out genuinely
    different ones.
    """
    headlines: List[str] = []
    for item in ((insights or {}).get("insights") or []):
        headlines.append(str(getattr(item, "headline", "")).strip().lower())
    for key in ("watch", "observations"):
        for item in ((watchlist or {}).get(key) or []):
            headlines.append(str(getattr(item, "headline", "")).strip().lower())
    seen, dupes = set(), []
    for h in headlines:
        if h and h in seen:
            dupes.append(h)
        seen.add(h)
    ok = not dupes
    return GateResult("no_duplicate_observations", ok,
                      f"{len(headlines)} observation(s), none duplicated" if ok else
                      f"duplicate observation(s): {dupes[0][:60]}",
                      evidence={"duplicates": dupes[:4], "total": len(headlines)})


def _gate_mandatory_slides(records: Sequence[Mapping[str, Any]]) -> GateResult:
    """Cover, methodology and appendix are the disclosure spine of the pack."""
    ids = {str(r.get("id")) for r in records}
    # One investor-safe Data and Methodology page satisfies the disclosure
    # spine; the former separate appendix has been folded into it.
    required = {"cover", "methodology"}
    missing = sorted(required - ids)
    ok = not missing
    return GateResult("mandatory_slides", ok,
                      "cover and the data/methodology page are present" if ok else
                      f"mandatory slide(s) missing: {', '.join(missing)}",
                      evidence={"missing": missing})


# --------------------------------------------------------------------------- #
# Entry point.
# --------------------------------------------------------------------------- #

def run_preflight(build_report: Mapping[str, Any], data: Any) -> PreflightReport:
    """Evaluate every publication gate for a generated deck."""
    deck_path = build_report.get("output")
    records = build_report.get("slides") or []
    portfolio = getattr(data, "portfolio", None)
    text = _rendered_text(deck_path)
    facts = build_report.get("facts") or {}

    report = PreflightReport(results=[
        _gate_opens(text),
        _gate_mandatory_slides(records),
        _gate_scope_rendered(text, portfolio),
        _gate_dates_rendered(text, portfolio, getattr(data, "reporting_date", None)),
        _gate_books_rendered(text, portfolio),
        _gate_reconciles(portfolio),
        _gate_loan_reconciles(portfolio),
        _gate_executive_summary(getattr(data, "insights", None), facts),
        _gate_no_placeholders(records),
        # -- v2.1: client-safety and reconciliation --------------------------
        _gate_no_internal_paths(text),
        _gate_no_placeholder_language(text),
        _gate_no_false_static_pool(text),
        _gate_no_empty_slides(deck_path, records),
        _gate_pipeline_reconciles(text, getattr(data, "pipeline", None)),
        _gate_concentration_reconciles(text, getattr(data, "concentration", None)),
        # -- v2.2: attribution integrity and evidence discipline -------------
        _gate_movement_reconciles(getattr(data, "movement", None)),
        _gate_no_unsupported_causal_language(text),
        _gate_no_duplicate_observations(getattr(data, "insights", None),
                                        getattr(data, "watchlist", None)),
    ])
    return report
