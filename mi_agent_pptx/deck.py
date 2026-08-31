"""mi_agent_pptx.deck — payload-driven deck assembly (dashboard-aligned).

Renders each slide directly from the MI API payloads (:mod:`mi_api`) so the pack
is a faithful export of the React dashboard: the Executive Summary is the funded
KPI tile grid, stratifications are the funded BarLists, the pipeline slide is the
pipeline snapshot, the forecast slide is the funded→forecast bridge, evolution
slides are the time series, geography/cohorts/risk mirror their tabs. Numbers are
taken verbatim from the payloads — never recomputed here.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .chart_resolver import render_bridge_waterfall
from .mi_api import DashboardData
from .placeholders import render_placeholder_png
from .pptx_theme import PptxTheme, THEME
from . import render as R

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
EMU_IN = 914400

_MONTHS = ("January", "February", "March", "April", "May", "June", "July",
           "August", "September", "October", "November", "December")


def _pretty_date(value: Any) -> str:
    """``2026-06-30`` -> ``30 June 2026``, for a client-facing page."""
    text = str(value or "").strip()
    parts = text.split("-")
    if len(parts) == 3 and len(parts[0]) == 4:
        try:
            return f"{int(parts[2])} {_MONTHS[int(parts[1]) - 1]} {parts[0]}"
        except (ValueError, IndexError):
            return text
    return text


@dataclass
class DeckContext:
    client_name: str
    as_of_date: str = ""
    run_dir: str = ""
    generated_by: str = "trakt MI Agent"
    footer: str = "trakt MI Agent · Confidential — for institutional funders/investors"
    deck_name: str = "Investor & Funder MI Pack"
    work_dir: str = ""
    logo_path: Optional[str] = None


class DeckBuilder:
    def __init__(self, data: DashboardData, ctx: DeckContext,
                 theme: PptxTheme = THEME):
        self.d = data
        self.ctx = ctx
        self.theme = theme
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._blank = self.prs.slide_layouts[6]
        self._page = 0
        self.work = Path(ctx.work_dir or (Path(ctx.run_dir) / "_pptx_charts"))
        self.work.mkdir(parents=True, exist_ok=True)
        self.appendix: List[str] = list(data.notes)
        #: What each renderer actually drew (see render.record_renders).
        self.rendered: List[Dict[str, Any]] = []
        self.records: List[Dict[str, Any]] = []
        #: Slides this portfolio did not justify, with reasons (rendered in the
        #: appendix so an omission is never silent).
        self.omissions: List[Any] = []
        self.facts: Dict[str, Any] = {}

    # ------------------------------------------------------------- pptx scaffold
    def _rgb(self, hx):
        r, g, b = self.theme.rgb(hx)
        return RGBColor(r, g, b)

    def _slide(self):
        s = self.prs.slides.add_slide(self._blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = self._rgb(self.theme.bg_page)
        return s

    def _text(self, slide, l, t, w, h, text, *, size=14, color=None, bold=False,
              align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, spacing=None):
        box = slide.shapes.add_textbox(l, t, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, 0)
        p = tf.paragraphs[0]
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = self.theme.font_sans
        run.font.color.rgb = self._rgb(color or self.theme.ink_100)
        return box

    def _panel(self, slide, l, t, w, h, *, fill=None, line=None, radius=True, lw=1.0):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h)
        try:
            shp.adjustments[0] = 0.045
        except Exception:
            pass
        shp.fill.solid()
        shp.fill.fore_color.rgb = self._rgb(fill or self.theme.bg_panel)
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = self._rgb(line)
            shp.line.width = Pt(lw)
        shp.shadow.inherit = False
        return shp

    def _header(self, slide, title, strap, *, accent=None):
        rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                      Inches(0.12), SLIDE_H)
        rail.fill.solid()
        rail.fill.fore_color.rgb = self._rgb(accent or self.theme.peri)
        rail.line.fill.background()
        rail.shadow.inherit = False
        self._text(slide, Inches(0.55), Inches(0.34), Inches(12.2), Inches(0.6),
                   title, size=25, bold=True)
        if strap:
            self._text(slide, Inches(0.57), Inches(1.0), Inches(12.4), Inches(0.5),
                       strap, size=12, color=self.theme.peri, italic=True)

    def scope_footnote(self) -> str:
        """The one-line scope + date stamp every slide carries.

        A slide read on its own — screenshotted, pasted into a memo — must still
        say which book it describes and as at when.
        """
        p = self.d.portfolio
        date = self.ctx.as_of_date or self.d.reporting_date
        if p is None:
            return f"Funded as at {date}" if date else ""
        label = f"{p.scope_label} portfolio"
        if p.has_mixed_reporting_dates:
            return f"{label} · mixed reporting dates (see cover)"
        return f"{label} · funded as at {date}" if date else label

    def _footer(self, slide):
        self._page += 1
        self._text(slide, Inches(0.55), Inches(7.08), Inches(6.6), Inches(0.3),
                   self.ctx.footer, size=8, color=self.theme.ink_500)
        stamp = self.scope_footnote()
        if stamp:
            self._text(slide, Inches(7.25), Inches(7.08), Inches(4.95), Inches(0.3),
                       stamp, size=8, color=self.theme.ink_400,
                       align=PP_ALIGN.RIGHT)
        self._text(slide, Inches(12.3), Inches(7.08), Inches(0.8), Inches(0.3),
                   str(self._page), size=8, color=self.theme.ink_500,
                   align=PP_ALIGN.RIGHT)

    def _place(self, slide, path, l, t, w_in, h_in):
        try:
            slide.shapes.add_picture(str(path), l, t, width=Inches(w_in),
                                     height=Inches(h_in))
        except Exception:
            pass

    def _card(self, slide, l, t, w, h, title):
        """A dashboard-style card: panel + title, returns the inner image box."""
        self._panel(slide, l, t, w, h, fill=self.theme.bg_panel, line=self.theme.line)
        self._text(slide, l + Inches(0.22), t + Inches(0.16), w - Inches(0.4),
                   Inches(0.34), title, size=12.5, bold=True)
        img_l = l + Inches(0.16)
        img_t = t + Inches(0.62)
        img_w = (int(w) - 2 * int(Inches(0.16))) / EMU_IN
        img_h = (int(h) - int(Inches(0.62)) - int(Inches(0.16))) / EMU_IN
        return img_l, img_t, img_w, img_h

    # ------------------------------------------------------------------- tiles
    #: Approximate widths, in ems, of the characters a KPI value is made of.
    #: A digit is not a full em and a full stop is not a digit; treating them
    #: as equal is what made a currency figure look too wide for its tile.
    _EM = {".": 0.30, ",": 0.30, " ": 0.30, "1": 0.55, "%": 0.95, "M": 0.90,
           "W": 0.90, "+": 0.60, "-": 0.35, "—": 0.90, "−": 0.60, "/": 0.35}

    @classmethod
    def _text_width_in(cls, text: str, size_pt: float, *, bold: bool = True) -> float:
        """Roughly how wide ``text`` draws at ``size_pt``, in inches."""
        ems = 0.0
        for ch in str(text or ""):
            if ch in cls._EM:
                ems += cls._EM[ch]
            elif ch.isdigit() or ch in "£$€":
                ems += 0.62
            elif ch.isupper():
                ems += 0.70
            else:
                ems += 0.55
        return ems * (size_pt / 72.0) * (1.06 if bold else 1.0)

    def _tile(self, slide, l, t, w, h, tile: Dict[str, Any]):
        self._panel(slide, l, t, w, h, fill=self.theme.bg_panel_alt,
                    line=self.theme.line_soft, lw=1.0)
        pad = Inches(0.16)
        iw = Emu(int(w) - 2 * int(pad))
        avail = bool(tile.get("available", True)) and tile.get("value") not in (None, "")
        self._text(slide, l + pad, t + Inches(0.14), iw, Inches(0.3),
                   str(tile.get("label", "")).upper(), size=8.5,
                   color=self.theme.ink_400, bold=True)
        val = str(tile.get("value") if avail else "—")
        # A KPI value may be a long label (an area name), not just a number. Step
        # the size down so it fits the tile instead of being clipped — a value the
        # reader cannot see is worse than a smaller one.
        #
        # Sized on the WIDTH THE STRING ACTUALLY DRAWS, not on how many
        # characters it has. Counting characters made "£78.4MM" — seven
        # characters, four of them narrow — step down twice, so the funded
        # balance printed at 12pt beside a 20pt "48.3%": the most important
        # number in the pack rendered as the least important thing on the page.
        width_in = int(iw) / EMU_IN
        size = 10
        for candidate in (20, 15, 12):
            if self._text_width_in(val, candidate) <= width_in * 0.98:
                size = candidate
                break
        self._text(slide, l + pad, t + Inches(0.44), iw, Inches(0.58), val,
                   size=size, bold=True,
                   color=self.theme.ink_100 if avail else self.theme.ink_500)
        y = t + Inches(1.02)
        delta, intent = tile.get("delta"), tile.get("deltaIntent")
        # THE MEASURE'S BASIS, on the measure. Two tiles on one page with
        # different weighting bases are only honest if each says which it used —
        # a reader who divides an unweighted mean by a balance-weighted one and
        # lands 15 points from the LTV tile has been misled by the layout, not
        # by any single number.
        #
        # It takes the HINT's slot rather than a line of its own: the tile has
        # room for one sub-line, and where a measure has both they say the same
        # thing ("balance-weighted current valuation" under "balance-weighted").
        # A delta still outranks both — a movement is news, a basis is a
        # standing property.
        basis = tile.get("basis")
        if delta:
            color = {"positive": self.theme.mint, "negative": self.theme.rose}.get(
                intent, self.theme.ink_400)
            arrow = {"positive": "▲ ", "negative": "▼ "}.get(intent, "")
            self._text(slide, l + pad, y, iw, Inches(0.3), f"{arrow}{delta}",
                       size=9.5, color=color, bold=True)
        elif basis:
            self._text(slide, l + pad, y, iw, Inches(0.3), str(basis),
                       size=7.5, color=self.theme.ink_500, italic=True)
        elif tile.get("hint"):
            self._text(slide, l + pad, y, iw, Inches(0.3), str(tile["hint"]),
                       size=9, color=self.theme.ink_400)

    def _tile_grid(self, slide, tiles: List[Dict[str, Any]], *, top=1.62, cols=5,
                   row_height: Optional[float] = None) -> float:
        """Lay out KPI tiles; returns the bottom edge of the block, in inches.

        ``row_height`` lets a slide that has to fit something BENEATH the tiles
        choose a more compact row rather than discovering afterwards that there
        is no room left for it.
        """
        rows = max(1, (len(tiles) + cols - 1) // cols)
        gx, gy = Inches(0.16), Inches(0.22)
        left0, top0 = Inches(0.55), Inches(top)
        tile_w = Emu(int((int(Inches(12.25)) - (cols - 1) * int(gx)) / cols))
        height_in = row_height if row_height is not None else (1.62 if rows <= 2 else 1.3)
        tile_h = Inches(height_in)
        for i, tile in enumerate(tiles):
            r, c = divmod(i, cols)
            l = Emu(int(left0) + c * (int(tile_w) + int(gx)))
            t = Emu(int(top0) + r * (int(tile_h) + int(gy)))
            self._tile(slide, l, t, tile_w, tile_h, tile)
        return top + rows * height_in + (rows - 1) * 0.22

    #: THE content band. Every tile row and every chart panel is laid out across
    #: these two edges, so a KPI strip and the charts beneath it cannot drift
    #: apart — which is exactly what happened when tiles used a 3.0in pitch and
    #: charts a 6.78in origin: the last tile ended at 12.47in and the right chart
    #: at 12.78in, a third of an inch of visible misalignment.
    CONTENT_L = 0.55
    CONTENT_R = 12.78

    #: The governed health checks the watchlist runs, named on the page. A
    #: reader cannot tell "nothing was flagged" from "nothing was checked"
    #: unless the pack says which checks ran, so the same list is printed
    #: whether or not any of them cleared its materiality threshold.
    GOVERNED_CHECKS = (
        "Concentration limits — current, expected and stress",
        "Reporting-date consistency across constituent books",
        "Portfolio-type balance movement",
        "Composition shift by region, channel, LTV and ticket band",
        "Weighted-average LTV movement",
        "Reporting-dimension coverage",
    )
    COLUMN_GAP = 0.28

    def _grid(self, n: int, *, gap: Optional[float] = None):
        """``[(left_in, width_in)]`` for *n* equal columns across the band."""
        gap = self.COLUMN_GAP if gap is None else gap
        span = self.CONTENT_R - self.CONTENT_L
        w = (span - gap * max(n - 1, 0)) / max(n, 1)
        return [(self.CONTENT_L + i * (w + gap), w) for i in range(n)]

    def _chart_boxes(self, n, *, top: float = 1.62, height: float = 4.95):
        return [(Inches(l), Inches(top), Inches(w), Inches(height))
                for l, w in self._grid(min(max(n, 1), 3))]

    #: Where the executive slide's risk strip sits, and the clearance the charts
    #: above it must leave. Without the clearance the trend card's border landed
    #: on the strip, and the one line a reader takes off the page read as part of
    #: the chart.
    RISK_STRIP_TOP = 6.58
    RISK_STRIP_CLEARANCE = 0.26

    #: Headline tiles on the executive page. One row, always full: a partial
    #: second row leaves a hole beside it, and the tiles are in priority order
    #: so filling the row keeps the measures that matter. Six is the widest a
    #: compact currency value still reads at across the content width.
    EXEC_MAX_TILES = 6

    #: Vertical room one bar-list row needs to stay readable, in inches. Derived
    #: from the renderer: the label and the mono value are ~9pt, and below this
    #: they begin to collide with the rows above and beneath them.
    ROW_PITCH_IN = 0.19

    @staticmethod
    def _fit_bars(rows, capacity: int, *, value_key: str = "balance"):
        """``rows`` reduced to ``capacity`` bars WITHOUT losing any value.

        NEVER DROP A ROW SILENTLY. Truncating the list left six bars under a
        heading a reader adds up, £11.8m short of the total they close on, with
        nothing on the page saying a row had been cut. The remainder is
        aggregated instead, so the bars still account for the whole.
        """
        if capacity < 2 or len(rows) <= capacity:
            return list(rows)
        rest = rows[capacity - 1:]
        return list(rows[:capacity - 1]) + [{
            "label": f"Other ({len(rest)})",
            value_key: sum(float(r.get(value_key) or 0.0) for r in rest)}]

    def _barlist_capacity(self, height_in: float, *, minimum: int = 3) -> int:
        """How many bars a panel of this height can carry legibly."""
        usable = max(0.0, height_in - 0.42)          # card title + padding
        return max(minimum, int(usable / self.ROW_PITCH_IN))

    def _matrix_boxes(self, n, *, top: float = 1.62, height: float = 4.95,
                      row_gap: float = 0.22):
        """Boxes for a 2 x 2 matrix of panels (or one row, for one or two).

        A four-panel matrix is the deck's standard stratification grammar: four
        governed dimensions, one visual language, readable at a glance. Falls
        back to a single row for fewer panels so the same handler draws both.
        """
        if n <= 2:
            return self._chart_boxes(n, top=top, height=height)
        cols = self._grid(2)
        rows = 2 if n <= 4 else (n + 1) // 2
        panel_h = (height - row_gap * (rows - 1)) / rows
        boxes = []
        for i in range(n):
            r, c = divmod(i, 2)
            left, width = cols[c]
            boxes.append((Inches(left), Inches(top + r * (panel_h + row_gap)),
                          Inches(width), Inches(panel_h)))
        return boxes

    #: A tile carrying a hint line needs this much height: the hint is drawn at
    #: 1.02in from the tile top and stands 0.30in tall, so anything shorter puts
    #: it outside its own panel and onto whatever the slide drew underneath.
    #: Enforced here rather than left to each caller, because two slides shipped
    #: with a short strip and a hint before anyone read the printed page.
    HINTED_TILE_HEIGHT = 1.34

    def _strip(self, tiles, *, top: float = 1.58, height: float = 1.45):
        """A KPI row on the same grid the charts use.

        Returns parallel sequences so a caller can ``zip`` them with its tiles;
        the columns are the grid's, never a per-slide constant.
        """
        if any((t or {}).get("hint") for t in tiles):
            height = max(height, self.HINTED_TILE_HEIGHT)
        cols = self._grid(len(tiles))
        return ([Inches(l) for l, _w in cols],
                [Inches(top)] * len(tiles),
                [Inches(w) for _l, w in cols],
                [Inches(height)] * len(tiles),
                list(tiles))

    #: A dimension is worth a panel when its balance is not all in one bucket.
    #: 99.5% rather than 100% because a handful of loans in a second band does
    #: not make a distribution either.
    SPREAD_FLOOR = 0.995

    @classmethod
    def _has_spread(cls, strat) -> bool:
        """Does this stratification actually distribute across its buckets?"""
        bars = [b for b in (strat or {}).get("bars") or ()
                if isinstance(b, dict)]
        values = [abs(float(b.get("balance") or 0.0)) for b in bars]
        total = sum(values)
        if len(values) < 2 or total <= 0:
            return len(values) > 1
        return (max(values) / total) < cls.SPREAD_FLOOR

    def _barlist_card(self, slide, box, title, rows, value_key, *, currency=True,
                      cid="bl", label_key="label", dimension=None):
        il, it, iw, ih = self._card(slide, *box, title)
        path = self.work / f"{cid}.png"
        if rows:
            R.draw_barlist(path, rows, value_key, iw, ih, theme=self.theme,
                           currency=currency, label_key=label_key,
                           chart_id=cid, dimension=dimension)
        else:
            render_placeholder_png(path, "", "No data for this run",
                                   theme=self.theme, width_in=iw, height_in=ih)
        self._place(slide, path, il, it, iw, ih)
        return bool(rows)

    # =====================================================================
    # SLIDE HANDLERS
    # =====================================================================
    def _record(self, sid, title, strap, *, placeholder=False):
        self.records.append({"id": sid, "title": title, "strapline": strap,
                             "placeholder": placeholder})

    def slide_cover(self, spec):
        """Cover — states, unambiguously, what this report is a report ABOUT.

        Scope, constituent books and every reporting date are rendered from the
        governed portfolio context, never from an operator-typed name, so a deck
        covering one book can never be mistaken for a total-portfolio deck.
        """
        s = self._slide()
        # Full-height left accent rail (replaces the old overlapping corner panel).
        rail = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                  Inches(0.16), SLIDE_H)
        rail.fill.solid()
        rail.fill.fore_color.rgb = self._rgb(self.theme.peri)
        rail.line.fill.background()
        rail.shadow.inherit = False
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.92),
                                 Inches(2.2), Inches(0.07))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(self.theme.peri)
        bar.line.fill.background()
        bar.shadow.inherit = False
        self._text(s, Inches(0.9), Inches(0.62), Inches(6), Inches(0.4),
                   "TRAKT · MI AGENT", size=12, color=self.theme.peri, bold=True)
        self._text(s, Inches(0.86), Inches(1.42), Inches(11.5), Inches(1.1),
                   self._entity_name(), size=38, bold=True)
        self._text(s, Inches(0.92), Inches(3.12), Inches(11), Inches(0.5),
                   self.ctx.deck_name, size=18, color=self.theme.peri)
        strap = self._cover_strapline()
        self._text(s, Inches(0.92), Inches(3.7), Inches(10.5), Inches(0.6),
                   strap, size=12.5, color=self.theme.ink_300, italic=True,
                   spacing=1.12)
        self._cover_scope_block(s)
        self._text(s, Inches(0.92), Inches(6.72), Inches(8), Inches(0.35),
                   f"Generated automatically by {self.ctx.generated_by}", size=10,
                   color=self.theme.ink_400)
        self._record("cover", self._entity_name(), strap)

    def _entity_name(self):
        """The reporting entity — the operator-supplied client name."""
        return self.ctx.client_name

    def _cover_scope_block(self, s):
        """Reporting scope + constituent books + every reporting date."""
        p = self.d.portfolio
        left, top = Inches(0.92), Inches(4.42)
        self._text(s, left, top, Inches(5.6), Inches(0.3), "REPORTING SCOPE",
                   size=9, color=self.theme.peri, bold=True)
        if p is None:
            self._text(s, left, top + Inches(0.3), Inches(6), Inches(0.4),
                       "Scope unavailable — no governed portfolio context resolved.",
                       size=11.5, color=self.theme.ink_300)
            return
        self._text(s, left, top + Inches(0.28), Inches(5.6), Inches(0.4),
                   f"{p.scope_label} portfolio", size=15, bold=True)
        # Constituent books, each with its type.
        y = top + Inches(0.68)
        for book in p.portfolios[:5]:
            self._text(s, left, y, Inches(5.6), Inches(0.3),
                       f"·  {book.label}  —  {book.type_label}", size=10.5,
                       color=self.theme.ink_300)
            y = Emu(int(y) + int(Inches(0.245)))
        if len(p.portfolios) > 5:
            self._text(s, left, y, Inches(5.6), Inches(0.3),
                       f"·  and {len(p.portfolios) - 5} further book(s)", size=10.5,
                       color=self.theme.ink_400)

        # Reporting dates — per type when they differ, else one line.
        rleft = Inches(7.1)
        self._text(s, rleft, top, Inches(5.3), Inches(0.3), "REPORTING DATES",
                   size=9, color=self.theme.peri, bold=True)
        ry = top + Inches(0.28)
        dates = p.type_reporting_dates
        if dates:
            for ptype, date in sorted(dates.items()):
                from .deck_context import type_label
                self._text(s, rleft, ry, Inches(5.3), Inches(0.32),
                           f"{type_label(ptype)}:  {date}", size=11.5,
                           color=self.theme.ink_100)
                ry = Emu(int(ry) + int(Inches(0.3)))
        else:
            self._text(s, rleft, ry, Inches(5.3), Inches(0.32),
                       f"Funded portfolio:  {self.ctx.as_of_date or 'n/a'}",
                       size=11.5)
            ry = Emu(int(ry) + int(Inches(0.3)))
        if p.has_mixed_reporting_dates:
            self._text(s, rleft, Emu(int(ry) + int(Inches(0.06))), Inches(5.3),
                       Inches(0.5),
                       "⚠  Constituent books are reported as at different dates; "
                       "the total combines them.",
                       size=9.5, color=self.theme.amber, italic=True)

    def _cover_strapline(self):
        p = self.d.portfolio
        kpis = {k.get("id"): k for k in self.d.funded.get("kpis", [])}
        bal = kpis.get("balance", {}).get("value")
        loans = kpis.get("loans", {}).get("value")
        ltv = kpis.get("wa_current_ltv", {}).get("value")
        if bal and p is not None and p.is_mixed:
            return (f"Funded book of {bal} across {loans} loans, reported as a total "
                    f"and split by portfolio type.")
        if bal:
            return (f"Funded book of {bal} across {loans} loans at {ltv} weighted "
                    f"current LTV.")
        return "Automated MI pack generated from the latest pipeline run."

    def slide_kpi_summary(self, spec):
        s = self._slide()
        self._header(s, spec.get("title", "Executive Summary"),
                     "Funded book snapshot" + (
                         f" · reporting {self.d.reporting_date}" if self.d.reporting_date else ""))
        tiles = list(self.d.funded.get("kpis", []))[:10]
        if not tiles:
            self._placeholder_body(s, "Funded book unavailable for this run.")
            self._footer(s)
            return self._record("executive_summary", spec.get("title"), "", placeholder=True)
        self._tile_grid(s, tiles, top=1.62, cols=5)
        self._footer(s)
        self._record("executive_summary", spec.get("title", "Executive Summary"),
                     "Funded KPIs (dashboard-aligned).")

    # ------------------------------------------------- executive dashboard
    def slide_executive(self, spec):
        """Slide 1 — where is the portfolio today, what is coming, and is
        anything approaching a limit?

        The three lenses on one page. Every figure is lifted from a governed
        payload that a later slide also renders, so the landing page can never
        disagree with the pack behind it: the funded tiles ARE the funded
        snapshot's KPI tiles, the pipeline tiles ARE the pipeline snapshot's, the
        forecast tile IS the forecast bridge's, and the risk strip IS the
        concentration evaluator's summary.

        There is no React executive landing page to mirror — the dashboard opens
        on the funded lens of a tabbed workspace. This composition is therefore
        new, and it is assembled from ``DashboardData`` alone (no extra compute
        call) precisely so it can be offered back to React as one payload later.
        """
        from .metric_resolver import compact_currency, compact_number

        s = self._slide()
        funded = self.d.funded or {}
        pipeline = self.d.pipeline or {}
        bridge = (self.d.forecast or {}).get("forecastBridge") or {}
        by_id = {k.get("id"): k for k in funded.get("kpis", [])}

        rd = self.d.reporting_date
        self._header(s, spec.get("title", "Executive Position"),
                     ("Funded, pipeline and forecast as at " + _pretty_date(rd))
                     if rd else "Funded, pipeline and forecast",
                     accent=self.theme.peri)

        def kpi_tile(kpi_id, label=None):
            """A funded KPI tile, verbatim from the governed snapshot."""
            k = by_id.get(kpi_id)
            if not k:
                return None
            return {"label": label or k.get("label"), "value": k.get("value"),
                    "delta": k.get("delta"), "deltaIntent": k.get("deltaIntent"),
                    "hint": None if k.get("delta") else k.get("hint"),
                    "available": k.get("available", True)}

        # SIX MEASURES, IN PRIORITY ORDER, AND NOTHING TO FILL A GRID.
        # The page used to carry every governed headline it could reach — seven
        # tiles, laid out four and three, with a hole beside the second row, and
        # three of them (pipeline balance, weighted expected, forecast funded)
        # restating one fact the Executive Summary then restated again in words
        # on the next page. What an opening page owes a reader is the position,
        # what is coming, and whether anything needs attention.
        tiles = [tile for tile in (
            kpi_tile("balance", "Funded balance"),
            kpi_tile("loans", "Loans funded"),
            kpi_tile("wa_current_ltv", "WA current LTV"),
        ) if tile]

        # PIPELINE — what is coming. Week-on-week deltas come from the governed
        # prior-week aggregates, never computed here.
        if pipeline:
            prior = pipeline.get("priorWeek") or {}
            amount = pipeline.get("pipelineAmount")
            cases = pipeline.get("pipelineRowCount")

            def wow(current, previous):
                if previous is None or current is None:
                    return None, None
                diff = float(current) - float(previous)
                intent = ("positive" if diff > 0 else
                          "negative" if diff < 0 else "neutral")
                return compact_currency(diff) + " vs prior wk", intent

            delta, intent = wow(amount, prior.get("pipelineAmount"))
            # ONE pipeline tile, carrying the case count as its hint. Two tiles
            # for one lens crowded out the forecast and the risk measure.
            tiles.append({"label": "Pipeline balance",
                          "value": compact_currency(amount),
                          "delta": delta, "deltaIntent": intent,
                          "hint": (None if delta else
                                   f"{compact_number(cases)} live cases"
                                   if cases else "current weekly extract")})

        # FORECAST — where the current book plus its pipeline lands. The
        # weighted-pipeline component is the forecast bridge's own subject and
        # is stated there; here the reader needs the destination.
        # A forecast that equals the funded balance is not a forecast — it is
        # the funded balance again, in a tile that claims to look forward. That
        # is what a book with no pipeline produced, and it is the same
        # duplication this page was redesigned to remove.
        forecast_balance = bridge.get("forecastFundedBalance")
        weighted = bridge.get("weightedExpectedFundedAmount") or 0.0
        if forecast_balance and weighted:
            tiles.append({"label": "Forecast funded",
                          "value": compact_currency(forecast_balance),
                          "hint": "funded + weighted pipeline"})

        # RISK, AS A MEASURE RATHER THAN ONLY A SENTENCE. The strip at the foot
        # of the page names the closest test; a reader scanning tiles should be
        # able to see there IS a limit position without reading a line of prose.
        tiles.append(self._headroom_tile())

        # Time to the nearest scale target the run-rate has not yet passed —
        # last, so it is the tile that gives way when the page is full.
        milestone = self._next_milestone()
        if milestone:
            tiles.append({"label": f"Time to {milestone['label']}",
                          "value": milestone["value"],
                          "hint": milestone.get("hint")})
        tiles = [t for t in tiles if t]

        if not tiles:
            self._placeholder_body(s, "No governed measures resolved for this run.")
            self._footer(s)
            return self._record("executive", spec.get("title"), "", placeholder=True)

        # ONE ROW, ALWAYS FULL. Seven tiles over four columns left a hole beside
        # the second row — the visual gap that made this page read as a
        # compressed web dashboard rather than an opening statement. The tiles
        # are in priority order, so filling the row keeps the measures that
        # matter and drops the ones that were there to fill a grid. Six is the
        # widest a compact currency value still reads at.
        shown = tiles[:self.EXEC_MAX_TILES]
        tiles_bottom = self._tile_grid(s, shown, top=1.58, cols=len(shown),
                                       row_height=1.44)
        self._executive_trends(s, top=tiles_bottom + 0.30)

        # Risk, last: the one line that says whether anything needs attention.
        self._executive_risk_strip(s)
        self._footer(s)
        self._record("executive", spec.get("title"),
                     "Funded, pipeline, forecast and risk on one page.")

    def _headroom_tile(self):
        """The closest approved limit, as a tile.

        Straight from ``concentration.summarise`` — the same evaluator the risk
        strip and the Concentration slide read, so the three cannot disagree.
        ``None`` where no operator-approved configuration is in force: an empty
        tile saying "no limits" would spend the page's scarcest space on an
        absence the strip below already states.
        """
        from . import concentration as C

        env = self.d.concentration or {}
        rows = C.adapt_tests(env)
        if not rows:
            return None
        summary = C.summarise(env, rows)
        closest = summary.get("closest")
        if not closest or closest.get("utilisation") is None:
            return None
        breaches, warnings = summary["breaches"], summary["warnings"]
        intent = ("negative" if breaches else
                  "neutral" if warnings else "positive")
        return {
            "label": "Closest limit",
            "value": f"{closest['utilisation']:.0f}%",
            "hint": f"{closest['label']} utilisation",
            "deltaIntent": intent,
        }

    def _next_milestone(self):
        """The nearest configured scale target the book has not yet reached.

        Straight from the governed extrapolation ladder
        (``forecast_extrapolation.build_extrapolation``) — the same milestones the
        Time to Scale slide tabulates. No projection is performed here.
        """
        extrap = self.d.extrapolation or {}
        for model_key in ("completionRunRateForecast", "kfiConversionForecast"):
            model = extrap.get(model_key) or {}
            if not model.get("available"):
                continue
            for row in model.get("milestones") or ():
                base = row.get("base") or row.get("expected") or {}
                period = base.get("period") if isinstance(base, dict) else None
                if not period:
                    for value in row.values():
                        if isinstance(value, dict) and value.get("period"):
                            period = value["period"]
                            break
                if period:
                    return {"label": row.get("thresholdLabel", "target"),
                            "value": str(period),
                            "hint": "central run-rate scenario"}
        return None

    @staticmethod
    def _period_labels(periods) -> List[str]:
        """The category axis for an evolution series.

        A WEEKLY series carries both ``week`` (the extract date) and ``period``
        (its calendar month). Labelling four weekly extracts by month drew an
        axis reading 2026-05, 2026-05, 2026-06, 2026-06 — four gridlines, two
        labels, each printed twice, which reads as a rendering fault and hides
        that the series has four observations. Where a distinct ``week`` exists
        on every point it is the label; otherwise the monthly period is.
        """
        weeks = [str(p.get("week") or "") for p in periods]
        if all(weeks) and len(set(weeks)) == len(weeks):
            return weeks
        return [str(p.get("period") or p.get("reporting_date") or p.get("run_id"))
                for p in periods]

    def _executive_trends(self, slide, *, top: float):
        """Up to two compact trends: funded balance, and weighted pipeline.

        Both are governed evolution series already resolved for later slides. A
        series with fewer than two periods is simply omitted — a single point is
        not a trend, and an empty chart frame on the landing page is worse than
        one fewer chart.
        """
        # The band is whatever is left between the tiles and the risk strip, and
        # all of it is used. A 2.30in cap left three quarters of an inch of dead
        # panel under the chart on a page whose whole job is to look deliberate.
        height = self.RISK_STRIP_TOP - self.RISK_STRIP_CLEARANCE - top
        if height < 1.4:
            return

        # ONE TRAJECTORY, FULL WIDTH. Two half-width trends competed for the
        # centre of the opening page and neither was legible enough to be the
        # thing a reader takes from it. The portfolio's own trajectory is the
        # visual that answers "where is this going"; the pipeline's own
        # trajectory is the subject of the pipeline pages, and its destination
        # is already on this page as the forecast tile.
        funded_evo = (self.d.funded_evolution or {}).get("periods") or []
        pipe_evo = (self.d.pipeline_evolution or {}).get("periods") or []
        if len(funded_evo) >= 2:
            cid, title, periods, metric = (
                "exec_funded", "Funded balance by period", funded_evo,
                "funded_balance")
        elif len(pipe_evo) >= 2:
            # A book with no funded history yet still has a story, and it is the
            # origination one.
            cid, title, periods, metric = (
                "exec_pipeline", "Weighted expected pipeline by week", pipe_evo,
                "weighted_expected_funded_amount")
        else:
            return

        values = [(p.get("metrics") or {}).get(metric) for p in periods]
        if sum(1 for v in values if v is not None) < 2:
            return
        box = (Inches(self.CONTENT_L), Inches(top),
               Inches(self.CONTENT_R - self.CONTENT_L), Inches(height))
        il, it, iw, ih = self._card(slide, *box, title)
        path = self.work / f"{cid}.png"
        R.draw_lines(path, self._period_labels(periods),
                     [{"name": title, "values": values}], iw, ih,
                     theme=self.theme, currency=True, area=True, chart_id=cid)
        self._place(slide, path, il, it, iw, ih)

    def _executive_risk_strip(self, slide):
        """One line on limits, from the governed concentration evaluator.

        The counts and the closest-to-breaching test are ``concentration.summarise``
        — the same figures the Concentration slide tabulates — so the landing page
        and the risk page cannot disagree. "No approved configuration" is itself a
        finding and is stated plainly rather than left blank.
        """
        from . import concentration as C

        env = self.d.concentration or {}
        rows = C.adapt_tests(env)
        top = self.RISK_STRIP_TOP
        width = Inches(self.CONTENT_R - self.CONTENT_L)
        if not rows:
            self._text(slide, Inches(self.CONTENT_L), Inches(top), width,
                       Inches(0.30),
                       "Concentration — no operator-approved limit configuration "
                       "is in force for this portfolio.",
                       size=10, color=self.theme.ink_400, italic=True)
            return

        summary = C.summarise(env, rows)
        breaches, warnings = summary["breaches"], summary["warnings"]
        within = max(summary["tests"] - breaches - warnings, 0)
        parts = [f"{within} within limit"]
        if warnings:
            parts.append(f"{warnings} approaching")
        if breaches:
            parts.append(f"{breaches} in breach")
        line = "Concentration — " + ", ".join(parts)

        # The tile above already names the closest test and its utilisation.
        # Repeating both here spends the page's one risk line restating a tile;
        # what the tile cannot carry is the DISTANCE to the limit, so that is
        # what this adds.
        closest = summary.get("closest")
        if closest and closest.get("utilisation") is not None:
            headroom = closest.get("headroom")
            if headroom is not None:
                line += (f". {C.format_headroom(headroom, closest.get('unit'))} "
                         f"of headroom on the closest, {closest['label']}")
            else:
                line += (f". Closest to its limit: {closest['label']} at "
                         f"{closest['utilisation']:.0f}% utilisation")
        if summary.get("expected_breaches"):
            line += f". {summary['expected_breaches']} forecast to breach"

        worst = (C.STATUS_BREACH if breaches else
                 C.STATUS_WARNING if warnings else C.STATUS_PASS)
        colour = {C.STATUS_BREACH: self.theme.rose,
                  C.STATUS_WARNING: self.theme.amber}.get(worst, self.theme.mint)
        self._panel(slide, Inches(self.CONTENT_L), Inches(top), width, Inches(0.40),
                    fill=self.theme.bg_panel_alt, line=self.theme.line_soft)
        self._text(slide, Inches(self.CONTENT_L + 0.18), Inches(top + 0.06),
                   Emu(int(width) - int(Inches(0.36))), Inches(0.28),
                   self._fit_label(line + ".", self.CONTENT_R - self.CONTENT_L - 0.4, 10),
                   size=10, color=colour, bold=True)

    # ------------------------------------------------- investor narrative
    def slide_exec_insights(self, spec):
        """Executive Summary — *what changed, and why*.

        Deterministic observations from :mod:`mi_agent_pptx.insights`. Every
        sentence is a template over a governed figure; nothing is generated.
        """
        s = self._slide()
        brief = self.d.insights or {}
        items = brief.get("insights") or []
        self._header(s, spec.get("title", "Executive Summary"),
                     "Governed observations for the period")
        if not items:
            self._placeholder_body(s, "No governed observations for this run.")
            self._footer(s)
            return self._record(spec.get("id", "executive_summary"),
                                spec.get("title"), "", placeholder=True)

        # Severity accent per observation — a caveat must not read as an aside.
        accent_for = {"concern": self.theme.rose, "attention": self.theme.amber}
        items = items[:8]
        top = 1.58
        # The card grid must FIT. Previously the row height was fixed, so a
        # seventh observation was laid out below the bottom of the slide and the
        # reader simply never saw it. Height is now derived from the available
        # band and the number of rows, so the grid always fits by construction.
        two_col = len(items) > 4
        col_w = Inches(6.0) if two_col else Inches(12.25)
        per_col = -(-len(items) // 2) if two_col else len(items)
        gap = 0.12
        band = 6.92 - top                      # bottom of the content area
        # Cards grow into the available band rather than sitting at a fixed
        # height: with two observations a 0.92in card clipped a long summary
        # while leaving most of the slide empty.
        row_h = max(0.62, min(1.62 if two_col else 1.55,
                              (band - gap * (per_col - 1)) / max(per_col, 1)))
        for i, ins in enumerate(items):
            col, row = (i // per_col, i % per_col) if two_col else (0, i)
            l = Inches(0.55) if col == 0 else Inches(6.78)
            t = Inches(top + row * (row_h + gap))
            h = Inches(row_h)
            accent = accent_for.get(getattr(ins, "severity", "info"), self.theme.peri)
            self._panel(s, l, t, col_w, h, fill=self.theme.bg_panel_alt,
                        line=self.theme.line_soft)
            chip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.045), h)
            chip.fill.solid()
            chip.fill.fore_color.rgb = self._rgb(accent)
            chip.line.fill.background()
            chip.shadow.inherit = False
            # The headline's box is sized to the lines it actually needs, and the
            # body starts below it. At a fixed 0.42in a two-line headline ran
            # straight into the summary beneath — which happens as soon as a
            # movement names two contributors, so it is not an edge case.
            headline = str(getattr(ins, "headline", ""))
            text_w = (col_w - Inches(0.46)) / EMU_IN
            per_line = max(16, int(text_w * 72 / (12.5 * 0.52)))
            head_lines = min(3, max(1, -(-len(headline) // per_line)))
            head_h = 0.26 + 0.24 * head_lines
            self._text(s, l + Inches(0.24), t + Inches(0.14),
                       col_w - Inches(0.46), Inches(head_h),
                       headline, size=12.5, bold=True)
            # Body size steps down when a summary is long, so a detailed
            # observation is never clipped by a fixed card height.
            summary = str(getattr(ins, "summary", ""))
            body_top = 0.14 + head_h + 0.04
            body_h = max(0.24, (h / EMU_IN) - body_top - 0.14)
            capacity = text_w * body_h * 210
            # Four tiers, not three: a runoff book narrates exits as well as
            # balances, and its observations are long enough that the old floor
            # of 8pt still overflowed a two-column card.
            body = 10 if len(summary) <= capacity else (
                9 if len(summary) <= capacity * 1.25 else
                8 if len(summary) <= capacity * 1.55 else 7)
            self._text(s, l + Inches(0.24), t + Inches(body_top),
                       col_w - Inches(0.46), Inches(body_h), summary,
                       size=body, color=self.theme.ink_300, spacing=1.06)
        self._footer(s)
        self._record(spec.get("id", "executive_summary"), spec.get("title"),
                     f"{len(items)} governed observations.")

    def slide_portfolio_composition(self, spec):
        """Portfolio Composition — *what do I own today, and how is it divided?*

        The anchor slide, and the second page an investor sees. It used to be a
        label column beside one value column per type — functionally complete and
        visually a spreadsheet. The measures are unchanged; the hierarchy is not:
        a summary strip, then the split itself as a proportional bar, then one
        card per portfolio type leading on the three figures that decide whether
        the reader keeps reading (balance, share, movement) with the rest beneath.

        Every portfolio type comes from the governed registry, so this renders a
        single book, a direct/acquired pair, or more than two types without any
        of them being named here.
        """
        s = self._slide()
        p = self.d.portfolio
        self._header(s, spec.get("title", "Portfolio Composition"),
                     "Funded book by portfolio type")
        if p is None:
            self._placeholder_body(s, "No governed portfolio context resolved.")
            self._footer(s)
            return self._record(spec.get("id", "portfolio_composition"),
                                spec.get("title"), "", placeholder=True)
        from .metric_resolver import compact_currency, compact_number

        total_bal = p.total_balance or 0.0
        slices = list(p.type_slices)

        # -- 1. summary strip -------------------------------------------------
        tiles = [
            {"label": "Total funded portfolio", "value": compact_currency(total_bal)},
            {"label": "Loans", "value": compact_number(p.loan_count or 0)},
            {"label": "Constituent books", "value": compact_number(p.portfolio_count),
             "hint": f"{len(slices) or 1} portfolio type"
                     f"{'s' if (len(slices) or 1) != 1 else ''}"},
            {"label": "Reporting date",
             "value": _pretty_date(self.d.reporting_date) or "—",
             "hint": ("dates differ by book" if p.has_mixed_reporting_dates
                      else "aligned across books")},
        ]
        for l, t, w, h, tile in zip(*self._strip(tiles, top=1.56, height=1.10)):
            self._tile(s, l, t, w, h, tile)

        if not slices:
            self._text(s, Inches(self.CONTENT_L), Inches(3.1),
                       Inches(self.CONTENT_R - self.CONTENT_L), Inches(0.5),
                       "Single-portfolio book — no type split applies.", size=12,
                       color=self.theme.ink_400, italic=True)
            self._footer(s)
            return self._record(spec.get("id", "portfolio_composition"),
                                spec.get("title"), "Single portfolio.")

        # -- 2. the split, as one proportional bar ----------------------------
        self._composition_bar(s, slices, p, top=3.06)

        # -- 3. one card per portfolio type -----------------------------------
        # Cards, not columns: with more than two types a shared-row table forces
        # every measure into a sliver, while cards simply get narrower and keep
        # their internal hierarchy.
        lead = [
            ("Balance", lambda sl: compact_currency(sl.balance)),
            ("Share", lambda sl: (f"{p.share_of(sl) * 100:.1f}%"
                                  if p.share_of(sl) is not None else "—")),
            ("Movement", lambda sl: self._signed_currency(sl.balance_movement)),
        ]
        rest = [
            ("Loans", lambda sl: compact_number(sl.loan_count or 0)),
            ("Average balance", lambda sl: compact_currency(sl.avg_balance)),
            ("WA current LTV", lambda sl: self._pct_display(sl, "wa_current_ltv")),
            ("WA interest rate", lambda sl: self._pct_display(sl, "wa_rate")),
            ("WA borrower age", lambda sl: self._pct_display(sl, "wa_age")),
            ("Loan movement", lambda sl: (f"{int(sl.loan_movement):+,}"
                                          if sl.loan_movement is not None else "—")),
            ("Reporting date", lambda sl: _pretty_date(sl.reporting_date) or "—"),
        ]
        top = 3.66
        height = 2.94
        for (l, w), sl in zip(self._grid(len(slices)), slices):
            self._type_card(s, sl, lead, rest, Inches(l), Inches(top),
                            Inches(w), Inches(height))
        self._footer(s)
        self._record(spec.get("id", "portfolio_composition"), spec.get("title"),
                     f"{len(slices)} portfolio type(s).")

    def _composition_bar(self, s, slices, ctx, *, top: float):
        """The split as ONE proportional bar.

        A restrained institutional visual rather than a donut: segment length is
        the only encoding, it reads left to right like the numbers beneath it,
        and it degrades to a single full-width segment for a one-type book.
        """
        left = self.CONTENT_L
        width = self.CONTENT_R - self.CONTENT_L
        colours = [self.theme.peri, self.theme.mint, self.theme.ink_400,
                   self.theme.rag.get("amber", self.theme.ink_300)]
        x = left
        for i, sl in enumerate(slices):
            # From the governed composition service. An equal split is the
            # fallback ONLY when there is no total to divide by.
            share = ctx.share_of(sl)
            if share is None:
                share = 1.0 / len(slices)
            seg = max(width * share, 0.06)     # a sliver must still be visible
            bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                                     Inches(top), Inches(seg), Inches(0.30))
            bar.fill.solid()
            bar.fill.fore_color.rgb = self._rgb(colours[i % len(colours)])
            bar.line.fill.background()
            bar.shadow.inherit = False
            # The label only goes INSIDE a segment wide enough to hold it;
            # otherwise the legend row beneath carries it.
            if seg >= 1.5:
                self._text(s, Inches(x + 0.12), Inches(top + 0.02),
                           Inches(seg - 0.24), Inches(0.26),
                           f"{sl.label} · {share * 100:.1f}%", size=9.5,
                           bold=True, color=self.theme.navy)
            x += seg
        # A legend row ONLY for the segments too narrow to carry their own
        # label. Listing every type again under a bar that already names them
        # spends a line of the slide saying the same thing twice.
        unlabelled = [(i, sl) for i, sl in enumerate(slices)
                      if width * (ctx.share_of(sl) if ctx.share_of(sl) is not None
                                  else 1.0) < 1.5]
        lx = left
        for i, sl in unlabelled:
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(lx), Inches(top + 0.40),
                                     Inches(0.10), Inches(0.10))
            dot.fill.solid()
            dot.fill.fore_color.rgb = self._rgb(colours[i % len(colours)])
            dot.line.fill.background()
            dot.shadow.inherit = False
            slice_share = ctx.share_of(sl)
            label = (f"{sl.label} — {slice_share * 100:.1f}%" if slice_share is not None
                     else str(sl.label))
            self._text(s, Inches(lx + 0.18), Inches(top + 0.34), Inches(3.2),
                       Inches(0.24), label, size=9, color=self.theme.ink_400)
            lx += 3.4

    def _type_card(self, s, sl, lead, rest, l, t, w, h):
        """One portfolio type: three headline measures, then the rest."""
        self._panel(s, l, t, w, h, fill=self.theme.bg_panel, line=self.theme.line)
        self._text(s, l + Inches(0.20), t + Inches(0.16), w - Inches(0.4),
                   Inches(0.30), sl.label.upper(), size=9.5, bold=True,
                   color=self.theme.peri)

        inner = (w / EMU_IN) - 0.40
        for i, (label, fn) in enumerate(lead):
            cw = inner / len(lead)
            lx = l + Inches(0.20 + i * cw)
            self._text(s, lx, t + Inches(0.56), Inches(cw - 0.08), Inches(0.24),
                       label.upper(), size=8, color=self.theme.ink_500, bold=True)
            self._text(s, lx, t + Inches(0.80), Inches(cw - 0.08), Inches(0.42),
                       str(self._cell(fn, sl)), size=15 if len(lead) <= 3 else 13,
                       bold=True)

        band = (h / EMU_IN) - 1.44
        row_h = min(0.30, band / max(len(rest), 1))
        size = 10 if row_h >= 0.27 else 9
        for i, (label, fn) in enumerate(rest):
            y = t + Inches(1.36 + i * row_h)
            self._text(s, l + Inches(0.20), y, Inches(inner * 0.58), Inches(0.26),
                       label, size=size, color=self.theme.ink_400)
            self._text(s, l + Inches(0.20 + inner * 0.56), y,
                       Inches(inner * 0.44), Inches(0.26),
                       str(self._cell(fn, sl)), size=size, bold=True,
                       align=PP_ALIGN.RIGHT)

    @staticmethod
    def _cell(fn, sl):
        try:
            return fn(sl)
        except Exception:  # noqa: BLE001 — one cell must not break the slide
            return "—"

    def _pct_display(self, sl, kpi_id):
        """A weighted-average percentage as the governed snapshot formatted it."""
        return sl.display(kpi_id) or "—"

    @staticmethod
    def _signed_currency(value):
        """A movement with an explicit sign — ``+£851K`` / ``-£296K``."""
        from .metric_resolver import compact_currency
        if value is None:
            return "—"
        text = compact_currency(value)
        return f"+{text}" if float(value) > 0 else text

    def slide_portfolio_comparison(self, spec):
        """Direct vs Acquired — *why did the total move?*

        Movement attribution plus the mix differences that a blended total hides.
        """
        s = self._slide()
        p = self.d.portfolio
        self._header(s, spec.get("title", "Direct vs Acquired"),
                     "Period movement attribution and portfolio differences")
        if p is None or len(p.type_slices) < 2:
            self._placeholder_body(s, "Only one portfolio type is in scope.")
            self._footer(s)
            return self._record(spec.get("id", "portfolio_comparison"),
                                spec.get("title"), "", placeholder=True)
        from .metric_resolver import compact_currency

        slices = list(p.type_slices)
        # Left: movement attribution waterfall (total = Σ type movements).
        il, it, iw, ih = self._card(s, Inches(0.55), Inches(1.58), Inches(6.0),
                                    Inches(4.98), "What moved the total")
        movers = [(sl, sl.balance_movement) for sl in slices
                  if sl.balance_movement is not None]
        path = self.work / "cmp_attrib.png"
        if movers:
            # ``total - Σ movements`` is the opening the waterfall must reach.
            # Both attribution slides derived it; one definition now serves both.
            from mi_agent_api.portfolio_context import opening_from_movement
            opening = opening_from_movement(p.total_balance,
                                            sum(v for _s, v in movers)) or 0.0
            steps = [("Opening", float(opening), "base")]
            for sl, v in movers:
                steps.append((sl.label.replace(" portfolio", "").replace(
                    " originations", ""), float(v), "add"))
            steps.append(("Closing", float(p.total_balance or 0), "total"))
            render_bridge_waterfall(path, steps, iw, ih, theme=self.theme)
        else:
            render_placeholder_png(path, "", "No prior reporting period to "
                                   "attribute movement against", theme=self.theme,
                                   width_in=iw, height_in=ih)
        self._place(s, path, il, it, iw, ih)

        # Right: side-by-side differences on identical governed measures.
        il, it, iw, ih = self._card(s, Inches(6.78), Inches(1.58), Inches(6.02),
                                    Inches(4.98), "How the books differ")
        cols = ["Measure"] + [sl.label for sl in slices]
        measures = [
            ("Funded balance", lambda sl: compact_currency(sl.balance)),
            ("Loans", lambda sl: f"{int(sl.loan_count or 0):,}"),
            ("Average balance", lambda sl: compact_currency(sl.avg_balance)),
            ("WA current LTV", lambda sl: self._pct_display(sl, "wa_current_ltv")),
            ("WA interest rate", lambda sl: self._pct_display(sl, "wa_rate")),
            ("WA borrower age", lambda sl: self._pct_display(sl, "wa_age")),
            ("Period movement", lambda sl: self._signed_currency(sl.balance_movement)),
            ("Loan movement", lambda sl: (f"{int(sl.loan_movement):+,}"
                                          if sl.loan_movement is not None else "—")),
        ]
        trows = []
        for label, fn in measures:
            row = [label]
            for sl in slices:
                try:
                    row.append(str(fn(sl)))
                except Exception:  # noqa: BLE001
                    row.append("—")
            trows.append(row)
        tpath = self.work / "cmp_table.png"
        R.draw_table(tpath, cols, trows, iw, ih, theme=self.theme)
        self._place(s, tpath, il, it, iw, ih)
        self._footer(s)
        self._record(spec.get("id", "portfolio_comparison"), spec.get("title"),
                     "Movement attribution by portfolio type.")

    def slide_movement_drivers(self, spec):
        """Portfolio Movement and Drivers — *why did funded AuM change?*

        Two governed views side by side: the movement by portfolio type (who
        moved it) and the movement by a chosen dimension (where it moved). Both
        reconcile exactly to the headline, because both come from governed
        decompositions whose parts sum to the whole by construction.
        """
        from . import movement as MV
        from .metric_resolver import compact_currency

        s = self._slide()
        p = self.d.portfolio
        bridges = self.d.movement or {}
        primary = next((bridges[k] for k in ("region", "broker", "ltv", "ticket")
                        if k in bridges and bridges[k].available), None)
        self._header(s, spec.get("title", "Portfolio Movement and Drivers"),
                     self._movement_window(primary), accent=self.theme.peri)

        # Left: movement by portfolio type (the attribution that matters most).
        il, it, iw, ih = self._card(s, Inches(0.55), Inches(1.62), Inches(6.02),
                                    Inches(3.62), "Movement by portfolio type")
        movers = [(sl, sl.balance_movement) for sl in (p.type_slices if p else ())
                  if sl.balance_movement is not None]
        path = self.work / "mv_type.png"
        if movers:
            # ``total - Σ movements`` is the opening the waterfall must reach.
            # Both attribution slides derived it; one definition now serves both.
            from mi_agent_api.portfolio_context import opening_from_movement
            opening = opening_from_movement(p.total_balance,
                                            sum(v for _s, v in movers)) or 0.0
            steps = [("Opening", float(opening), "base")]
            for sl, v in movers:
                steps.append((sl.label.replace(" portfolio", "")
                              .replace(" originations", ""), float(v), "add"))
            steps.append(("Closing", float(p.total_balance or 0), "total"))
            render_bridge_waterfall(path, steps, iw, ih, theme=self.theme)
        elif primary is not None:
            steps = [("Opening", float(primary.opening or 0), "base"),
                     ("Net change", float(primary.total_delta or 0), "add"),
                     ("Closing", float(primary.closing or 0), "total")]
            render_bridge_waterfall(path, steps, iw, ih, theme=self.theme)
        self._place(s, path, il, it, iw, ih)

        # Right: movement by the leading governed dimension.
        if primary is not None:
            il, it, iw, ih = self._card(s, Inches(6.78), Inches(1.62), Inches(6.02),
                                        Inches(3.62), f"Movement by {primary.label.lower()}")
            rows = [{"category": c.category, "delta": c.delta, "is_other": c.is_other}
                    for c in sorted(primary.contributors,
                                    key=lambda c: -abs(c.delta))[:7]]
            dpath = self.work / "mv_dim.png"
            R.draw_diverging(dpath, rows, iw, ih, theme=self.theme)
            self._place(s, dpath, il, it, iw, ih)

        # Deterministic takeaways across the governed dimensions.
        lines: List[str] = []
        if p is not None and len(p.type_slices) > 1 and movers:
            ups = [(sl, v) for sl, v in movers if v > 0]
            downs = [(sl, v) for sl, v in movers if v < 0]
            if ups and downs:
                lines.append(
                    f"{ups[0][0].label} added {compact_currency(abs(ups[0][1]))}, "
                    f"partly offset by a {compact_currency(abs(downs[0][1]))} "
                    f"reduction in the {downs[0][0].label.lower()}.")
        for key in ("region", "broker", "ticket", "ltv"):
            b = bridges.get(key)
            if b is not None and b.available:
                head = MV.headline(b)
                if head:
                    lines.append(head)
            if len(lines) >= 3:
                break
        self._takeaway_strip(s, lines[:3], top=5.42)
        self._footer(s)
        self._record(spec.get("id", "movement_drivers"), spec.get("title"),
                     "Governed movement attribution.")

    def _movement_window(self, bridge) -> str:
        """The period a movement was measured over — never left implicit."""
        if bridge is None or not getattr(bridge, "available", False):
            return "Period movement attribution"
        if bridge.start_date and bridge.end_date:
            return (f"Movement from {_pretty_date(bridge.start_date)} to "
                    f"{_pretty_date(bridge.end_date)}")
        if bridge.start_period and bridge.end_period:
            return f"Movement from {bridge.start_period} to {bridge.end_period}"
        return "Period movement attribution"

    def _takeaway_strip(self, slide, lines, *, top: float, width: float = 12.25):
        """The slide's conclusions — the sentence an investor actually keeps."""
        lines = [l for l in lines if l]
        if not lines:
            return
        self._panel(slide, Inches(0.55), Inches(top), Inches(width),
                    Inches(0.34 + 0.26 * len(lines)),
                    fill=self.theme.bg_panel_alt, line=self.theme.line_soft)
        for i, line in enumerate(lines):
            self._text(slide, Inches(0.78), Inches(top + 0.14 + i * 0.26),
                       Inches(width - 0.46), Inches(0.26), f"·  {line}",
                       size=10, color=self.theme.ink_300)

    def slide_strat(self, spec):
        """Stratifications — composition, PLUS what changed.

        The composition view is legitimate investor content and is kept. What it
        could not answer is whether the shape of the book is moving, so where a
        comparable prior period exists the slide pairs each dimension with its
        governed marginal change and states the movement in words.
        """
        from . import movement as MV

        s = self._slide()
        # LENS. ``funded`` (default) reads the funded snapshot; ``pipeline``
        # reads the governed pipeline stratifications, which carry the SAME
        # bands from the SAME bucket registry. One handler, one visual grammar,
        # so a pipeline LTV band and a funded LTV band are read the same way.
        lens = spec.get("lens", "funded")
        source = self.d.pipeline if lens == "pipeline" else self.d.funded
        strats = (source or {}).get("stratifications", []) or []
        # ``keys`` IS A PREFERENCE, NOT A FILTER. It names the cuts a reader asks
        # for first; where one of them has nothing to distribute on this book —
        # pipeline broker/channel at 100% Direct is the standing example — the
        # slot goes to the next most informative dimension the book DOES
        # support, rather than to a panel with one full-width bar in it.
        #
        # The judgement is the shared one in ``mi_agent_api.presentation``, so a
        # dimension React declines to chart and a dimension the pack declines to
        # draw are the same dimension.
        from mi_agent_api import presentation as _sel

        strats = [st for st in strats if st.get("bars")]
        chosen = _sel.select_dimensions(strats, want=4, value_key="balance",
                                        preferred=tuple(spec.get("keys") or ()))
        rejected = chosen["rejected"]
        # On a book where EVERY dimension is concentrated, that is the finding —
        # and an empty page is not an improvement on a flat one.
        strats = chosen["selected"] or strats[:4]
        flat = [st for st in strats if not self._has_spread(st)]

        # Movement attribution is a FUNDED concept: it compares two funded
        # reporting periods. A pipeline stratification has no such bridge.
        bridges = (self.d.movement or {}) if lens == "funded" else {}
        moved = [bridges[st.get("key")] for st in strats
                 if bridges.get(st.get("key")) is not None
                 and bridges[st.get("key")].available]
        window = self._movement_window(moved[0]) if moved else "Balance by dimension"
        default_strap = ("Pipeline balance by dimension" if lens == "pipeline"
                         else "Balance by dimension")

        # A 2 x 2 matrix leaves no room for a takeaway strip; two panels do.
        has_takeaways = bool(moved) and len(strats) <= 2
        # THE STRAPLINE DESCRIBES THE PAGE, NOT THE DATA BEHIND IT. A four-panel
        # matrix suppresses the movement strip for room, and a subtitle that
        # still promised "period movement" would send a reader hunting the page
        # for a view that is not on it. Movement is available either way — it is
        # decomposed on Funded Balance Movement — so the honest strapline here
        # is the one that names what these panels actually show.
        self._header(s, spec.get("title", "Stratifications"),
                     ("Composition and period movement" if has_takeaways
                      else default_strap), accent=self.theme.peri)
        chart_h = 3.62 if has_takeaways else 4.95
        ph = True
        if len(strats) == 1:
            boxes = [(Inches(self.CONTENT_L), Inches(1.62),
                      Inches(self.CONTENT_R - self.CONTENT_L), Inches(chart_h))]
        else:
            boxes = self._matrix_boxes(len(strats), height=chart_h)

        for st, box in zip(strats, boxes):
            key = st.get("key")
            rows = st.get("bars", [])
            ok = self._barlist_card(s, box, st.get("label", key or ""), rows,
                                    "balance", cid=f"strat_{lens}_{key}",
                                    dimension=key)
            ph = ph and not ok

        # A single marginal-change panel beneath, for the dimension that moved
        # most — one clear change view rather than a grid of small ones.
        lines: List[str] = []
        # NEVER SILENT. A reader must be able to tell a dimension dropped for
        # being uniform from one the tape does not carry.
        note = ""
        dropped = [r for r in rejected if "one category" in (r.get("reason") or "")
                   or "single" in (r.get("reason") or "")]
        if dropped:
            names = ", ".join(str(r.get("label") or r.get("key"))
                              for r in dropped[:3])
            note = (f"{names}: the whole book sits in a single band, so the "
                    f"distribution is not charted.")
        elif flat:
            names = ", ".join(str(st.get("label") or st.get("key")) for st in flat)
            note = (f"{names}: the whole book sits in a single band, so the "
                    f"distribution is not charted.")
        if has_takeaways:
            if note:
                lines.append(note)
            for b in moved:
                lines.extend(MV.takeaways(b, limit=1))
            self._takeaway_strip(s, lines[:3], top=5.42)
        elif note:
            self._text(s, Inches(self.CONTENT_L), Inches(6.62),
                       Inches(self.CONTENT_R - self.CONTENT_L), Inches(0.28),
                       note, size=9, color=self.theme.ink_500, italic=True)
        if not strats:
            self._placeholder_body(
                s, f"No {lens} stratifications for this run.")
        self._footer(s)
        self._record(spec.get("id", "strat"), spec.get("title"),
                     window if moved else "", placeholder=ph)

    # ------------------------------------------- economic funded movement
    def slide_balance_movement(self, spec):
        """Why did funded balance change? — the ECONOMIC bridge.

        Opening, plus the loans that arrived, less the loans that left, plus what
        the loans present throughout did. Every figure is
        ``evolution.funded_balance_movement``, which composes the governed
        reconciled bridge with the governed evidence-based exit split; the deck
        computes none of it and draws only what reconciled.

        The continuing-loan leg is deliberately NOT labelled interest. It is the
        movement on loans present at both dates, and separating accretion from
        repayment needs per-loan period movement the canonical model does not
        carry — so the slide says what it measured.
        """
        from .metric_resolver import compact_currency, compact_number

        s = self._slide()
        bm = self.d.balance_movement or {}
        if not bm.get("available"):
            self._header(s, spec.get("title", "Funded Balance Movement"),
                         "Opening to closing", accent=self.theme.peri)
            self._placeholder_body(s, str(bm.get("reason") or
                                          "No reconciled movement for this period."))
            self._footer(s)
            return self._record("balance_movement", spec.get("title"), "",
                                placeholder=True)

        opening = float(bm["openingBalance"])
        closing = float(bm["closingBalance"])
        net = float(bm["netChange"])
        window = f"{bm.get('openingPeriod')} to {bm.get('closingPeriod')}"

        # FINDING-LED SUBTITLE, from the reconciled figures alone.
        direction = "increase" if net >= 0 else "reduction"
        strap = (f"{compact_currency(abs(net))} {direction} over {window} — "
                 f"{self._movement_finding(bm)}")
        self._header(s, spec.get("title", "Funded Balance Movement"),
                     self._fit_label(strap, self.CONTENT_R - self.CONTENT_L, 11),
                     accent=self.theme.peri)

        # The waterfall. Exits are shown split where the evidence classified
        # them, as one bar per reason, so the reader sees WHY loans left.
        steps = [("Opening", opening, "base"),
                 ("+ New funding", float(bm["newLoanBalance"]), "add")]
        components = bm.get("exitComponents") or []
        if components and bm.get("exitsReconcile"):
            for comp in components:
                steps.append((f"− {comp['label']}", -float(comp["balance"]), "sub"))
        else:
            steps.append(("− Exits", -float(bm["exitedLoanBalance"]), "sub"))
        movement = float(bm["continuingMovement"])
        steps.append((("+ " if movement >= 0 else "− ") + "Continuing book",
                      movement, "add" if movement >= 0 else "sub"))
        steps.append(("Closing", closing, "total"))

        # The tile strip below carries a HINT line, which needs 1.34in of tile —
        # at 1.10 the hint rendered outside its own panel and landed on the
        # disclosure. The chart gives up the difference rather than the caption.
        box = (Inches(self.CONTENT_L), Inches(1.72),
               Inches(self.CONTENT_R - self.CONTENT_L), Inches(3.34))
        il, it, iw, ih = self._card(s, *box, f"Funded balance movement, {window}")
        path = self.work / "econ_bridge.png"
        render_bridge_waterfall(path, steps, iw, ih, theme=self.theme,
                                chart_id="balance_movement")
        self._place(s, path, il, it, iw, ih)

        # The counts beneath, and the disclosure the identity depends on.
        tiles = [
            {"label": "Loans added", "value": compact_number(bm.get("newLoanCount")),
             "hint": compact_currency(bm["newLoanBalance"])},
            {"label": "Loans exited", "value": compact_number(bm.get("exitedLoanCount")),
             "hint": compact_currency(bm["exitedLoanBalance"])},
            {"label": "Loans throughout",
             "value": compact_number(bm.get("continuingLoanCount")),
             "hint": compact_currency(movement) + " movement"},
            {"label": "Net change", "value": compact_currency(net),
             "deltaIntent": "positive" if net >= 0 else "negative"},
        ]
        for l, t2, w, h, tile in zip(*self._strip(tiles, top=5.18, height=1.34)):
            self._tile(s, l, t2, w, h, tile)

        self._text(s, Inches(self.CONTENT_L), Inches(6.58),
                   Inches(self.CONTENT_R - self.CONTENT_L), Inches(0.44),
                   self._movement_disclosure(bm), size=8.5,
                   color=self.theme.ink_500, italic=True, spacing=1.04)
        self._footer(s)
        self._record("balance_movement", spec.get("title"), strap)

    #: A bridge leg must hold this much of the gross movement before the page
    #: names it. Deliberately STRICTER than ``materiality.DOMINANCE_SHARE``,
    #: which governs contributions competing within one dimension: the three
    #: legs here are structural — every bridge has all of them — so a leg that
    #: merely leads is not a driver, and only one that carries most of the
    #: movement is worth naming as the reason the balance changed.
    LEG_DOMINANCE_SHARE = 0.45

    def _movement_finding(self, bm) -> str:
        """The one clause that says where the movement came from.

        Chosen from the reconciled legs by magnitude, so it cannot claim a driver
        the bridge does not show.
        """
        legs = [("new funding", abs(float(bm.get("newLoanBalance") or 0.0))),
                ("exits", abs(float(bm.get("exitedLoanBalance") or 0.0))),
                ("movement on the continuing book",
                 abs(float(bm.get("continuingMovement") or 0.0)))]
        total = sum(v for _, v in legs)
        if not total:
            return "no loan-level movement in the period"
        label, value = max(legs, key=lambda x: x[1])
        share = value / total
        if share < self.LEG_DOMINANCE_SHARE:
            return "movement spread across new funding, exits and the continuing book"
        return f"driven primarily by {label}"

    @staticmethod
    def _movement_disclosure(bm) -> str:
        """What the identity rests on, stated on the page rather than assumed."""
        parts = [f"Reconciled on {bm.get('identifierField') or 'loan identity'}; "
                 f"residual within tolerance."]
        if bm.get("exitsClassified") and bm.get("exitsReconcile"):
            evidence = ", ".join(bm.get("exitEvidenceFields") or ()) or "governed exit evidence"
            parts.append(f"Exit reasons from {evidence}.")
        elif bm.get("exitsClassified") is False:
            parts.append("Exit reasons not evidenced on this tape; exits shown in total.")
        parts.append("Continuing-book movement is the change on loans present at "
                     "both dates; it is not split into interest, repayment or "
                     "further advance.")
        return " ".join(parts)

    # ------------------------------------------------ stock by constituent book
    def slide_funded_stock(self, spec):
        """Where does funded exposure sit, and how has that moved?

        A stacked area over ``funded_evolution.breakdowns['portfolio']`` — the
        governed period x book series that has always been computed and never
        drawn. The stack reconciles to the period total by construction: the
        breakdown routes every blank to an explicit Unknown / Missing bucket
        precisely so it does.

        With ONE book the stack conveys nothing a single line does not, so the
        slide falls back to the total series rather than drawing a one-colour
        stack of itself.
        """
        s = self._slide()
        evo = self.d.funded_evolution or {}
        periods = evo.get("periods") or []
        books = self._book_series(evo)

        if len(periods) < 2:
            self._header(s, spec.get("title", "Funded Stock"), "Funded balance over time")
            self._placeholder_body(s, "Funded stock over time needs at least two "
                                      "reporting periods.")
            self._footer(s)
            return self._record("funded_stock", spec.get("title"), "", placeholder=True)

        x = self._period_labels(periods)
        totals = [(p.get("metrics") or {}).get("funded_balance") for p in periods]
        multi = len(books) > 1
        # FINDING-LED SUBTITLE, from the series itself. The title is stable; the
        # subtitle says what the series shows, so a reader who reads only the
        # headings still gets the direction of the book.
        strap = self._stock_strap(totals, x, books, multi)
        self._header(s, spec.get("title", "Funded Stock"), strap, accent=self.theme.peri)

        box = (Inches(self.CONTENT_L), Inches(1.72),
               Inches(self.CONTENT_R - self.CONTENT_L), Inches(4.10))
        il, it, iw, ih = self._card(
            s, *box, "Funded balance by constituent book" if multi
            else "Funded balance by reporting period")
        path = self.work / "funded_stock.png"
        if multi:
            series = [{"name": name, "values": vals} for name, vals in books]
            R.draw_lines(path, x, series, iw, ih, theme=self.theme, currency=True,
                         stack=True, chart_id="funded_stock")
        else:
            R.draw_lines(path, x, [{"name": "Funded balance", "values": totals}],
                         iw, ih, theme=self.theme, currency=True, area=True,
                         chart_id="funded_stock")
        self._place(s, path, il, it, iw, ih)

        self._stock_takeaway(s, books, totals, x, multi)
        self._footer(s)
        self._record("funded_stock", spec.get("title"), strap)

    def _stock_strap(self, totals, x, books, multi) -> str:
        """What the stock series shows, in one clause."""
        from .metric_resolver import compact_currency

        opening = next((v for v in totals if v is not None), None)
        closing = next((v for v in reversed(totals) if v is not None), None)
        prefix = f"{len(books)} constituent books · " if multi else ""
        if opening is None or closing is None or not opening:
            return prefix + "Funded balance over time"
        pct = (closing - opening) / abs(opening) * 100.0
        direction = "up" if pct >= 0 else "down"
        return (f"{prefix}{compact_currency(closing)} at {x[-1]}, "
                f"{direction} {abs(pct):.1f}% since {x[0]}")

    def _book_series(self, evo):
        """``[(book, [value per period])]`` from the governed breakdown, ordered
        largest-closing-first so the stack reads top-down by size.

        Every period is filled — a book absent from a period contributes zero
        there rather than breaking the stack — so the series sum equals the
        period total exactly.
        """
        rows = ((evo.get("breakdowns") or {}).get("portfolio")) or []
        periods = [str(p.get("period") or p.get("reporting_date") or p.get("run_id"))
                   for p in (evo.get("periods") or [])]
        if not rows or not periods:
            return []
        index = {p: i for i, p in enumerate(periods)}
        by_book = {}
        for row in rows:
            per, key = str(row.get("period")), str(row.get("key"))
            if per not in index:
                continue
            by_book.setdefault(key, [0.0] * len(periods))[index[per]] = float(
                row.get("value") or 0.0)
        return sorted(by_book.items(), key=lambda kv: kv[1][-1], reverse=True)

    def _stock_takeaway(self, slide, books, totals, x, multi):
        """One line: what the stack says. Materiality decides whether a book may
        be named as the one that moved."""
        from mi_agent_api import materiality as MAT
        from .metric_resolver import compact_currency

        lines, largest = [], ""
        opening = next((v for v in totals if v is not None), None)
        closing = next((v for v in reversed(totals) if v is not None), None)
        if opening is not None and closing is not None:
            delta = closing - opening
            lines.append(
                f"Funded balance moved from {compact_currency(opening)} at {x[0]} "
                f"to {compact_currency(closing)} at {x[-1]} "
                f"({'+' if delta >= 0 else '−'}{compact_currency(abs(delta))}).")
        if multi and len(totals) >= 2:
            moves = [{"label": name, "value": (vals[-1] - vals[0])}
                     for name, vals in books]
            outcome = MAT.classify(moves, base=opening)
            sentence = MAT.describe(outcome, dimension="constituent book",
                                    money=compact_currency)
            if sentence:
                lines.append(sentence)
            share = closing and books and books[0][1][-1] / closing
            if share:
                largest = (f"{books[0][0]} is the largest book at "
                           f"{share * 100:.0f}% of closing balance.")
        # ONE STORY, NOT TWO. Where the movement page is also in the deck, this
        # page names the number they share and hands the reader on to it. Stated
        # only when the two engines actually agree — a pointer to a page that
        # closes on a different figure would be worse than no pointer.
        bm = self.d.balance_movement or {}
        if bm.get("available") and closing is not None:
            bridge_close = float(bm.get("closingBalance") or 0.0)
            if abs(bridge_close - float(closing)) <= max(0.01, abs(bridge_close) * 1e-9):
                lines.append(
                    f"The same {compact_currency(bridge_close)} closing balance is "
                    f"decomposed loan by loan on Funded Balance Movement.")
        # The hand-off to the movement page outranks the largest-book line: the
        # strip holds three, and connecting the two pages is the point.
        if largest and len(lines) < 3:
            lines.append(largest)
        self._takeaway_strip(slide, lines[:3], top=5.96)

    # -------------------------------------------------- per-book forward view
    def slide_portfolio_projections(self, spec):
        """Which constituent book is expected to drive the portfolio?

        ``forecast_bridge.portfolio_projections`` — current balance, expected
        originations under each book's governed forecast treatment, and a
        retention factor applied ONLY where the client supplied an approved
        run-off curve. Trakt models no run-off of its own; where none was
        supplied the balance is held flat and this slide says so, because a
        projection that quietly assumes a book never redeems is worse than one
        that admits it does not know.
        """
        from .metric_resolver import compact_currency

        s = self._slide()
        pp = self.d.portfolio_projections or {}
        books = pp.get("portfolios") or []
        if len(books) < 1:
            self._header(s, spec.get("title", "Forward View by Book"), "")
            self._placeholder_body(s, "No constituent-book projection for this scope.")
            self._footer(s)
            return self._record("portfolio_projections", spec.get("title"), "",
                                placeholder=True)

        horizon = pp.get("horizonMonths")
        strap = (f"Current and projected balance by book, {horizon}-month horizon"
                 if horizon else "Current and projected balance by book")
        self._header(s, spec.get("title", "Forward View by Book"), strap,
                     accent=self.theme.mint)

        ordered = sorted(books, key=lambda b: float(b.get("projectedBalance") or 0.0),
                         reverse=True)
        rows = [[str(b.get("label") or b.get("portfolioId")),
                 compact_currency(b.get("currentBalance")),
                 compact_currency(b.get("expectedNewOriginations")),
                 ("—" if b.get("balanceRetentionFactor") is None
                  else f"{float(b['balanceRetentionFactor']) * 100:.1f}%"),
                 compact_currency(b.get("projectedBalance"))]
                for b in ordered]
        # THE TOTAL MUST BE THE SUM OF THE ROWS ABOVE IT. Where the governed
        # pipeline cannot be attributed to an individual book, the engine holds
        # the whole weighted amount OUTSIDE the per-book rows and adds it to the
        # total — correctly, and with its own disclosure. Omitting that row put a
        # £111.8MM total over rows summing to £104.8MM, which is the first thing
        # a funder checks and the last thing a pack can afford to get wrong.
        unattributed = float(pp.get("unattributedExpectedOriginations") or 0.0)
        if unattributed:
            rows.append(["Expected originations — not attributed to a book", "",
                         compact_currency(unattributed), "",
                         compact_currency(unattributed)])
        rows.append(["Total", compact_currency(pp.get("totalCurrentBalance")), "",
                     "", compact_currency(pp.get("totalProjectedBalance"))])

        box = (Inches(self.CONTENT_L), Inches(1.72),
               Inches(self.CONTENT_R - self.CONTENT_L), Inches(2.95))
        il, it, iw, ih = self._card(s, *box, "Projection by constituent book")
        path = self.work / "book_projection.png"
        # "Run-off retained" rather than "Retention": this is the share of the
        # existing balance the CLIENT's approved run-off curve keeps over the
        # horizon, not a survival rate Trakt observed.
        R.draw_table(path, ["Book", "Current", "Expected additions",
                            "Run-off retained", "Projected"], rows, iw, ih,
                     theme=self.theme, chart_id="book_projection")
        self._place(s, path, il, it, iw, ih)

        # Current vs projected, side by side, so the shape is visible not read.
        chart = (Inches(self.CONTENT_L), Inches(4.90),
                 Inches(self.CONTENT_R - self.CONTENT_L), Inches(1.62))
        cl, ct, cw, ch = self._card(s, *chart, "Projected balance by book")
        p2 = self.work / "book_projection_bars.png"
        self._barlist_rows(p2, ordered, cw, ch)
        self._place(s, p2, cl, ct, cw, ch)

        self._text(s, Inches(self.CONTENT_L), Inches(6.62),
                   Inches(self.CONTENT_R - self.CONTENT_L), Inches(0.30),
                   self._projection_disclosure(pp), size=9,
                   color=self.theme.ink_500, italic=True)
        self._footer(s)
        self._record("portfolio_projections", spec.get("title"), strap)

    def _barlist_rows(self, path, books, w, h):
        rows = [{"label": str(b.get("label") or b.get("portfolioId")),
                 "balance": float(b.get("projectedBalance") or 0.0)}
                for b in books]
        R.draw_barlist(path, rows, "balance", w, h, theme=self.theme,
                       chart_id="book_projection_bars", dimension="portfolio")

    @staticmethod
    def _projection_disclosure(pp) -> str:
        """The run-off disclosure, verbatim in substance from the governed
        payload — never dropped, because it is the caveat that makes the
        projection honest."""
        not_modelled = pp.get("runoffNotModelled") or []
        modelled = pp.get("runoffModelled") or []
        parts = []
        if pp.get("unattributedExpectedOriginations"):
            # The engine's own words on why the pipeline is not split by book.
            attribution = next(
                (d for d in (pp.get("disclosures") or ())
                 if "not attributed to an individual portfolio" in str(d)), None)
            parts.append(str(attribution) if attribution else
                         "Expected originations are not attributed to an "
                         "individual book and are shown for the originating "
                         "group.")
        if modelled:
            parts.append(f"Run-off applied from the client's approved curve for "
                         f"{', '.join(str(x) for x in modelled)}.")
        if not_modelled:
            parts.append(f"No approved run-off curve for "
                         f"{', '.join(str(x) for x in not_modelled)}; those "
                         f"balances are held flat, not projected to decay.")
        parts.append("Trakt generates no mortality, decay or run-off assumption.")
        return " ".join(parts)

    def slide_geo(self, spec):
        s = self._slide()
        self._header(s, spec.get("title", "Geographic Exposure"),
                     "Funded exposure by ITL3 area")
        geo = self.d.geo
        areas = sorted(geo.get("areas", []), key=lambda a: a.get("balance", 0),
                       reverse=True)
        # concentration tiles
        total = geo.get("total") or sum(a.get("balance", 0) for a in areas)
        from .metric_resolver import compact_currency, format_percent
        top = areas[0] if areas else {}
        top5 = sum(a.get("balance", 0) for a in areas[:5])
        tiles = [
            {"label": "Total funded exposure", "value": compact_currency(total),
             "hint": f"{geo.get('areaCount', len(areas))} ITL3 areas"},
            {"label": "Top area", "value": top.get("itl3_name", "—"),
             "hint": compact_currency(top.get("balance", 0)) if top else ""},
            {"label": "Top-5 concentration",
             "value": format_percent(top5 / total) if total else "—"},
            {"label": "Postcode coverage",
             "value": (f"{geo.get('coveragePct')}%" if geo.get("coveragePct") is not None else "—")},
        ]
        # 4 tiles across the top
        tw = Inches(2.98)
        for i, tile in enumerate(tiles):
            l = Emu(int(Inches(0.55)) + i * int(Inches(3.08)))
            self._tile(s, l, Inches(1.6), tw, Inches(1.5), tile)
        # region BarList below
        box = (Inches(0.55), Inches(3.35), Inches(12.25), Inches(3.25))
        rows = [{"label": a.get("itl3_name", a.get("itl3_code", "")),
                 "balance": a.get("balance", 0)} for a in areas[:12]]
        self._barlist_card(s, box, "Top areas by funded exposure", rows, "balance",
                           cid="geo_bars")
        self._footer(s)
        self._record("geography", spec.get("title"), "", placeholder=not areas)

    def _evolution_lines(self, s, spec, evo, chart_specs, accent=None,
                         height=None):
        """Render N line-chart cards from an evolution payload's periods[].

        A time series needs ≥2 reporting periods; the dashboard flags a single cut
        with ``singlePeriod`` and shows an insufficient-history state rather than a
        lone point — so do the same (a one-dot 'trend' reads as broken)."""
        periods = evo.get("periods", [])
        single = bool(evo.get("singlePeriod")) or len(periods) < 2
        x = self._period_labels(periods)
        # Only measures the governed series actually carries. A tape without an
        # interest rate should lose that panel, not gain an empty one.
        chart_specs = [cs for cs in chart_specs
                       if any((p.get("metrics") or {}).get(ser["key"]) is not None
                              for p in periods for ser in cs["series"])] or chart_specs[:1]
        boxes = (self._matrix_boxes(len(chart_specs)) if len(chart_specs) > 2
                 else self._chart_boxes(len(chart_specs), **(
                     {"height": height} if height else {})))
        for cs, box in zip(chart_specs, boxes):
            il, it, iw, ih = self._card(s, *box, cs["title"])
            series = [{"name": ser.get("name", ""),
                       "values": [(p.get("metrics") or {}).get(ser["key"]) for p in periods],
                       "color": ser.get("color")}
                      for ser in cs["series"]]
            path = self.work / f"{cs['id']}.png"
            if not single:
                R.draw_lines(path, x, series, iw, ih, theme=self.theme,
                             currency=cs.get("currency", True),
                             percent=cs.get("percent", False),
                             area=cs.get("area", False), chart_id=cs["id"])
            else:
                render_placeholder_png(path, "", "Insufficient reporting history "
                                       "(needs ≥2 periods)", theme=self.theme,
                                       width_in=iw, height_in=ih)
            self._place(s, path, il, it, iw, ih)
        return single

    def slide_funded_evolution(self, spec):
        s = self._slide()
        self._header(s, spec.get("title", "Funded Evolution"), "Funded book over time")
        # The four measures the dashboard's funded Evolution tab plots. The deck
        # used to show two of them, so the same tab and the same slide described
        # the book differently.
        ph = self._evolution_lines(s, spec, self.d.funded_evolution, [
            {"id": "evo_bal", "title": "Funded balance by month",
             "series": [{"name": "Funded balance", "key": "funded_balance"}],
             "currency": True},
            {"id": "evo_count", "title": "Funded loan count by month",
             "series": [{"name": "Loan count", "key": "loan_count"}],
             "currency": False},
            {"id": "evo_ltv", "title": "WA current LTV by month",
             "series": [{"name": "WA LTV", "key": "wa_ltv"}],
             "currency": False, "percent": True},
            {"id": "evo_rate", "title": "WA interest rate by month",
             "series": [{"name": "WA rate", "key": "wa_interest_rate"}],
             "currency": False, "percent": True},
        ])
        self._footer(s)
        self._record("funded_evolution", spec.get("title"), "", placeholder=ph)

    def slide_cohorts(self, spec):
        """Vintage Formation — *how much business entered each vintage?*

        The governed cohort table (``cohorts.cohort_analysis``, the same service
        behind the React Cohorts composition table), presented as formation
        rather than as a bare cross-section. Where the scope spans more than one
        portfolio type this reports them separately, because a ten-year acquired
        book and a two-year direct book averaged into one vintage series
        describe neither.
        """
        from . import cohorts as CO
        from .metric_resolver import compact_currency, compact_number

        s = self._slide()
        formation = CO.adapt_formation(self.d.cohorts)
        basis = "Origination vintage" if (formation.cohort_basis or "").find(
            "origination") >= 0 else "Cohort"
        self._header(s, spec.get("title", "Vintage Formation"),
                     f"{basis} composition of the funded book"
                     + (f" · {formation.span}" if formation.span else ""))
        if not CO.formation_is_meaningful(formation):
            self._placeholder_body(
                s, formation.reason or "No governed cohort composition for this book.")
            self._footer(s)
            return self._record("cohorts", spec.get("title"), "", placeholder=True)

        rows = [r for r in formation.rows if r.vintage != "Unknown"]
        # -- summary strip ----------------------------------------------------
        oldest, newest = rows[0], rows[-1]
        tiles = [
            {"label": "Vintages", "value": compact_number(len(rows))},
            {"label": "Oldest vintage", "value": oldest.vintage,
             "hint": compact_currency(oldest.balance)},
            {"label": "Newest vintage", "value": newest.vintage,
             "hint": compact_currency(newest.balance)},
            {"label": "Largest vintage",
             "value": max(rows, key=lambda r: r.balance or 0).vintage,
             "hint": compact_currency(max(r.balance or 0 for r in rows))},
        ]
        for l, t, w, h, tile in zip(*self._strip(tiles, top=1.58, height=1.16)):
            self._tile(s, l, t, w, h, tile)

        # -- balance by vintage, then the governed per-vintage measures -------
        boxes = self._chart_boxes(2, top=3.02, height=3.26)
        il, it, iw, ih = self._card(s, *boxes[0], "Funded balance by vintage")
        p1 = self.work / "cohort_balance.png"
        R.draw_barlist(p1, [{"label": r.vintage, "balance": r.balance}
                            for r in sorted(rows, key=lambda r: -(r.balance or 0))],
                       "balance", iw, ih, theme=self.theme)
        self._place(s, p1, il, it, iw, ih)

        il, it, iw, ih = self._card(s, *boxes[1], "Vintage measures")
        self._vintage_table(s, rows, formation, il, it, iw, ih)
        self._footer(s)
        self._record("cohorts", spec.get("title"),
                     f"{len(rows)} governed vintages.")

    def _vintage_table(self, s, rows, formation, il, it, iw, ih):
        """The governed per-vintage measures, as native text.

        Only columns the service reported in ``metricsAvailable`` are drawn: a
        tape without an interest rate must lose the rate column, not gain a
        column of dashes.
        """
        from .metric_resolver import compact_currency, compact_number

        # Columns are declared with RELATIVE weights and laid out across the card
        # that exists, so a tape carrying every optional measure cannot push the
        # last column past the panel edge — which a fixed offset table did.
        have = set(formation.metrics_available)
        spec_cols = [("Vintage", 1.05, PP_ALIGN.LEFT, lambda r: r.vintage)]
        if "balance" in have:
            spec_cols += [
                ("Balance", 1.15, PP_ALIGN.RIGHT,
                 lambda r: compact_currency(r.balance)),
                ("Share", 0.72, PP_ALIGN.RIGHT,
                 lambda r: f"{r.share_pct:.1f}%" if r.share_pct is not None else "—")]
        if "loanCount" in have:
            spec_cols.append(("Loans", 0.72, PP_ALIGN.RIGHT,
                              lambda r: compact_number(r.loan_count)))
        if "balance" in have and "loanCount" in have:
            spec_cols.append(("Avg", 0.80, PP_ALIGN.RIGHT,
                              lambda r: compact_currency(r.average_balance)))
        if "waLtv" in have:
            spec_cols.append(("WA LTV", 0.80, PP_ALIGN.RIGHT,
                              lambda r: f"{r.wa_ltv:.1f}%" if r.wa_ltv is not None else "—"))
        if "waRate" in have:
            spec_cols.append(("WA rate", 0.82, PP_ALIGN.RIGHT,
                              lambda r: f"{r.wa_rate:.1f}%" if r.wa_rate is not None else "—"))
        if "waMonthsOnBook" in have:
            spec_cols.append(("Months", 0.72, PP_ALIGN.RIGHT,
                              lambda r: compact_number(r.wa_months_on_book)))

        x0 = (il / EMU_IN) + 0.14
        scale = (iw - 0.28) / sum(c[1] for c in spec_cols)   # iw is the WIDTH
        cols, dx = [], 0.0
        for label, weight, align, fn in spec_cols:
            cols.append((label, dx, weight * scale, align, fn))
            dx += weight * scale
        band = ih - 0.46                      # ih is already INCHES from _card
        shown = rows[: max(1, int(band / 0.26))]
        row_h = min(0.34, band / max(len(shown), 1))
        head_y = (it / EMU_IN) + 0.16
        for label, dx, cw, align, _fn in cols:
            self._text(s, Inches(x0 + dx), Inches(head_y + pad), Inches(cw),
                       Inches(0.24), label, size=8.5, color=self.theme.ink_400,
                       bold=True, align=align)
        size = 9.5 if len(shown) <= 8 else 8.5
        for i, r in enumerate(shown):
            y = Inches(head_y + 0.30 + i * row_h)
            for label, dx, cw, align, fn in cols:
                self._text(s, Inches(x0 + dx), y, Inches(cw), Inches(0.26),
                           str(fn(r)), size=size, align=align,
                           color=self.theme.ink_100 if label == "Vintage"
                           else self.theme.ink_300)
        if len(shown) < len(rows):
            self._text(s, Inches(x0), Inches(head_y + 0.30 + len(shown) * row_h),
                       Inches(5.4), Inches(0.24),
                       f"{len(rows) - len(shown)} further vintages not shown.",
                       size=8, color=self.theme.ink_500)

    def slide_cohort_progression(self, spec):
        """Cohort Progression — *how have cohorts seasoned since they formed?*

        The governed static pool (``evolution.funded_cohort_progression``), one
        call per cohort, plotted on a SEASONING axis: periods since formation,
        not calendar dates, so a 2019 vintage and a 2024 vintage can be compared
        at the same age.
        """
        from . import cohorts as CO
        from .metric_resolver import compact_currency

        s = self._slide()
        payload = self.d.cohort_series or {}
        series = [CO.adapt_progression(p, v)
                  for v, p in sorted((payload.get("series") or {}).items())]
        live = CO.plottable(series)
        declined = CO.rejected(series)
        self._header(s, spec.get("title", "Cohort Progression"),
                     "Static-pool seasoning — the same cohort tracked across "
                     "reporting periods")
        if not live:
            reason = next((x.reason for x in series if x.reason), None)
            self._placeholder_body(
                s, "Static-pool seasoning needs at least two reporting periods "
                   "in which a cohort holds loans." + (f" {reason}." if reason else ""))
            self._footer(s)
            return self._record("cohort_progression", spec.get("title"), "",
                                placeholder=True)

        # -- balance curves, indexed to formation ----------------------------
        boxes = self._chart_boxes(2, top=1.62, height=3.62)
        il, it, iw, ih = self._card(s, *boxes[0],
                                    "Funded balance by reporting periods since formation")
        p1 = self.work / "cohort_prog_balance.png"
        longest = max(len(x.live) for x in live)
        R.draw_lines(p1, [str(i) for i in range(longest)],
                     [{"name": x.vintage,
                       "values": [x.value("funded_balance", i) if i < len(x.live)
                                  else None for i in range(longest)]}
                      for x in live],
                     iw, ih, theme=self.theme, currency=True)
        self._place(s, p1, il, it, iw, ih)

        # -- how each pool has changed since formation ------------------------
        # Titled for what the table measures rather than for "retention": the
        # balance column is a net-of-everything ratio, and on a roll-up book it
        # exceeds 100% while loans are leaving. Only the loan column is survival.
        il, it, iw, ih = self._card(s, *boxes[1], "Change since formation")
        self._cohort_change_table(s, live, il, it, iw, ih)

        overflow = payload.get("overflow") or []
        # Deliberately avoids naming what a balance movement WOULD be attributed
        # to. The publication gate that bans causal vocabulary is a substring
        # check and cannot read a negation, and it is the more valuable of the
        # two properties — so the sentence is written without the word.
        note = ("Cohorts are the governed static pool: a vintage fixed at "
                "formation and tracked across reporting periods. Loan survival "
                "is surviving loans as a percentage of the loans at formation; "
                "exits are loans in the pool at formation that are no longer in "
                "it. Balance vs formation is the latest balance as a percentage "
                "of the balance at formation — a net figure that is not "
                "decomposed here, and which can exceed 100%. Periods is the "
                "number of reporting periods the cohort has been observed for "
                "since it formed.")
        if declined:
            note += (f" {len(declined)} cohort"
                     f"{'s were' if len(declined) != 1 else ' was'} not plotted "
                     f"because the governed series does not hold the pool fixed.")
        if overflow:
            note += (f" {len(overflow)} smaller vintage"
                     f"{'s are' if len(overflow) != 1 else ' is'} not plotted.")
        self._text(s, Inches(0.57), Inches(5.46), Inches(12.2), Inches(0.5), note,
                   size=9, color=self.theme.ink_500, spacing=1.06)
        self._footer(s)
        self._record("cohort_progression", spec.get("title"),
                     f"{len(live)} governed cohort series.")

    def _cohort_change_table(self, s, live, il, it, iw, ih):
        from .metric_resolver import compact_currency, compact_number

        # COLUMN NAMES ARE THE MEASURE, NOT A HOUSE WORD. "Retention" is
        # reserved for the count-based survival ratio; the balance ratio is
        # named for what it is, because calling a >100% balance figure
        # "retention" invites a reader to conclude the pool grew.
        # Headers must fit their column on ONE line. "Loan survival" and
        # "Seasoning" each wrapped, and the wrap clipped mid-word — the meaning
        # of every one of these columns is carried by the note beneath the card,
        # so the headers are short and the note does the explaining.
        spec_cols = [("Cohort", 0.68, PP_ALIGN.LEFT),
                     ("At formation", 1.00, PP_ALIGN.RIGHT),
                     ("Latest", 0.86, PP_ALIGN.RIGHT),
                     ("Balance vs formation", 1.30, PP_ALIGN.RIGHT),
                     ("Loans", 0.64, PP_ALIGN.RIGHT),
                     ("Survival", 0.78, PP_ALIGN.RIGHT),
                     ("Exits", 0.56, PP_ALIGN.RIGHT),
                     ("Periods", 0.66, PP_ALIGN.RIGHT)]
        scale = (iw - 0.28) / sum(c[1] for c in spec_cols)
        cols, dx = [], 0.0
        for label, weight, align in spec_cols:
            cols.append((label, dx, weight * scale, align))
            dx += weight * scale
        x0 = (il / EMU_IN) + 0.14
        head_y = (it / EMU_IN) + 0.16
        for label, dx, cw, align in cols:
            self._text(s, Inches(x0 + dx), Inches(head_y), Inches(cw), Inches(0.24),
                       label, size=8.5, color=self.theme.ink_400, bold=True,
                       align=align)
        band = ih - 0.52                      # ih is already INCHES from _card
        row_h = min(0.44, band / max(len(live), 1))
        # Centred in the band. Four rows pinned to the top of a 3.6in card leave
        # an inch and a half of empty panel, which reads as a rendering fault
        # rather than as a short list.
        pad = max(0.0, (band - len(live) * row_h) / 2)
        for i, x in enumerate(live):
            y = Inches(head_y + 0.32 + pad + i * row_h)
            bal_ret = x.balance_vs_formation
            survival = x.loan_survival
            exits = x.exits
            values = [
                (x.vintage, self.theme.ink_100),
                (compact_currency(x.value("funded_balance", 0)), self.theme.ink_300),
                (compact_currency(x.value("funded_balance", -1)), self.theme.ink_100),
                (f"{bal_ret:.0f}%" if bal_ret is not None else "—",
                 self.theme.ink_100),
                (f"{x.surviving_count:,}/{x.formation_count:,}"
                 if x.formation_count is not None else "—", self.theme.ink_300),
                (f"{survival:.0f}%" if survival is not None else "—",
                 self.theme.ink_100),
                (f"{exits:,}" if exits is not None else "—",
                 self.theme.rag.get("amber") if exits else self.theme.ink_300),
                (str(len(x.live) - 1), self.theme.ink_300),
            ]
            for (value, colour), (_label, dx, cw, align) in zip(values, cols):
                self._text(s, Inches(x0 + dx), y, Inches(cw), Inches(0.28),
                           value, size=10, color=colour, align=align)

    def slide_pipeline(self, spec):
        s = self._slide()
        p = self.d.pipeline
        self._header(s, spec.get("title", "Pipeline"),
                     "Origination pipeline (pre-funded)", accent=self.theme.peri)
        if not p:
            self._placeholder_body(s, "No pipeline source resolved for this run.")
            self._footer(s)
            return self._record("pipeline", spec.get("title"), "", placeholder=True)
        from .metric_resolver import compact_currency, compact_number
        pw = p.get("priorWeek") or {}
        def delta(cur, prv, cur_key):
            if not pw or prv is None:
                return None, None
            diff = (cur or 0) - (prv or 0)
            intent = "positive" if diff > 0 else ("negative" if diff < 0 else "neutral")
            return (compact_currency(diff) if cur_key == "amt" else
                    ("+" if diff >= 0 else "−") + compact_number(abs(diff))) + " vs prior wk", intent
        d1, i1 = delta(p.get("pipelineAmount"), pw.get("pipelineAmount"), "amt")
        d2, i2 = delta(p.get("pipelineRowCount"), pw.get("pipelineRowCount"), "cnt")
        cases = p.get("pipelineRowCount") or 0
        amount = p.get("pipelineAmount") or 0
        avg_case = (amount / cases) if cases else 0
        tiles = [
            {"label": "Pipeline cases", "value": compact_number(cases),
             "delta": d2, "deltaIntent": i2},
            {"label": "Total pipeline amount", "value": compact_currency(amount),
             "delta": d1, "deltaIntent": i1},
            {"label": "Average case amount", "value": compact_currency(avg_case),
             "hint": "total ÷ cases"},
            {"label": "Weighted expected funded",
             "value": compact_currency(p.get("weightedExpectedFundedAmount")),
             "hint": "probability-weighted"},
        ]
        # Tiles and charts share ONE grid, so the strip and the panels beneath
        # it line up at both outer edges.
        for l, t, w, h, tile in zip(*self._strip(tiles, top=1.60, height=1.45)):
            self._tile(s, l, t, w, h, tile)
        # two BarLists: stage + broker
        box1, box2 = self._chart_boxes(2, top=3.28, height=3.35)
        self._barlist_card(s, box1, "Pipeline amount by stage",
                           self._stage_rows(p.get("stageBreakdown", [])), "pipelineAmount",
                           cid="pipe_stage")
        # THE SECOND CHART EARNS ITS PLACE. It was broker/channel, falling back
        # to region — and on a direct-only book that drew one bar labelled
        # "Direct", which is the pipeline total already in the tile above it,
        # redrawn as a chart. The strongest dimension this pipeline actually
        # distributes across takes the panel instead, judged by the shared rule.
        second = self._pipeline_second_cut(p)
        if second:
            self._barlist_card(s, box2, second["title"], second["rows"],
                               second["value_key"], cid="pipe_second",
                               label_key=second["label_key"],
                               dimension=second.get("dimension"))
        else:
            # Nothing distributes. Rather than a meaningless chart, the panel
            # carries the pipeline facts a second chart would have competed
            # with — the expected-completion profile the tiles only summarise.
            self._pipeline_secondary_facts(s, box2, p)
        self._footer(s)
        self._record("pipeline", spec.get("title"), "", placeholder=False)

    def _pipeline_second_cut(self, pipeline):
        """The strongest informative pipeline dimension, or ``None``.

        Reads the governed pipeline stratifications — the SAME payload the
        Pipeline Stratifications page draws — so the two cannot disagree about
        which cuts this pipeline supports, and applies the shared
        informativeness rule to pick the one worth a panel here.
        """
        from mi_agent_api import presentation as _sel

        strats = [st for st in (pipeline.get("stratifications") or ())
                  if isinstance(st, dict) and st.get("bars")]
        if strats:
            chosen = _sel.select_dimensions(
                strats, want=1, value_key="balance",
                preferred=("product", "region", "ltv", "ticket", "age", "rate"))
            if chosen["selected"]:
                st = chosen["selected"][0]
                return {"title": f"Pipeline amount {str(st.get('label', '')).lower()}",
                        "rows": st["bars"], "value_key": "balance",
                        "label_key": "label", "dimension": st.get("key")}

        # No governed stratifications on this payload (an older pipeline
        # source): fall back to the flat breakdowns, still judged on shape.
        for key, title in (("brokerBreakdown", "Pipeline amount by broker / channel"),
                           ("regionBreakdown", "Pipeline amount by region")):
            rows = list(pipeline.get(key) or ())
            if not rows:
                continue
            rows.sort(key=lambda r: r.get("pipelineAmount", 0), reverse=True)
            if _sel.is_informative(rows, value_key="pipelineAmount"):
                return {"title": title, "rows": rows,
                        "value_key": "pipelineAmount", "label_key": "key"}
        return None

    def _pipeline_secondary_facts(self, slide, box, pipeline):
        """Governed pipeline facts, where no dimension earns a chart.

        Every figure is lifted from the pipeline snapshot the tiles above read;
        nothing is computed here.
        """
        from .metric_resolver import compact_currency, compact_number

        il, it, iw, ih = self._card(slide, *box, "Expected completion profile")
        rows = []
        nxt = pipeline.get("nextExpectedCompletionMonth")
        if nxt:
            rows.append(("Next expected completion month", str(nxt)))
        cur = pipeline.get("currentMonthExpectedCompletionCount")
        if cur is not None:
            rows.append(("Cases expected to complete this month",
                         compact_number(cur)))
        overdue = pipeline.get("overdueExpectedCompletionCount")
        if overdue is not None:
            rows.append(("Cases past their expected completion date",
                         compact_number(overdue)))
        overdue_amt = pipeline.get("overdueExpectedCompletionWeightedAmount")
        if overdue_amt:
            rows.append(("Weighted amount past expected completion",
                         compact_currency(overdue_amt)))
        stages = pipeline.get("pipelineLiveStages") or ()
        if stages:
            rows.append(("Live stages", ", ".join(
                str(x).title() for x in stages)))
        if not rows:
            self._text(slide, il, it + Inches(0.2), iw, Inches(0.3),
                       "No further governed pipeline detail for this book.",
                       size=10, color=self.theme.ink_500, italic=True)
            return
        y = float(it) / EMU_IN + 0.24
        for label, value in rows[:6]:
            self._text(slide, il, Inches(y), Inches(iw * 0.66),
                       Inches(0.3), label, size=10, color=self.theme.ink_400)
            self._text(slide, il, Inches(y), Inches(iw),
                       Inches(0.3), str(value), size=10.5,
                       color=self.theme.ink_100, align=PP_ALIGN.RIGHT, bold=True)
            y += 0.42

    def _stage_rows(self, rows, value_key="pipelineAmount"):
        order = {"KFI": 0, "APPLICATION": 1, "OFFER": 2, "COMPLETED": 3, "WITHDRAWN": 4}
        pretty = {"KFI": "KFI", "APPLICATION": "Application", "OFFER": "Offer",
                  "COMPLETED": "Completed", "WITHDRAWN": "Withdrawn", "UNKNOWN": "Other"}
        rows = sorted(rows, key=lambda r: order.get(str(r.get("stage", "")).upper(), 9))
        return [{"label": pretty.get(str(r.get("stage", "")).upper(),
                                     str(r.get("stage", "")).title()),
                 "pipelineAmount": r.get("pipelineAmount", 0),
                 "caseCount": r.get("caseCount", 0)} for r in rows]

    _STAGE_PRETTY = {"KFI": "KFI", "APPLICATION": "Application", "OFFER": "Offer",
                     "COMPLETED": "Completed", "WITHDRAWN": "Withdrawn"}
    _STAGE_COLOR = {"APPLICATION": "#7c9cf0", "OFFER": "#5ec6b8",
                    "COMPLETED": "#e0a458", "WITHDRAWN": "#eb6f6f"}

    def slide_pipeline_evolution(self, spec):
        s = self._slide()
        self._header(s, spec.get("title", "Pipeline Evolution"),
                     "Pipeline stock over time", accent=self.theme.peri)
        evo = self.d.pipeline_evolution or {}
        periods = evo.get("periods", [])
        single = bool(evo.get("singlePeriod")) or len(periods) < 2
        boxes = self._chart_boxes(2)
        x = [str(p.get("week") or p.get("period")) for p in periods]

        il, it, iw, ih = self._card(s, *boxes[0], "Pipeline amount by week")
        p1 = self.work / "pevo_amt.png"
        if not single:
            R.draw_lines(p1, x, [{"name": "Pipeline amount",
                                  "values": [(p.get("metrics") or {}).get("pipeline_amount")
                                             for p in periods]}],
                         iw, ih, theme=self.theme, currency=True, area=True)
        else:
            render_placeholder_png(p1, "", "Insufficient reporting history (needs ≥2 weeks)",
                                   theme=self.theme, width_in=iw, height_in=ih)
        self._place(s, p1, il, it, iw, ih)

        # Pipeline by stage over time — EXCLUDING the KFI line (dashboard view).
        il, it, iw, ih = self._card(s, *boxes[1], "Pipeline by stage over time")
        p2 = self.work / "pevo_stage.png"
        by_stage = evo.get("byStage", [])
        if not single and by_stage:
            weeks = sorted({str(r.get("week") or r.get("period")) for r in by_stage})
            lut = {(str(r.get("stage", "")).upper(), str(r.get("week") or r.get("period"))):
                   r.get("value") for r in by_stage}
            series = []
            for st in ("APPLICATION", "OFFER", "COMPLETED", "WITHDRAWN"):
                vals = [lut.get((st, w)) for w in weeks]
                if any(v for v in vals):
                    series.append({"name": self._STAGE_PRETTY[st], "values": vals,
                                   "color": self._STAGE_COLOR[st]})
            R.draw_lines(p2, weeks, series, iw, ih, theme=self.theme, currency=True)
        else:
            render_placeholder_png(p2, "", "Insufficient reporting history (needs ≥2 weeks)",
                                   theme=self.theme, width_in=iw, height_in=ih)
        self._place(s, p2, il, it, iw, ih)
        self._footer(s)
        self._record("pipeline_evolution", spec.get("title"), "", placeholder=single)

    def slide_pipeline_movement(self, spec):
        """Pipeline Stage Movement — *what happened to cases between periods?*

        For each live stage, on counts AND amounts::

            opening live + arrivals - departures +/- amount change on stayers
                = closing live

        Two aligned reads, because a funder asks both: a stage table carrying
        the identity in cases and money, and — where cases left a stage — where
        they actually went. "Left the stage" and "left the pipeline" are
        different events, and a completion is not attrition.

        Nothing is computed here. The reconciliation is
        ``evolution.pipeline_stage_movement``, which ``/mi/evolution/pipeline-
        movement`` serves to the dashboard, so both surfaces read one result.
        """
        from .metric_resolver import compact_currency, compact_number

        s = self._slide()
        mv = self.d.pipeline_movement or {}
        stages = [st for st in (mv.get("stages") or ())
                  if st.get("openingCaseCount") or st.get("closingCaseCount")]
        window = ""
        if mv.get("openingWeek") and mv.get("closingWeek"):
            window = f"{mv['openingWeek']} to {mv['closingWeek']}"
        self._header(s, spec.get("title", "Pipeline Stage Movement"),
                     (f"Case and balance movement by stage, {window}" if window
                      else "Case and balance movement by stage"),
                     accent=self.theme.peri)
        if not stages:
            self._placeholder_body(
                s, mv.get("reason") or "No governed stage movement for this book.")
            self._footer(s)
            return self._record("pipeline_movement", spec.get("title"), "",
                                placeholder=True)

        # -- the identity, stage by stage ------------------------------------
        cols = ["Stage", "Opening", "Arrived", "Departed", "On stayers", "Closing"]
        rows = []
        for st in stages:
            def _leg(count, amount, sign):
                # A zero leg is a dash. "−0  £0" reads as a rendering fault.
                if not count and not amount:
                    return "—"
                return (f"{sign}{compact_number(count)}  "
                        f"{compact_currency(amount)}")

            rows.append([
                self._STAGE_PRETTY.get(st["stage"], st["stage"]),
                f"{compact_number(st['openingCaseCount'])}  "
                f"{compact_currency(st['openingAmount'])}",
                _leg(st["arrivalCaseCount"], st["arrivalAmount"], "+"),
                _leg(st["departureCaseCount"], st["departureAmount"], "−"),
                self._signed_currency(st.get("amountChangeOnPersisting")),
                f"{compact_number(st['closingCaseCount'])}  "
                f"{compact_currency(st['closingAmount'])}",
            ])
        # -- where the departures went ---------------------------------------
        destinations: Dict[str, Dict[str, Any]] = {}
        for st in stages:
            for dest in st.get("departuresByDestination") or ():
                bucket = destinations.setdefault(
                    str(dest.get("stage", "")).upper(),
                    {"cases": 0, "amount": 0.0})
                bucket["cases"] += int(dest.get("caseCount") or 0)
                bucket["amount"] += float(dest.get("amount") or 0.0)
        bars = [{"label": ("Left the extract" if key == "ABSENT"
                           else self._STAGE_PRETTY.get(key, key.title())),
                 "balance": info["amount"], "count": info["cases"]}
                for key, info in destinations.items() if info["amount"]]
        bars.sort(key=lambda b: -b["balance"])

        # THE TABLE TAKES THE PAGE WHEN NOTHING DEPARTED. Reserving a second
        # panel for "no case left a stage" spends half the slide saying nothing
        # happened; the sentence carries that, and the table gets the room.
        table_h = 2.78 if bars else 3.92
        il, it, iw, ih = self._card(
            s, Inches(self.CONTENT_L), Inches(1.62),
            Inches(self.CONTENT_R - self.CONTENT_L), Inches(table_h),
            "Live stock by stage")
        path = self.work / "pipe_move_table.png"
        R.draw_table(path, cols, rows, iw, ih, theme=self.theme)
        self._place(s, path, il, it, iw, ih)

        if bars:
            self._barlist_card(
                s, (Inches(self.CONTENT_L), Inches(4.58),
                    Inches(self.CONTENT_R - self.CONTENT_L), Inches(1.94)),
                "Where departing cases went", bars, "balance",
                cid="pipe_move_dest")
        else:
            self._text(s, Inches(self.CONTENT_L), Inches(5.76),
                       Inches(self.CONTENT_R - self.CONTENT_L), Inches(0.32),
                       "No case left a live stage between these two extracts.",
                       size=11, color=self.theme.ink_400, italic=True)

        note = mv.get("lineage", {}).get("identity") or ""
        if note:
            self._text(s, Inches(self.CONTENT_L), Inches(6.60),
                       Inches(self.CONTENT_R - self.CONTENT_L), Inches(0.30),
                       ("Reconciled on the governed case identifier "
                        f"({mv.get('identifierField')}): {note}. An amount "
                        "amendment is a movement on the same case, not an exit "
                        "and an arrival."),
                       size=8.5, color=self.theme.ink_500, italic=True)
        self._footer(s)
        self._record("pipeline_movement", spec.get("title"),
                     f"{len(stages)} live stage(s); "
                     f"{'reconciles' if mv.get('reconciles') else 'residual outside tolerance'}.")

    def slide_origination_flow(self, spec):
        """KFI and Completion weekly-flow panels (bars) with a cumulative line —
        the dashboard's pipeline→origination flow view."""
        s = self._slide()
        self._header(s, spec.get("title", "Origination Flow — KFIs & Completions"),
                     "Weekly flow with cumulative build", accent=self.theme.peri)
        f = self.d.funnel or {}
        series = f.get("series", {}) or {}
        summary = f.get("summary", {}) or {}
        single = bool(f.get("singlePeriod")) or not series
        boxes = self._chart_boxes(2)
        for box, stage in zip(boxes, ("KFI", "COMPLETED")):
            label = self._STAGE_PRETTY.get(stage, stage)
            il, it, iw, ih = self._card(s, *box, f"{label}s · weekly flow")
            path = self.work / f"flow_{stage.lower()}.png"
            pts = series.get(stage) or []
            if not single and pts:
                weeks = [str(pt.get("week")) for pt in pts]
                vals = [pt.get("value") for pt in pts]
                cum, run = [], 0.0
                for v in vals:
                    run += float(v or 0)
                    cum.append(run)
                avg = (summary.get(stage) or {}).get("fiveWeekAvgFlowValue")
                R.draw_bars_with_line(path, weeks, vals, cum, iw, ih, theme=self.theme,
                                      avg=avg, line_label="Cumulative")
            else:
                render_placeholder_png(path, "", "Insufficient reporting history "
                                       "(needs ≥2 weeks)", theme=self.theme,
                                       width_in=iw, height_in=ih)
            self._place(s, path, il, it, iw, ih)
        self._footer(s)
        self._record("origination_flow", spec.get("title"), "", placeholder=single)

    def slide_multidim(self, spec):
        """Funded balance across paired dimensions — one visual grammar.

        LTV × borrower age used to be a bubble matrix while LTV × borrower type
        and LTV × region were heatmaps, though all three answer the same
        question: how much balance sits in each cell of a two-way band cross-tab.
        Bubble AREA encodes the value but bubble POSITION encodes nothing, so a
        reader compares circles by eye where a heatmap lets them read the number.
        All three are now heatmaps on one colour scale methodology.

        The underlying cross-tabs are unchanged — the bubble was already drawing
        the same matrix, via its ``points`` projection.
        """
        s = self._slide()
        self._header(s, spec.get("title", "Multi-Dimensional Risk Analytics"),
                     "Funded balance across paired dimensions", accent=self.theme.peri)
        md = self.d.multidim or {}
        # WHAT THIS BOOK SUPPORTS, not a fixed three. The pairs are chosen by
        # the engine's governed selection — both dimensions present, both axes
        # real, the matrix dense enough to read, and no crossing repeating a
        # story a crossing above already told. Only panels that resolved are
        # drawn; the composition guard omits the slide when none do.
        wanted = spec.get("pairs")
        panels = [(key, str(entry.get("label", key)).replace(" x ", " × "))
                  for key, entry in md.items()
                  if not key.startswith("_")
                  and isinstance(entry, dict) and entry.get("matrix")
                  and (not wanted or key in wanted)]
        panels = panels[:4]
        if not panels:
            self._placeholder_body(s, "No paired funded dimensions resolved.")
            self._footer(s)
            return self._record("multidim", spec.get("title"), "", placeholder=True)

        # A crossing with long row labels (region, product) takes the full width
        # when it would otherwise share a row and be squeezed to unreadable.
        _wide_dims = ("region", "product", "status")
        if len(panels) == 3:
            wide = [p for p in panels
                    if any(d in p[0] for d in _wide_dims)][:1]
            narrow = [p for p in panels if p not in wide]
            boxes = [(Inches(l), Inches(1.62), Inches(w), Inches(2.42))
                     for l, w in self._grid(2)]
            boxes += [(Inches(self.CONTENT_L), Inches(4.20),
                       Inches(self.CONTENT_R - self.CONTENT_L), Inches(2.38))]
            ordered = (narrow + wide) if wide else panels
        elif len(panels) == 4:
            ordered = panels
            boxes = self._matrix_boxes(4, top=1.62, height=4.96)
        else:
            boxes = self._chart_boxes(len(panels))
            ordered = panels

        for box, (key, title) in zip(boxes, ordered):
            il, it, iw, ih = self._card(s, *box, title)
            hm = md[key]
            path = self.work / f"md_{key}.png"
            R.draw_heatmap(path, hm["xLabels"], hm["yLabels"], hm["matrix"],
                           iw, ih, theme=self.theme, chart_id=f"multidim_{key}",
                           x_dimension=hm.get("xDimension"),
                           y_dimension=hm.get("yDimension"))
            self._place(s, path, il, it, iw, ih)
        self._footer(s)
        self._record("multidim", spec.get("title"),
                     f"{len(panels)} paired dimensions.", placeholder=False)

    def slide_funnel(self, spec):
        s = self._slide()
        self._header(s, spec.get("title", "Origination Funnel"),
                     "Weekly origination funnel by stage", accent=self.theme.peri)
        summary = self.d.funnel.get("summary", {}) or {}
        stages = self.d.funnel.get("stages", []) or ["KFI", "APPLICATION", "OFFER", "COMPLETED"]
        pretty = {"KFI": "KFI", "APPLICATION": "Application", "OFFER": "Offer",
                  "COMPLETED": "Completed", "WITHDRAWN": "Withdrawn"}
        rows = [{"label": pretty.get(st, st),
                 "v": (summary.get(st) or {}).get("latestFlowValue", 0)}
                for st in stages]
        title = "Latest weekly origination flow by stage"
        # Weekly flow needs ≥2 pipeline extracts. With a single extract, fall back to
        # the CURRENT pipeline funnel — case counts by stage — so the slide still
        # carries real data (matching the dashboard's single-period funnel).
        # Weekly FLOW is a balance; the single-extract fallback is a CASE COUNT.
        # They are different measures and must not be formatted the same way —
        # the fallback used to render amounts unlabelled, so a reader could not
        # tell which they were looking at.
        as_currency = True
        if not any(r["v"] for r in rows):
            stage_rows = self._stage_rows(self.d.pipeline.get("stageBreakdown", []),
                                          value_key="caseCount")
            rows = [{"label": r["label"], "v": r.get("caseCount", 0)} for r in stage_rows]
            title = "Current pipeline cases by stage"
            as_currency = False
        # The flow chart takes the upper band; the governed conversion rates sit
        # beneath it. The deck used to drop the conversion block entirely, which
        # is the single most-asked question of a growing book — and it is
        # already computed, with its lag and its sufficiency flag, by
        # ``evolution.pipeline_funnel_evolution``. Nothing is derived here.
        conv_rows = self._conversion_rows(summary)
        box = (Inches(self.CONTENT_L), Inches(1.62),
               Inches(self.CONTENT_R - self.CONTENT_L),
               Inches(3.30 if conv_rows else 4.95))
        ok = self._barlist_card(s, box, title, [r for r in rows if r.get("v")], "v",
                                currency=as_currency, cid="funnel")
        if conv_rows:
            self._conversion_strip(s, conv_rows,
                                   self.d.funnel.get("conversionLagWeeks"))
        self._footer(s)
        self._record("funnel", spec.get("title"), "", placeholder=not ok)

    def _conversion_rows(self, summary):
        """Governed forward conversion per stage, read straight off the funnel.

        ``weeklyRateValue`` is the evaluator's own forward conversion — average
        weekly flow into the stage over the lagged KFI stock. It is never
        recomputed here, and a stage the evaluator marks INSUFFICIENT is carried
        with that mark rather than silently presented as a rate to plan from.
        """
        rows = []
        for stage, block in (summary or {}).items():
            conv = (block or {}).get("conversion")
            if not conv or conv.get("weeklyRateValue") is None:
                continue
            rows.append({
                "label": (block.get("label") or stage),
                "rate": float(conv["weeklyRateValue"]),
                "sufficient": bool(conv.get("sufficient")),
                "weeks": conv.get("weeksInWindow"),
                "min_weeks": conv.get("minWeeks"),
                # WHAT THE CALCULATION ACTUALLY DID, per stage. The slide used
                # to word the basis line from a deck-level lag field and every
                # tile from the fixed phrase "of lagged KFI stock" — so a deck
                # whose rates were computed UNLAGGED said "(unlagged)" once and
                # "of lagged KFI stock" four times on the same page.
                "lag_applied": bool(conv.get("lagApplied")),
                "lag_weeks": conv.get("lagWeeks"),
            })
        return rows

    def _conversion_strip(self, slide, rows, lag_weeks):
        """Conversion rates as a labelled strip beneath the funnel.

        The basis is read from what the EVALUATOR did, per stage, not from a
        deck-level field — so the sentence at the top of the strip and the note
        under each rate can never describe two different calculations.
        """
        top = 5.10
        width = self.CONTENT_R - self.CONTENT_L
        lags = {(r.get("lag_weeks") if r.get("lag_applied") else None)
                for r in rows}
        basis = ("Forward conversion — average weekly flow into each stage over "
                 "the KFI stock")
        if lags == {None}:
            basis += ", unlagged: no KFI-to-completion lag was estimable."
        elif len(lags) == 1:
            basis += f", lagged {lags.pop()} week(s)."
        else:
            # Mixed is a real state — a stage with too little history gets no
            # lag — and stating one number for the page would misdescribe it.
            basis += "; the lag applied is stated per stage below."
        self._text(slide, Inches(self.CONTENT_L), Inches(top), Inches(width),
                   Inches(0.26), basis, size=9.5, color=self.theme.ink_400,
                   italic=True)
        cells = self._grid(min(len(rows), 4))
        for (left, w), row in zip(cells, rows[:4]):
            self._panel(slide, Inches(left), Inches(top + 0.32), Inches(w),
                        Inches(1.05), fill=self.theme.bg_panel_alt,
                        line=self.theme.line_soft)
            self._text(slide, Inches(left + 0.16), Inches(top + 0.42),
                       Inches(w - 0.32), Inches(0.26),
                       str(row["label"]).upper(), size=8.5,
                       color=self.theme.ink_400, bold=True)
            self._text(slide, Inches(left + 0.16), Inches(top + 0.66),
                       Inches(w - 0.32), Inches(0.34),
                       f"{row['rate']:.1f}%/wk", size=17, bold=True,
                       color=self.theme.ink_100 if row["sufficient"]
                       else self.theme.ink_500)
            if not row["sufficient"]:
                note = f"provisional — {row['weeks']} of {row['min_weeks']}+ weeks"
            elif row.get("lag_applied"):
                note = f"of KFI stock {row.get('lag_weeks')}wk earlier"
            else:
                note = "of current KFI stock (unlagged)"
            self._text(slide, Inches(left + 0.16), Inches(top + 1.00),
                       Inches(w - 0.32), Inches(0.24), note, size=8.5,
                       color=self.theme.ink_500)

    def slide_forecast_bridge(self, spec):
        s = self._slide()
        self._header(s, spec.get("title", "Forecast Bridge"),
                     "Funded + weighted pipeline → forecast funded",
                     accent=self.theme.mint)
        fb = self.d.forecast.get("forecastBridge") or {}
        if not fb or fb.get("forecastFundedBalance") is None:
            self._placeholder_body(s, "Forecast bridge requires funded + pipeline data.")
            self._footer(s)
            return self._record("forecast_bridge", spec.get("title"), "", placeholder=True)
        # Clarify the forecast is the CURRENT book's expected completions only.
        rd = self.d.reporting_date or "the reporting date"
        self._text(s, Inches(0.57), Inches(1.54), Inches(12.4), Inches(0.30),
                   f"Expected completions from the current book only (pipeline as of "
                   f"{rd}), weighted by historical stage conversion — not future new business.",
                   size=10, color=self.theme.ink_400, italic=True)
        # Full-width waterfall: split the weighted-pipeline block across expected
        # completion months (byCompletionMonth), Funded → +months → Forecast.
        funded = float(fb.get("fundedBalance") or 0)
        brk = (self.d.forecast.get("forecastBreakdowns") or {})
        months = [(str(m.get("month")), float(m.get("weightedExpectedFundedAmount") or 0))
                  for m in (brk.get("byCompletionMonth") or [])
                  if m.get("weightedExpectedFundedAmount")]
        months.sort()
        steps = [("Funded", funded, "base")]
        head, tail = months[:8], months[8:]
        for mth, val in head:
            steps.append((f"+ {mth}", val, "add"))
        if tail:
            steps.append(("+ Later", sum(v for _, v in tail), "add"))
        if len(steps) == 1:  # no monthly breakdown — single weighted block
            steps.append(("+ Weighted Pipeline",
                          float(fb.get("weightedExpectedFundedAmount") or 0), "add"))
        steps.append(("Forecast Funded", float(fb.get("forecastFundedBalance") or 0), "total"))
        # The bridge, then WHERE the forecast lands. ``forecast_breakdowns`` was
        # already being resolved for this deck and thrown away, while the
        # dashboard's Forecast view rendered both cuts — so the pack answered
        # "how much" and the screen also answered "where".
        by_region = self._forecast_breakdown_rows(brk, "byRegionCapped", "byRegion")
        by_ltv = self._forecast_breakdown_rows(brk, "byLtvBucketCapped", "byLtvBucket")
        cuts = [(k, lbl, rows, dim) for k, lbl, rows, dim in (
            ("fc_region", "Forecast balance by region", by_region, "region"),
            ("fc_ltv", "Forecast balance by LTV band", by_ltv, "ltv"),
        ) if rows]

        # The cuts beneath need room for every band a governed stratification
        # carries; the bridge gives up a quarter-inch rather than the cuts
        # dropping a region off the bottom of the panel.
        bridge_h = 2.46 if cuts else 4.64
        box = (Inches(self.CONTENT_L), Inches(1.92),
               Inches(self.CONTENT_R - self.CONTENT_L), Inches(bridge_h))
        il, it, iw, ih = self._card(s, *box,
                                    "Funded + weighted pipeline (by expected completion month) → Forecast")
        path = self.work / "bridge.png"
        render_bridge_waterfall(path, steps, iw, ih, theme=self.theme)
        self._place(s, path, il, it, iw, ih)

        if cuts:
            boxes = self._chart_boxes(len(cuts), top=1.92 + bridge_h + 0.20,
                                      height=6.55 - (1.92 + bridge_h + 0.20))
            for (cid, label, rows, dim), cbox in zip(cuts, boxes):
                # A bar list needs roughly a fifth of an inch per row to stay
                # legible. Show the bands the panel can actually carry rather
                # than compressing seven rows into an inch, which produced
                # overlapping labels.
                capacity = self._barlist_capacity(float(cbox[3]) / EMU_IN)
                self._barlist_card(s, cbox, label,
                                   self._fit_bars(rows, capacity), "balance",
                                   cid=cid, dimension=dim)
        self._footer(s)
        self._record("forecast_bridge", spec.get("title"), "", placeholder=False)

    @staticmethod
    def _forecast_breakdown_rows(breakdowns, capped_key, full_key):
        """Forecast-by-dimension rows in the deck's bar-list shape.

        ``workspace.forecast_breakdowns`` is the SAME payload the dashboard's
        Forecast view renders. The capped form (top 10 + Other) is preferred so
        a long region list stays legible, exactly as it does on screen.
        """
        rows = (breakdowns or {}).get(capped_key) or (breakdowns or {}).get(full_key) or []
        out = []
        for row in rows:
            value = row.get("forecastAmount")
            if value is None:
                value = row.get("pipelineAmount")
            if value is None:
                continue
            out.append({"label": str(row.get("key", "")), "balance": float(value)})
        return out

    def slide_forecast_projection(self, spec):
        s = self._slide()
        self._header(s, spec.get("title", "Forecast Projection"),
                     "Run-rate scale-up (downside / base / upside)",
                     accent=self.theme.mint)
        ex = self.d.extrapolation or {}
        model = {}
        for key in ("completionRunRateForecast", "kfiConversionForecast"):
            cand = ex.get(key) or {}
            if cand.get("available") and cand.get("projectedBalances"):
                model = cand
                break
        proj = model.get("projectedBalances", [])
        # Projection chart (top ~62%) + milestone table (bottom) when available.
        chart_box = ((Inches(0.55), Inches(1.62), Inches(12.25), Inches(3.05)) if proj
                     else (Inches(0.55), Inches(1.62), Inches(12.25), Inches(4.95)))
        il, it, iw, ih = self._card(s, *chart_box, "Projected funded balance")
        path = self.work / "projection.png"
        band_note = ""
        if proj:
            x = [str(p.get("month")) for p in proj]
            series = [
                {"name": "Downside", "values": [p.get("downside") for p in proj], "color": "#eb6f6f"},
                {"name": "Base", "values": [p.get("base") for p in proj], "color": "#7c9cf0"},
                {"name": "Upside", "values": [p.get("upside") for p in proj], "color": "#5ec6b8"},
            ]
            # INDISTINGUISHABLE SCENARIOS. Three lines that sit on top of each
            # other read as a rendering fault, and invite a reader to look for a
            # difference between scenarios that this book's run-rate history does
            # not produce. Where the band is immaterial the chart carries the
            # base case alone and the band is stated in words instead.
            series, band_note = self._scenario_series(proj, series)
            R.draw_lines(path, x, series, iw, ih, theme=self.theme, currency=True)
        else:
            render_placeholder_png(path, "", "Insufficient run-rate history for a "
                                   "scale-up projection", theme=self.theme,
                                   width_in=iw, height_in=ih)
        self._place(s, path, il, it, iw, ih)
        # Milestone dates to funding thresholds.
        milestones = model.get("milestones", [])
        if proj and milestones:
            box2 = (Inches(0.55), Inches(4.85), Inches(12.25), Inches(1.7))
            il, it, iw, ih = self._card(s, *box2, "Milestone dates to funding thresholds")
            cols = ["Threshold", "Downside", "Base", "Upside"]
            def _d(m, k):
                v = m.get(f"{k}Date")
                return "reached" if v == "reached" else (str(v) if v else "—")
            trows = [[m.get("thresholdLabel", ""), _d(m, "downside"), _d(m, "base"),
                      _d(m, "upside")] for m in milestones[:6]]
            tpath = self.work / "milestones.png"
            R.draw_table(tpath, cols, trows, iw, ih, theme=self.theme)
            self._place(s, tpath, il, it, iw, ih)
        if band_note:
            self._text(s, Inches(0.57), Inches(6.62), Inches(12.2), Inches(0.28),
                       band_note, size=9, color=self.theme.ink_500, italic=True)
        self._footer(s)
        self._record("forecast_projection", spec.get("title"), band_note,
                     placeholder=not proj)

    #: A scenario band narrower than this share of the base terminal balance is
    #: not a band a reader can act on, and three lines drawn through it overlap.
    SCENARIO_BAND_FLOOR = 0.03

    def _scenario_series(self, proj, series):
        """The scenario lines worth drawing, and the band stated in words.

        Returns the series unchanged where downside and upside genuinely
        separate from base by the horizon. Where they do not, returns the base
        case alone plus a sentence carrying the range — which is the same
        information, legibly, and does not imply a spread the run-rate history
        did not produce.
        """
        from .metric_resolver import compact_currency

        terminal = proj[-1] if proj else {}
        base = terminal.get("base")
        low, high = terminal.get("downside"), terminal.get("upside")
        if base in (None, 0) or low is None or high is None:
            return series, ""
        try:
            spread = (float(high) - float(low)) / abs(float(base))
        except (TypeError, ValueError, ZeroDivisionError):
            return series, ""
        if spread >= self.SCENARIO_BAND_FLOOR:
            return series, ""
        base_only = [s for s in series if s.get("name") == "Base"]
        return (base_only or series), (
            f"Downside and upside sit within {spread * 100:.1f}% of the base "
            f"case at the horizon ({compact_currency(low)} to "
            f"{compact_currency(high)}), so the scenarios are not separately "
            f"plotted. The milestone table below carries all three.")

    def slide_forecast_evolution(self, spec):
        """Was the prior forecast right? — the credibility page.

        The two charts are secondary here. What a funder wants from this page is
        a number: how far off has this forecaster been, and does it lean. Both
        are arithmetic over figures the governed evolution service has already
        reconciled — ``prior_forecast`` at period N IS the forecast period N-1
        published — so nothing is modelled and nothing is projected.
        """
        from . import forecast_accuracy as FA

        s = self._slide()
        accuracy = FA.measure(self.d.forecast_evolution)
        # FINDING-LED SUBTITLE. The title is stable; the subtitle states the
        # track record where one exists, and says nothing where one does not.
        strap = "Forecast funded balance across reporting runs"
        if accuracy.available:
            lean = accuracy.lean
            strap = (f"Typically {accuracy.error_pct:.1f}% from the outturn "
                     f"across {accuracy.observations} periods"
                     + (f", {lean}stated on average" if lean else
                        ", with no consistent lean"))
        self._header(s, spec.get("title", "Forecast Evolution"), strap,
                     accent=self.theme.mint)

        # THE VARIANCE CHART LEADS. The forecast's own travel is the supporting
        # view; whether it held is the question the page is titled with.
        charts = [
            {"id": "fvar", "title": "Actual funded vs the prior run's forecast",
             "series": [
                 {"name": "Prior-run forecast", "key": "prior_forecast", "color": "#e0a458"},
                 {"name": "Actual funded", "key": "funded_balance", "color": "#7c9cf0"}],
             "currency": True},
            {"id": "fevo", "title": "Forecast funded balance by run",
             "series": [
                 {"name": "Funded actual", "key": "funded_balance", "color": "#7c9cf0"},
                 {"name": "Weighted pipeline", "key": "weighted_expected_pipeline", "color": "#5ec6b8"},
                 {"name": "Forecast", "key": "forecast_funded_balance", "color": "#e0a458"}],
             "currency": True}]
        ph = self._evolution_lines(s, spec, self.d.forecast_evolution, charts,
                                   height=3.72)
        # The sentence the reader keeps — or, where there is no track record
        # yet, why there is not. Silence would read as an accurate forecast.
        self._text(s, Inches(self.CONTENT_L), Inches(5.60),
                   Inches(self.CONTENT_R - self.CONTENT_L), Inches(0.46),
                   FA.describe(accuracy), size=10, color=self.theme.ink_300,
                   italic=True, spacing=1.06)
        self._text(s, Inches(self.CONTENT_L), Inches(6.16),
                   Inches(self.CONTENT_R - self.CONTENT_L), Inches(0.44),
                   "Error is the actual funded balance against the forecast the "
                   "PRIOR run published, as a percentage of that forecast. Bias "
                   "is the mean signed error; a negative bias means the forecast "
                   "was high. No forecast is restated after the fact.",
                   size=8.5, color=self.theme.ink_500, spacing=1.06)
        self._footer(s)
        self._record("forecast_evolution", spec.get("title"), strap, placeholder=ph)

    def slide_risk(self, spec):
        s = self._slide()
        self._header(s, spec.get("title", "Risk Limits"),
                     "Concentration versus limits", accent=self.theme.peri)
        risk = self.d.risk or {}
        summary = risk.get("summary") or {}
        if not risk.get("available", True) and not risk.get("tests"):
            self._placeholder_body(s, risk.get("limitsReason")
                                   or "Risk-limit artifacts not present for this run.")
            self._footer(s)
            self.appendix.append("Risk monitor: no risk-limit artifact for this run.")
            return self._record("risk", spec.get("title"), "", placeholder=True)
        tiles = [
            {"label": "Tests passed", "value": summary.get("testsPassed", "—")},
            {"label": "Warnings", "value": summary.get("warnings", "—")},
            {"label": "Breaches", "value": summary.get("breaches", "—")},
            {"label": "Needs review", "value": summary.get("needsReview", "—")},
        ]
        tw = Inches(2.98)
        colors = [self.theme.rag["green"], self.theme.rag["amber"], self.theme.rag["red"],
                  self.theme.ink_400]
        for i, (tile, col) in enumerate(zip(tiles, colors)):
            l = Emu(int(Inches(0.55)) + i * int(Inches(3.08)))
            self._panel(s, l, Inches(1.6), tw, Inches(1.4), fill=self.theme.bg_panel_alt,
                        line=col, lw=1.2)
            self._text(s, l + Inches(0.2), Inches(1.82), tw, Inches(0.7),
                       str(tile["value"]), size=26, bold=True, color=col)
            self._text(s, l + Inches(0.2), Inches(2.62), tw, Inches(0.3),
                       tile["label"].upper(), size=9.5, color=self.theme.ink_400, bold=True)
        # tests table
        tests = risk.get("tests", [])[:10]
        box = (Inches(0.55), Inches(3.2), Inches(12.25), Inches(3.4))
        il, it, iw, ih = self._card(s, *box, "Limit tests")
        if tests:
            cols = ["Limit", "Actual", "Limit", "Headroom", "Status"]
            trows = [[str(t.get("label", "")), str(t.get("actualValue", "")),
                      str(t.get("limitValue", "")), str(t.get("headroom", "")),
                      str(t.get("status", ""))] for t in tests]
            path = self.work / "risk.png"
            R.draw_table(path, cols, trows, iw, ih, theme=self.theme, status_col=4)
            self._place(s, path, il, it, iw, ih)
        self._footer(s)
        self._record("risk", spec.get("title"), "", placeholder=False)

    def slide_watchlist(self, spec):
        """Portfolio Health and Watch Items — *what needs attention next?*

        At most five watch items and three observations, ranked by governed
        severity. No management actions are proposed: recommending what to do
        about a breach is a decision, and nothing in the evidence authorises the
        deck to make it.
        """
        s = self._slide()
        wl = self.d.watchlist or {}
        watch = wl.get("watch") or []
        observations = wl.get("observations") or []
        self._header(s, spec.get("title", "Portfolio Health and Watch Items"),
                     "Governed items requiring attention before the next period",
                     accent=self.theme.amber if watch else self.theme.peri)

        colour = {"concern": self.theme.rag["red"],
                  "attention": self.theme.rag["amber"]}
        # ONE band, shared by both columns, so the watch stack and the
        # observations panel start and finish on the same two lines whatever the
        # item count. A fixed 0.86in pitch left the stack ending above the panel
        # beside it for every count except five.
        BAND_TOP, BAND_BOTTOM = 1.62, 6.52
        band = BAND_BOTTOM - BAND_TOP
        top = BAND_TOP
        # The two columns are sized to what there is to say. With no watch items
        # a narrow left column left roughly 60% of the slide empty, which reads
        # as a rendering failure rather than as a clean bill of health.
        left_w = 7.7 if watch else 6.02
        obs_l = 8.5 if watch else 6.78
        obs_w = (self.CONTENT_R - 8.5) if watch else 6.02
        # ... and to what there ISN'T. The reverse case was not handled: watch
        # items beside an observations panel whose only content was the words
        # "None recorded." left a reader looking at four inches of empty box on
        # the page that is supposed to say what needs attention. With nothing to
        # observe there is no second column, and the watch stack takes the page.
        show_obs = bool(observations) or not watch
        if watch and not show_obs:
            left_w = self.CONTENT_R - self.CONTENT_L
        if watch:
            items = watch[:5]
            # The pitch is capped so a short list stays a short list. Spreading
            # one item over the whole band drew a card at the top of four inches
            # of nothing; the group is sized to the items and centred in the
            # band instead. At four and five items the cap does not bind and the
            # layout is unchanged.
            pitch = min(band / len(items), 1.42)
            row_h = min(1.30, pitch - 0.10)
            used = (len(items) - 1) * pitch + row_h
            # WHAT ELSE WAS TESTED. A page headed "requiring attention" that
            # shows one item tells a reader what was flagged and nothing about
            # how much was checked to flag it. Where the stack leaves room, the
            # governed checks are named underneath — the same list the all-clear
            # branch prints, because it is the same set of checks either way.
            # Where it does not, the stack is centred rather than left hanging
            # at the top of the band.
            checks_room = (band - used) >= 2.50
            top = BAND_TOP if checks_room else BAND_TOP + max(0.0, (band - used) / 2)
            for i, item in enumerate(items):
                t = Inches(top + i * pitch)
                accent = colour.get(item.severity, self.theme.ink_400)
                self._panel(s, Inches(0.55), t, Inches(left_w), Inches(row_h),
                            fill=self.theme.bg_panel_alt, line=self.theme.line_soft)
                chip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), t,
                                          Inches(0.05), Inches(row_h))
                chip.fill.solid()
                chip.fill.fore_color.rgb = self._rgb(accent)
                chip.line.fill.background()
                chip.shadow.inherit = False
                # A taller row (few items) gets a larger headline and more room
                # for the body, rather than the same 0.76in card with white space
                # underneath it.
                head_pt = 12.5 if row_h >= 1.0 else 11
                self._text(s, Inches(0.78), t + Inches(0.10), Inches(left_w - 0.5),
                           Inches(0.30), item.headline, size=head_pt, bold=True)
                self._text(s, Inches(0.78), t + Inches(0.10 + 0.32),
                           Inches(left_w - 0.5), Inches(max(0.28, row_h - 0.52)),
                           item.summary, size=9.5 if row_h >= 1.0 else 8.5,
                           color=self.theme.ink_400, spacing=1.06)
        else:
            # An empty list is a finding, and must be stated rather than leaving
            # the reader to wonder whether the check ran at all.
            self._panel(s, Inches(0.55), Inches(BAND_TOP), Inches(left_w),
                        Inches(band),
                        fill=self.theme.bg_panel_alt, line=self.theme.line_soft)
            self._text(s, Inches(0.85), Inches(1.94), Inches(left_w - 0.6),
                       Inches(0.72),
                       "No material watch items identified for this period.",
                       size=15, bold=True, color=self.theme.rag["green"])
            self._text(s, Inches(0.85), Inches(2.56), Inches(left_w - 0.6),
                       Inches(0.6),
                       "Every governed check ran and none cleared its "
                       "materiality threshold.", size=10,
                       color=self.theme.ink_300)
            self._text(s, Inches(0.85), Inches(3.16), Inches(left_w - 0.6),
                       Inches(0.28), "CHECKS PERFORMED", size=8.5,
                       color=self.theme.peri, bold=True)
            # Naming the checks is what separates "all clear" from "nothing ran".
            for i, line in enumerate(self.GOVERNED_CHECKS):
                self._text(s, Inches(0.95), Inches(3.48 + i * 0.34),
                           Inches(left_w - 0.7), Inches(0.3), f"·  {line}",
                           size=9.5, color=self.theme.ink_400)

        if watch and checks_room:
            cy = BAND_TOP + used + 0.34
            self._text(s, Inches(0.85), Inches(cy), Inches(left_w - 0.6),
                       Inches(0.28), "CHECKS PERFORMED", size=8.5,
                       color=self.theme.peri, bold=True)
            self._text(s, Inches(0.85), Inches(cy + 0.28), Inches(left_w - 0.6),
                       Inches(0.28),
                       "Every governed check below ran for this period. Those "
                       "that cleared their materiality threshold are listed "
                       "above; the rest did not.",
                       size=9, color=self.theme.ink_500, italic=True)
            for i, line in enumerate(self.GOVERNED_CHECKS):
                self._text(s, Inches(0.95), Inches(cy + 0.64 + i * 0.30),
                           Inches(left_w - 0.7), Inches(0.28), f"·  {line}",
                           size=9.5, color=self.theme.ink_400)

        # Observations column.
        if not show_obs:
            self._footer(s)
            self._record(spec.get("id", "watchlist"), spec.get("title"),
                         f"{len(watch)} watch item(s).")
            return
        self._panel(s, Inches(obs_l), Inches(BAND_TOP), Inches(obs_w), Inches(band),
                    fill=self.theme.bg_panel, line=self.theme.line)
        self._text(s, Inches(obs_l + 0.22), Inches(1.78), Inches(obs_w - 0.4),
                   Inches(0.3), "OBSERVATIONS", size=9, bold=True,
                   color=self.theme.peri)
        # Observations are distributed down the SAME band rather than stacked at
        # the top of it, so a single observation does not sit above four inches
        # of empty panel.
        shown_obs = observations[:3]
        obs_pitch = (band - 0.60) / max(len(shown_obs), 1)
        y = BAND_TOP + 0.52
        for item in shown_obs:
            self._text(s, Inches(obs_l + 0.22), Inches(y), Inches(obs_w - 0.4),
                       Inches(0.46), item.headline, size=10, bold=True)
            self._text(s, Inches(obs_l + 0.22), Inches(y + 0.5), Inches(obs_w - 0.4),
                       Inches(min(1.30, obs_pitch - 0.60)), item.summary[:220],
                       size=9 if obs_pitch >= 1.3 else 8.5,
                       color=self.theme.ink_400, spacing=1.06)
            y += obs_pitch
        if not observations:
            self._text(s, Inches(obs_l + 0.22), Inches(BAND_TOP + 0.52),
                       Inches(obs_w - 0.4),
                       Inches(0.3), "None recorded.", size=10,
                       color=self.theme.ink_400)
        self._footer(s)
        self._record(spec.get("id", "watchlist"), spec.get("title"),
                     f"{len(watch)} watch item(s).")

    def slide_concentration(self, spec):
        """Concentration Tests and Headroom — *am I within my limits?*

        Three states, kept visually and verbally distinct: CURRENT funded (the
        only actual), EXPECTED forecast, and the ALL-PIPELINE-CONVERTS stress.
        Presenting the stress as an expectation would misstate the risk, so it
        is labelled a stress everywhere it appears.
        """
        from . import concentration as C

        s = self._slide()
        env = self.d.concentration or {}
        rows = C.adapt_tests(env)
        self._header(s, spec.get("title", "Concentration Tests and Headroom"),
                     "Utilisation of contractual limits — current, expected and stress",
                     accent=self.theme.peri)
        if not rows:
            self._placeholder_body(s, "No governed concentration tests configured.")
            self._footer(s)
            return self._record(spec.get("id", "concentration"), spec.get("title"),
                                "", placeholder=True)

        # DIRECTION OF TRAVEL. The prior governed value comes from the history
        # service, which re-evaluates today's approved configuration against
        # each historical frame — so "moved toward the limit" is a statement
        # about the book, not about a changed definition.
        rows = C.attach_stress(C.attach_history(rows, self.d.concentration_history))
        summary = C.summarise(env, rows)
        top = C.select_tests(rows)
        forward = C.forward_states_available(env)
        historic = any(r.get("prior_value") is not None for r in top)

        # -- summary strip --------------------------------------------------
        tiles = [
            ("TESTS", str(summary["tests"]), self.theme.ink_400),
            ("IN BREACH", str(summary["breaches"]),
             self.theme.rag["red"] if summary["breaches"] else self.theme.rag["green"]),
            ("WARNING", str(summary["warnings"]),
             self.theme.rag["amber"] if summary["warnings"] else self.theme.ink_400),
            ("FORECAST BREACH", str(summary["expected_breaches"]) if forward else "—",
             self.theme.rag["amber"] if summary["expected_breaches"] else self.theme.ink_400),
            ("STRESS BREACH", str(summary["stress_breaches"]) if forward else "—",
             self.theme.ink_400),
        ]
        tw = Inches(2.35)
        for i, (label, value, colour) in enumerate(tiles):
            l = Emu(int(Inches(0.55)) + i * int(Inches(2.45)))
            self._panel(s, l, Inches(1.56), tw, Inches(0.92),
                        fill=self.theme.bg_panel_alt, line=self.theme.line_soft)
            self._text(s, l + Inches(0.18), Inches(1.66), tw - Inches(0.3),
                       Inches(0.28), label, size=8, color=self.theme.ink_400,
                       bold=True)
            self._text(s, l + Inches(0.18), Inches(1.92), tw - Inches(0.3),
                       Inches(0.42), value, size=19, bold=True, color=colour)

        # -- utilisation bars ------------------------------------------------
        bars = [{"label": r["label"], "utilisation": r["utilisation"] or 0,
                 "status": r["status"],
                 "expectedUtilisation": r["expected_utilisation"] if forward else None,
                 "stressUtilisation": r["stress_utilisation"] if forward else None}
                for r in top]
        il, it, iw, ih = self._card(s, Inches(0.55), Inches(2.66), Inches(6.5),
                                    Inches(3.62), "Utilisation of limit")
        path = self.work / "conc_util.png"
        R.draw_utilisation_tests(path, bars, iw, ih, theme=self.theme)
        self._place(s, path, il, it, iw, ih)

        # -- the numbers behind the bars -------------------------------------
        # Rendered as NATIVE PowerPoint text, not a chart image: a covenant table
        # is the page an investor reads closely, copies figures out of and zooms
        # into, and an image is none of those things.
        self._panel(s, Inches(7.28), Inches(2.66), Inches(5.52), Inches(3.62),
                    fill=self.theme.bg_panel, line=self.theme.line)
        self._text(s, Inches(7.5), Inches(2.82), Inches(5.1), Inches(0.34),
                   "Current position against limit", size=12.5, bold=True)
        # Column offsets carry an explicit WIDTH so the rightmost column always
        # ends inside the panel. Deriving the width from whether dx was zero
        # pushed the fifth column 0.2in off the slide once the Expected column
        # appeared, which only happens when forward states exist.
        # The table reads left to right as the sequence a covenant actually
        # moves through: where it was, where it is, where it is expected to go,
        # and the limit it is measured against.
        # Six columns would squeeze the test name below the width at which a
        # governed limit name is still legible, so where prior AND expected are
        # both present the HEADROOM column gives way: it is limit less current,
        # both of which are on the row, and the detail line states it in words.
        if forward and historic:
            cols = [("Test", 0.0, 1.98, PP_ALIGN.LEFT),
                    ("Prior", 2.04, 0.68, PP_ALIGN.RIGHT),
                    ("Current", 2.78, 0.70, PP_ALIGN.RIGHT),
                    ("Expected", 3.54, 0.78, PP_ALIGN.RIGHT),
                    ("Limit", 4.38, 0.72, PP_ALIGN.RIGHT)]
        elif forward:
            cols = [("Test", 0.0, 2.00, PP_ALIGN.LEFT),
                    ("Current", 2.06, 0.70, PP_ALIGN.RIGHT),
                    ("Expected", 2.82, 0.72, PP_ALIGN.RIGHT),
                    ("Limit", 3.60, 0.66, PP_ALIGN.RIGHT),
                    ("Headroom", 4.32, 0.76, PP_ALIGN.RIGHT)]
        elif historic:
            cols = [("Test", 0.0, 1.96, PP_ALIGN.LEFT),
                    ("Prior", 2.02, 0.78, PP_ALIGN.RIGHT),
                    ("Current", 2.86, 0.78, PP_ALIGN.RIGHT),
                    ("Limit", 3.70, 0.72, PP_ALIGN.RIGHT),
                    ("Headroom", 4.48, 0.84, PP_ALIGN.RIGHT)]
        else:
            cols = [("Test", 0.0, 2.20, PP_ALIGN.LEFT),
                    ("Current", 2.35, 0.85, PP_ALIGN.RIGHT),
                    ("Limit", 3.30, 0.80, PP_ALIGN.RIGHT),
                    ("Headroom", 4.20, 0.90, PP_ALIGN.RIGHT)]
        # Pitch derives from how many tests there ARE, so one test does not sit
        # in a sliver at the top of an empty card and five do not collide. When
        # tests overflow the slide, the last line of the band is reserved for
        # saying so — the note belongs to this table, and putting it below the
        # panel would run it into the takeaway.
        rest = C.overflow(rows)
        band = (6.20 - 3.62) - (0.30 if rest else 0.0)
        row_h = min(0.80, band / max(len(top), 1))
        value_pt = 10 if len(top) > 2 else 11
        # With few tests the rows would otherwise sit in the top fifth of the
        # card above a large blank, which reads as a rendering fault rather than
        # as a short list. They are given a detail line and centred in the band.
        detail = row_h >= 0.72
        row_span = row_h + (0.22 if detail else 0.0)
        top_y = 3.62 + max(0.0, (band - len(top) * row_span) / 2)
        # The column header travels WITH the rows. Pinned at a constant it would
        # be orphaned at the top of the card whenever a short list is centred.
        for label, dx, cw, align in cols:
            self._text(s, Inches(7.5 + dx), Inches(top_y - 0.32), Inches(cw),
                       Inches(0.26), label, size=8.5, color=self.theme.ink_400,
                       bold=True, align=align)
        for i, r in enumerate(top):
            y = Inches(top_y + i * row_span)
            status_colour = self.theme.rag.get(
                {"breach": "red", "warning": "amber"}.get(r["status"], "green"),
                self.theme.ink_300)
            values = [(self._fit_label(r["label"], cols[0][2]), self.theme.ink_100)]
            if historic:
                values.append((C.format_measure(r["prior_value"], r["unit"])
                               if r.get("prior_value") is not None else "—",
                               self.theme.ink_500))
            values.append((C.format_measure(r["value"], r["unit"]), status_colour))
            if forward:
                values.append((C.format_measure(r["expected_value"], r["unit"])
                               if r["expected_value"] is not None else "—",
                               self.theme.peri))
            values.append((C.format_measure(r["limit"], r["unit"]),
                           self.theme.ink_300))
            if not (forward and historic):
                values.append((C.format_headroom(r["headroom"], r["unit"])
                               if r["headroom"] is not None else "—",
                               self.theme.ink_300))
            for i, ((value, colour), (_label, dx, cw, align)) in enumerate(
                    zip(values, cols)):
                self._text(s, Inches(7.5 + dx), y, Inches(cw), Inches(0.3),
                           str(value), size=9.5 if i == 0 else value_pt,
                           color=colour, align=align,
                           bold=(align == PP_ALIGN.RIGHT and colour is status_colour))
            # A test that passes today but is forecast to cross says BOTH, and
            # says which is which — "PASS · breaches 2026-07" reads as a
            # contradiction rather than as a forward-looking warning.
            status_line = r["status"].upper()
            moved = C.travel(r)
            if moved:
                status_line += f" · {moved} since {r.get('prior_date') or 'the prior period'}"
            # HEADROOM MUST APPEAR SOMEWHERE. Its column gives way when prior and
            # expected are both present, and the detail line below only renders
            # when the rows are tall enough — so on a four-test page it would
            # otherwise vanish entirely from the one slide about headroom.
            if not (forward and historic) or r["headroom"] is None:
                pass
            elif not detail:
                status_line += (
                    f" · {C.format_headroom(abs(r['headroom']), r['unit'])} "
                    + ("of headroom" if r["headroom"] >= 0 else "beyond the limit"))
            if r.get("expected_breach") and r.get("breach_horizon"):
                status_line += f" now · forecast breach {r['breach_horizon']}"
            elif r.get("expected_breach"):
                status_line += " now · forecast breach"
            elif r.get("stress_breach"):
                status_line += " now · breaches under stress only"
            self._text(s, Inches(7.5), Emu(int(y) + int(Inches(0.28))),
                       Inches(5.0), Inches(0.22), status_line,
                       size=7.5, color=status_colour)
            if detail:
                # Only where the evaluator produced them. A deployment that
                # evaluates neither forward state gets the current position
                # restated in words, not two invented ones.
                bits = []
                if forward and r["expected_value"] is not None:
                    bits.append(
                        f"Expected {C.format_measure(r['expected_value'], r['unit'])}"
                        + (f" ({r['expected_utilisation']:.0f}% of limit)"
                           if r["expected_utilisation"] is not None else ""))
                if forward and r["stress_value"] is not None:
                    # A stress that eases the test, or moves it not at all, is
                    # explained rather than printed bare — an "under stress"
                    # figure BELOW the current one reads as a fault.
                    explained = C.stress_note(r)
                    bits.append(explained or (
                        "under the all-pipeline-converts stress "
                        f"{C.format_measure(r['stress_value'], r['unit'])}"))
                if (not bits or (forward and historic)) and r["headroom"] is not None:
                    bits.append(
                        f"{C.format_headroom(abs(r['headroom']), r['unit'])} "
                        + ("of headroom remaining" if r["headroom"] >= 0
                           else "beyond the limit"))
                if bits:
                    self._text(s, Inches(7.5), Emu(int(y) + int(Inches(0.50))),
                               Inches(5.1), Inches(0.22), " · ".join(bits),
                               size=7.5, color=self.theme.ink_500)

        if rest:
            # Never assert the hidden tests are within limit. They are ranked
            # BELOW the shown ones, which on a book with five breaches does not
            # mean they pass — and "N further tests within limit" would then be
            # a false statement on the one slide that must not make one.
            adverse = sum(1 for r in rest
                          if r.get("status") in (C.STATUS_BREACH, C.STATUS_WARNING))
            note = (f"{len(rest)} further test{'s' if len(rest) != 1 else ''} "
                    f"ranked below these")
            note += (f", of which {adverse} in breach or warning." if adverse
                     else f"; nearest is {self._fit_label(rest[0]['label'], 2.2, 8)}.")
            self._text(s, Inches(7.5), Inches(3.62 + len(top) * row_h + 0.04),
                       Inches(5.1), Inches(0.24), note,
                       size=8, color=(self.theme.rag.get("amber") if adverse
                                      else self.theme.ink_500))

        # -- deterministic takeaway + source disclosure -----------------------
        takeaway = self._concentration_takeaway(summary, top, forward)
        # The conclusion can run to several clauses when many states are in
        # play, so it gets the full content width and steps down rather than
        # clipping — this is the sentence the reader keeps.
        self._text(s, Inches(0.57), Inches(6.38), Inches(12.2), Inches(0.44),
                   takeaway, size=10 if len(takeaway) < 190 else 9,
                   color=self.theme.ink_300, italic=True, spacing=1.06)
        disclosure = C.source_disclosure(env)
        if disclosure:
            self._text(s, Inches(0.57), Inches(6.86), Inches(9.4), Inches(0.26),
                       disclosure, size=8, color=self.theme.ink_500)
        self._footer(s)
        self._record(spec.get("id", "concentration"), spec.get("title"),
                     f"{summary['tests']} governed tests.")

    @staticmethod
    def _fit_label(text: str, width_in: float, size_pt: float = 9.5) -> str:
        """Trim a label to one line at *width_in*, so it cannot wrap into the
        status line beneath it."""
        capacity = max(6, int(width_in * 72 / (size_pt * 0.55)))
        text = str(text)
        return text if len(text) <= capacity else text[: capacity - 1].rstrip() + "…"

    def _concentration_takeaway(self, summary, rows, forward) -> str:
        """One deterministic sentence over the governed evidence.

        Composed from what the configuration actually produced: it never assumes
        a direction, a scenario is available, or that any particular state
        exists. Where a state was not evaluated it is simply not mentioned,
        rather than being reported as "none".
        """
        from . import concentration as C
        parts: List[str] = []
        if summary["breaches"]:
            worst = next((r for r in rows if r["status"] == C.STATUS_BREACH), None)
            lead = (f"{summary['breaches']} test"
                    f"{'s are' if summary['breaches'] != 1 else ' is'} in breach "
                    f"at the reporting date")
            if worst:
                lead += f", led by {worst['label']}"
            parts.append(lead + ".")
        else:
            lead = "All current tests remain within limit."
            # "Nothing is in breach" is a weak conclusion on its own. Where the
            # evaluator identified the tightest test, name it and its position:
            # an investor's next question is always how much room is left.
            # Only when nothing at all is adverse. With a test in warning range
            # the warning clause below is the answer, and naming the closest as
            # well would say the same thing twice.
            closest = summary.get("closest") if not summary["warnings"] else None
            if closest and closest.get("utilisation") is not None:
                lead += (f" The closest is {closest['label']} at "
                         f"{C.format_measure(closest['value'], closest['unit'])} "
                         f"against a "
                         f"{C.format_measure(closest['limit'], closest['unit'])} "
                         f"limit, {closest['utilisation']:.0f}% utilised.")
            parts.append(lead)
        if summary["warnings"]:
            parts.append(f"{summary['warnings']} test"
                         f"{'s are' if summary['warnings'] != 1 else ' is'} in "
                         f"warning range.")
        if forward and summary["expected_breaches"]:
            horizon = next((r.get("breach_horizon") for r in rows
                            if r.get("expected_breach") and r.get("breach_horizon")),
                           None)
            parts.append(
                f"{summary['expected_breaches']} "
                f"{'is' if summary['expected_breaches'] == 1 else 'are'} expected "
                f"to breach on the governed forecast"
                + (f", crossing around {horizon}" if horizon else "") + ".")
        if forward and summary["stress_breaches"]:
            only_stress = summary.get("stress_only_breaches",
                                      summary["stress_breaches"])
            if only_stress > 0:
                further = "further " if summary["expected_breaches"] else ""
                parts.append(
                    f"{only_stress} {further}test"
                    f"{'s' if only_stress != 1 else ''} breach"
                    f"{'' if only_stress != 1 else 'es'} only under the "
                    f"all-pipeline-converts stress, which is a stress rather "
                    f"than the expected outcome.")
        if not forward:
            parts.append("Forward-looking states were not evaluated for this "
                         "book, so only the current position is shown.")
        return " ".join(parts)

    def slide_methodology(self, spec):
        """Kept as an alias of the investor-safe Data and Methodology page.

        The former version of this slide listed endpoint paths, internal compute
        function names and resolved source filenames. Those are implementation
        details, not methodology, and had no place in a client-facing pack — the
        single page below states the basis of preparation in business language.
        """
        return self.slide_appendix(spec)

    def slide_appendix(self, spec):
        """Data and Methodology — investor-safe.

        The previous version of this page printed the generator's own
        diagnostics: discovery roots (absolute filesystem paths), resolved
        source filenames and internal compute-function names. None of that
        belongs in a client-facing document. This page states the same facts in
        business language — what the report covers, as at when, what was
        excluded and why — and never where the bytes live.
        """
        s = self._slide()
        self._header(s, spec.get("title", "Data and Methodology"),
                     "Scope, source dates, coverage and basis of preparation")
        p = self.d.portfolio
        d = self.d.diagnostics or {}

        left = []
        left.append("REPORTING SCOPE")
        if p is not None:
            left.append(f"   {p.scope_label} portfolio")
            for book in p.portfolios[:4]:
                left.append(f"   ·  {book.label} — {book.type_label}")
            if len(p.portfolios) > 4:
                left.append(f"   ·  and {len(p.portfolios) - 4} further book(s)")
        else:
            left.append("   Scope unavailable for this run.")

        left.append("")
        left.append("SOURCE REPORTING DATES")
        if p is not None and p.type_reporting_dates:
            from .deck_context import type_label
            for ptype, date in sorted(p.type_reporting_dates.items()):
                left.append(f"   {type_label(ptype)} funded data as at {_pretty_date(date)}.")
        elif self.d.reporting_date:
            left.append(f"   Funded portfolio data as at {_pretty_date(self.d.reporting_date)}.")
        pipe_date = (self.d.pipeline or {}).get("pipelineAsOfDate")
        if pipe_date:
            left.append(f"   Pipeline data as at {_pretty_date(pipe_date)}.")
        if p is not None and p.has_mixed_reporting_dates:
            left.append("   Constituent books are reported as at different dates;")
            left.append("   the total combines them.")

        # THE CLAIM MATCHES THE EVIDENCE. The pack used to assert that every
        # figure was "identical to the management dashboard". That is true of the
        # tiles, stratifications, cross-tabs, cohort series, concentration tests
        # and the balance bridge — and was NOT true while economic values were
        # still derived independently downstream of the engine. The composition
        # shares, forecast accuracy and limit direction have since moved to
        # shared owners; a handful of derivations remain, so the claim states
        # what is provable today rather than what will be provable when they do.
        right = ["BASIS OF PREPARATION",
                 "   Figures are generated deterministically from governed MI",
                 "   outputs using shared reporting definitions, so this pack and",
                 "   the management dashboard read the same measures for the same",
                 "   portfolio and reporting date.",
                 "   Commentary is generated deterministically from those figures.",
                 "   No language model is used in its production.",
                 "   Averages are stated with their weighting basis. Measures on",
                 "   different bases are not intended to divide into one another."]
        conc = self.d.concentration or {}
        if conc.get("tests"):
            from . import concentration as C
            disclosure = C.source_disclosure(conc)
            if disclosure:
                right.append(f"   Concentration limits: {disclosure.lower()}.")
        right.append("")
        right.append("COVERAGE")
        # Prefer the diagnostic; fall back to the periods the funded history
        # actually resolved. The diagnostic is not populated on every path, and
        # a coverage block that states the pipeline extract count and nothing
        # about the funded book is the wrong half of the answer.
        cuts = d.get("fundedCutsFound") or len(
            (getattr(self.d, "funded_evolution", {}) or {}).get("periods") or ())
        if cuts:
            right.append(f"   {cuts} funded reporting period(s) available.")
        else:
            right.append("   One funded reporting period available.")
        snaps = d.get("pipelineSnapshotsFound") or 0
        right.append(f"   {snaps} weekly pipeline extract(s) available."
                     if snaps else "   No weekly pipeline extracts available.")

        right.extend(self._capability_lines())

        # THE OMISSIONS GO ON THE LEFT. They are a statement about SCOPE, which
        # is what the left column is, and the right column already carries the
        # basis, the coverage and every measure this book cannot report — it ran
        # off the bottom of the slide and over the footer, so the sections a
        # reader most needs to see listed were the ones printed past the edge of
        # the page. The left column was two-thirds empty throughout.
        if self.omissions:
            left.append("")
            left.append("SECTIONS NOT INCLUDED")
            # GROUPED BY REASON. Three consecutive lines repeating "the pipeline
            # is small relative to the funded book" spend three of the six lines
            # this block has room for on one fact, and push other sections behind
            # "and N further" where the reader cannot see them at all.
            grouped: List[tuple] = []
            for o in self.omissions:
                for i, (reason, titles) in enumerate(grouped):
                    if reason == o.reason:
                        grouped[i][1].append(o.title)
                        break
                else:
                    grouped.append((o.reason, [o.title]))
            shown, hidden = grouped[:6], grouped[6:]
            for reason, titles in shown:
                left.append(f"   {', '.join(titles)}: {reason}.")
            if hidden:
                count = sum(len(t) for _r, t in hidden)
                left.append(f"   and {count} further section(s).")

        self._column_text(s, left, Inches(0.6), Inches(6.0))
        self._column_text(s, right, Inches(6.95), Inches(5.85))
        self._footer(s)
        self._record(spec.get("id", "appendix"), spec.get("title"), "")

    #: Why a measure is absent, in the language a funder reads. The registry
    #: distinguishes these deliberately: "we lack a field" is a data request to
    #: the client, "this book has no such thing" is a property of the asset,
    #: and "that needs a model we do not run" is a boundary Trakt has drawn.
    #: Collapsing them into "not available" is what makes a pack look evasive.
    _CAPABILITY_WORDING = {
        "NOT_APPLICABLE": "not applicable to this portfolio",
        "UNAVAILABLE": "required data not supplied",
        "ASSUMPTION_REQUIRED": "would require an assumption Trakt does not make",
        "MODEL_REQUIRED": "would require behavioural modelling Trakt does not perform",
        "METHODOLOGY_NOT_APPROVED": "methodology not yet approved",
    }
    #: Grouped in the order a reader can act on: what they can fix, what is
    #: inherent, what Trakt has chosen not to do.
    _CAPABILITY_ORDER = ("UNAVAILABLE", "NOT_APPLICABLE", "ASSUMPTION_REQUIRED",
                         "MODEL_REQUIRED", "METHODOLOGY_NOT_APPROVED")

    def _capability_lines(self):
        """"Measures not reported for this book", and the reason for each.

        Read from the published capability registry's own resolution against
        this portfolio's canonical shape — the same catalogue the API and the
        agent tools answer from. Nothing here branches on what the book IS: a
        capability declares the economic conditions it needs, and this page
        reports which of them this tape did not meet.
        """
        from trakt_core import capability as cap

        resolved = self.d.capabilities or {}
        if not resolved:
            return []
        registry = cap.load_registry()
        grouped = {}
        for metric, availability in resolved.items():
            status = getattr(availability, "status", cap.AVAILABLE)
            if status == cap.AVAILABLE:
                continue
            entry = registry.get(metric)
            grouped.setdefault(status, []).append(
                str(getattr(entry, "name", None) or metric))
        if not grouped:
            return []

        lines = ["", "MEASURES NOT REPORTED FOR THIS BOOK"]
        budget = 4
        for status in self._CAPABILITY_ORDER:
            names = sorted(grouped.get(status) or ())
            if not names or budget <= 0:
                continue
            shown = ", ".join(names[:3])
            if len(names) > 3:
                shown += f" and {len(names) - 3} other measure(s)"
            wording = self._CAPABILITY_WORDING.get(status, "not available")
            lines.append(f"   {shown} — {wording}.")
            budget -= 1
        return lines

    #: The methodology column's box, in inches. Text does not shrink to fit a
    #: PowerPoint textbox: it simply draws past the bottom, over the footer and
    #: off the slide, and python-pptx reports nothing. The height is a real
    #: limit and has to be treated as one.
    _COLUMN_TOP, _COLUMN_HEIGHT = 1.6, 5.3

    @staticmethod
    def _column_extent(lines, width_in: float, size: float) -> float:
        """Roughly how tall this column will render, in inches.

        Wrapping is estimated rather than measured — the renderer is the only
        thing that knows for certain — so it is deliberately pessimistic: a
        column judged slightly too tall loses half a point of type, while one
        judged too short runs off the page.
        """
        per_line = max(10.0, width_in * 125.0 / size)   # chars that fit on a line
        total = 0.0
        for line in lines:
            heading = bool(line) and line == line.upper() and not line.startswith(" ")
            pt = (size - 1) if heading else size
            wraps = max(1, -(-len(line) // int(per_line)))
            total += wraps * (pt * 1.25) / 72.0 + (5 if heading else 2) / 72.0
        return total

    def _column_text(self, slide, lines, left, width):
        """A column of the methodology page; section headings pick up the accent."""
        width_in = int(width) / EMU_IN
        size = 10.0
        for candidate in (10.0, 9.5, 9.0, 8.5, 8.0):
            size = candidate
            if self._column_extent(lines, width_in, candidate) <= self._COLUMN_HEIGHT:
                break
        box = slide.shapes.add_textbox(left, Inches(self._COLUMN_TOP), width,
                                       Inches(self._COLUMN_HEIGHT))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = para.add_run()
            run.text = line
            heading = bool(line) and line == line.upper() and not line.startswith(" ")
            run.font.size = Pt(size - 1 if heading else size)
            run.font.bold = heading
            run.font.name = self.theme.font_sans
            run.font.color.rgb = self._rgb(
                self.theme.peri if heading else self.theme.ink_300)
            para.space_after = Pt(5 if heading else 2)

    # ------------------------------------------------------------------ helpers
    def _placeholder_body(self, slide, msg):
        path = self.work / f"ph_{self._page}.png"
        render_placeholder_png(path, "", msg, theme=self.theme, width_in=12.2,
                               height_in=4.9)
        self._place(slide, path, Inches(0.55), Inches(1.62), 12.2, 4.9)

    def _bullets(self, slide, lines, *, size=12):
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.62), Inches(12.1), Inches(5.2))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(size)
            run.font.name = self.theme.font_sans
            run.font.color.rgb = self._rgb(self.theme.ink_300)
            p.space_after = Pt(6)

    # ------------------------------------------------------------------- build
    _DISPATCH = {
        "cover": "slide_cover", "kpi_summary": "slide_kpi_summary",
        "executive": "slide_executive",
        "exec_insights": "slide_exec_insights",
        "portfolio_composition": "slide_portfolio_composition",
        "portfolio_comparison": "slide_portfolio_comparison",
        "movement_drivers": "slide_movement_drivers",
        "balance_movement": "slide_balance_movement",
        "funded_stock": "slide_funded_stock",
        "portfolio_projections": "slide_portfolio_projections",
        "watchlist": "slide_watchlist",
        "strat_barlists": "slide_strat", "multidim": "slide_multidim", "geo": "slide_geo",
        "funded_evolution": "slide_funded_evolution", "cohorts": "slide_cohorts",
        "cohort_progression": "slide_cohort_progression",
        "pipeline_summary": "slide_pipeline", "pipeline_evolution": "slide_pipeline_evolution",
        "funnel": "slide_funnel", "origination_flow": "slide_origination_flow",
        "pipeline_movement": "slide_pipeline_movement",
        "forecast_bridge": "slide_forecast_bridge",
        "forecast_projection": "slide_forecast_projection",
        "forecast_evolution": "slide_forecast_evolution", "risk": "slide_risk",
        "concentration": "slide_concentration",
        "methodology": "slide_methodology", "appendix": "slide_appendix",
    }

    def build(self, slides: List[Dict[str, Any]], output: str | Path) -> Dict[str, Any]:
        """Render the slides this portfolio justifies, and record the rest.

        Composition runs BEFORE any rendering, so a slide that would have had
        nothing to show never reaches the deck — an investor pack contains no
        "no data available" pages. Everything dropped is carried into the
        appendix with its reason, because a silent omission is indistinguishable
        from a book that had nothing to report.
        """
        from .composition import build_facts, select_slides

        facts = build_facts(self.d)
        selected, omissions = select_slides(slides, self.d, facts)
        self.omissions = list(omissions)
        self.facts = dict(facts)

        # Record what each renderer actually draws. A bar list becomes a PNG, so
        # its category order is not recoverable from the finished file; the
        # record is how a publication gate — and a parity test — can see it.
        with R.record_renders() as drawn:
            for spec in selected:
                handler = getattr(self, self._DISPATCH.get(spec.get("type"), ""), None)
                if handler is None:
                    continue
                handler(spec)
            self.rendered = list(drawn)
        out = Path(output)
        out.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(out))
        return {"output": str(out), "slides": self.records,
                "rendered": self.rendered,
                "currency_code": getattr(self.d, "currency_code", None),
                "coverage_notes": self.appendix,
                "omitted_slides": [o.to_dict() for o in self.omissions],
                "facts": self.facts,
                "portfolio_context": (self.d.portfolio.to_dict()
                                      if self.d.portfolio else None),
                "insights": self._insight_records()}

    def _insight_records(self) -> Dict[str, Any]:
        """The executive summary in a serialisable form, for the run manifest."""
        brief = self.d.insights or {}
        return {
            "insight_version": brief.get("insight_version"),
            "status": brief.get("status"),
            "count": brief.get("count", 0),
            "headlines": [str(getattr(i, "headline", ""))
                          for i in (brief.get("insights") or [])],
            "omitted": [o.to_dict() for o in (brief.get("omitted") or [])],
        }
