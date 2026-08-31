#!/usr/bin/env python3
"""scripts/pptx_visual_qa.py — generate representative investor packs and QA them.

Builds a deck for each representative book THROUGH THE REACT ROUTE
(``POST /mi/decks/generate`` → ``GET /mi/decks/download``), then inspects the
downloaded PowerPoint for the defects a test suite does not see: overlapping or
clipped shapes, unreadable type, orphaned titles, empty chart frames, duplicated
slide titles, a foreign currency symbol, and bar lists drawn out of the governed
bucket order.

It is a QA instrument, not a test: it reports, it does not assert. Run it, read
the report, look at the decks it wrote.

    python scripts/pptx_visual_qa.py [--out artifacts/pptx_qa]
"""

from __future__ import annotations

import argparse
import io
import json
import os
import shutil
import sys
import tempfile
import time
from pathlib import Path

_REPO = Path(__file__).resolve().parents[1]
if str(_REPO) not in sys.path:
    sys.path.insert(0, str(_REPO))

import pandas as pd  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

_CENTRAL = "18_central_lender_tape.csv"
SLIDE_W, SLIDE_H = Emu(12192000), Emu(6858000)   # 13.333in x 7.5in


# --------------------------------------------------------------------------- #
# Representative books.
# --------------------------------------------------------------------------- #

def _loan(i, pid, ptype, *, ltv, balance, region, age, rate, cut, origination):
    return {
        "unique_identifier": f"{pid}_L{i:05d}",
        "source_portfolio_id": pid, "source_portfolio_type": ptype,
        "source_portfolio_label": pid.replace("_", " ").title(),
        "current_outstanding_balance": balance,
        "current_principal_balance": balance,
        "original_principal_balance": balance * 1.05,
        "current_valuation_amount": balance / (ltv / 100.0),
        "original_valuation_amount": balance / (ltv / 100.0),
        "current_loan_to_value": ltv,
        "original_loan_to_value": (ltv / 100.0) - 0.02,
        "current_interest_rate": rate,
        "youngest_borrower_age": age,
        "geographic_region_collateral": region,
        "collateral_geography": region,
        "origination_channel": "Direct", "broker_channel": "Direct",
        "product_type": "Lifetime Mortgage",
        "origination_date": origination, "data_cut_off_date": cut,
    }


_REGIONS = ["London", "South East", "Wales", "Scotland", "North West",
            "Yorkshire and The Humber", "East of England"]
_LTVS = [22.0, 27.0, 34.0, 45.0, 52.0, 63.0, 74.0, 86.0]
_AGES = [58, 63, 67, 71, 76, 81, 87]
_RATES = [5.4, 6.2, 6.8, 7.1, 7.6]


#: Constituent books, as (portfolio id, type, share of the whole). A funder pack
#: for a warehouse or a forward-flow facility is nearly always reporting on more
#: than one book, and the pages that matter most there — the stock stack, the
#: per-book forward view, the composition table — are exactly the ones a
#: single-book fixture cannot exercise.
_SINGLE = (("direct_001", "direct", 1.0),)
_MULTI = (("direct_001", "direct", 0.52),
          ("acquired_001", "acquired", 0.31),
          ("direct_002", "direct", 0.17))


def _pid_for(j, books):
    """The constituent book a loan belongs to — a function of the LOAN, so a
    loan keeps its book for life and the per-book stack is a real series."""
    r = ((j * 7919) % 1000) / 1000.0
    running = 0.0
    for pid, ptype, share in books:
        running += share
        if r < running:
            return pid, ptype
    return books[-1][0], books[-1][1]


def _book(first, last, cut, *, vintages, scale=1.0, seed=7, books=_SINGLE,
          exit_below=None):
    """Loans ``first..last`` as at ``cut``.

    A loan's IDENTITY and its static attributes — book, region, borrower age,
    LTV band, origination vintage — are functions of the loan number alone, so
    the same loan is the same loan in every period. Only its BALANCE moves, with
    ``scale``. That is what makes the economic bridge meaningful: the continuing
    leg is a real movement on real continuing loans rather than an artefact of
    regenerating the book each period.

    ``exit_below`` marks the loans that will be absent from the NEXT period.
    They carry redemption, default and maturity evidence here — on the opening
    frame, which is where the classifier reads it from — so the exit leg splits
    by reason instead of collapsing into the unevidenced bucket.
    """
    rows = []
    for j in range(first, last + 1):
        pid, ptype = _pid_for(j, books)
        ltv = _LTVS[(j * 3) % len(_LTVS)]
        balance = (85_000 + ((j * 37_000 + seed * 11_000) % 420_000)) * scale
        row = _loan(
            j, pid, ptype, ltv=ltv, balance=float(balance),
            region=_REGIONS[(j * 5) % len(_REGIONS)],
            age=_AGES[(j * 2) % len(_AGES)],
            rate=_RATES[j % len(_RATES)], cut=cut,
            origination=vintages[(j * 7) % len(vintages)])
        if exit_below is not None and j < exit_below:
            kind = j % 3
            if kind == 0:
                row["loan_redemption_flag"] = "Y"
            elif kind == 1:
                row["default_date"] = cut
            else:
                row["maturity_date"] = cut
        rows.append(row)
    return rows


#: (name, [(run_id, reporting_date, first_loan, last_loan, vintages)], books,
#:  description). Loans leave from the bottom of the range and arrive at the
#: top, so a period-on-period pair has both a new leg and an exit leg.
_SEASONED_VINTAGES = ["2019-03-01", "2020-07-01", "2021-05-01",
                      "2022-09-01", "2023-06-01", "2024-04-01"]
_GROWING_VINTAGES = ["2022-05-01", "2023-08-01", "2024-06-01",
                     "2025-09-01", "2026-03-01"]

BOOKS = {
    # A — a newly originated single book. One period: no history, no movement.
    "new_book": (
        [("mi_2026_06", "2026-06-30", 1, 90,
          ["2026-02-01", "2026-04-01", "2026-06-01"])],
        _SINGLE,
        "A newly originated book: one reporting period, one origination year."),
    # B — a seasoned single book. Five periods, six vintages, real exits.
    "seasoned_book": (
        [("mi_2026_02", "2026-02-28", 1, 300, _SEASONED_VINTAGES),
         ("mi_2026_03", "2026-03-31", 9, 312, _SEASONED_VINTAGES),
         ("mi_2026_04", "2026-04-30", 18, 325, _SEASONED_VINTAGES),
         ("mi_2026_05", "2026-05-31", 26, 338, _SEASONED_VINTAGES),
         ("mi_2026_06", "2026-06-30", 33, 350, _SEASONED_VINTAGES)],
        _SINGLE,
        "A seasoned book: five reporting periods, six vintages, loans leaving "
        "with evidence each period."),
    # C — the same seasoning across THREE constituent books, which is what a
    # warehouse or forward-flow facility actually reports on.
    "multi_seasoned": (
        [("mi_2026_02", "2026-02-28", 1, 300, _SEASONED_VINTAGES),
         ("mi_2026_03", "2026-03-31", 9, 312, _SEASONED_VINTAGES),
         ("mi_2026_04", "2026-04-30", 18, 325, _SEASONED_VINTAGES),
         ("mi_2026_05", "2026-05-31", 26, 338, _SEASONED_VINTAGES),
         ("mi_2026_06", "2026-06-30", 33, 350, _SEASONED_VINTAGES)],
        _MULTI,
        "A seasoned multi-book portfolio: three constituent books, five "
        "periods, evidenced exits."),
    # D — a growing multi-book portfolio: origination dominates the story.
    "multi_growing": (
        [("mi_2026_04", "2026-04-30", 1, 180, _GROWING_VINTAGES),
         ("mi_2026_05", "2026-05-31", 5, 205, _GROWING_VINTAGES),
         ("mi_2026_06", "2026-06-30", 10, 260, _GROWING_VINTAGES)],
        _MULTI,
        "A growing multi-book portfolio: three periods, still originating "
        "hard, three constituent books."),
    # Retained under its old name so an existing invocation still resolves.
    "mixed_book": (
        [("mi_2026_04", "2026-04-30", 1, 180, _GROWING_VINTAGES),
         ("mi_2026_05", "2026-05-31", 5, 205, _GROWING_VINTAGES),
         ("mi_2026_06", "2026-06-30", 10, 240, _GROWING_VINTAGES)],
        _SINGLE,
        "A growing single book: three periods, five vintages, still originating."),
}


def write_book(root: Path, client: str, key: str) -> Path:
    periods, books, _ = BOOKS[key]
    for idx, (run_id, date, first, last, vintages) in enumerate(periods):
        central = root / client / run_id / "central"
        central.mkdir(parents=True, exist_ok=True)
        # The loans that will be gone next period are evidenced HERE.
        nxt = periods[idx + 1][2] if idx + 1 < len(periods) else None
        pd.DataFrame(_book(first, last, date, vintages=vintages,
                           scale=1.0 + 0.03 * idx, seed=7 + idx,
                           books=books, exit_below=nxt)).to_csv(
            central / _CENTRAL, index=False)
    return root


# --------------------------------------------------------------------------- #
# The pipeline and the approved limits — so the pipeline, forecast and
# concentration slides are actually exercised rather than composed away.
# --------------------------------------------------------------------------- #

#: A live origination pipeline: several weeks of dated extracts, spread across
#: the same bands and regions as the funded book so the two lenses are
#: comparable. Written in the governed weekly-extract layout that
#: ``pipeline_contract.discover_pipeline_sources`` walks.
_PIPE_WEEKS = ("2026-05-15", "2026-05-29", "2026-06-12", "2026-06-26")
_STAGES = ("KFI", "APPLICATION", "OFFER", "COMPLETED")


def _pipeline_case(i, week_index, *, region, ltv, balance, stage, prob):
    """One live pipeline case.

    ``Account Number`` is the contract's own alias for
    ``pipeline_case_identifier``, and it is what a real lender extract carries
    (see ``tests/fixtures/pipeline_history_5w``, which maps it cleanly). The
    fixture used to write ``unique_identifier`` — the ESMA RREL1 name for a
    FUNDED underlying exposure, which is not a pre-funding case key and is not
    an alias here. That made the prepared frame case-anonymous, which is a
    property of this fixture and was never a property of the contract.

    The id is a function of the CASE, not of the week, so a case persists across
    extracts and its amount can move without changing who it is.
    """
    return {
        "Account Number": f"ACC{i:05d}",
        "current_outstanding_balance": balance,
        "current_principal_balance": balance,
        "current_valuation_amount": balance / (ltv / 100.0),
        "current_loan_to_value": ltv,
        "youngest_borrower_age": _AGES[i % len(_AGES)],
        "current_interest_rate": _RATES[i % len(_RATES)],
        "collateral_geography": region,
        "geographic_region_obligor": region,
        "broker_channel": "Direct",
        "product_type": "Lifetime Mortgage",
        "pipeline_status": "pipeline",
        "pipeline_stage": stage,
        "completion_probability": prob,
    }


def write_pipeline(root: Path, client: str, weeks=_PIPE_WEEKS):
    """Dated weekly extracts under the client's pipeline tree."""
    for w, week in enumerate(weeks):
        folder = root / client / "pipeline" / week
        folder.mkdir(parents=True, exist_ok=True)
        rows = []
        for i in range(28 + w * 4):
            # A CASE PROGRESSES, AND NEW ONES ARRIVE. The stage used to be a
            # function of the case index alone, so every case sat in the same
            # stage for ever: a pipeline in which nothing moves, which is not a
            # pipeline, and on which the funnel, the conversion rates and the
            # stage-movement reconciliation all had nothing to measure.
            #
            # Each case now joins in its own week and walks the ladder from
            # there, so every extract carries a mix of stages, cases move
            # between extracts, and some complete and leave the live stock.
            joined = i % 5
            if joined > w:
                continue                       # not originated yet
            stage = _STAGES[min(w - joined, len(_STAGES) - 1)]
            rows.append(_pipeline_case(
                i + 1, w, region=_REGIONS[(i * 3 + w) % len(_REGIONS)],
                ltv=_LTVS[(i * 2 + w) % len(_LTVS)],
                balance=float(95_000 + ((i * 29_000 + w * 7_000) % 360_000)),
                stage=stage,
                prob={"KFI": 0.25, "APPLICATION": 0.5,
                      "OFFER": 0.85, "COMPLETED": 1.0}[stage]))
        pd.DataFrame(rows).to_csv(
            folder / f"M2L_KFI_and_Pipeline_{week.replace('-', '_')}.csv",
            index=False)
    return root


def approve_limits(client: str):
    """Commit an operator-approved concentration configuration, through the
    real store — the same path the operations console writes."""
    from apps.blob_trigger_app.storage import open_storage
    from mi_agent.concentration_tests.models import (
        ActiveConfiguration, ActiveTest, ApprovalRecord, SourceEvidence)
    from mi_agent.concentration_tests.store import ConcentrationStore

    limits = (("London concentration", 30.0, ["London"]),
              ("South East concentration", 45.0, ["South East"]),
              ("Wales concentration", 40.0, ["Wales"]),
              ("Scotland concentration", 50.0, ["Scotland"]))
    tests = [
        ActiveTest(metric_id="geo_region_share", threshold=thr, operator="max",
                   display_name=name, parameters={"regions": regions},
                   evidence=SourceEvidence(source_reference="facility_agreement.pdf",
                                           source_text="Schedule 8, clause 3.1"),
                   approval=ApprovalRecord(decision="approved", operator="Operator",
                                           decided_at="2026-01-10T09:00:00+00:00"))
        for name, thr, regions in limits]
    config = ActiveConfiguration(client_id=client, version=1,
                                 activated_by="Operator",
                                 library_version="1.0.0", tests=tests)
    ConcentrationStore(open_storage(),
                       container="operations-control").commit_configuration(
        client, config)


# --------------------------------------------------------------------------- #
# Generation, through the React route.
# --------------------------------------------------------------------------- #

def generate(tmp: Path, client: str, book: str, currency: str, out: Path, *,
             pipeline: bool = True, limits: bool = True):
    from fastapi.testclient import TestClient

    root = write_book(tmp / "runs", client, book)
    if pipeline:
        write_pipeline(root, client)
    config = tmp / f"cfg_{client}.yaml"
    config.write_text(f"portfolio:\n  base_currency: {currency}\n", encoding="utf-8")

    os.environ.update({
        "MI_AGENT_ONBOARDING_OUTPUT_ROOT": str(root),
        "MI_AGENT_CLIENT_ID": client,
        "MI_AGENT_PIPELINE_ROOT": str(root),
        "TRAKT_LOCAL_BLOB_ROOT": str(tmp / "blob"),
        "TRAKT_INVESTOR_PPTX_PERSIST": "true",
        "TRAKT_INVESTOR_PPTX_ON_DEMAND": "true",
        "MI_AGENT_AUTH_ENABLED": "false",
        "TRAKT_MI_CLIENT_CONFIG": str(config),
    })
    for key in ("MI_AGENT_DECK_ROOT", "AZURE_STORAGE_CONNECTION_STRING",
                "TRAKT_BLOB_CONNECTION"):
        os.environ.pop(key, None)

    os.environ["TRAKT_STORAGE_BACKEND"] = "file"
    os.environ["TRAKT_RUNTIME_MODE"] = "test"
    if limits:
        approve_limits(client)

    from mi_agent_api import currency as currency_mod, datasets, deck_generation
    currency_mod._load_client_config.cache_clear()
    datasets._CLIENT_CURRENCY_CACHE.clear()
    from mi_agent_api import data_source
    data_source.reset_cache()
    deck_generation.reset_jobs()

    from mi_agent_api.app import app
    api = TestClient(app)
    accepted = api.post("/mi/decks/generate", json={})
    if accepted.status_code != 202:
        return None, {"error": f"generate -> {accepted.status_code}: {accepted.text}"}
    job_id = accepted.json()["jobId"]
    deadline, body = time.time() + 600, None
    while time.time() < deadline:
        body = api.get(f"/mi/decks/generate/{job_id}").json()
        if body["state"] in ("completed", "blocked", "failed"):
            break
        time.sleep(0.5)
    if not body or body["state"] != "completed":
        return None, {"error": f"generation {body and body['state']}", "job": body}

    got = api.get("/mi/decks/download")
    if got.status_code != 200:
        return None, {"error": f"download -> {got.status_code}"}
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_bytes(got.content)

    sidecars = sorted(root.glob("**/investor_pack.pptx.preflight.json"))
    side = json.loads(sidecars[-1].read_text(encoding="utf-8")) if sidecars else {}
    return got.content, {"job": body, "sidecar": side}


# --------------------------------------------------------------------------- #
# Visual checks against the rendered file.
# --------------------------------------------------------------------------- #

_SYMBOLS = {"GBP": "£", "EUR": "€", "USD": "$"}
_MIN_PT = 7.5


def inspect(content: bytes, currency: str, side: dict):
    deck = Presentation(io.BytesIO(content))
    findings, titles = [], []

    for n, slide in enumerate(deck.slides, start=1):
        boxes, texts = [], []
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            l, t2 = int(shape.left), int(shape.top)
            r, b = l + int(shape.width or 0), t2 + int(shape.height or 0)
            if l < -1000 or t2 < -1000 or r > int(SLIDE_W) + 1000 or b > int(SLIDE_H) + 1000:
                findings.append(f"slide {n}: shape extends off-canvas ({l},{t2})-({r},{b})")
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size.pt < _MIN_PT:
                            findings.append(
                                f"slide {n}: {run.font.size.pt}pt type — below "
                                f"{_MIN_PT}pt: {run.text[:40]!r}")
            if shape.shape_type == 13:                      # PICTURE
                boxes.append((l, t2, r, b))
        # Charts must not overlap each other.
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                a, c = boxes[i], boxes[j]
                if a[0] < c[2] and c[0] < a[2] and a[1] < c[3] and c[1] < a[3]:
                    findings.append(f"slide {n}: two chart images overlap")
        if not texts:
            findings.append(f"slide {n}: no text at all — orphaned slide")
        else:
            titles.append(texts[0].split("\n")[0])

    duplicates = sorted({t for t in titles if titles.count(t) > 1})
    if duplicates:
        findings.append(f"duplicated slide titles: {duplicates}")

    body = "\n".join(sh.text_frame.text for s in deck.slides for sh in s.shapes
                     if sh.has_text_frame)
    expected = _SYMBOLS.get(currency, "")
    for code, sym in _SYMBOLS.items():
        if code != currency and sym in body:
            findings.append(f"foreign currency symbol {sym} in a {currency} deck")
    if expected and expected not in body:
        findings.append(f"no {currency} amounts ({expected}) anywhere in the deck")

    # Bucket order, from the render record the drawing functions wrote.
    from mi_agent_api import presentation as P
    for entry in side.get("rendered") or ():
        if entry.get("kind") != "barlist" or not entry.get("dimension"):
            continue
        dim = entry["dimension"]
        if not P.governed_ladder(dim):
            continue
        drawn = [str(c) for c in entry.get("categories") or ()]
        if drawn != P.order_categories(drawn, dimension=dim):
            findings.append(f"{dim} bar list drawn out of the governed order: {drawn}")
        if not drawn:
            findings.append(f"{dim} bar list drawn with no categories")

    # Bar lists that cannot fit their own rows. The renderer scales its type to
    # the row band and stops at a floor; below that floor the labels collide,
    # which is a defect no text inspection can see because the bars are an image.
    pictures = {}
    for n, slide in enumerate(deck.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == 13 and shape.height:
                pictures.setdefault(n, []).append(int(shape.height) / 914400)
    for entry in side.get("rendered") or ():
        if entry.get("kind") != "barlist":
            continue
        rows = len(entry.get("categories") or ())
        if rows > 14:
            findings.append(
                f"{entry.get('chart')}: {rows} bars in one list — beyond what a "
                f"slide panel can carry legibly")

    gates = side.get("preflight", {}) or {}
    for failed in gates.get("failed_gates") or ():
        findings.append(f"publication gate FAILED: {failed}")

    return {"slides": len(deck.slides), "titles": titles, "findings": findings}


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="artifacts/pptx_qa")
    args = ap.parse_args()
    out_root = _REPO / args.out
    out_root.mkdir(parents=True, exist_ok=True)

    # The five shapes the pack has to be right for:
    #   A  single-book, newly originated, GBP   — the shortest honest deck
    #   B  single-book, seasoned, GBP           — history, cohorts, exits
    #   C  multi-book, seasoned, GBP            — stack, per-book forward view
    #   D  multi-book, growing, GBP             — origination is the story
    #   E  seasoned, EUR, no pipeline           — currency, and composition in
    #                                             the other direction
    cases = [("new_book", "GBP", True, True),
             ("seasoned_book", "GBP", True, True),
             ("multi_seasoned", "GBP", True, True),
             ("multi_growing", "GBP", True, True),
             ("seasoned_book", "EUR", False, True)]
    report = {}
    for book, ccy, has_pipeline, has_limits in cases:
        name = f"{book}_{ccy.lower()}"
        client = f"qa{abs(hash(name)) % 9973}"
        tmp = Path(tempfile.mkdtemp(prefix=f"qa_{name}_"))
        try:
            path = out_root / f"{name}.pptx"
            content, meta = generate(tmp, client, book, ccy, path,
                                     pipeline=has_pipeline, limits=has_limits)
            if content is None:
                report[name] = {"ok": False, **meta}
                print(f"[{name}] FAILED: {meta.get('error')}")
                continue
            result = inspect(content, ccy, meta.get("sidecar") or {})
            report[name] = {"ok": not result["findings"], "deck": str(path), **result}
            flag = "clean" if not result["findings"] else f"{len(result['findings'])} finding(s)"
            print(f"[{name}] {result['slides']} slides — {flag}  -> {path}")
            for f in result["findings"]:
                print(f"    · {f}")
        finally:
            shutil.rmtree(tmp, ignore_errors=True)

    (out_root / "qa_report.json").write_text(
        json.dumps(report, indent=2, default=str), encoding="utf-8")
    print(f"\nReport: {out_root / 'qa_report.json'}")
    return 0 if all(v.get("ok") for v in report.values()) else 1


if __name__ == "__main__":
    sys.exit(main())
